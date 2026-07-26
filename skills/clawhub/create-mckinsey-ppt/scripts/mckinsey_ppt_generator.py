#!/usr/bin/env python3
"""
McKinsey-style consulting PPT Generator
========================================
Usage:
  python mckinsey_ppt_generator.py --title "报告标题" --subtitle "副标题" \\
    --sections "章节1,章节2,章节3" \\
    --content <json_file_or_json_string> \\
    --output output.pptx

The --content JSON should follow this schema:
{
  "cover_subtitle": "副标题",
  "cover_extra": "额外信息行",
  "date": "2026年7月",
  "sections": [
    {
      "title": "第一章 市场概览",
      "type": "cards",        // "cards" | "list" | "table" | "text"
      "items": [
        {"title": "标题1", "desc": "描述内容…", "color": "navy"},
        ...
      ]
    },
    {
      "title": "第二章 分析框架",
      "type": "text",
      "content": "详细文本内容…"
    },
    {
      "title": "第三章 数据对比",
      "type": "table",
      "headers": ["指标","基准值","目标值"],
      "rows": [
        ["OEE","65%","85%"],
        ...
      ],
      "col_widths": [3.0, 3.0, 3.0]
    }
  ]
}
"""

import json
import sys
import argparse
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ── Color Palette ──
NAVY = RGBColor(0x00, 0x2B, 0x5C)
DARK_BLUE = RGBColor(0x00, 0x3D, 0x7A)
MID_BLUE = RGBColor(0x00, 0x6D, 0xBA)
LIGHT_BLUE = RGBColor(0x4D, 0xA8, 0xE0)
TEAL = RGBColor(0x00, 0x96, 0x88)
PURPLE = RGBColor(0x6A, 0x1B, 0x9A)
DARK_GRAY = RGBColor(0x33, 0x33, 0x33)
MED_GRAY = RGBColor(0x66, 0x66, 0x66)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
RED_ACCENT = RGBColor(0xC0, 0x28, 0x2E)
GREEN_ACC = RGBColor(0x2E, 0x7D, 0x32)
ORANGE_ACC = RGBColor(0xE6, 0x7E, 0x22)
GOLD = RGBColor(0xD4, 0xA0, 0x17)

COLOR_MAP = {
    "navy": NAVY, "dark_blue": DARK_BLUE, "mid_blue": MID_BLUE,
    "light_blue": LIGHT_BLUE, "teal": TEAL, "purple": PURPLE,
    "red": RED_ACCENT, "green": GREEN_ACC, "orange": ORANGE_ACC,
    "gold": GOLD, "dark_gray": DARK_GRAY, "med_gray": MED_GRAY,
}

def resolve_color(c):
    if isinstance(c, RGBColor):
        return c
    if isinstance(c, str) and c.lower() in COLOR_MAP:
        return COLOR_MAP[c.lower()]
    # Try hex
    if isinstance(c, str) and c.startswith("#"):
        h = c.lstrip("#")
        return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))
    return DARK_BLUE  # fallback

# ── Slide Helpers ──

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

def add_header_bar(slide, title_text, subtitle_text=None):
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, Inches(0.95))
    bar.fill.solid(); bar.fill.fore_color.rgb = DARK_BLUE; bar.line.fill.background()
    a = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(0.95), SLIDE_W, Inches(0.04))
    a.fill.solid(); a.fill.fore_color.rgb = LIGHT_BLUE; a.line.fill.background()
    tf = bar.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = title_text
    p.font.size = Pt(26); p.font.color.rgb = WHITE; p.font.bold = True; p.font.name = 'Arial'
    p.alignment = PP_ALIGN.LEFT; tf.margin_left = Inches(0.6); tf.margin_top = Inches(0.12)
    if subtitle_text:
        b2 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(1.0), SLIDE_W, Inches(0.45))
        b2.fill.solid(); b2.fill.fore_color.rgb = RGBColor(0xE8,0xEF,0xF5); b2.line.fill.background()
        t2 = b2.text_frame; t2.margin_left = Inches(0.6)
        p2 = t2.paragraphs[0]; p2.text = subtitle_text
        p2.font.size = Pt(14); p2.font.color.rgb = MED_GRAY; p2.font.name = 'Arial'

def add_footer(slide, page_num):
    f = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(7.15), SLIDE_W, Inches(0.35))
    f.fill.solid(); f.fill.fore_color.rgb = RGBColor(0xF0,0xF2,0xF5); f.line.fill.background()
    tf = f.text_frame; tf.margin_left = Inches(0.6); tf.margin_top = Inches(0.05)
    p = tf.paragraphs[0]; p.text = f"CONFIDENTIAL AND PROPRIETARY — {page_num}"
    p.font.size = Pt(9); p.font.color.rgb = MED_GRAY; p.font.name = 'Arial'

def tb(slide, l, t, w, h, text, sz=12, bold=False, color=DARK_GRAY, align=PP_ALIGN.LEFT, fn='Arial'):
    box = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = box.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = text
    p.font.size = Pt(sz); p.font.bold = bold; p.font.color.rgb = color; p.font.name = fn; p.alignment = align
    return box

def add_table(slide, l, t, w, h, rows, cols, data, cw=None):
    ts = slide.shapes.add_table(rows, cols, Inches(l), Inches(t), Inches(w), Inches(h))
    table = ts.table
    if cw:
        for i, v in enumerate(cw):
            table.columns[i].width = Inches(v)
    for r in range(rows):
        for c in range(cols):
            cell = table.cell(r, c); cell.text = ""
            tf = cell.text_frame; tf.word_wrap = True
            p = tf.paragraphs[0]; p.text = str(data[r][c])
            p.font.size = Pt(10); p.font.name = 'Arial'; p.alignment = PP_ALIGN.CENTER
            if r == 0:
                cell.fill.solid(); cell.fill.fore_color.rgb = NAVY
                p.font.color.rgb = WHITE; p.font.bold = True
                cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            else:
                cell.vertical_anchor = MSO_ANCHOR.MIDDLE
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(0xF0,0xF4,0xF8) if r%2==0 else WHITE
                p.font.color.rgb = DARK_GRAY
                if c == 0: p.alignment = PP_ALIGN.LEFT
    return ts

def card(slide, l, t, w, h, fc, text="", sz=11, fg=WHITE, bold=False, border=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(l), Inches(t), Inches(w), Inches(h))
    shape.fill.solid(); shape.fill.fore_color.rgb = fc
    if border: shape.line.color.rgb = border; shape.line.width = Pt(1)
    else: shape.line.fill.background()
    tf = shape.text_frame; tf.word_wrap = True
    tf.margin_left = Inches(0.08); tf.margin_right = Inches(0.08)
    tf.margin_top = Inches(0.04); tf.margin_bottom = Inches(0.04)
    p = tf.paragraphs[0]; p.text = text
    p.font.size = Pt(sz); p.font.color.rgb = fg; p.font.bold = bold; p.font.name = 'Arial'
    p.alignment = PP_ALIGN.CENTER
    return shape


# ════════════════════════════════════════════════════════════════
# SLIDE BUILDERS
# ════════════════════════════════════════════════════════════════

def build_cover(prs, title, subtitle, cover_extra="", date_str="", client=""):
    """Build a McKinsey-style cover slide."""
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    bg = sl.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
    bg.fill.solid(); bg.fill.fore_color.rgb = NAVY; bg.line.fill.background()
    accent_line = sl.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(2.8), SLIDE_W, Inches(0.06))
    accent_line.fill.solid(); accent_line.fill.fore_color.rgb = LIGHT_BLUE; accent_line.line.fill.background()
    tb(sl, 0.8, 1.0, 11.7, 0.6, client, sz=18, bold=False, color=RGBColor(0x99,0xAA,0xBB))
    tb(sl, 0.8, 1.5, 11.7, 1.2, title, sz=36, bold=True, color=WHITE)
    tb(sl, 0.8, 3.2, 11.7, 0.8, subtitle, sz=24, bold=True, color=LIGHT_BLUE)
    if cover_extra:
        tb(sl, 0.8, 4.2, 11.7, 0.6, cover_extra, sz=16, color=RGBColor(0xBB,0xCC,0xDD))
    tb(sl, 0.8, 5.5, 5, 0.5, "CONFIDENTIAL AND PROPRIETARY", sz=11, color=MED_GRAY)
    if date_str:
        tb(sl, 0.8, 5.9, 5, 0.5, date_str, sz=14, color=WHITE)
    tb(sl, 9.5, 5.9, 3, 0.5, "麦肯锡风格 · 战略咨询版", sz=12, color=MED_GRAY, align=PP_ALIGN.RIGHT)


def build_executive_summary(prs, data, page_num):
    """Build executive summary slide with card-style insights."""
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    add_header_bar(sl, data.get("title", "执行摘要"), data.get("subtitle", ""))
    add_footer(sl, page_num)

    items = data.get("items", [])
    for i, item in enumerate(items):
        color = resolve_color(item.get("color", DARK_BLUE))
        x = 0.25 + i * 2.58
        width = 2.45
        card(sl, x, 1.3, width, 5.5, WHITE, border=RGBColor(0xDD,0xDD,0xDD))
        tb_bar = sl.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(1.3), Inches(width), Inches(0.5))
        tb_bar.fill.solid(); tb_bar.fill.fore_color.rgb = color; tb_bar.line.fill.background()
        tf = tb_bar.text_frame; tf.margin_left = Inches(0.08)
        p = tf.paragraphs[0]; p.text = f" {item.get('title','')}"; p.font.size = Pt(11)
        p.font.bold = True; p.font.color.rgb = WHITE; p.font.name = 'Arial'
        tb(sl, x + 0.1, 2.0, width - 0.2, 4.5, item.get("desc", ""), sz=10.5, color=DARK_GRAY)


def build_cards_slide(prs, section, page_num):
    """Build a slide with card-style items in a grid."""
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    add_header_bar(sl, section.get("title", ""), section.get("subtitle", ""))
    add_footer(sl, page_num)

    items = section.get("items", [])
    cols = section.get("cols", 3)
    rows = section.get("rows", 2)

    cw = min(12.7 / cols, 4.0)
    ch = min(5.5 / rows, 2.5)
    start_x = (12.7 - cw * cols) / 2 + 0.3
    start_y = 1.3
    gap = 0.1

    for i, item in enumerate(items):
        col = i % cols
        row = i // cols
        if row >= rows: break
        x = start_x + col * (cw + gap)
        y = start_y + row * (ch + gap + 0.2)

        color = resolve_color(item.get("color", MID_BLUE))
        # Card body
        card(sl, x, y, cw, ch, WHITE, border=color)
        # Header bar
        hdr = sl.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(cw), Inches(0.5))
        hdr.fill.solid(); hdr.fill.fore_color.rgb = color; hdr.line.fill.background()
        tf = hdr.text_frame; p = tf.paragraphs[0]; p.text = f" {item.get('title','')}"
        p.font.size = Pt(13); p.font.bold = True; p.font.color.rgb = WHITE; p.font.name = 'Arial'
        tf.margin_top = Inches(0.06)
        tb(sl, x + 0.1, y + 0.6, cw - 0.2, ch - 0.7, item.get("desc", ""), sz=10, color=DARK_GRAY)


def build_text_slide(prs, section, page_num):
    """Build a text-heavy slide."""
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    add_header_bar(sl, section.get("title", ""), section.get("subtitle", ""))
    add_footer(sl, page_num)

    content = section.get("content", "")
    tb(sl, 0.6, 1.5, 12.1, 5.3, content, sz=13, color=DARK_GRAY)


def build_table_slide(prs, section, page_num):
    """Build a table slide."""
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    add_header_bar(sl, section.get("title", ""), section.get("subtitle", ""))
    add_footer(sl, page_num)

    headers = section.get("headers", [])
    rows = section.get("rows", [])
    col_widths = section.get("col_widths", None)
    data = [headers] + rows

    add_table(sl, 0.3, 1.3, 12.7, min(5.5, 0.4 * len(data)), len(data), len(headers), data, col_widths)

    # Optional footnote
    footnote = section.get("footnote", "")
    if footnote:
        tb(sl, 0.5, 6.2, 12, 0.7, footnote, sz=11, color=MED_GRAY)


def build_list_slide(prs, section, page_num):
    """Build a list/bullet-point slide."""
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    add_header_bar(sl, section.get("title", ""), section.get("subtitle", ""))
    add_footer(sl, page_num)

    items = section.get("items", [])
    for i, item in enumerate(items):
        y = 1.3 + i * 0.85
        color = resolve_color(item.get("color", RED_ACCENT))
        bar = sl.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(y), Inches(0.08), Inches(0.65))
        bar.fill.solid(); bar.fill.fore_color.rgb = color; bar.line.fill.background()
        tb(sl, 0.8, y, 3.5, 0.3, item.get("title", ""), sz=15, bold=True, color=color)
        tb(sl, 4.5, y, 8.3, 0.3, item.get("desc", ""), sz=12, color=DARK_GRAY)


def build_ending(prs, page_num):
    """Build the ending slide: '谢谢，请指导' """
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    bg = sl.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
    bg.fill.solid(); bg.fill.fore_color.rgb = NAVY; bg.line.fill.background()
    al = sl.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(3.2), SLIDE_W, Inches(0.06))
    al.fill.solid(); al.fill.fore_color.rgb = LIGHT_BLUE; al.line.fill.background()
    tb(sl, 0.8, 2.0, 11.7, 1.2, "谢谢，请指导", sz=44, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    tb(sl, 9.5, 6.2, 3, 0.5, f"CONFIDENTIAL AND PROPRIETARY — {page_num}", sz=9, color=MED_GRAY, align=PP_ALIGN.RIGHT)
    return sl


# ════════════════════════════════════════════════════════════════
# MAIN
# ════════════════════════════════════════════════════════════════

def generate_pptx(title, subtitle, content_data, output_path, client="", date_str=""):
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    # ── Parse content ──
    if isinstance(content_data, str):
        data = json.loads(content_data)
    else:
        data = content_data

    page = 1

    # Cover
    build_cover(
        prs,
        title=title,
        subtitle=subtitle,
        cover_extra=data.get("cover_extra", ""),
        date_str=date_str or data.get("date", ""),
        client=client or data.get("client", "")
    )
    page += 1

    sections = data.get("sections", [])

    for sec in sections:
        stype = sec.get("type", "text")
        if stype == "executive_summary":
            build_executive_summary(prs, sec, page)
        elif stype == "cards":
            build_cards_slide(prs, sec, page)
        elif stype == "table":
            build_table_slide(prs, sec, page)
        elif stype == "list":
            build_list_slide(prs, sec, page)
        elif stype == "text":
            build_text_slide(prs, sec, page)
        elif stype == "cover":
            pass  # Already handled
        else:
            # fallback: text
            build_text_slide(prs, sec, page)
        page += 1

    # ── Ending slide ──
    build_ending(prs, page)
    page += 1

    prs.save(output_path)
    print(f"✅ PPT saved: {output_path}")
    print(f"   Slides: {len(prs.slides)}")
    return output_path


if __name__ == "__main__":
    parser = argparse.ArgumentParser(description="McKinsey-style PPT Generator")
    parser.add_argument("--title", required=True, help="PPT main title")
    parser.add_argument("--subtitle", default="", help="PPT subtitle")
    parser.add_argument("--client", default="", help="Client name")
    parser.add_argument("--date", default="", help="Date string")
    parser.add_argument("--content", required=True, help="JSON content file path or JSON string")
    parser.add_argument("--output", "-o", default="output.pptx", help="Output file path")
    args = parser.parse_args()

    # Load content
    try:
        with open(args.content, "r", encoding="utf-8") as f:
            content = f.read()
    except (FileNotFoundError, OSError):
        content = args.content  # treat as JSON string

    generate_pptx(
        title=args.title,
        subtitle=args.subtitle,
        content_data=content,
        output_path=args.output,
        client=args.client,
        date_str=args.date,
    )
