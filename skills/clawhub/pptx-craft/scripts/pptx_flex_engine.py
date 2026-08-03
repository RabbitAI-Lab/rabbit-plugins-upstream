# -*- coding: utf-8 -*-
"""
方法3 PPTX 出码引擎 (pptx_flex_engine.py)  —  pptx-craft 技能核心【纯原语库 + 通用版式库】
==========================================================================================

文本优先多遍布局 (Text-First Multi-Pass)：把"数据模型 / HTML 解析结果"渲染成可编辑 PPT，
保证 不重叠 / 不裁切 / 有气口。

⚠️ 本文件**不再内嵌任何"某看板的 demo"**（旧的 KPI/批次硬编码已移除）。
   - demo 已移至 examples/dashboard_demo.py（独立运行：python examples/dashboard_demo.py）
   - 引擎只提供可复用原语 + 通用版式组件；真实内容由调用方（专家团 / html2ppt / 用户）注入

架构层（可被 configure() 配置，支持多页 + 文档级令牌覆盖）：
  ① 设计令牌(DEFAULT_TOKENS) + 间距令牌(SP) + 密度(DENSITY) —— 可被文档级 :root 令牌覆盖
  ② 虚拟画布 + 坐标映射（configurable）
  ③ 容器系统 Row/Column + flex 权重（仅用于"区域框"层级）
  ④ 文本优先布局(layout_texts) + CJK 字体度量
  ⑤ 校验层(validate)：越界 + 跨组重叠 + 留白下限
  ⑥ Deck 渲染原语（rect/ellipse/text → python-pptx shapes）+ SVG/PNG 预览
  ⑦ 通用版式组件（KpiCard / CompareCard 等，内容为调用方注入；见底部 builders 分区）

运行 demo:  python examples/dashboard_demo.py
"""

from pptx import Presentation
from pptx.util import Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
import os, copy, tempfile, shutil
from PIL import Image, ImageDraw, ImageFont

# ---------------------------------------------------------------------------
# ① 设计令牌 (可被文档级 :root CSS 变量覆盖)
# ---------------------------------------------------------------------------
DEFAULT_TOKENS = {
    "color": {
        "bg":      "F4F5F8",   # 页面底
        "ink":     "0F172A",   # 主文字/深底
        "muted":   "64748B",   # 次级文字
        "sub":     "94A3B8",   # 辅助文字
        "border":  "E2E8F0",   # 描边
        "card":    "FFFFFF",   # 卡片底
        "blue":    "2563EB",   # 主色(电光蓝)
        "green":   "16A34A",   # 成功
        "amber":   "D97706",   # 警示
        "greenBg": "DCFCE7",   # 状态芯片底
        "greenTx": "86EFAC",   # 状态芯片字
        "bar":     "CBD5E1",   # 灰色条/基线
        "red":     "B5483D",   # 警示红
        "gold":    "B8954A",   # 金(玛奇朵)
        "goldSoft":"EBDFC4",
    },
    "font": "Microsoft YaHei",
    "radius": {"sm": 4, "md": 10, "lg": 12, "pill": 14},
}

# ---------------------------------------------------------------------------
# 引擎配置（configure() 覆盖；支持多页 + 文档级令牌）
# ---------------------------------------------------------------------------
CFG = {
    "vw": 1440, "vh": 680,
    "slide_w_emu": int(33.867 * 360000),
    "slide_h_emu": int(19.05 * 360000),
    "scale": 1.0, "offx": 0.0, "offy": 0.0,
    "density": 1.20,
    "tokens": DEFAULT_TOKENS,
    "g_ts": 1.0,
}

def configure(vw=1440, vh=680, slide_w_cm=33.867, slide_h_cm=19.05, density=1.20, tokens=None):
    """配置虚拟画布与目标尺寸；文档级令牌可整体覆盖（如 HTML :root 调色板）。"""
    CFG["vw"], CFG["vh"] = vw, vh
    CFG["slide_w_emu"] = int(slide_w_cm * 360000)
    CFG["slide_h_emu"] = int(slide_h_cm * 360000)
    CFG["scale"] = min(CFG["slide_w_emu"] / vw, CFG["slide_h_emu"] / vh)
    CFG["offx"] = (CFG["slide_w_emu"] - vw * CFG["scale"]) / 2
    CFG["offy"] = 0.0
    CFG["density"] = density
    if tokens:
        # 合并：以传入令牌覆盖默认，缺省项回退默认，避免 KeyError
        t = {k: dict(DEFAULT_TOKENS[k]) if isinstance(DEFAULT_TOKENS[k], dict) else DEFAULT_TOKENS[k]
             for k in DEFAULT_TOKENS}
        for k, v in tokens.items():
            if isinstance(v, dict) and isinstance(t.get(k), dict):
                t[k].update(v)
            else:
                t[k] = v
        CFG["tokens"] = t

def set_g_ts(v):
    CFG["g_ts"] = v

# ---------------------------------------------------------------------------
# ②b 模板克隆（Branch B：模板为尺寸锚点）
# ---------------------------------------------------------------------------
def load_presentation_safe(path):
    """加载模板; 非 ASCII 路径(lxml 可能 segfault)先复制到临时 ascii 副本再加载"""
    if any(ord(c) > 127 for c in path):
        tmp = os.path.join(tempfile.gettempdir(), "pptx_tpl_tmp.pptx")
        if os.path.exists(tmp):
            os.remove(tmp)
        shutil.copyfile(path, tmp)
        return Presentation(tmp)
    return Presentation(path)

def compute_available_area(prs, slide_idx, margin_cm):
    """扫描目标页(含版式继承)已有形状, 算可用区; 同时扫 slide + slide_layout。"""
    SW, SH = prs.slide_width, prs.slide_height
    M = int(margin_cm * 360000)
    header_bottom, footer_top = 0, SH
    has_footer = False
    def scan(shapes):
        nonlocal header_bottom, footer_top, has_footer
        for sh in shapes:
            try:
                L, T, R, B = sh.left, sh.top, sh.left + sh.width, sh.top + sh.height
            except Exception:
                continue
            if R - L <= 0 or B - T <= 0:
                continue
            # 排除整页背景/铺满图形
            if L <= 0 and T <= 0 and R >= SW - 1000 and B >= SH - 1000:
                continue
            cy = (T + B) / 2
            if cy < SH / 2:
                header_bottom = max(header_bottom, B)
            else:
                footer_top = min(footer_top, T)
                has_footer = True
    scan(prs.slides[slide_idx].shapes)
    scan(prs.slides[slide_idx].slide_layout.shapes)
    left = M
    right = SW - M
    top = max(M, header_bottom + M)
    bottom = (footer_top - M) if has_footer else (SH - M)
    return {"left": left, "top": top, "cw": right - left, "ch": bottom - top,
            "has_footer": has_footer, "header_bottom": header_bottom / 360000}

def configure_from_template(prs, slide_idx, margin_cm):
    """Branch B: 把虚拟画布 1440×VH 直接映射模板可用区(比例精确, 无拉伸, 原生生成)。"""
    a = compute_available_area(prs, slide_idx, margin_cm)
    VW = 1440
    VH = round(a["ch"] / a["cw"] * VW)
    CFG["vw"], CFG["vh"] = VW, VH
    CFG["slide_w_emu"], CFG["slide_h_emu"] = a["cw"], a["ch"]
    CFG["scale"] = a["cw"] / VW
    CFG["offx"], CFG["offy"] = a["left"], a["top"]
    CFG["_tpl_idx"] = slide_idx
    CFG["_tpl_layout"] = prs.slides[slide_idx].slide_layout   # 缓存版式, 删原页后仍可克隆
    CFG["_tpl_src"] = prs.slides[slide_idx]                   # 缓存源幻灯片(含真实标题文本框)
    return a

def delete_all_slides(prs):
    """清空模板自带的全部幻灯片(含物理 slide part), 仅保留版式用于克隆。
    仅移除 sldIdLst 不够——模板的 slideN.xml / slideN.xml.rels 仍残留在 package 中,
    保存时新幻灯片会与它们重名(Duplicate name)。这里把它们从 package._parts 一并剔除。"""
    ns_uri = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"
    sldIdLst = prs.slides._sldIdLst
    for sldId in list(sldIdLst):
        rId = sldId.get("{%s}id" % ns_uri)
        if rId:
            try:
                prs.part.drop_rel(rId)
            except Exception:
                pass
        sldIdLst.remove(sldId)
    pkg = prs.part.package
    parts = getattr(pkg, "_parts", None)
    if parts is None:
        return
    def _orphan(p):
        n = str(p.partname)
        return n.startswith("/ppt/slides/slide") or n.startswith("/ppt/slides/_rels/slide")
    if isinstance(parts, list):
        parts[:] = [p for p in parts if not _orphan(p)]
    else:
        for p in list(parts):
            if _orphan(p):
                try:
                    parts.discard(p)
                except Exception:
                    pass

def _max_shape_id(slide):
    mx = 0
    for sh in slide.shapes:
        try:
            i = sh.shape_id
            if i and i > mx:
                mx = i
        except Exception:
            pass
    return mx

def add_slide_from_template(prs, deck, title="", badge="", src_els=None, bg=None):
    """克隆模板【源幻灯片】(保留其真实标题文本框与装饰), 回填标题/编号, 再把 deck 渲染进可用区。
    注意：克隆源幻灯片而非版式——模板标题是落在幻灯片上的真实文本框, 版式本身常无占位符,
    直接 add_slide(layout) 会得到空白幻灯片, 导致标题丢失。"""
    C = TOK()["color"]
    src = CFG.get("_tpl_src") or prs.slides[CFG.get("_tpl_idx", 0)]
    layout = src.slide_layout
    new = prs.slides.add_slide(layout)
    spTree = new.shapes._spTree
    extLst = spTree.find(qn("p:extLst"))
    mx = _max_shape_id(new)
    # 深拷贝源幻灯片中"非内容占位符"的形状(标题/编号/装饰条), 保留模板观感
    for sh in src.shapes:
        if sh.is_placeholder and _is_content_placeholder(sh):
            continue
        c = copy.deepcopy(sh._element)
        cp = c.find(".//" + qn("p:cNvPr"))
        if cp is not None:
            mx += 1
            cp.set("id", str(mx))
        if extLst is not None:
            extLst.addprevious(c)
        else:
            spTree.append(c)
    # 回填标题/编号：深拷贝来的文本框中, 最大宽度=标题, 最小=编号/副标题
    tbs = [sh for sh in new.shapes if sh.has_text_frame]
    tbs.sort(key=lambda s: (s.width or 0))
    if tbs:
        tbs[-1].text_frame.text = title            # 标题(最大文本框)
        if badge and len(tbs) >= 2:
            tbs[0].text_frame.text = badge         # 编号/副标题(仅在有 badge 时回填)
    if bg:
        new.background.fill.solid()
        new.background.fill.fore_color.rgb = rgb(bg)
    deck.render(new)
    return new


def _is_content_placeholder(sh):
    """判断是否为内容/正文占位符(应由 deck 内容覆盖, 不深拷贝, 避免旧文字透出)。"""
    if not sh.is_placeholder:
        return False
    try:
        from pptx.enum.shapes import MSO_PLACEHOLDER as PH
        return sh.placeholder_format.type in (PH.CONTENT, PH.BODY)
    except Exception:
        return False

# ---------------------------------------------------------------------------
# ② 间距令牌 + 坐标映射
# ---------------------------------------------------------------------------
SP_RAW = {"xs": 4, "sm": 8, "md": 12, "lg": 16, "xl": 20, "2xl": 28}

def sp(k):
    """间距令牌 -> 经密度系数放大的虚拟像素值"""
    return round(SP_RAW[k] * CFG["density"])

def PAD_CARD():
    return sp("lg")

def PAD_CARD_LG():
    return sp("xl")

def X(v):  return int(round(CFG["offx"] + v * CFG["scale"]))
def Y(v):  return int(round(CFG["offy"] + v * CFG["scale"]))
def W_(v): return int(round(v * CFG["scale"]))
def H_(v): return int(round(v * CFG["scale"]))
def FS(px): return Pt(px * 0.75)        # px -> pt (@96dpi)
def rgb(h): return RGBColor.from_string(str(h).replace("#", ""))
def TOK(): return CFG["tokens"]

# ---------------------------------------------------------------------------
# ④ 字体度量（估算）—— 不需真实渲染，用保守近似保证"宁多勿少"，文字绝不裁切
#   CJK 全角字符宽 ≈ fs*0.98；拉丁 ≈ fs*0.55；空格 ≈ fs*0.30
#   行高 ≈ fs*1.25（含行内上下留白，模拟 PPT 真实包络）
# ---------------------------------------------------------------------------
def _char_w(ch, fs):
    o = ord(ch)
    if o > 0x2E7F:                 # CJK + 全角标点
        return fs * 1.05           # Pitfall#8: CJK 系数 ≥1.0(保守估计), 防换行爆炸
    if ch == ' ':
        return fs * 0.30
    return fs * 0.55

def text_w(s, fs):
    return sum(_char_w(c, fs) for c in s)

def text_h(s, fs, max_w):
    if max_w <= 0:
        return fs * 1.25
    cur = 0.0
    lines = 1
    for ch in s:
        cw = _char_w(ch, fs)
        if cur + cw > max_w and cur > 0:
            lines += 1
            cur = cw
        else:
            cur += cw
    return lines * fs * 1.25

# ---------------------------------------------------------------------------
# ③ Box —— 虚拟坐标矩形
# ---------------------------------------------------------------------------
class Box:
    __slots__ = ("x", "y", "w", "h")
    def __init__(self, x, y, w, h):
        self.x, self.y, self.w, self.h = x, y, w, h

# ---------------------------------------------------------------------------
# ③ 容器系统 (flex) —— 仅用于"区域框"层级
# ---------------------------------------------------------------------------
def hbox(deck, box, items, gap=16, pad=0):
    ix, iw = box.x + pad, box.w - 2 * pad
    n = len(items)
    tg = gap * (n - 1) if n > 1 else 0
    fixed = sum(it.get("size") or 0 for it in items)
    flex_total = sum(it.get("flex") or 0 for it in items)
    unit = (iw - tg - fixed) / flex_total if flex_total else 0
    cx = ix
    for it in items:
        cw = it["flex"] * unit if it.get("flex") else it.get("size", 0)
        cb = Box(cx, box.y + pad, cw, box.h - 2 * pad)
        it["build"](deck, cb, it.get("group", "g"))
        cx += cw + gap

def vbox(deck, box, items, gap=16, pad=0):
    iy, ih = box.y + pad, box.h - 2 * pad
    n = len(items)
    tg = gap * (n - 1) if n > 1 else 0
    fixed = sum(it.get("size") or 0 for it in items)
    flex_total = sum(it.get("flex") or 0 for it in items)
    unit = (ih - tg - fixed) / flex_total if flex_total else 0
    cy = iy
    for it in items:
        ch = it["flex"] * unit if it.get("flex") else it.get("size", 0)
        cb = Box(box.x + pad, cy, box.w - 2 * pad, ch)
        it["build"](deck, cb, it.get("group", "g"))
        cy += ch + gap

# ---------------------------------------------------------------------------
# ⑥ Deck —— 渲染 + 校验登记
# ---------------------------------------------------------------------------
class VRec:
    __slots__ = ("x", "y", "w", "h", "group", "kind")
    def __init__(self, x, y, w, h, group, kind):
        self.x, self.y, self.w, self.h, self.group, self.kind = x, y, w, h, group, kind

class Deck:
    def __init__(self):
        self.ops = []
        self.recs = []
    def rect(self, b, fill, line=None, radius=0, group="g", gradient=None, shadow=None):
        self.ops.append(("rect", b, fill, line, radius, group, gradient, shadow))
        self.recs.append(VRec(b.x, b.y, b.w, b.h, group, "rect"))
    def ellipse(self, b, fill, group="g"):
        self.ops.append(("ellipse", b, fill, group))
        self.recs.append(VRec(b.x, b.y, b.w, b.h, group, "ellipse"))
    def text(self, b, content, fs, color, bold=False, align="left", group="g", runs=None):
        self.ops.append(("text", b, content, fs, color, bold, align, group, runs))
        self.recs.append(VRec(b.x, b.y, b.w, b.h, group, "text"))
    def render(self, slide):
        for op in self.ops:
            if op[0] in ("rect", "ellipse", "text"):
                b = op[1]
                if b.w <= 0 or b.h <= 0:
                    raise AssertionError(
                        f"[Pitfall#1 负尺寸] {op[0]} w={b.w} h={b.h} —— 会触发 PowerPoint 修复弹窗, 已拦截")
            if op[0] == "rect":
                _, b, fill, line, radius, _, gradient, shadow = op
                sp_ = slide.shapes.add_shape(
                    MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE,
                    X(b.x), Y(b.y), W_(b.w), H_(b.h))
                if gradient:
                    try:
                        sp_.fill.gradient()
                        stops = sp_.fill.gradient_stops
                        stops[0].color.rgb = rgb(gradient[0])
                        stops[1].color.rgb = rgb(gradient[1])
                    except Exception:
                        sp_.fill.solid(); sp_.fill.fore_color.rgb = rgb(fill)
                else:
                    sp_.fill.solid(); sp_.fill.fore_color.rgb = rgb(fill)
                if line:
                    sp_.line.color.rgb = rgb(line[0]); sp_.line.width = Pt(line[1])
                else:
                    sp_.line.fill.background()
                if shadow:
                    try:
                        sp_.shadow.inherit = False
                        sp_.shadow.visible = True
                        sp_.shadow.blur = Pt(7)
                        sp_.shadow.offset = Pt(3)
                        sp_.shadow.angle = 90
                        sp_.shadow.color.rgb = RGBColor(0x1A, 0x20, 0x2E)
                    except Exception:
                        sp_.shadow.inherit = False
                else:
                    sp_.shadow.inherit = False
            elif op[0] == "ellipse":
                _, b, fill, _ = op
                sp_ = slide.shapes.add_shape(MSO_SHAPE.OVAL, X(b.x), Y(b.y), W_(b.w), H_(b.h))
                sp_.fill.solid(); sp_.fill.fore_color.rgb = rgb(fill); sp_.line.fill.background()
                sp_.shadow.inherit = False
            elif op[0] == "text":
                _, b, content, fs, color, bold, align, _, runs = op
                tb = slide.shapes.add_textbox(X(b.x), Y(b.y), W_(b.w), H_(b.h))
                tf = tb.text_frame
                tf.word_wrap = True
                tf.vertical_anchor = MSO_ANCHOR.MIDDLE
                tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
                p = tf.paragraphs[0]
                p.alignment = {"left": PP_ALIGN.LEFT, "center": PP_ALIGN.CENTER, "right": PP_ALIGN.RIGHT}[align]
                if runs:
                    for rn in runs:
                        r = p.add_run(); r.text = rn["text"]
                        r.font.size = FS(rn.get("fs", fs)); r.font.bold = rn.get("bold", False)
                        r.font.color.rgb = rgb(rn.get("color", color)); r.font.name = TOK()["font"]
                else:
                    r = p.add_run(); r.text = content
                    r.font.size = FS(fs); r.font.bold = bold
                    r.font.color.rgb = rgb(color); r.font.name = TOK()["font"]

# ---------------------------------------------------------------------------
# ⑤ 校验层 (Validator)
# ---------------------------------------------------------------------------
MIN_BREATH = 10   # 跨组最小间隙（虚拟px），低于即报"拥挤"

def _overlap(a, b, tol=1):
    return not (a.x + a.w <= b.x + tol or b.x + b.w <= a.x + tol or
                a.y + a.h <= b.y + tol or b.y + b.h <= a.y + tol)

def _contains(o, i, tol=1):
    return (i.x >= o.x - tol and i.y >= o.y - tol and
            i.x + i.w <= o.x + o.w + tol and i.y + i.h <= o.y + o.h + tol)

def _gap_between(a, b):
    ox = min(a.x + a.w, b.x + b.w) - max(a.x, b.x)
    oy = min(a.y + a.h, b.y + b.h) - max(a.y, b.y)
    if ox <= 0 and oy <= 0:
        return max(ox, oy)
    if ox <= 0:
        return ox
    if oy <= 0:
        return oy
    return min(ox, oy)

def validate(recs):
    errors, warnings = [], []
    vw, vh = CFG["vw"], CFG["vh"]
    for r in recs:
        if r.x < -1 or r.y < -1 or r.x + r.w > vw + 1 or r.y + r.h > vh + 1:
            errors.append(f"[越界:{r.group}] ({r.x:.0f},{r.y:.0f},{r.w:.0f},{r.h:.0f}) 超出 {vw}x{vh}")
    for i in range(len(recs)):
        for j in range(i + 1, len(recs)):
            a, b = recs[i], recs[j]
            if a.group == b.group:
                continue
            if _contains(a, b) or _contains(b, a):
                continue
            if _overlap(a, b):
                errors.append(f"[重叠:{a.group}↔{b.group}] 几何相交 -> 元素重复/错位")
    for i in range(len(recs)):
        for j in range(i + 1, len(recs)):
            a, b = recs[i], recs[j]
            if a.group == b.group:
                continue
            if _contains(a, b) or _contains(b, a):
                continue
            ox = min(a.x + a.w, b.x + b.w) - max(a.x, b.x)
            oy = min(a.y + a.h, b.y + b.h) - max(a.y, b.y)
            # 仅当两者在"分离轴"方向上间隙过小才算拥挤:
            #   x 范围交叠(ox>0)→上下堆叠, 看纵向间隙 -oy;
            #   y 范围交叠(oy>0)→左右并列, 看横向间隙 -ox;
            #   对角分离(ox<=0 且 oy<=0)不算相邻, 跳过(避免并列元素被误报)。
            if ox > 0:
                gap = -oy
            elif oy > 0:
                gap = -ox
            else:
                continue
            if 0 <= gap < MIN_BREATH:
                # 仅当两者同属"结构框"(rect↔rect)或同属"文字"(text↔text)才报拥挤；
                # text↔rect(文字贴其卡片或邻卡)属常态，不报，避免稀释 QA 信号。
                txt_mix = (a.kind == "text") ^ (b.kind == "text")
                if txt_mix:
                    continue
                warnings.append(f"[拥挤:{a.group}↔{b.group}] 间隙 {gap:.0f}px < {MIN_BREATH}px")
    return errors, warnings

# ============================================================================
# ④ 文本优先布局核心 (Pass1 + Pass2)
# ============================================================================
def layout_texts(deck, items, box, pad, breath, align="left", group="g", ts=None):
    """
    在 box 内"文本优先"放置一组文字（纵栈）：
      - 上下左右留 pad 气口；多文字间留 breath 气口
      - 若总高超出 box.h，按比例缩所有字号（各自保 min_fs，绝不删内容）
    返回 leftover：文字块下方剩余空间 Box（供装饰填充）
    """
    ts = CFG["g_ts"] if ts is None else ts
    avail_w = max(1, box.w - 2 * pad)

    def block_h(scale):
        h = 2 * pad
        for it in items:
            fs = max(it.get("min_fs", 9), it["fs"] * ts * scale)
            h += text_h(it["content"], fs, avail_w)
            h += breath
        return h - breath

    scale = 1.0
    if block_h(1.0) > box.h:
        lo, hi = 0.40, 1.0
        for _ in range(22):
            mid = (lo + hi) / 2
            if block_h(mid) > box.h:
                hi = mid
            else:
                lo = mid
        scale = lo

    y = box.y + pad
    for it in items:
        fs = max(it.get("min_fs", 9), it["fs"] * ts * scale)
        th = text_h(it["content"], fs, avail_w)
        tb = Box(box.x + pad, y, avail_w, th)
        deck.text(tb, it["content"], fs, it["color"], bold=it.get("bold", False),
                  align=align, group=group)
        y += th + breath
    used = (y - breath) - box.y
    leftover = Box(box.x, y - breath, box.w, max(0, box.y + box.h - (y - breath)))
    return leftover, scale

# ---------------------------------------------------------------------------
# 视觉增强原语 (#9)：强调条 / 分隔线 —— 纯装饰，不登记为内容、不改变几何校验
# ---------------------------------------------------------------------------
def accent_bar(deck, box, color, thickness=6, vertical=True, group="accent"):
    """卡片左侧/顶部强调条。vertical=True 画竖条(默认左侧)，False 画顶部分隔条。"""
    cb = Box(box.x, box.y, thickness, box.h) if vertical else Box(box.x, box.y, box.w, thickness)
    deck.rect(cb, color, radius=0, group=group)

def divider(deck, x, y, w, color, thickness=2, group="div"):
    deck.rect(Box(x, y, w, thickness), color, radius=0, group=group)

# ============================================================================
# ⑦ 通用版式组件（内容为调用方注入；内部文本优先，装饰填空）
# ============================================================================
def KpiCard(deck, box, group, kpi):
    C = TOK()["color"]
    deck.rect(box, C["card"], line=(C["border"], 1), radius=TOK()["radius"]["lg"], group=group)
    items = [
        {"content": kpi["label"], "fs": 13, "min_fs": 10, "color": C["muted"]},
        {"content": kpi["value"], "fs": 30, "bold": True, "min_fs": 20, "color": kpi.get("color", C["ink"])},
        {"content": kpi["sub"],   "fs": 12, "min_fs": 9,  "color": C["sub"]},
    ]
    layout_texts(deck, items, box, pad=PAD_CARD(), breath=sp("sm"), align="left", group=group)

def CompareCard(deck, box, group, cmp):
    C = TOK()["color"]
    deck.rect(box, C["card"], line=(C["border"], 1), radius=TOK()["radius"]["md"], group=group)
    items = [
        {"content": cmp["label"], "fs": 12, "min_fs": 10, "color": C["muted"]},
        {"content": cmp["value"], "fs": 24, "bold": True, "min_fs": 16, "color": C["ink"]},
        {"content": cmp["sub"],   "fs": 11, "min_fs": 9,  "color": C.get("green", C["ink"])},
    ]
    layout_texts(deck, items, box, pad=sp("md"), breath=sp("xs"), align="left", group=group)

# ============================================================================
# ⑥ SVG 预览（虚拟坐标 1:1）
# ============================================================================
def emit_svg(deck, path):
    C = TOK()["color"]
    vw, vh = CFG["vw"], CFG["vh"]
    parts = [f'<svg xmlns="http://www.w3.org/2000/svg" viewBox="0 0 {vw} {vh}" font-family="{TOK()["font"]}">']
    parts.append(f'<rect x="0" y="0" width="{vw}" height="{vh}" fill="#{C["bg"]}"/>')
    gi = [0]
    for op in deck.ops:
        if op[0] == "rect":
            _, b, fill, line, radius, _, gradient, shadow = op
            r = radius if radius else 0
            stroke = f' stroke="#{str(line[0]).replace("#","")}" stroke-width="{line[1]}"' if line else ''
            fill_attr = f'fill="#{str(fill).replace("#","")}"'
            extra = ''
            if gradient:
                gid = f"grad{gi[0]}"; gi[0] += 1
                parts.append(f'<defs><linearGradient id="{gid}" x1="0" y1="0" x2="0" y2="1">'
                             f'<stop offset="0%" stop-color="#{str(gradient[0]).replace("#","")}"/>'
                             f'<stop offset="100%" stop-color="#{str(gradient[1]).replace("#","")}"/></linearGradient></defs>')
                fill_attr = f'fill="url(#{gid})"'
            if shadow:
                fid = f"sh{gi[0]}"; gi[0] += 1
                parts.append(f'<defs><filter id="{fid}" x="-20%" y="-20%" width="140%" height="140%">'
                             f'<feDropShadow dx="2" dy="3" stdDeviation="4" flood-color="#1A202E" flood-opacity="0.18"/></filter></defs>')
                extra = f' filter="url(#{fid})"'
            parts.append(f'<rect x="{b.x:.0f}" y="{b.y:.0f}" width="{b.w:.0f}" height="{b.h:.0f}" rx="{r}" {fill_attr}{stroke}{extra}/>')
        elif op[0] == "ellipse":
            _, b, fill, _ = op
            cx, cy, rx, ry = b.x + b.w/2, b.y + b.h/2, b.w/2, b.h/2
            parts.append(f'<ellipse cx="{cx:.0f}" cy="{cy:.0f}" rx="{rx:.0f}" ry="{ry:.0f}" fill="#{str(fill).replace("#","")}"/>')
        elif op[0] == "text":
            _, b, content, fs, color, bold, align, _, runs = op
            if runs:
                content = "".join(rn["text"] for rn in runs)
                color = runs[0].get("color", color)
                bold = runs[0].get("bold", bold)
            anc = {"left": "start", "center": "middle", "right": "end"}[align]
            tx = {"left": b.x, "center": b.x + b.w/2, "right": b.x + b.w}[align]
            ty = b.y + b.h/2 + fs*0.35
            wt = ' font-weight="700"' if bold else ''
            esc = (content.replace("&","&amp;").replace("<","&lt;").replace(">","&gt;"))
            parts.append(f'<text x="{tx:.0f}" y="{ty:.0f}" font-size="{fs}" fill="#{str(color).replace("#","")}" text-anchor="{anc}"{wt}>{esc}</text>')
    parts.append('</svg>')
    with open(path, "w", encoding="utf-8") as f:
        f.write("\n".join(parts))
    return path

# ============================================================================
# ⑥ PNG 预览（供 QA Layer2 逐页目检；与 PPTX 同源坐标，最忠实）
# ============================================================================
_FONT_CACHE = {}
def _pil_font(size, font_path=None, index=0):
    key = (size, font_path, index)
    if key in _FONT_CACHE:
        return _FONT_CACHE[key]
    try:
        if font_path:
            f = ImageFont.truetype(font_path, size, index=index)
        else:
            f = ImageFont.load_default()
    except Exception:
        f = ImageFont.load_default()
    _FONT_CACHE[key] = f
    return f

def _hex2rgb(h):
    s = str(h).replace("#", "")
    if len(s) == 3:
        s = "".join(c * 2 for c in s)
    try:
        return (int(s[0:2], 16), int(s[2:4], 16), int(s[4:6], 16))
    except Exception:
        return (0, 0, 0)

def render_deck_png(deck, path, scale_px=1.4, font_path=None):
    """把 deck 渲染成 PNG（虚拟坐标 * scale_px），用于 QA 逐页视觉校验。"""
    vw, vh = CFG["vw"], CFG["vh"]
    W, H = int(vw * scale_px), int(vh * scale_px)
    C = TOK()["color"]
    img = Image.new("RGB", (W, H), "#" + str(C["bg"]).replace("#", ""))
    d = ImageDraw.Draw(img)

    # 先画形状，再画文字（文字在上层）
    def sx(v): return int(v * scale_px)
    for op in deck.ops:
        if op[0] == "rect":
            _, b, fill, line, radius, _, gradient, shadow = op
            rad = int(radius*scale_px) if radius else 0
            # 阴影近似：先画偏移的整块深色圆角，再在原始位置画填充 -> 仅右下露出阴影
            if shadow:
                d.rounded_rectangle([sx(b.x)+3, sx(b.y)+4, sx(b.x+b.w)+3, sx(b.y+b.h)+4],
                                    radius=rad, fill="#1A202E")
            if gradient:
                c1 = _hex2rgb(gradient[0]); c2 = _hex2rgb(gradient[1])
                y0, y1 = int(sx(b.y)), int(sx(b.y+b.h))
                for yy in range(y0, y1 + 1):
                    t = (yy - y0) / max(1, (y1 - y0))
                    col = tuple(int(c1[k] + (c2[k] - c1[k]) * t) for k in range(3))
                    d.line([(sx(b.x), yy), (sx(b.x+b.w), yy)], fill=col, width=1)
            else:
                d.rounded_rectangle([sx(b.x), sx(b.y), sx(b.x+b.w), sx(b.y+b.h)],
                                    radius=rad, fill="#" + str(fill).replace("#", ""),
                                    outline=None, width=1)
            if line:
                d.rounded_rectangle([sx(b.x), sx(b.y), sx(b.x+b.w), sx(b.y+b.h)],
                                    radius=rad, outline="#" + str(line[0]).replace("#", ""),
                                    width=int(line[1]*scale_px) if line else 1)
        elif op[0] == "ellipse":
            _, b, fill, _ = op
            d.ellipse([sx(b.x), sx(b.y), sx(b.x+b.w), sx(b.y+b.h)],
                      fill="#" + str(fill).replace("#", ""))
    for op in deck.ops:
        if op[0] == "text":
            _, b, content, fs, color, bold, align, _, runs = op
            if runs:
                content = "".join(rn["text"] for rn in runs)
            fsize = max(8, int(fs * 0.75 * scale_px))
            font = _pil_font(fsize, font_path)
            # 计算换行（与 text_h 同算法）
            avail = b.w
            lines, cur = [], ""
            for ch in content:
                test = cur + ch
                if text_w(test, fs) > avail and cur:
                    lines.append(cur); cur = ch
                else:
                    cur = test
            if cur:
                lines.append(cur)
            line_h = fsize * 1.25
            total_h = line_h * len(lines)
            ty = sx(b.y) + (sx(b.h) - total_h) / 2
            anchor = {"left": "la", "center": "ma", "right": "ra"}[align]
            for ln in lines:
                tw = d.textlength(ln, font=font)
                if align == "left":
                    tx = sx(b.x)
                elif align == "center":
                    tx = sx(b.x) + (sx(b.w) - tw) / 2
                else:
                    tx = sx(b.x) + sx(b.w) - tw
                d.text((tx, ty), ln, font=font, fill="#" + str(color).replace("#", ""))
                ty += line_h
    img.save(path)
    return path

# ============================================================================
# 多页辅助：把 deck 渲染进一个新建/复用 slide
# ============================================================================
def add_slide_from_deck(prs, deck, bg=None):
    """在 prs 中新增一页空白 slide，背景填 bg（默认令牌 bg），渲染 deck。"""
    C = TOK()["color"]
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bgc = bg or C["bg"]
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = rgb(bgc)
    deck.render(slide)
    return slide

# ---------------------------------------------------------------------------
# 便捷：创建一个标准 16:9 的空白 Presentation
# ---------------------------------------------------------------------------
def new_presentation(slide_w_cm=None, slide_h_cm=None):
    prs = Presentation()
    w = (slide_w_cm if slide_w_cm else CFG["slide_w_emu"]/360000)
    h = (slide_h_cm if slide_h_cm else CFG["slide_h_emu"]/360000)
    prs.slide_width = Emu(int(w * 360000))
    prs.slide_height = Emu(int(h * 360000))
    return prs

if __name__ == "__main__":
    print("pptx_flex_engine 是纯原语库，不直接运行。请运行 examples/dashboard_demo.py 查看 demo。")
