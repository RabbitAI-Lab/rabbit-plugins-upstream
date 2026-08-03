# -*- coding: utf-8 -*-
"""
qa_geometric.py — pptx-craft L1 几何校验工具（严过关专用）
==================================================================================
独立、确定性地检测生成的 .pptx，捕获「一定弹修复 / 一定顶出画面」的硬伤：

  · negative_size   负尺寸（width/height < 0 → PowerPoint 必弹修复，致命）
  · out_of_canvas   越界（元素超出幻灯片画布）
  · usable_overflow 可用区溢出（模板基底路线：元素底/右超出模板可用区死线）

这是 pptx-craft「实战踩坑经验库」里 P0 级问题的自动化闸门：
  #1 负高度（动态高度算成负值）  #3 可用区测量 + 死线

输出 JSON：{passed, slide_size, usable_area, issue_count, issues[]}

用法:
  python qa_geometric.py <in.pptx> \
      [--usable-x 0 --usable-y 0.88 --usable-w 13.333 --usable-h 5.52] \
      [--out qa.json]
依赖: python-pptx
"""
import sys
import json
import argparse
from pptx import Presentation

EMU_PER_IN = 914400.0


def inch(v):
    return v / EMU_PER_IN


def main():
    ap = argparse.ArgumentParser(description="pptx-craft L1 几何校验")
    ap.add_argument("pptx", help="待校验的 .pptx 路径")
    ap.add_argument("--usable-x", type=float, default=0.0, help="可用区左上 X(英寸)")
    ap.add_argument("--usable-y", type=float, default=0.0, help="可用区左上 Y(英寸)")
    ap.add_argument("--usable-w", type=float, default=13.333, help="可用区宽(英寸)")
    ap.add_argument("--usable-h", type=float, default=7.5, help="可用区高(英寸)")
    ap.add_argument("--out", default=None, help="QA 报告输出路径(json)")
    args = ap.parse_args()

    prs = Presentation(args.pptx)
    SW = inch(prs.slide_width)
    SH = inch(prs.slide_height)
    ux, uy, uw, uh = args.usable_x, args.usable_y, args.usable_w, args.usable_h
    uR = ux + uw
    uB = uy + uh

    issues = []
    for si, slide in enumerate(prs.slides, 1):
        for shape in slide.shapes:
            try:
                l = inch(shape.left)
                t = inch(shape.top)
                w = inch(shape.width)
                h = inch(shape.height)
            except Exception:
                continue
            r = l + w
            b = t + h
            name = getattr(shape, "name", "?") or "?"
            txt = ""
            if hasattr(shape, "text_frame"):
                txt = shape.text_frame.text[:30] if shape.text_frame.text else ""

            # 负尺寸（致命，必弹修复）
            if w < 0 or h < 0:
                issues.append({
                    "slide": si, "shape": name, "type": "negative_size",
                    "detail": "负尺寸 w=%.2f h=%.2f [%s]" % (w, h, txt),
                })
                continue

            # 越界（画布）
            if l < -0.01 or t < -0.01 or r > SW + 0.01 or b > SH + 0.01:
                issues.append({
                    "slide": si, "shape": name, "type": "out_of_canvas",
                    "detail": "越界 L=%.2f T=%.2f R=%.2f B=%.2f / 画布 %.2fx%.2f [%s]"
                              % (l, t, r, b, SW, SH, txt),
                })

            # 可用区溢出（模板基底路线死线）
            if b > uB + 0.02 or r > uR + 0.02:
                issues.append({
                    "slide": si, "shape": name, "type": "usable_overflow",
                    "detail": "超可用区 B=%.2f(限%.2f) R=%.2f(限%.2f) [%s]"
                              % (b, uB, r, uR, txt),
                })

    passed = len(issues) == 0
    report = {
        "passed": passed,
        "slide_size": [round(SW, 3), round(SH, 3)],
        "usable_area": [ux, uy, uw, uh],
        "issue_count": len(issues),
        "issues": issues,
    }
    if args.out:
        with open(args.out, "w", encoding="utf-8") as f:
            json.dump(report, f, ensure_ascii=False, indent=2)
    print(json.dumps(report, ensure_ascii=False, indent=2))


if __name__ == "__main__":
    main()
