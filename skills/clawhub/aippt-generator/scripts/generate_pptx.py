#!/usr/bin/env python3
"""
AIPPT Generator — Markdown to PowerPoint converter.

Converts AI-generated Markdown PPT content (with chart data blocks) into
a professional .pptx file with native charts, color themes, and font configs.

Usage:
    python generate_pptx.py \
        --title "Q4季度营销预算申请" \
        --content-file slides.md \
        --theme business \
        --font-standard \
        --output output.pptx

    # Or pipe content via stdin
    cat slides.md | python generate_pptx.py --title "My PPT" --output output.pptx

    # Custom colors
    python generate_pptx.py \
        --title "My PPT" \
        --content-file slides.md \
        --custom-colors '{"primary":"#6F42C1","secondary":"#EC4899","accent":"#F59E0B"}' \
        --output output.pptx
"""

import argparse
import json
import re
import sys
import os

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE


# ==================== Color Themes ====================

THEMES = {
    "business": {
        "name": "商务经典",
        "colors": {
            "primary": RGBColor(0, 82, 164),
            "secondary": RGBColor(237, 125, 49),
            "accent": RGBColor(16, 137, 62),
            "dark": RGBColor(33, 37, 41),
            "light": RGBColor(248, 249, 250),
            "white": RGBColor(255, 255, 255),
            "gray": RGBColor(108, 117, 125),
        },
    },
    "tech": {
        "name": "科技蓝调",
        "colors": {
            "primary": RGBColor(0, 120, 212),
            "secondary": RGBColor(0, 188, 242),
            "accent": RGBColor(16, 110, 190),
            "dark": RGBColor(28, 28, 28),
            "light": RGBColor(243, 249, 255),
            "white": RGBColor(255, 255, 255),
            "gray": RGBColor(119, 119, 119),
        },
    },
    "minimal": {
        "name": "简约黑白",
        "colors": {
            "primary": RGBColor(33, 37, 41),
            "secondary": RGBColor(108, 117, 125),
            "accent": RGBColor(173, 181, 189),
            "dark": RGBColor(33, 37, 41),
            "light": RGBColor(248, 249, 250),
            "white": RGBColor(255, 255, 255),
            "gray": RGBColor(108, 117, 125),
        },
    },
    "vibrant": {
        "name": "活力创新",
        "colors": {
            "primary": RGBColor(111, 66, 193),
            "secondary": RGBColor(236, 72, 153),
            "accent": RGBColor(245, 158, 11),
            "dark": RGBColor(33, 37, 41),
            "light": RGBColor(250, 245, 255),
            "white": RGBColor(255, 255, 255),
            "gray": RGBColor(108, 117, 125),
        },
    },
}

FONT_PRESETS = {
    "compact": {"title": 24, "h2": 18, "body": 13},
    "standard": {"title": 28, "h2": 22, "body": 16},
    "loose": {"title": 32, "h2": 26, "body": 19},
}

XL_MAP = {
    "pie": XL_CHART_TYPE.PIE,
    "bar": XL_CHART_TYPE.COLUMN_CLUSTERED,
    "xychart": XL_CHART_TYPE.LINE,
}

CHART_LABELS = {"pie": "饼图", "bar": "柱状图", "xychart": "折线图"}


# ==================== Helpers ====================

def hex_to_rgb(hex_color):
    hex_color = hex_color.lstrip("#")
    if len(hex_color) == 6:
        return RGBColor(int(hex_color[0:2], 16), int(hex_color[2:4], 16), int(hex_color[4:6], 16))
    return RGBColor(0, 0, 0)


def get_theme_colors(theme_key="business", custom_colors=None):
    if custom_colors:
        colors = {}
        for key, hex_value in custom_colors.items():
            colors[key] = hex_to_rgb(hex_value)
        defaults = THEMES["business"]["colors"]
        for key in ["primary", "secondary", "accent", "dark", "light", "white", "gray"]:
            if key not in colors:
                colors[key] = defaults[key]
        return colors
    return THEMES.get(theme_key, THEMES["business"])["colors"]


def resolve_font(font_config):
    base = FONT_PRESETS.get(font_config.get("theme", "standard"), FONT_PRESETS["standard"])
    out = dict(base)
    for k in ("title", "h2", "body"):
        v = font_config.get(k)
        if v:
            try:
                out[k] = int(v)
            except (ValueError, TypeError):
                pass
    font_name = (font_config.get("font_name") or "").strip()
    if font_name:
        out["name"] = font_name
    return out


def split_segments(content):
    """Split markdown into [(kind, payload)] segments: 'chart' or 'text'.

    For chart segments, payload is a tuple (chart_type, data_text).
    """
    segments = []
    lines = (content or "").split("\n")
    i, buf = 0, []

    def flush():
        if buf:
            segments.append(("text", "\n".join(buf)))
            buf.clear()

    while i < len(lines):
        head = lines[i].strip().lower()
        if head.startswith("```chart"):
            flush()
            # Extract chart type from the fence header (e.g. ```chart: pie)
            ctype_match = re.match(r"```\s*chart\s*:\s*(\w+)", lines[i].strip(), re.IGNORECASE)
            ctype = ctype_match.group(1).lower() if ctype_match else "pie"
            code, i = [], i + 1
            while i < len(lines) and not lines[i].strip().startswith("```"):
                code.append(lines[i])
                i += 1
            if i < len(lines):
                i += 1
            segments.append(("chart", (ctype, "\n".join(code))))
        else:
            buf.append(lines[i])
            i += 1
    flush()
    return segments


def parse_chart_data(text):
    cats, vals = [], []
    for raw in (text or "").splitlines():
        line = raw.strip().lstrip("-*").strip()
        if not line:
            continue
        m = re.match(r"^(.+?)\s*[:：]?\s*(-?\d+(?:\.\d+)?)\s*[%万千亿元]*\s*$", line)
        if m:
            cats.append(m.group(1).strip())
            vals.append(float(m.group(2)))
    return cats, vals


def split_by_h2(content):
    """Split content by ## headings. Returns [(title, body), ...]."""
    segments = []
    cur_title = None
    cur_body = []

    def flush():
        body = "\n".join(cur_body).strip()
        if cur_title is not None or body:
            segments.append((cur_title, body))

    for line in (content or "").split("\n"):
        if line.startswith("## "):
            flush()
            cur_title = line[3:].strip()
            cur_body = []
        elif line.startswith("# "):
            flush()
            cur_title = line[2:].strip()
            cur_body = []
        else:
            cur_body.append(line)
    flush()
    return [(t, b) for t, b in segments if (t is not None and t) or b]


# ==================== Text Rendering ====================

def _add_runs(paragraph, text, size, color, bold=False, italic=False, font_name=None):
    for part in re.split(r"(\*\*[^*]+?\*\*)", text or ""):
        if not part:
            continue
        run = paragraph.add_run()
        if part.startswith("**") and part.endswith("**"):
            run.text = part[2:-2]
            run.font.bold = True
        else:
            run.text = part.replace("**", "")
            if bold:
                run.font.bold = True
        run.font.size = Pt(size)
        run.font.color.rgb = color
        if italic:
            run.font.italic = True
        if font_name:
            run.font.name = font_name


def _set_para_font(p, size, color, bold=False, font_name=None):
    """Set font properties on both paragraph and all existing runs.

    python-pptx paragraph-level font is a fallback; many PPT readers
    only read run-level font. Setting both ensures consistent rendering.
    """
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    if font_name:
        p.font.name = font_name
    for run in p.runs:
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color
        if font_name:
            run.font.name = font_name


def _estimate_text_height(text, width_in=8.4):
    h = 0.0
    for ln in (text or "").split("\n"):
        if not ln.strip():
            h += 0.12
        elif ln.startswith("##"):
            h += 0.5
        elif ln.startswith("-") or ln.startswith("*"):
            h += 0.34
        elif ln.startswith(">"):
            h += 0.3
        else:
            chars_per_line = max(20, int(width_in * 13))
            n = max(1, (len(ln) + chars_per_line - 1) // chars_per_line)
            h += 0.28 * n
    return max(0.5, h)


# Max usable height for content area (slide_height 5.625 - top0 1.5 - bottom_margin 0.3)
MAX_CONTENT_HEIGHT = 3.8


def _estimate_sub_height(sub, width_in=8.4):
    """Estimate rendered height of a single subsegment."""
    kind, payload = sub
    if kind == "code":
        nlines = max(1, (payload or "").count("\n") + 1)
        return min(3.0, 0.26 * nlines + 0.2) + 0.1
    if kind == "table":
        header, rows = payload
        nrows = len(rows) + 1
        return min(3.2, 0.4 * nrows + 0.2) + 0.1
    if kind == "chart":
        return 3.2 + 0.1
    # text
    return _estimate_text_height(payload, width_in) + 0.1


def _split_text_sub_by_lines(text, max_h, width_in=8.4):
    """Split a large text block into smaller text subs by grouping lines.

    Groups lines into chunks whose estimated height stays under max_h.
    """
    lines = (text or "").split("\n")
    chunks = []
    cur_lines = []
    cur_h = 0.0
    for ln in lines:
        if not ln.strip():
            ln_h = 0.12
        elif ln.startswith("##"):
            ln_h = 0.5
        elif ln.startswith("-") or ln.startswith("*"):
            ln_h = 0.34
        elif ln.startswith(">"):
            ln_h = 0.3
        else:
            chars_per_line = max(20, int(width_in * 13))
            n = max(1, (len(ln) + chars_per_line - 1) // chars_per_line)
            ln_h = 0.28 * n

        if cur_lines and cur_h + ln_h > max_h:
            chunks.append(("text", "\n".join(cur_lines)))
            cur_lines = []
            cur_h = 0.0
        cur_lines.append(ln)
        cur_h += ln_h

    if cur_lines:
        chunks.append(("text", "\n".join(cur_lines)))
    return chunks


def _split_subs_by_height(subs, max_h=MAX_CONTENT_HEIGHT, width_in=8.4):
    """Split subsegments into batches that fit within max_h.

    Returns a list of batches, each a list of subsegments.
    If a single text subsegment exceeds max_h, its lines are split into
    smaller text subsegments before batching.
    """
    # Pre-split oversized text subs into line-grouped chunks
    expanded = []
    for sub in subs:
        kind, payload = sub
        if kind == "text" and _estimate_text_height(payload, width_in) > max_h:
            expanded.extend(_split_text_sub_by_lines(payload, max_h, width_in))
        else:
            expanded.append(sub)

    batches = []
    current = []
    current_h = 0.0
    for sub in expanded:
        h = _estimate_sub_height(sub, width_in)
        if current and current_h + h > max_h:
            batches.append(current)
            current = []
            current_h = 0.0
        current.append(sub)
        current_h += h
    if current:
        batches.append(current)
    return batches


def _render_text_block(slide, text, top, height, colors, font=None, left=0.8, width=8.4):
    font = font or {}
    h2_sz = font.get("h2", 22)
    body_sz = font.get("body", 16)
    font_name = font.get("name")
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    current_paragraph = None
    for line in (text or "").split("\n"):
        line = line.rstrip()
        if not line:
            current_paragraph = None
            continue
        if line.startswith("##"):
            current_paragraph = None
            p = tf.add_paragraph()
            _add_runs(p, line.lstrip("#").strip(), h2_sz, colors["primary"], bold=True, font_name=font_name)
            p.space_before = Pt(12)
            current_paragraph = p
        elif line.startswith("-") or line.startswith("*"):
            p = tf.add_paragraph()
            _add_runs(p, "• " + line.lstrip("-*").strip(), body_sz, colors["dark"], font_name=font_name)
            p.level = 1
            p.space_before = Pt(6)
        elif line.startswith(">"):
            p = tf.add_paragraph()
            _add_runs(p, line.lstrip(">").strip(), max(12, body_sz - 2), colors["gray"], italic=True, font_name=font_name)
        else:
            if not current_paragraph:
                current_paragraph = tf.add_paragraph()
                _add_runs(current_paragraph, line, body_sz, colors["dark"], font_name=font_name)
                current_paragraph.space_after = Pt(8)
            else:
                _add_runs(current_paragraph, " " + line, body_sz, colors["dark"], font_name=font_name)


# ==================== Table & Code Rendering ====================

def _is_table_row(line):
    s = line.strip()
    return s.startswith("|") and s.endswith("|") and s.count("|") >= 2


def _is_table_sep(line):
    s = line.strip()
    return s.startswith("|") and "-" in s and re.match(r"^\|[\s:|\-]+$", s) is not None


def _parse_table_row(line):
    return [c.strip() for c in line.strip().strip("|").split("|")]


def _split_text_subsegments(text):
    subs = []
    lines = (text or "").split("\n")
    i = 0
    buf = []

    def flush():
        if buf:
            subs.append(("text", "\n".join(buf)))
            buf.clear()

    while i < len(lines):
        line = lines[i]
        if line.strip().startswith("```"):
            flush()
            code_lines = []
            i += 1
            while i < len(lines) and not lines[i].strip().startswith("```"):
                code_lines.append(lines[i])
                i += 1
            if i < len(lines):
                i += 1
            subs.append(("code", "\n".join(code_lines)))
        elif _is_table_row(line) and i + 1 < len(lines) and _is_table_sep(lines[i + 1]):
            flush()
            header = _parse_table_row(line)
            i += 2
            rows = []
            while i < len(lines) and _is_table_row(lines[i]):
                rows.append(_parse_table_row(lines[i]))
                i += 1
            subs.append(("table", (header, rows)))
        else:
            buf.append(line)
            i += 1
    flush()
    return subs


def _render_code_block(slide, code, top, colors, left=0.8, width=8.4):
    code = code or ""
    nlines = max(1, code.count("\n") + 1)
    h = min(3.0, 0.26 * nlines + 0.2)
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    first = True
    for line in code.split("\n"):
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.text = line
        for r in p.runs:
            r.font.size = Pt(11)
            r.font.name = "Consolas"
            r.font.color.rgb = colors["dark"]
    return top + h + 0.1


def _render_table(slide, header, rows, top, colors, left=0.8, width=8.4):
    ncols = max(1, len(header))
    nrows = len(rows) + 1
    h = min(3.2, 0.4 * nrows + 0.2)
    tbl_shape = slide.shapes.add_table(nrows, ncols, Inches(left), Inches(top), Inches(width), Inches(h))
    tbl = tbl_shape.table
    for j in range(ncols):
        cell = tbl.cell(0, j)
        cell.text = header[j] if j < len(header) else ""
        cell.fill.solid()
        cell.fill.fore_color.rgb = colors["primary"]
        for p in cell.text_frame.paragraphs:
            for r in p.runs:
                r.font.bold = True
                r.font.size = Pt(12)
                r.font.color.rgb = colors["white"]
    for i, row in enumerate(rows):
        for j in range(ncols):
            cell = tbl.cell(i + 1, j)
            cell.text = row[j] if j < len(row) else ""
            for p in cell.text_frame.paragraphs:
                for r in p.runs:
                    r.font.size = Pt(11)
                    r.font.color.rgb = colors["dark"]
    return top + h + 0.1


# ==================== Chart Rendering ====================

def _smart_chart_type(ctype, cats, vals):
    """Auto-correct chart type based on data characteristics.

    - Single-value pie is meaningless → bar
    - Categories that look like time series → xychart (line)
    - Categories that look like A vs B comparison → bar
    """
    if ctype != "pie":
        return ctype

    # Single data point pie is meaningless
    if len(vals) <= 1:
        return "bar"

    # Check if categories look like time series
    time_patterns = ["月", "季", "周", "天", "年", "Q1", "Q2", "Q3", "Q4",
                     "week", "month", "year", "day", "jan", "feb", "mar",
                     "apr", "may", "jun", "jul", "aug", "sep", "oct", "nov", "dec"]
    if any(any(p.lower() in c.lower() for p in time_patterns) for c in cats):
        return "xychart"

    # If exactly 2 categories and they look like a comparison (not %)
    if len(cats) == 2 and not any("%" in c for c in cats):
        return "bar"

    return ctype


def add_native_chart(slide, chart_type, categories, values, colors, top, height=3.2, left=1.0, width=8.0):
    max_h = max(1.8, 5.42 - top)
    height = min(height, max_h)
    cd = CategoryChartData()
    cd.categories = categories
    cd.add_series("数据", values)
    gf = slide.shapes.add_chart(
        XL_MAP[chart_type], Inches(left), Inches(top), Inches(width), Inches(height), cd
    )
    chart = gf.chart
    chart.has_title = True
    chart.chart_title.text_frame.text = CHART_LABELS.get(chart_type, "图表")
    palette = [colors["primary"], colors["secondary"], colors["accent"]]
    try:
        series = chart.plots[0].series[0]
        for idx, pt in enumerate(series.points):
            pt.format.fill.solid()
            pt.format.fill.fore_color.rgb = palette[idx % len(palette)]
    except Exception:
        pass
    chart.has_legend = chart_type == "pie"
    if chart.has_legend:
        try:
            chart.legend.position = 2
        except Exception:
            pass
    return top + height


def _render_one(slide, sub, top, colors, font, left, width):
    kind, payload = sub
    if kind == "code":
        return _render_code_block(slide, payload, top, colors, left, width)
    if kind == "table":
        header, rows = payload
        return _render_table(slide, header, rows, top, colors, left, width)
    if kind == "chart":
        ctype, data_text = payload
        if ctype not in XL_MAP:
            ctype = "pie"
        cats, vals = parse_chart_data(data_text)
        if cats and vals:
            ctype = _smart_chart_type(ctype, cats, vals)
            return add_native_chart(slide, ctype, cats, vals, colors, top, left=left + 0.2, width=max(2.0, width - 0.4))
        return top + 0.6
    est_h = _estimate_text_height(payload, width_in=width)
    _render_text_block(slide, payload, top, est_h, colors, font, left, width)
    return top + est_h + 0.1


def _render_column(slide, subs, left, width, top, colors, font):
    for sub in subs:
        top = _render_one(slide, sub, top, colors, font, left, width)
    return top


# ==================== Slide Builders ====================

def _add_cover_slide(prs, title, subtitle, colors, font):
    font_name = font.get("name")
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    background = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    background.fill.solid()
    background.fill.fore_color.rgb = colors["primary"]
    background.line.fill.background()

    accent_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(1.2), prs.slide_width, Inches(0.15))
    accent_bar.fill.solid()
    accent_bar.fill.fore_color.rgb = colors["secondary"]
    accent_bar.line.fill.background()

    title_box = slide.shapes.add_textbox(Inches(1), Inches(2.2), Inches(8), Inches(1))
    tf = title_box.text_frame
    tf.text = title
    p = tf.paragraphs[0]
    _set_para_font(p, 48, colors["white"], bold=True, font_name=font_name)
    p.alignment = PP_ALIGN.CENTER

    if subtitle:
        sub_box = slide.shapes.add_textbox(Inches(1), Inches(3.5), Inches(8), Inches(1.2))
        stf = sub_box.text_frame
        stf.word_wrap = True
        p = stf.paragraphs[0]
        p.text = subtitle
        _set_para_font(p, 20, colors["white"], font_name=font_name)
        p.alignment = PP_ALIGN.CENTER


def _add_toc_slide(prs, page_titles, colors, font):
    font_name = font.get("name")
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    background = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    background.fill.solid()
    background.fill.fore_color.rgb = colors["light"]
    background.line.fill.background()

    accent_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(0.3))
    accent_bar.fill.solid()
    accent_bar.fill.fore_color.rgb = colors["primary"]
    accent_bar.line.fill.background()

    title_bg = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(0.5), Inches(9), Inches(0.9))
    title_bg.fill.solid()
    title_bg.fill.fore_color.rgb = colors["primary"]
    title_bg.line.color.rgb = colors["primary"]

    title_box = slide.shapes.add_textbox(Inches(0.7), Inches(0.55), Inches(8.6), Inches(0.8))
    tf = title_box.text_frame
    tf.text = "目录"
    p = tf.paragraphs[0]
    _set_para_font(p, 32, colors["white"], bold=True, font_name=font_name)

    top = Inches(1.8)
    n_items = len(page_titles)
    # Dynamic spacing: fit all items within slide height (5.625 - 1.8 - 0.3 bottom margin)
    available_h = 5.625 - 1.8 - 0.3
    spacing = min(0.65, available_h / max(n_items, 1)) if n_items > 0 else 0.65
    for i, pt_title in enumerate(page_titles):
        item_top = top + Inches(spacing * i)
        num_bg = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.3), item_top, Inches(0.4), Inches(0.4))
        num_bg.fill.solid()
        num_bg.fill.fore_color.rgb = colors["primary"]
        num_bg.line.color.rgb = colors["primary"]

        num_box = slide.shapes.add_textbox(Inches(1.3), item_top, Inches(0.4), Inches(0.4))
        ntf = num_box.text_frame
        ntf.text = str(i + 1)
        p = ntf.paragraphs[0]
        _set_para_font(p, 14, colors["white"], bold=True, font_name=font_name)
        p.alignment = PP_ALIGN.CENTER

        item_box = slide.shapes.add_textbox(Inches(1.8), item_top, Inches(7), Inches(0.5))
        itf = item_box.text_frame
        itf.text = pt_title
        p = itf.paragraphs[0]
        _set_para_font(p, 18, colors["dark"], font_name=font_name)
        p.space_before = Pt(6)


def _add_content_slide(prs, title, content, colors, font, layout="auto"):
    subs = []
    for kind, payload in split_segments(content):
        if kind == "chart":
            subs.append(("chart", payload))
        else:
            subs.extend(_split_text_subsegments(payload))

    RIGHT_KINDS = ("code", "table", "chart")
    has_right = any(s[0] in RIGHT_KINDS for s in subs)
    if layout == "auto":
        layout = "horizontal" if has_right else "vertical"
    top0 = 1.5

    if layout == "vertical" or not has_right:
        # Split subs into batches that fit within slide height
        batches = _split_subs_by_height(subs, MAX_CONTENT_HEIGHT)
        if not batches:
            batches = [[]]
        # First batch on the main slide
        slide = _create_content_slide_base(prs, title, colors, font)
        _render_column(slide, batches[0], 0.8, 8.4, top0, colors, font)
        # Continuation slides for overflow
        for i, batch in enumerate(batches[1:], 2):
            cont_title = f"{title}（续{i-1}）" if title else f"（续{i-1}）"
            cont_slide = _create_content_slide_base(prs, cont_title, colors, font)
            _render_column(cont_slide, batch, 0.8, 8.4, top0, colors, font)
    elif layout == "horizontal":
        slide = _create_content_slide_base(prs, title, colors, font)
        left_subs = [s for s in subs if s[0] == "text"]
        right_subs = [s for s in subs if s[0] in RIGHT_KINDS]
        # Split left column if too tall
        left_batches = _split_subs_by_height(left_subs, MAX_CONTENT_HEIGHT)
        _render_column(slide, left_batches[0] if left_batches else [], 0.6, 4.5, top0, colors, font)
        _render_column(slide, right_subs, 5.3, 4.3, top0, colors, font)
        # Continuation for left column overflow
        for i, batch in enumerate(left_batches[1:], 2):
            cont_title = f"{title}（续{i-1}）" if title else f"（续{i-1}）"
            cont_slide = _create_content_slide_base(prs, cont_title, colors, font)
            _render_column(cont_slide, batch, 0.6, 4.5, top0, colors, font)
    else:
        slide = _create_content_slide_base(prs, title, colors, font)
        first_text = next((s for s in subs if s[0] == "text"), None)
        rest = [s for s in subs if s is not first_text]
        top = top0
        if first_text:
            top = _render_one(slide, first_text, top, colors, font, 0.8, 8.4)
        left_subs = [s for s in rest if s[0] == "text"]
        right_subs = [s for s in rest if s[0] in RIGHT_KINDS]
        _render_column(slide, left_subs, 0.6, 4.5, top + 0.1, colors, font)
        _render_column(slide, right_subs, 5.3, 4.3, top + 0.1, colors, font)

    return slide


def _create_content_slide_base(prs, title, colors, font):
    """Create a content slide with background, accent bar, and title."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    background = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    background.fill.solid()
    background.fill.fore_color.rgb = colors["white"]
    background.line.fill.background()

    accent_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(0.2))
    accent_bar.fill.solid()
    accent_bar.fill.fore_color.rgb = colors["primary"]
    accent_bar.line.fill.background()

    title_bg = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.4), Inches(0.3), Inches(9.2), Inches(0.8))
    title_bg.fill.solid()
    title_bg.fill.fore_color.rgb = colors["primary"]
    title_bg.line.color.rgb = colors["primary"]
    title_box = slide.shapes.add_textbox(Inches(0.6), Inches(0.35), Inches(8.8), Inches(0.7))
    tf = title_box.text_frame
    tf.word_wrap = True
    _add_runs(tf.paragraphs[0], title or "", font.get("title", 28), colors["white"], bold=True, font_name=font.get("name"))

    return slide


def _add_summary_slide(prs, content, colors, font):
    font_name = font.get("name")
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    background = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    background.fill.solid()
    background.fill.fore_color.rgb = colors["light"]
    background.line.fill.background()

    accent_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(0.3))
    accent_bar.fill.solid()
    accent_bar.fill.fore_color.rgb = colors["accent"]
    accent_bar.line.fill.background()

    title_bg = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(0.5), Inches(9), Inches(0.9))
    title_bg.fill.solid()
    title_bg.fill.fore_color.rgb = colors["accent"]
    title_bg.line.color.rgb = colors["accent"]

    title_box = slide.shapes.add_textbox(Inches(0.7), Inches(0.55), Inches(8.6), Inches(0.8))
    tf = title_box.text_frame
    tf.text = "总结"
    p = tf.paragraphs[0]
    _set_para_font(p, 32, colors["white"], bold=True, font_name=font_name)

    content_bg = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.8), Inches(8.4), Inches(3.3))
    content_bg.fill.solid()
    content_bg.fill.fore_color.rgb = colors["white"]
    content_bg.line.color.rgb = colors["accent"]
    content_bg.line.width = Pt(2)

    # Parse content into subsegments (text, table, chart) and render properly
    subs = []
    for kind, payload in split_segments(content):
        if kind == "chart":
            subs.append(("chart", payload))
        else:
            subs.extend(_split_text_subsegments(payload))

    _render_column(slide, subs, 1.0, 8.0, 2.0, colors, font)


# ==================== Main Export ====================

def export_pptx(title, content, theme_key="business", custom_colors=None, font_config=None, subtitle=""):
    """
    Export markdown content to a PowerPoint presentation.

    Args:
        title: PPT title (cover page)
        content: Markdown content with ## page separators and chart data blocks
        theme_key: Color theme (business/tech/minimal/vibrant)
        custom_colors: Custom color dict {primary, secondary, accent} as hex strings
        font_config: Font configuration dict {theme, font_name, title, h2, body}
        subtitle: Cover page subtitle

    Returns:
        pptx.Presentation object
    """
    colors = get_theme_colors(theme_key, custom_colors)
    font = resolve_font(font_config or {})

    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(5.625)

    # 1. Cover slide
    _add_cover_slide(prs, title, subtitle, colors, font)

    # 2. Split content into pages by ##
    pages = split_by_h2(content) if content else []

    # 3. TOC slide (from page titles)
    page_titles = [t for t, _ in pages if t]
    _add_toc_slide(prs, page_titles, colors, font)

    # 4. Content slides
    summary_titles_exact = ("总结", "summary", "Summary", "结语", "结论")
    summary_keywords = ["总结", "结论", "行动", "下一步", "计划"]

    # Pre-detect: should the last page be treated as summary?
    last_is_summary = False
    if pages:
        last_title = pages[-1][0] or ""
        if last_title.lower() in (t.lower() for t in summary_titles_exact):
            last_is_summary = True
        elif any(kw in last_title for kw in summary_keywords):
            last_is_summary = True

    summary_content = None
    for i, (h2_title, body) in enumerate(pages):
        slide_title = h2_title or ""
        is_last = (i == len(pages) - 1)

        if slide_title.lower() in (t.lower() for t in summary_titles_exact):
            summary_content = body
            _add_summary_slide(prs, body, colors, font)
        elif is_last and last_is_summary:
            # Last page auto-detected as summary — render as summary, not content slide
            summary_content = body
            _add_summary_slide(prs, body, colors, font)
        else:
            _add_content_slide(prs, slide_title, body, colors, font)

    return prs


# ==================== CLI ====================

def main():
    parser = argparse.ArgumentParser(
        description="Convert AI-generated Markdown PPT content to .pptx file",
        formatter_class=argparse.RawDescriptionHelpFormatter,
    )
    parser.add_argument("--title", required=True, help="PPT title (cover page)")
    parser.add_argument("--subtitle", default="", help="Cover page subtitle")
    parser.add_argument("--content-file", "-", help="Path to markdown content file (use - for stdin)")
    parser.add_argument("--theme", default="business", choices=["business", "tech", "minimal", "vibrant"], help="Color theme")
    parser.add_argument("--custom-colors", help='Custom colors JSON, e.g. \'{"primary":"#0052A4","secondary":"#ED7D31","accent":"#10893E"}\'')
    parser.add_argument("--font-standard", action="store_true", help="Standard font size (28/22/16)")
    parser.add_argument("--font-compact", action="store_true", help="Compact font size (24/18/13)")
    parser.add_argument("--font-loose", action="store_true", help="Loose font size (32/26/19)")
    parser.add_argument("--font-name", default="", help="Font family name")
    parser.add_argument("--output", required=True, help="Output .pptx file path")

    args = parser.parse_args()

    # Read content
    if not args.content_file or args.content_file == "-":
        content = sys.stdin.read()
    else:
        with open(args.content_file, "r", encoding="utf-8") as f:
            content = f.read()

    # Parse custom colors
    custom_colors = None
    if args.custom_colors:
        try:
            custom_colors = json.loads(args.custom_colors)
        except json.JSONDecodeError as e:
            print(f"Error: Invalid custom-colors JSON: {e}", file=sys.stderr)
            sys.exit(1)

    # Determine font theme
    font_theme = "standard"
    if args.font_compact:
        font_theme = "compact"
    elif args.font_loose:
        font_theme = "loose"

    font_config = {"theme": font_theme}
    if args.font_name:
        font_config["font_name"] = args.font_name

    # Generate PPTX
    prs = export_pptx(
        title=args.title,
        content=content,
        theme_key=args.theme,
        custom_colors=custom_colors,
        font_config=font_config,
        subtitle=args.subtitle,
    )

    # Save
    prs.save(args.output)
    print(f"✅ PPT generated: {args.output}")
    print(f"   Slides: {len(prs.slides)}")
    print(f"   Theme: {args.theme}")
    print(f"   Font: {font_theme}")


if __name__ == "__main__":
    main()
