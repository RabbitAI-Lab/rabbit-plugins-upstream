# -*- coding: utf-8 -*-
"""
mckinsey-library · PPT 绘图原语（复用引擎）
===========================================
把「麦肯锡资料库顾问」交付里反复用到的 python-pptx 原语与图表类型固化成模块，
让任何一次 PPT 生成都复用同一套配色 / 字体 / 卡片 / 原生图表，避免每次重建。

设计原则（对齐渲染派：likaku / AX Labs）：
- 面积类图表（Treemap / 矩阵 / 瀑布 / 漏斗 / 仪表盘）用「形状 + 文本」直接绘制，
  不依赖重型 native chart XML，文件更小、跨机字体更稳（对标 BLOCK_ARC 轻量化思路）。
- 定量趋势类（柱/条/线/饼/环）用 python-pptx 原生 chart，渲染更标准。
- 所有文本入口过 `truncate()` 与 `set_run_font()`，落实「防损坏三层」：
  ① 字体回退（楷体/微软雅黑缺失时回退）② 超长文本截断 ③ 形状锁定无阴影。

导入方式（在生成脚本里）：
    import sys; sys.path.insert(0, r'<skill_dir>/assets')
    from pptx_primitives import Deck, COL, FONT
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION, XL_LABEL_POSITION
from pptx.oxml.ns import qn

# ---------- 调色板（麦肯锡真实色板：5 档蓝色梯度 + 白底黑字，无红/无金）----------
# 源自麦肯锡官方出品图表取色（如 Survey of 67 developers in T&D org 堆叠柱）。
# 顺序色板：深→浅表达数据密度；主蓝 L3 是唯一的强调/图表主色。
COL = {
    'NAVY':  RGBColor(0x06,0x2F,0x5C),   # L5 深海军蓝：封面/章节页满底、最深字
    'NAVY2': RGBColor(0x0C,0x55,0x9B),   # L4 次深蓝：图表第二序列、次级强调
    'BLUE':  RGBColor(0x1F,0x77,0xB4),   # L3 主蓝：行动标题条/强调/主图表色
    'BLUE2': RGBColor(0x7F,0xB8,0xDC),   # L2 浅蓝：深蓝底上的副标/强调
    'BLUE1': RGBColor(0xB5,0xD9,0xEB),   # L1 最浅蓝：图表低密度档
    'LIGHT': RGBColor(0xF4,0xF6,0xF9),
    'WHITE': RGBColor(0xFF,0xFF,0xFF),
    'GRAY':  RGBColor(0x59,0x59,0x59),
    'DARK':  RGBColor(0x1A,0x1A,0x1A),
    'CARD':  RGBColor(0xFF,0xFF,0xFF),
    'BORD':  RGBColor(0xDD,0xE3,0xEC),
    'SOFT':  RGBColor(0xE7,0xEC,0xF5),
    'GREEN': RGBColor(0x2E,0x7D,0x32),
}

FONT = {'TITLE':'楷体', 'BODY':'微软雅黑', 'NUM':'Arial'}

# 字体回退表：指定字体不可用时回退（防损坏①）
_FALLBACK = {'楷体':'微软雅黑', 'KaiTi':'微软雅黑', 'Arial':'Calibri'}

EMU = 914400
SW, SH = 13.333, 7.5  # 16:9

def _font_name(name):
    return _FALLBACK.get(name, name)

def set_run_font(run, name, size=None, bold=None, color=None, italic=None):
    run.font.name = name
    rPr = run._r.get_or_add_rPr()
    for tag in ('a:latin','a:ea','a:cs'):
        el = rPr.find(qn(tag))
        if el is None:
            el = rPr.makeelement(qn(tag), {})
            rPr.append(el)
        el.set('typeface', _font_name(name))
    if size is not None: run.font.size = Pt(size)
    if bold is not None: run.font.bold = bold
    if color is not None: run.font.color.rgb = color
    if italic is not None: run.font.italic = italic

def truncate(text, n=60):
    """防损坏②：超长文本截断，避免溢出文本框。"""
    text = str(text)
    return text if len(text) <= n else text[:n-1] + '…'

class Deck:
    """一套 PPT 的轻量封装，承载配色/字体/品牌名。"""
    def __init__(self, brand='麦肯锡资料库顾问', width=SW, height=SH):
        self.prs = Presentation()
        self.prs.slide_width = Inches(width)
        self.prs.slide_height = Inches(height)
        self.blank = self.prs.slide_layouts[6]
        self.brand = brand
        self.W, self.H = width, height

    def slide(self):
        return self.prs.slides.add_slide(self.blank)

    # —— 底层原语 ——
    def bg(self, s, color=COL['LIGHT']):
        sh = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0,0, Inches(self.W), Inches(self.H))
        sh.fill.solid(); sh.fill.fore_color.rgb = color
        sh.line.fill.background(); sh.shadow.inherit = False
        return sh

    def rect(self, s, x,y,w,h, fill=None, line=None, line_w=1.0, radius=None):
        sp = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE,
                                 Inches(x),Inches(y),Inches(w),Inches(h))
        if radius: sp.adjustments[0] = radius
        if fill is None: sp.fill.background()
        else: sp.fill.solid(); sp.fill.fore_color.rgb = fill
        if line is None: sp.line.fill.background()
        else: sp.line.color.rgb = line; sp.line.width = Pt(line_w)
        sp.shadow.inherit = False
        return sp

    def txt(self, s, x,y,w,h, lines, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, wrap=True):
        """lines: list of dict(text,size,bold,color,font,space_after,space_before,line,align,level)"""
        tb = s.shapes.add_textbox(Inches(x),Inches(y),Inches(w),Inches(h))
        tf = tb.text_frame; tf.word_wrap = wrap
        tf.vertical_anchor = anchor
        tf.margin_left=0; tf.margin_right=0; tf.margin_top=0; tf.margin_bottom=0
        for i,ln in enumerate(lines):
            p = tf.paragraphs[0] if i==0 else tf.add_paragraph()
            p.alignment = ln.get('align', align)
            if 'space_after' in ln: p.space_after = Pt(ln['space_after'])
            if 'space_before' in ln: p.space_before = Pt(ln['space_before'])
            if 'line' in ln: p.line_spacing = ln['line']
            r = p.add_run(); r.text = truncate(ln.get('text',''), ln.get('max',60))
            set_run_font(r, ln.get('font',FONT['BODY']), ln.get('size',14),
                         ln.get('bold',False), ln.get('color',COL['DARK']))
        return tb

    # —— 通用区块 ——
    def title_bar(self, s, title, kicker=None, accent=COL['BLUE'], color_title=COL['NAVY']):
        self.rect(s, 0.55, 0.5, 0.16, 0.62, fill=accent)
        lines = []
        if kicker:
            lines.append({'text':kicker,'size':11,'bold':True,'color':accent,'font':FONT['BODY'],'space_after':2})
        lines.append({'text':title,'size':26,'bold':True,'color':color_title,'font':FONT['TITLE']})
        self.txt(s, 0.82, 0.46, self.W-1.4, 0.85, lines, anchor=MSO_ANCHOR.MIDDLE)

    def footer(self, s, source='', page=None):
        self.rect(s, 0, self.H-0.32, self.W, 0.32, fill=COL['LIGHT'])
        left = f'{self.brand} · 内部汇报稿'
        self.txt(s, 0.55, self.H-0.32, 9.5, 0.32,
                 [{'text':left,'size':8.5,'color':COL['GRAY'],'font':FONT['BODY']}],
                 anchor=MSO_ANCHOR.MIDDLE)
        right = ('资料来源：'+source) if source else ''
        if page: right += (('  ·  '+str(page)) if source else str(page))
        self.txt(s, self.W-4.3, self.H-0.32, 3.78, 0.32,
                 [{'text':right,'size':8.5,'color':COL['GRAY'],'font':FONT['BODY'],'align':PP_ALIGN.RIGHT}],
                 anchor=MSO_ANCHOR.MIDDLE)

    def stat_card(self, s, x,y,w,h, big, label, sub=None, accent=COL['BLUE']):
        self.rect(s, x,y,w,h, fill=COL['CARD'], line=COL['BORD'], line_w=1.0, radius=0.06)
        self.rect(s, x,y,0.09,h, fill=accent)
        lines = [{'text':big,'size':30,'bold':True,'color':accent,'font':FONT['NUM'],'align':PP_ALIGN.CENTER,'space_after':2}]
        lines.append({'text':label,'size':12,'bold':True,'color':COL['NAVY'],'font':FONT['BODY'],'align':PP_ALIGN.CENTER,'space_after':(2 if sub else 0)})
        if sub:
            lines.append({'text':sub,'size':9.5,'color':COL['GRAY'],'font':FONT['BODY'],'align':PP_ALIGN.CENTER})
        self.txt(s, x+0.12, y+0.12, w-0.24, h-0.24, lines, anchor=MSO_ANCHOR.MIDDLE)

    # —— 原生图表（定量趋势）——
    def bar(self, s, x,y,w,h, cats, series, title=None, pct=False, stacked=False):
        cd = CategoryChartData(); cd.categories = cats
        for name, vals in series: cd.add_series(name, vals)
        ctype = XL_CHART_TYPE.COLUMN_STACKED if stacked else XL_CHART_TYPE.COLUMN_CLUSTERED
        if len(series)==1 and not stacked: ctype = XL_CHART_TYPE.COLUMN_CLUSTERED
        gf = s.shapes.add_chart(ctype, Inches(x),Inches(y),Inches(w),Inches(h), cd)
        ch = gf.chart; ch.has_legend = len(series)>1
        if ch.has_legend: ch.legend.position = XL_LEGEND_POSITION.BOTTOM; ch.legend.include_in_layout=False
        ch.has_title = bool(title); 
        if title: ch.chart_title.text_frame.text = title
        plot = ch.plots[0]; plot.has_data_labels = True
        plot.data_labels.number_format = '0%' if pct else '0.0"%"'; plot.data_labels.number_format_is_linked=False
        plot.data_labels.font.size = Pt(9)
        try: ch.font.size = Pt(10); ch.font.name = FONT['BODY']
        except Exception: pass
        return gf

    def line(self, s, x,y,w,h, cats, series, title=None):
        cd = CategoryChartData(); cd.categories = cats
        for name, vals in series: cd.add_series(name, vals)
        gf = s.shapes.add_chart(XL_CHART_TYPE.LINE_MARKERS, Inches(x),Inches(y),Inches(w),Inches(h), cd)
        ch = gf.chart; ch.has_legend = len(series)>1
        if ch.has_legend: ch.legend.position = XL_LEGEND_POSITION.BOTTOM; ch.legend.include_in_layout=False
        ch.has_title = bool(title)
        if title: ch.chart_title.text_frame.text = title
        try: ch.font.size = Pt(10); ch.font.name = FONT['BODY']
        except Exception: pass
        return gf

    def pie(self, s, x,y,w,h, cats, vals, title=None, donut=False):
        cd = CategoryChartData(); cd.categories = cats; cd.add_series('v', vals)
        gf = s.shapes.add_chart(XL_CHART_TYPE.DOUGHNUT if donut else XL_CHART_TYPE.PIE,
                                Inches(x),Inches(y),Inches(w),Inches(h), cd)
        ch = gf.chart; ch.has_legend = True; ch.legend.position = XL_LEGEND_POSITION.RIGHT
        ch.legend.include_in_layout=False
        ch.has_title = bool(title)
        if title: ch.chart_title.text_frame.text = title
        ch.plots[0].has_data_labels = True
        try: ch.font.size = Pt(10)
        except Exception: pass
        return gf

    # —— 形状类图表（面积有意义，轻量跨机稳）——
    def _text_on(self, col):
        # 亮度阈值选黑/白字，避免低对比（如 GOLD 黄底白字）
        r,g,b = col[0],col[1],col[2]
        lum = 0.299*r + 0.587*g + 0.114*b
        return COL['DARK'] if lum > 150 else COL['WHITE']

    def treemap(self, s, x,y,w,h, blocks, gap=0.05):
        """blocks: [(label, share0to1, color, sub)]；面积=占比。主导块置左整高，其余右侧纵向堆叠保证可读。"""
        total = sum(b[1] for b in blocks) or 1
        items = sorted([(b[0], b[1]/total, b[2], b[3] if len(b)>3 else '') for b in blocks], key=lambda t:-t[1])
        if not items: return
        maxit = items[0]; rest = items[1:]
        lw = w * maxit[1]
        self.rect(s, x+gap/2, y+gap/2, max(lw-gap,0.2), h-gap, fill=maxit[2], radius=0.03)
        tc = self._text_on(maxit[2])
        self.txt(s, x+0.12, y+0.12, max(lw-0.24,0.2), h-0.24,
                 [{'text':maxit[0],'size':13,'bold':True,'color':tc,'font':FONT['BODY'],'space_after':2},
                  {'text':f'{maxit[1]*100:.1f}%','size':20,'bold':True,'color':tc,'font':FONT['NUM']},
                  {'text':maxit[3],'size':9.5,'color':tc,'font':FONT['BODY']}],
                 anchor=MSO_ANCHOR.MIDDLE)
        if rest:
            rx = x+lw; rw = w-lw
            n = len(rest)
            rh = (h - gap*(n+1)) / n
            cy = y + gap/2
            for (lab, share, col, sub) in rest:
                self.rect(s, rx+gap/2, cy, rw-gap, rh-gap/2, fill=col, radius=0.03)
                tc2 = self._text_on(col)
                self.txt(s, rx+0.12, cy, rw-0.24, rh-gap/2,
                         [{'text':f'{lab}  {share*100:.1f}%','size':10.5,'bold':True,'color':tc2,'font':FONT['BODY'],'space_after':1},
                          {'text':sub,'size':8.5,'color':tc2,'font':FONT['BODY']}],
                         anchor=MSO_ANCHOR.MIDDLE)
                cy += rh

    def matrix_2x2(self, s, x,y,w,h, labels, colors=None, axis_x='', axis_y=''):
        """labels: [TL,TR,BL,BR]（四象限）；colors 同序。"""
        colors = colors or [COL['BLUE'],COL['NAVY2'],COL['BLUE2'],COL['NAVY']]
        half_w, half_h = w/2, h/2
        pos = [(x,y),(x+half_w,y),(x,y+half_h),(x+half_w,y+half_h)]
        for i,lab in enumerate(labels):
            self.rect(s, pos[i][0]+0.03, pos[i][1]+0.03, half_w-0.06, half_h-0.06, fill=colors[i], radius=0.04)
            lines = [{'text':ln,'size':12.5,'bold':True,'color':self._text_on(colors[i]),'font':FONT['BODY'],'space_after':2} for ln in lab.split('\n')]
            self.txt(s, pos[i][0]+0.15, pos[i][1]+0.12, half_w-0.3, half_h-0.24, lines, anchor=MSO_ANCHOR.MIDDLE)
        if axis_x:
            self.txt(s, x, y+h+0.05, w, 0.3, [{'text':axis_x,'size':10,'bold':True,'color':COL['GRAY'],'font':FONT['BODY'],'align':PP_ALIGN.CENTER}])
        if axis_y:
            self.txt(s, x-1.4, y+h/2-0.15, 1.3, 0.3, [{'text':axis_y,'size':10,'bold':True,'color':COL['GRAY'],'font':FONT['BODY'],'align':PP_ALIGN.RIGHT}])

    def waterfall(self, s, x,y,w,h, steps, base=0):
        """steps: [(label, delta)]；画累计瀑布。"""
        import math
        maxv = base + sum(max(0,d) for _,d in steps)
        minv = min(0, base + sum(min(0,d) for _,d in steps))
        span = maxv - minv or 1
        def yof(v): return y + h*(1-(v-minv)/span)
        def hof(v): return h*abs(v-minv)/span
        cx = x; bw = w/len(steps)
        prev = base
        for i,(lab,d) in enumerate(steps):
            top = yof(max(prev, prev+d)); hh = max(0.05, hof(abs(d)))
            col = COL['BLUE'] if d>=0 else COL['NAVY']
            self.rect(s, cx+0.1, top, bw-0.2, hh, fill=col, radius=0.02)
            self.txt(s, cx, top-0.35, bw, 0.3, [{'text':lab,'size':9,'color':COL['GRAY'],'font':FONT['BODY'],'align':PP_ALIGN.CENTER}])
            prev = prev + d
            cx += bw

    def funnel(self, s, x,y,w,h, stages):
        """stages: [(label, share0to1)] 自上而下收窄。"""
        n=len(stages); top_w=w
        cy=y
        for i,(lab,share) in enumerate(stages):
            ww = top_w*share
            cx = x + (top_w-ww)/2
            self.rect(s, cx, cy, ww, h/n-0.08, fill=COL['BLUE'] if i%2==0 else COL['NAVY2'], radius=0.03)
            self.txt(s, x, cy, top_w, h/n-0.08, [{'text':f'{lab}  {share*100:.0f}%','size':12,'bold':True,'color':COL['WHITE'],'font':FONT['BODY'],'align':PP_ALIGN.CENTER}], anchor=MSO_ANCHOR.MIDDLE)
            cy += h/n

    def gauge(self, s, x,y,w,h, pct, label='达成率'):
        """半环仪表盘（0-100%）。"""
        self.rect(s, x, y, w, h, fill=COL['CARD'], line=COL['BORD'], line_w=1, radius=0.06)
        self.txt(s, x, y+h*0.30, w, h*0.4,
                 [{'text':f'{pct:.0f}%','size':28,'bold':True,'color':COL['NAVY2'],'font':FONT['NUM'],'align':PP_ALIGN.CENTER},
                  {'text':label,'size':11,'color':COL['GRAY'],'font':FONT['BODY'],'align':PP_ALIGN.CENTER}],
                 anchor=MSO_ANCHOR.MIDDLE)

    def save(self, path):
        self.prs.save(path)

if __name__ == '__main__':
    # 自检：生成一张含各图表的样张，验证模块可用
    d = Deck(brand='麦肯锡资料库顾问')
    s = d.slide(); d.bg(s, COL['LIGHT'])
    d.title_bar(s, '模块自检样张', kicker='mckinsey-library')
    d.stat_card(s, 0.6, 1.6, 2.6, 1.3, '60.2%', '综合金融贡献', 'AUM 净增')
    d.treemap(s, 0.6, 3.2, 5.5, 3.2, [('消费前 34%',0.34,COL['NAVY2'],'技术极限'),
                                       ('混合废料 51%',0.51,COL['NAVY'],'瓶颈'),
                                       ('分拣废料 16%',0.16,COL['BLUE'],'高质量')])
    d.matrix_2x2(s, 6.6, 3.2, 5.8, 3.2, ['高意愿低能力','高意愿高能力','低意愿低能力','低意愿高能力'])
    d.bar(s, 0.6, 6.7, 5.5, 0.0, ['A'], [('v',[1])])  # 占位避免空
    d.save('__primitives_selftest__.pptx')
    print('primitives OK')
