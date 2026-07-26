#!/usr/bin/env python3
"""
Knowledge-base driven PPTX generator with Apple-style minimalist design.

Input: JSON file defining slide structure
Output: .pptx file with speaker notes, chart placeholders, and quality metadata

Usage:
    python generate_pptx.py --input slides.json --output report.pptx
    python generate_pptx.py --input slides.json --output report.pptx --template custom.pptx
"""

import argparse
import json
import os
import sys
from datetime import datetime
from pathlib import Path

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu, Cm
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
    from pptx.enum.shapes import MSO_SHAPE
except ImportError:
    print("Error: python-pptx is required. Install with: pip install python-pptx")
    sys.exit(1)

# ── Apple-style Design Constants ──────────────────────────────────────────

# Color palette (light, clean, minimal)
COLORS = {
    "bg": RGBColor(0xFF, 0xFF, 0xFF),       # White
    "text_primary": RGBColor(0x1D, 0x1D, 0x1F),  # Near black
    "text_secondary": RGBColor(0x86, 0x86, 0x8B),  # Gray
    "accent": RGBColor(0x00, 0x7A, 0xFF),    # Apple Blue
    "accent_alt": RGBColor(0x34, 0xC7, 0x59),  # Green (positive)
    "accent_warn": RGBColor(0xFF, 0x95, 0x00),  # Orange (warning)
    "accent_danger": RGBColor(0xFF, 0x3B, 0x30),  # Red (danger)
    "line_light": RGBColor(0xE5, 0xE5, 0xEA),  # Light border
    "tag_bg": RGBColor(0xF2, 0xF2, 0xF7),   # Light gray tag background
    "tag_text": RGBColor(0x00, 0x7A, 0xFF),  # Blue tag text
}

# Typography
FONT_PRIMARY = "PingFang SC"
FONT_FALLBACK = "Microsoft YaHei"
FONT_MONO = "SF Mono"

# Sizes
SLIDE_WIDTH = Inches(13.333)  # 16:9 widescreen
SLIDE_HEIGHT = Inches(7.5)

TITLE_SIZE = Pt(36)
SUBTITLE_SIZE = Pt(20)
BODY_SIZE = Pt(16)
SMALL_SIZE = Pt(12)
CAPTION_SIZE = Pt(10)

# Margins and padding
MARGIN_LEFT = Inches(1.2)
MARGIN_RIGHT = Inches(1.2)
MARGIN_TOP = Inches(0.8)
CONTENT_TOP = Inches(2.0)

# ── Helper Functions ──────────────────────────────────────────────────────

def get_font_family():
    """Return available font, falling back gracefully."""
    return FONT_PRIMARY


def set_slide_bg(slide, color=COLORS["bg"]):
    """Set slide background to solid color."""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_textbox(slide, left, top, width, height, text="",
                font_size=BODY_SIZE, color=COLORS["text_primary"],
                bold=False, alignment=PP_ALIGN.LEFT, font_name=None):
    """Add a text box to a slide."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = font_size
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name or get_font_family()
    p.alignment = alignment
    return txBox, tf


def add_rich_textbox(slide, left, top, width, height, lines,
                     font_size=BODY_SIZE, color=COLORS["text_primary"],
                     line_spacing=Pt(24)):
    """
    Add a text box with multiple lines.
    Each line is a dict: {text, bold, color, size} or a plain string.
    """
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True

    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()

        if isinstance(line, str):
            p.text = line
            p.font.size = font_size
            p.font.color.rgb = color
            p.font.bold = False
        else:
            p.text = line.get("text", "")
            p.font.size = line.get("size", font_size)
            p.font.color.rgb = line.get("color", color)
            p.font.bold = line.get("bold", False)

        p.font.name = get_font_family()
        p.space_after = line_spacing

    return txBox, tf


def add_placeholder_shape(slide, left, top, width, height, label="",
                          color=COLORS["line_light"]):
    """Add a dashed-border placeholder rectangle for charts/images."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0xF9, 0xF9, 0xFB)

    # Dashed line
    shape.line.color.rgb = color
    shape.line.width = Pt(1)
    shape.line.dash_style = 2  # dash

    if label:
        tf = shape.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = label
        p.font.size = Pt(11)
        p.font.color.rgb = COLORS["text_secondary"]
        p.font.name = get_font_family()
        p.alignment = PP_ALIGN.CENTER

    return shape


def add_notes(slide, notes_text):
    """Add speaker notes to a slide."""
    notes_slide = slide.notes_slide
    tf = notes_slide.notes_text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = notes_text
    p.font.size = Pt(10)
    p.font.name = get_font_family()


def add_page_number(slide, page_num, total_pages):
    """Add page number at bottom right."""
    add_textbox(
        slide, SLIDE_WIDTH - Inches(1.5), SLIDE_HEIGHT - Inches(0.5),
        Inches(1.0), Inches(0.3),
        f"{page_num} / {total_pages}",
        font_size=CAPTION_SIZE, color=COLORS["text_secondary"],
        alignment=PP_ALIGN.RIGHT
    )


# ── Page Layouts ──────────────────────────────────────────────────────────

def create_cover_slide(prs, slide_data):
    """Create a cover/title slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    set_slide_bg(slide)

    title = slide_data.get("title", "Untitled")
    subtitle = slide_data.get("subtitle", "")
    author = slide_data.get("author", "")
    date_str = slide_data.get("date", datetime.now().strftime("%Y.%m.%d"))

    # Decorative top accent line
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        MARGIN_LEFT, Inches(1.8),
        Inches(0.8), Pt(4)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = COLORS["accent"]
    line.line.fill.background()

    # Title
    add_textbox(
        slide, MARGIN_LEFT, Inches(2.2), Inches(10.0), Inches(1.2),
        title, font_size=Pt(44), bold=True, color=COLORS["text_primary"]
    )

    # Subtitle
    if subtitle:
        add_textbox(
            slide, MARGIN_LEFT, Inches(3.4), Inches(10.0), Inches(0.8),
            subtitle, font_size=SUBTITLE_SIZE, color=COLORS["text_secondary"]
        )

    # Author + date
    add_textbox(
        slide, MARGIN_LEFT, Inches(5.2), Inches(5.0), Inches(0.4),
        f"{author}  |  {date_str}",
        font_size=SMALL_SIZE, color=COLORS["text_secondary"]
    )

    if slide_data.get("notes"):
        add_notes(slide, slide_data["notes"])

    return slide


def create_section_slide(prs, slide_data):
    """Create a section divider slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)

    section_title = slide_data.get("title", "")
    section_desc = slide_data.get("subtitle", "")

    # Large section number or title centered
    add_textbox(
        slide, MARGIN_LEFT, Inches(2.8), Inches(10.0), Inches(1.0),
        section_title, font_size=Pt(40), bold=True, color=COLORS["text_primary"]
    )

    if section_desc:
        add_textbox(
            slide, MARGIN_LEFT, Inches(3.8), Inches(10.0), Inches(0.6),
            section_desc, font_size=SUBTITLE_SIZE, color=COLORS["text_secondary"]
        )

    if slide_data.get("notes"):
        add_notes(slide, slide_data["notes"])

    return slide


def create_content_slide(prs, slide_data):
    """Create a standard content slide with title, subtitle, and bullet points."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)

    title = slide_data.get("title", "")
    subtitle = slide_data.get("subtitle", "")
    bullets = slide_data.get("bullets", [])
    chart_config = slide_data.get("chart", None)
    notes = slide_data.get("notes", "")
    chart_label = slide_data.get("chart_label", "[ 图表占位: 请在 PPT 中替换 ]")

    # Title
    add_textbox(
        slide, MARGIN_LEFT, Inches(0.6), Inches(10.0), Inches(0.7),
        title, font_size=TITLE_SIZE, bold=True, color=COLORS["text_primary"]
    )

    # Thin separator line under title
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        MARGIN_LEFT, Inches(1.35), Inches(10.8), Pt(1)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = COLORS["line_light"]
    line.line.fill.background()

    # Subtitle / one-line conclusion
    if subtitle:
        add_textbox(
            slide, MARGIN_LEFT, Inches(1.5), Inches(10.0), Inches(0.5),
            subtitle, font_size=SUBTITLE_SIZE, color=COLORS["accent"],
            bold=True
        )
        bullet_top = Inches(2.3)
    else:
        bullet_top = Inches(1.7)

    # Bullet points or chart: choose layout
    if chart_config:
        # Left side: bullets, Right side: chart placeholder
        content_width = Inches(4.8)

        # Bullets
        if bullets:
            add_rich_textbox(
                slide, MARGIN_LEFT, bullet_top, content_width, Inches(4.5),
                bullets, font_size=BODY_SIZE
            )

        # Chart placeholder on the right
        chart_left = Inches(7.2)
        add_placeholder_shape(
            slide, chart_left, bullet_top,
            Inches(5.0), Inches(4.5),
            label=chart_label
        )
    else:
        # Full width bullets
        add_rich_textbox(
            slide, MARGIN_LEFT, bullet_top, Inches(10.8), Inches(4.8),
            bullets, font_size=BODY_SIZE
        )

    if notes:
        add_notes(slide, notes)

    return slide


def create_conclusion_slide(prs, slide_data):
    """Create a conclusion/closing slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)

    title = slide_data.get("title", "总结与下一步")
    bullets = slide_data.get("bullets", [])

    # Accent line
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        MARGIN_LEFT, Inches(1.8), Inches(0.8), Pt(4)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = COLORS["accent"]
    line.line.fill.background()

    # Title
    add_textbox(
        slide, MARGIN_LEFT, Inches(2.2), Inches(10.0), Inches(1.0),
        title, font_size=Pt(40), bold=True, color=COLORS["text_primary"]
    )

    if bullets:
        add_rich_textbox(
            slide, MARGIN_LEFT, Inches(3.4), Inches(10.8), Inches(3.5),
            bullets, font_size=Pt(18)
        )

    if slide_data.get("notes"):
        add_notes(slide, slide_data["notes"])

    return slide


def create_blank_slide(prs, slide_data):
    """Create a blank slide for custom content."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)

    if slide_data.get("notes"):
        add_notes(slide, slide_data["notes"])

    return slide


LAYOUT_MAP = {
    "cover": create_cover_slide,
    "section": create_section_slide,
    "content": create_content_slide,
    "conclusion": create_conclusion_slide,
    "blank": create_blank_slide,
    "内容": create_content_slide,
    "封面": create_cover_slide,
    "分割": create_section_slide,
    "总结": create_conclusion_slide,
    "空白": create_blank_slide,
}


# Generation credit (appended to last page speaker notes)
GENERATION_CREDIT = "此文件由小朱AI养成skill生成，有问题请联系VX：zyc19892233"


# ── Main Generation Logic ─────────────────────────────────────────────────

def _add_generation_credit(slide, existing_notes=""):
    """Append generation credit to the last page's speaker notes."""
    notes_slide = slide.notes_slide
    tf = notes_slide.notes_text_frame

    if existing_notes:
        # Append credit below existing notes
        p = tf.add_paragraph()
        p.text = ""
        p.font.size = Pt(10)
        p.font.name = get_font_family()

        p2 = tf.add_paragraph()
        p2.text = "—"
        p2.font.size = Pt(10)
        p2.font.name = get_font_family()

        p3 = tf.add_paragraph()
        p3.text = GENERATION_CREDIT
        p3.font.size = Pt(10)
        p3.font.name = get_font_family()
    else:
        p = tf.paragraphs[0]
        p.text = GENERATION_CREDIT
        p.font.size = Pt(10)
        p.font.name = get_font_family()


def generate_pptx(slides_data, output_path):
    """
    Generate PPTX from slides_data.

    slides_data structure:
    {
        "meta": {
            "title": "Presentation Title",
            "author": "Author Name",
            "date": "2024.01.15"
        },
        "slides": [
            {
                "layout": "cover",           # cover | content | section | conclusion | blank
                "title": "Page Title",
                "subtitle": "One-line conclusion",
                "bullets": [
                    "Bullet point 1",
                    {"text": "Bullet 2 with formatting", "bold": true, "color": "#007AFF"},
                    ...
                ],
                "chart": "chart-type",        # Optional: indicates a chart placeholder
                "chart_label": "Custom chart placeholder text",
                "notes": "Speaker notes for this slide"
            },
            ...
        ]
    }
    """
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT

    meta = slides_data.get("meta", {})
    slides = slides_data.get("slides", [])
    total = len(slides)

    # Track metadata for output report
    chart_slides = []
    missing_data_slides = []
    placeholder_slides = []

    for i, slide_data in enumerate(slides):
        layout = slide_data.get("layout", "content")
        create_func = LAYOUT_MAP.get(layout, create_content_slide)

        # Auto-set cover meta if not provided
        if layout == "cover":
            slide_data.setdefault("title", meta.get("title", ""))
            slide_data.setdefault("subtitle", meta.get("subtitle", ""))
            slide_data.setdefault("author", meta.get("author", ""))
            slide_data.setdefault("date", meta.get("date", ""))

        slide = create_func(prs, slide_data)

        # Track chart placeholders
        if slide_data.get("chart"):
            chart_slides.append({
                "page": i + 1,
                "title": slide_data.get("title", ""),
                "chart_type": slide_data["chart"],
            })

        # Track missing data
        missing = slide_data.get("missing_data", [])
        if missing:
            missing_data_slides.append({
                "page": i + 1,
                "items": missing if isinstance(missing, list) else [missing],
            })

        # Track placeholder content
        placeholders = slide_data.get("placeholders", [])
        if placeholders:
            placeholder_slides.append({
                "page": i + 1,
                "items": placeholders if isinstance(placeholders, list) else [placeholders],
            })

        # Add page numbers (skip cover)
        if layout != "cover":
            add_page_number(slide, i + 1, total)

        # Add generation credit to the last page's speaker notes
        if i == total - 1:
            _add_generation_credit(slide, slide_data.get("notes", ""))

    # Save
    output_path = Path(output_path)
    prs.save(str(output_path))

    # Generate report
    report = generate_report(
        output_path, meta, chart_slides, missing_data_slides,
        placeholder_slides
    )

    return output_path, report


def generate_report(output_path, meta, chart_slides, missing_data_slides,
                    placeholder_slides):
    """Generate a summary report of the generated PPTX."""
    report_path = output_path.with_suffix(".report.txt")

    lines = []
    lines.append("=" * 60)
    lines.append("PPTX 生成报告")
    lines.append("=" * 60)
    lines.append(f"文件: {output_path}")
    lines.append(f"标题: {meta.get('title', 'N/A')}")
    lines.append(f"生成时间: {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}")
    lines.append("")

    # Chart placeholders
    lines.append("-" * 40)
    lines.append("📊 图表占位清单 (需用实际图表替换):")
    lines.append("-" * 40)
    if chart_slides:
        for item in chart_slides:
            lines.append(f"  第 {item['page']} 页 [{item['title']}]: {item['chart_type']}")
    else:
        lines.append("  (无)")

    # Missing data
    lines.append("")
    lines.append("-" * 40)
    lines.append("📋 需要用户补充的数据清单:")
    lines.append("-" * 40)
    if missing_data_slides:
        for item in missing_data_slides:
            for data_item in item["items"]:
                lines.append(f"  第 {item['page']} 页: {data_item}")
    else:
        lines.append("  (无)")

    # Replaceable content
    lines.append("")
    lines.append("-" * 40)
    lines.append("🔄 可替换占位内容清单:")
    lines.append("-" * 40)
    if placeholder_slides:
        for item in placeholder_slides:
            for ph in item["items"]:
                lines.append(f"  第 {item['page']} 页: {ph}")
    else:
        lines.append("  (无)")

    lines.append("")
    lines.append("=" * 60)
    lines.append("Speaker Notes 清单:")
    lines.append("=" * 60)
    lines.append("(每页 Speaker Notes 已嵌入 PPTX 文件中，可通过 PowerPoint 的「备注」视图查看)")
    lines.append("")

    report_text = "\n".join(lines)
    with open(report_path, "w", encoding="utf-8") as f:
        f.write(report_text)

    return report_text


# ── CLI ───────────────────────────────────────────────────────────────────

def main():
    parser = argparse.ArgumentParser(
        description="Generate Apple-style PPTX from JSON slide definitions"
    )
    parser.add_argument(
        "--input", "-i", required=True,
        help="Path to JSON file with slide definitions"
    )
    parser.add_argument(
        "--output", "-o", required=True,
        help="Output PPTX file path"
    )
    args = parser.parse_args()

    # Load input
    with open(args.input, "r", encoding="utf-8") as f:
        slides_data = json.load(f)

    # Validate
    if "slides" not in slides_data:
        print("Error: JSON must contain a 'slides' array")
        sys.exit(1)

    # Generate
    output_path, report = generate_pptx(slides_data, args.output)

    print(f"✅ PPTX 已生成: {output_path}")
    print(f"📄 报告已生成: {output_path.with_suffix('.report.txt')}")
    print()
    print(report)


if __name__ == "__main__":
    main()
