#!/usr/bin/env python3
"""Render a ClawVision summary as HTML, PNG, Markdown, or PowerPoint."""

import argparse
import json
import re
import sys
from pathlib import Path

from playwright.sync_api import sync_playwright
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

PRESETS = {
    "minimal": {
        "font_family": "Inter", "pptx_font": "Inter", "title_size": "22px",
        "radius": "14px", "shadow": "0 2px 8px rgba(0,0,0,.04)", "pptx_radius": 0.12, "pptx_shadow": False,
        "light": {
            "bg": "#f5f6f8",
            "card": "#ffffff",
            "text": "#1a1a1a",
            "muted": "#666666",
            "accent": "#2a9df4",
            "green": "#4cd964",
            "orange": "#ff9500",
            "red": "#ff3b30",
            "border": "#e5e7eb",
            "card_alt": "#f8f9fa",
        },
        "dark": {
            "bg": "#0f1115",
            "card": "#1a1d23",
            "text": "#e8e8e8",
            "muted": "#9aa0a6",
            "accent": "#4aa8ff",
            "green": "#5dd877",
            "orange": "#ffae33",
            "red": "#ff6659",
            "border": "#2c3038",
            "card_alt": "#22262e",
        },
    },
    "editorial": {
        "font_family": "Georgia", "pptx_font": "Georgia", "title_size": "24px",
        "radius": "6px", "shadow": "0 1px 3px rgba(0,0,0,.05)", "pptx_radius": 0.03, "pptx_shadow": False,
        "light": {
            "bg": "#faf9f6",
            "card": "#ffffff",
            "text": "#2a1f1d",
            "muted": "#7d6f66",
            "accent": "#7d2c2c",
            "green": "#3a7d44",
            "orange": "#b86d29",
            "red": "#a83232",
            "border": "#e6e2db",
            "card_alt": "#f4f1ec",
        },
        "dark": {
            "bg": "#1a1714",
            "card": "#26221d",
            "text": "#efe8df",
            "muted": "#a89b8c",
            "accent": "#d17b7b",
            "green": "#7bc47f",
            "orange": "#e5a15c",
            "red": "#d66a6a",
            "border": "#3d362e",
            "card_alt": "#2f2923",
        },
    },
    "retro": {
        "font_family": "Space Grotesk", "pptx_font": "Arial Rounded MT Bold", "title_size": "24px",
        "radius": "22px", "shadow": "4px 4px 0 rgba(0,0,0,.18)", "pptx_radius": 0.2, "pptx_shadow": True,
        "light": {
            "bg": "#fff6e5",
            "card": "#fffdf5",
            "text": "#3d2b1f",
            "muted": "#8c7b66",
            "accent": "#e67300",
            "green": "#5a8a4b",
            "orange": "#c47b00",
            "red": "#c13b3b",
            "border": "#d9c9a8",
            "card_alt": "#ffefc8",
        },
        "dark": {
            "bg": "#2a2118",
            "card": "#3b2f22",
            "text": "#ffeecd",
            "muted": "#c4b598",
            "accent": "#ff9a3c",
            "green": "#8bd17a",
            "orange": "#ffbf47",
            "red": "#ff7b7b",
            "border": "#5c4b36",
            "card_alt": "#4a3a2a",
        },
    },
    "luxury": {
        "font_family": "Playfair Display", "pptx_font": "Playfair Display", "title_size": "26px",
        "radius": "3px", "shadow": "0 6px 24px rgba(0,0,0,.1)", "pptx_radius": 0.01, "pptx_shadow": False,
        "light": {
            "bg": "#f7f5f2",
            "card": "#ffffff",
            "text": "#1a1a1a",
            "muted": "#7a736a",
            "accent": "#bfa35a",
            "green": "#3d6e4e",
            "orange": "#9a6b2e",
            "red": "#8a3333",
            "border": "#e0d8cd",
            "card_alt": "#faf8f5",
        },
        "dark": {
            "bg": "#12100e",
            "card": "#1d1a16",
            "text": "#f2efe9",
            "muted": "#a89f93",
            "accent": "#e3c87c",
            "green": "#7fc493",
            "orange": "#d4a65a",
            "red": "#c77a7a",
            "border": "#3a332b",
            "card_alt": "#252019",
        },
    },
    "playful": {
        "font_family": "Nunito", "pptx_font": "Nunito", "title_size": "24px",
        "radius": "26px", "shadow": "0 8px 0 rgba(0,0,0,.08)", "pptx_radius": 0.28, "pptx_shadow": True,
        "light": {
            "bg": "#f3f6ff",
            "card": "#ffffff",
            "text": "#2d2a45",
            "muted": "#6e6a8a",
            "accent": "#7c4dff",
            "green": "#00c853",
            "orange": "#ff9100",
            "red": "#ff1744",
            "border": "#d6dcff",
            "card_alt": "#e8edff",
        },
        "dark": {
            "bg": "#1a1830",
            "card": "#252244",
            "text": "#eae8ff",
            "muted": "#9d99bd",
            "accent": "#b388ff",
            "green": "#69f0ae",
            "orange": "#ffd180",
            "red": "#ff8a80",
            "border": "#3f3a6b",
            "card_alt": "#2f2b55",
        },
    },
}

def _css_var(preset: dict) -> str:
    def _hex6(c: str) -> str:
        c = c.lstrip("#")
        if len(c) == 3:
            c = "".join(ch * 2 for ch in c)
        return c
    def _rgb(c: str) -> str:
        h = _hex6(c)
        return f"{int(h[0:2],16)},{int(h[2:4],16)},{int(h[4:6],16)}"
    light = preset["light"]
    dark = preset["dark"]
    return f"""
:root{{
  --bg:{light['bg']};--card:{light['card']};--text:{light['text']};--muted:{light['muted']};
  --accent:{light['accent']};--accent-rgb:{_rgb(light['accent'])};
  --green:{light['green']};--orange:{light['orange']};--red:{light['red']};
  --border:{light['border']};--card-bg-alt:{light.get('card_alt','#f8f9fa')};
}}
:root.dark{{
  --bg:{dark['bg']};--card:{dark['card']};--text:{dark['text']};--muted:{dark['muted']};
  --accent:{dark['accent']};--accent-rgb:{_rgb(dark['accent'])};
  --green:{dark['green']};--orange:{dark['orange']};--red:{dark['red']};
  --border:{dark['border']};--card-bg-alt:{dark.get('card_alt','#22262e')};
}}
"""

def _base_css(preset: dict) -> str:
    radius = preset["radius"]
    shadow = preset["shadow"]
    font_stack = f"{preset['font_family']}, -apple-system, BlinkMacSystemFont, 'Segoe UI', Roboto, Arial, sans-serif"
    css = """
*{box-sizing:border-box}
body{margin:0;font-family:__FONT_STACK__;background:var(--bg);color:var(--text);line-height:1.45}
.wrap{max-width:820px;margin:0 auto;padding:24px}
.topbar{display:flex;justify-content:space-between;align-items:center;gap:12px;flex-wrap:wrap;margin-bottom:22px}
.controls{display:flex;gap:8px;flex-wrap:wrap}
.export-btn,.lang-btn,.theme-btn{border:none;border-radius:20px;padding:8px 14px;font-size:12px;font-weight:600;cursor:pointer;transition:.15s}
.export-btn{background:var(--accent);color:#fff;box-shadow:__SHADOW__}
.export-btn:hover{filter:brightness(1.1)}
.lang-switch,.theme-switch{display:flex;gap:6px;background:var(--card);border:1px solid var(--border);border-radius:20px;padding:4px}
.lang-btn,.theme-btn{background:transparent;color:var(--muted)}
.lang-btn.active,.theme-btn.active{background:var(--accent);color:#fff}
header{margin-bottom:20px}
header h1{font-size:__TITLE_SIZE__;font-weight:800;margin:0 0 6px;letter-spacing:-0.5px}
header p{color:var(--muted);margin:0;font-size:15px}
.badge{display:inline-block;background:rgba(var(--accent-rgb),.12);color:var(--accent);font-size:12px;font-weight:700;padding:5px 12px;border-radius:14px;margin-bottom:14px;letter-spacing:.5px;text-transform:uppercase}
.tabs{display:flex;gap:8px;flex-wrap:wrap;margin-bottom:20px}
.tab{background:var(--card);border:1px solid var(--border);border-radius:__RADIUS__;padding:9px 15px;font-size:13px;cursor:pointer;transition:.15s}
.tab.active{background:var(--accent);color:#fff;border-color:var(--accent)}
.panel{display:none;background:var(--card);border-radius:__RADIUS__;padding:22px;box-shadow:__SHADOW__}
.panel.active{display:block}
.lead{color:var(--muted);font-size:16px;margin-bottom:18px;line-height:1.5}
.flow{display:flex;align-items:center;gap:10px;flex-wrap:wrap;margin:16px 0}
.flow-item{background:var(--card-bg-alt);border:1px solid var(--border);border-radius:__RADIUS__;padding:12px 16px;min-width:110px;text-align:center}
.flow-item strong{display:block;font-size:14px;margin-bottom:3px}
.flow-item small{color:var(--muted);font-size:12px}
.arrow{color:var(--muted);font-size:18px}
.grid{display:grid;grid-template-columns:repeat(auto-fit,minmax(220px,1fr));gap:12px;margin-top:14px}
.card{background:var(--card-bg-alt);border-radius:__RADIUS__;padding:16px}
.card h3{margin:0 0 6px;font-size:17px;color:var(--accent)}
.card p{margin:0;color:var(--muted);font-size:13px}
.checklist{margin:0;padding:0;list-style:none}
.checklist li{display:flex;justify-content:space-between;align-items:center;padding:12px 0;border-bottom:1px solid var(--border);font-size:14px}
.checklist li:last-child{border-bottom:none}
.status{font-size:12px;font-weight:700;color:var(--green);background:rgba(76,217,100,.12);padding:3px 10px;border-radius:10px}
.status.pending{color:var(--orange);background:rgba(255,149,0,.12)}
.status.blocked{color:var(--red);background:rgba(255,59,48,.12)}
.two{display:grid;grid-template-columns:1fr 1fr;gap:12px}
ul.clean{padding-left:18px;margin:0}
ul.clean li{margin-bottom:6px;font-size:13px}
@media(max-width:600px){.two{grid-template-columns:1fr}.topbar{flex-direction:column;align-items:flex-start}}
"""
    css = css.replace('__FONT_STACK__', font_stack)
    css = css.replace('__RADIUS__', radius)
    css = css.replace('__SHADOW__', shadow)
    css = css.replace('__TITLE_SIZE__', preset['title_size'])
    return _css_var(preset) + css

TEMPLATE_HTML = r"""
<!DOCTYPE html>
<html lang="{{lang}}">
<head>
<meta charset="UTF-8">
<meta name="viewport" content="width=device-width, initial-scale=1.0">
<title>{{title}}</title>
<style>
{{css}}
</style>
</head>
<body>
<div class="wrap">
  <div class="topbar">
    <div class="controls">
      <a class="export-btn" href="{{md_file}}" download>Export Markdown</a>
      <a class="export-btn" href="{{pptx_file}}" download>Export PowerPoint</a>
    </div>
    <div class="controls">
      <div class="theme-switch"><button class="theme-btn" data-theme="light">☀</button><button class="theme-btn" data-theme="dark">🌙</button></div>
      <div class="lang-switch"><button class="lang-btn" data-lang="en">EN</button><button class="lang-btn" data-lang="ru">RU</button><button class="lang-btn" data-lang="zh">中文</button></div>
    </div>
  </div>

  <span class="badge">{{badge}}</span>
  <header>
    <h1>{{title}}</h1>
    <p>{{subtitle}}</p>
  </header>

  <div class="tabs">
    <div class="tab active" data-tab="main"><span id="main-title">{{tab_main}}</span></div>
    <div class="tab" data-tab="format"><span id="format-title">{{tab_format}}</span></div>
    <div class="tab" data-tab="built"><span id="built-title">{{tab_built}}</span></div>
    <div class="tab" data-tab="next"><span id="next-title">{{tab_next}}</span></div>
  </div>

  <div id="main" class="panel active">
    <p class="lead">{{main_takeaway}}</p>
    <div class="flow">{{flow_items}}</div>
    <div class="grid">{{metric_cards}}</div>
  </div>

  <div id="format" class="panel">
    <p class="lead">{{format_takeaway}}</p>
    <div class="two">
      <div class="card">
        <h3 id="dos-title">{{dos_title}}</h3>
        <ul class="clean">{{dos}}</ul>
      </div>
      <div class="card">
        <h3 id="donts-title">{{donts_title}}</h3>
        <ul class="clean">{{donts}}</ul>
      </div>
    </div>
  </div>

  <div id="built" class="panel">
    <ul class="checklist">{{checklist}}</ul>
  </div>

  <div id="next" class="panel">
    <p class="lead">{{next_takeaway}}</p>
    <ul class="clean">{{next_steps}}</ul>
  </div>
</div>
<script>
const LANGS={{langs_json}};
function setLang(l){
  document.body.dataset.lang=l;
  localStorage.setItem('cv-lang',l);
  document.querySelectorAll('.lang-btn').forEach(b=>b.classList.toggle('active',b.dataset.lang===l));
  const labels=LANGS[l]||LANGS.en;
  ['main','format','built','next'].forEach(k=>document.getElementById(k+'-title').textContent=labels[k]);
  document.getElementById('dos-title').textContent=labels.dos;
  document.getElementById('donts-title').textContent=labels.donts;
}
function setTheme(t){
  document.documentElement.classList.toggle('dark',t==='dark');
  localStorage.setItem('cv-theme',t);
  document.querySelectorAll('.theme-btn').forEach(b=>b.classList.toggle('active',b.dataset.theme===t));
}
document.querySelectorAll('.tab').forEach(t=>t.addEventListener('click',()=>{
  document.querySelectorAll('.tab').forEach(x=>x.classList.remove('active'));
  document.querySelectorAll('.panel').forEach(x=>x.classList.remove('active'));
  t.classList.add('active');
  document.getElementById(t.dataset.tab).classList.add('active');
}));
document.querySelectorAll('.lang-btn').forEach(b=>b.addEventListener('click',()=>setLang(b.dataset.lang)));
document.querySelectorAll('.theme-btn').forEach(b=>b.addEventListener('click',()=>setTheme(b.dataset.theme)));
setLang(localStorage.getItem('cv-lang')||'{{lang}}');
setTheme(localStorage.getItem('cv-theme')||'light');
</script>
</body>
</html>
"""

LANG_LABELS = {
    "en": {"main": "Main takeaway", "format": "Format", "built": "What we built", "next": "Next steps", "dos": "Do", "donts": "Don't", "export_md": "Export Markdown", "export_pptx": "Export PowerPoint", "badge": "ClawVision · Summary"},
    "ru": {"main": "Главный вывод", "format": "Формат", "built": "Что построено", "next": "Что дальше", "dos": "Нормально", "donts": "Риски", "export_md": "Экспорт Markdown", "export_pptx": "Экспорт PowerPoint", "badge": "ClawVision · Сводка"},
    "zh": {"main": "主要结论", "format": "形式", "built": "已完成", "next": "下一步", "dos": "建议", "donts": "风险", "export_md": "导出 Markdown", "export_pptx": "导出 PowerPoint", "badge": "ClawVision · 摘要"},
}

def _h(text: str) -> str:
    return ((text or '').replace('&', '&amp;').replace('<', '&lt;').replace('>', '&gt;').replace('"', '&quot;'))

def render_html(summary: dict, slug: str, md_file: str, pptx_file: str, lang: str = 'en', preset_name: str = 'minimal') -> str:
    preset = PRESETS.get(preset_name, PRESETS['minimal'])
    labels = LANG_LABELS.get(lang, LANG_LABELS['en'])
    template = TEMPLATE_HTML
    template = template.replace('{{css}}', _base_css(preset))
    template = template.replace('{{lang}}', _h(lang))
    template = template.replace('{{langs_json}}', json.dumps(LANG_LABELS, ensure_ascii=False))
    template = template.replace('{{md_file}}', _h(md_file))
    template = template.replace('{{pptx_file}}', _h(pptx_file))
    template = template.replace('{{title}}', _h(summary.get('title', 'ClawVision summary')))
    template = template.replace('{{subtitle}}', _h(summary.get('subtitle', '')))
    template = template.replace('{{badge}}', _h(labels['badge']))
    template = template.replace('{{tab_main}}', labels['main'])
    template = template.replace('{{tab_format}}', labels['format'])
    template = template.replace('{{tab_built}}', labels['built'])
    template = template.replace('{{tab_next}}', labels['next'])
    template = template.replace('{{main_takeaway}}', _h(summary.get('main_takeaway', '')))
    template = template.replace('{{format_takeaway}}', _h(summary.get('format_takeaway', '')))
    template = template.replace('{{next_takeaway}}', _h(summary.get('next_takeaway', '')))
    template = template.replace('{{dos_title}}', labels['dos'])
    template = template.replace('{{donts_title}}', labels['donts'])
    flow_items = ''
    for item in summary.get('flow', []):
        label = item.get('label', '')
        sub = item.get('sub', '')
        if label in ('→', '->'):
            flow_items += f"<div class=\"arrow\">{_h(label)}</div>\n"
        else:
            flow_items += f"<div class=\"flow-item\"><strong>{_h(label)}</strong><small>{_h(sub)}</small></div>\n"
    template = template.replace('{{flow_items}}', flow_items)
    metric_cards = ''
    for m in summary.get('metrics', []):
        metric_cards += f"<div class=\"card\"><h3>{_h(m.get('title', ''))}</h3><p>{_h(m.get('text', ''))}</p></div>\n"
    template = template.replace('{{metric_cards}}', metric_cards)
    dos = ''.join(f'<li>{_h(d)}</li>\n' for d in summary.get('dos', []))
    template = template.replace('{{dos}}', dos)
    donts = ''.join(f'<li>{_h(d)}</li>\n' for d in summary.get('donts', []))
    template = template.replace('{{donts}}', donts)
    checklist = ''
    for c in summary.get('checklist', []):
        status = c.get('status', 'pending')
        checklist += f"<li>{_h(c.get('text', ''))} <span class=\"status {status}\">{_h(status)}</span></li>\n"
    template = template.replace('{{checklist}}', checklist)
    next_steps = ''.join(f'<li>{_h(s)}</li>\n' for s in summary.get('next_steps', []))
    template = template.replace('{{next_steps}}', next_steps)
    return template

def render_md(summary: dict, lang: str = 'en', preset_name: str = 'minimal') -> str:
    labels = LANG_LABELS.get(lang, LANG_LABELS['en'])
    lines = [f"# {summary.get('title', 'ClawVision summary')}", f"_{summary.get('subtitle', '')}_", "", f"**Preset:** {preset_name}", "", f"## {labels['main']}", summary.get("main_takeaway", ""), "", f"## {labels['format']}", summary.get("format_takeaway", ""), "", f"### {labels['dos']}"]
    lines.extend(f'- {d}' for d in summary.get('dos', []))
    lines.append('')
    lines.append(f"### {labels['donts']}")
    lines.extend(f'- {d}' for d in summary.get('donts', []))
    lines.append('')
    lines.append(f"## {labels['built']}")
    for c in summary.get('checklist', []):
        mark = '[x]' if c.get('status') == 'ready' else '[ ]'
        lines.append(f"- {mark} {c.get('text', '')} ({c.get('status', 'pending')})")
    lines.append('')
    lines.append(f"## {labels['next']}")
    lines.append(summary.get('next_takeaway', ''))
    lines.append('')
    lines.extend(f'- {s}' for s in summary.get('next_steps', []))
    return '\n'.join(lines)

def _rgb(color: str):
    color = color.lstrip("#")
    if len(color) == 3:
        color = ''.join(c * 2 for c in color)
    return RGBColor(int(color[0:2], 16), int(color[2:4], 16), int(color[4:6], 16))

def render_pptx(summary: dict, pptx_path: Path, lang: str = 'en', preset_name: str = 'minimal'):
    preset = PRESETS.get(preset_name, PRESETS['minimal'])
    labels = LANG_LABELS.get(lang, LANG_LABELS['en'])
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    BG = _rgb(preset['light']['bg'])
    CARD = _rgb(preset['light']['card'])
    TEXT = _rgb(preset['light']['text'])
    MUTED = _rgb(preset['light']['muted'])
    ACCENT = _rgb(preset['light']['accent'])
    GREEN = _rgb(preset['light']['green'])
    ORANGE = _rgb(preset['light']['orange'])
    RED = _rgb(preset['light']['red'])
    BORDER = _rgb(preset['light']['border'])
    FONT = preset['pptx_font']
    radius = preset['pptx_radius']
    WHITE = RGBColor(255, 255, 255)

    def set_slide_bg(slide, color):
        fill = slide.background.fill
        fill.solid()
        fill.fore_color.rgb = color

    def rounded_rect(slide, left, top, width, height, fill):
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
        shape.line.color.rgb = BORDER
        shape.line.width = Pt(1)
        if radius:
            shape.adjustments[0] = radius
        return shape

    def set_text(shape, text, font_size=14, bold=False, color=TEXT, align=PP_ALIGN.LEFT):
        tf = shape.text_frame
        tf.word_wrap = True
        tf.clear()
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(font_size)
        p.font.bold = bold
        p.font.color.rgb = color
        p.font.name = FONT
        p.alignment = align
        return tf

    def textbox(slide, left, top, width, height, text, font_size=14, bold=False, color=TEXT, align=PP_ALIGN.LEFT):
        box = slide.shapes.add_textbox(left, top, width, height)
        set_text(box, text, font_size, bold, color, align)
        return box

    slide = prs.slides.add_slide(blank)
    set_slide_bg(slide, BG)
    badge = rounded_rect(slide, Inches(0.6), Inches(0.5), Inches(2.8), Inches(0.45), ACCENT)
    set_text(badge, labels['badge'], font_size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    title_size = int(preset['title_size'].replace('px', ''))
    textbox(slide, Inches(0.6), Inches(1.2), Inches(12), Inches(1.2), summary.get('title', 'ClawVision summary'), font_size=max(30, int(title_size * 0.65)), bold=True, color=TEXT)
    textbox(slide, Inches(0.6), Inches(2.4), Inches(12), Inches(0.8), summary.get('subtitle', ''), font_size=22, color=MUTED)

    slide = prs.slides.add_slide(blank)
    set_slide_bg(slide, BG)
    textbox(slide, Inches(0.6), Inches(0.4), Inches(12), Inches(0.7), labels['main'], font_size=32, bold=True, color=TEXT)
    lead = rounded_rect(slide, Inches(0.6), Inches(1.2), Inches(12), Inches(0.9), CARD)
    set_text(lead, summary.get('main_takeaway', ''), font_size=18, color=TEXT)
    flow = summary.get('flow', [])
    y_flow = Inches(2.35)
    x = Inches(0.6)
    for item in flow:
        label = item.get('label', '')
        sub = item.get('sub', '')
        if label in ('→', '->'):
            textbox(slide, x, y_flow + Inches(0.1), Inches(0.5), Inches(0.5), '→', font_size=24, bold=True, color=MUTED, align=PP_ALIGN.CENTER)
            x += Inches(0.55)
        else:
            width = min(max(Inches(1.7), Inches(0.32 * max(len(label), len(sub) or 0))), Inches(2.8))
            card = rounded_rect(slide, x, y_flow, width, Inches(0.9), CARD)
            set_text(card, f'{label}\n{sub}', font_size=13, bold=True, color=TEXT, align=PP_ALIGN.CENTER)
            x += width + Inches(0.12)
    metrics = summary.get('metrics', [])
    if metrics:
        y_metrics = Inches(3.55)
        total_width = Inches(12)
        gap = Inches(0.25)
        n = len(metrics)
        card_w = (total_width - gap * (n - 1)) / n
        for i, m in enumerate(metrics):
            left = Inches(0.6) + i * (card_w + gap)
            card = rounded_rect(slide, left, y_metrics, card_w, Inches(1.85), CARD)
            tf = card.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = m.get('title', '')
            p.font.size = Pt(18)
            p.font.bold = True
            p.font.color.rgb = ACCENT
            p.font.name = FONT
            p.alignment = PP_ALIGN.LEFT
            p2 = tf.add_paragraph()
            p2.text = m.get('text', '')
            p2.font.size = Pt(15)
            p2.font.color.rgb = TEXT
            p2.font.name = FONT
            p2.space_before = Pt(8)

    slide = prs.slides.add_slide(blank)
    set_slide_bg(slide, BG)
    textbox(slide, Inches(0.6), Inches(0.4), Inches(12), Inches(0.7), labels['format'], font_size=32, bold=True, color=TEXT)
    lead = rounded_rect(slide, Inches(0.6), Inches(1.2), Inches(12), Inches(0.9), CARD)
    set_text(lead, summary.get('format_takeaway', ''), font_size=18, color=TEXT)
    col_w = (Inches(12) - Inches(0.3)) / 2
    do_card = rounded_rect(slide, Inches(0.6), Inches(2.35), col_w, Inches(4.2), CARD)
    tf = do_card.text_frame
    p = tf.paragraphs[0]
    p.text = labels['dos']
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = GREEN
    p.font.name = FONT
    for d in summary.get('dos', []):
        p = tf.add_paragraph()
        p.text = f'• {d}'
        p.font.size = Pt(16)
        p.font.color.rgb = TEXT
        p.font.name = FONT
        p.space_before = Pt(6)
    dont_card = rounded_rect(slide, Inches(0.6) + col_w + Inches(0.3), Inches(2.35), col_w, Inches(4.2), CARD)
    tf = dont_card.text_frame
    p = tf.paragraphs[0]
    p.text = labels['donts']
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = RED
    p.font.name = FONT
    for d in summary.get('donts', []):
        p = tf.add_paragraph()
        p.text = f'• {d}'
        p.font.size = Pt(16)
        p.font.color.rgb = TEXT
        p.font.name = FONT
        p.space_before = Pt(6)

    slide = prs.slides.add_slide(blank)
    set_slide_bg(slide, BG)
    textbox(slide, Inches(0.6), Inches(0.4), Inches(12), Inches(0.7), labels['built'], font_size=32, bold=True, color=TEXT)
    y = Inches(1.3)
    for c in summary.get('checklist', []):
        status = c.get('status', 'pending')
        status_color = GREEN if status == 'ready' else ORANGE if status == 'pending' else RED
        icon = '✓' if status == 'ready' else '…' if status == 'pending' else '✕'
        row = rounded_rect(slide, Inches(0.6), y, Inches(12), Inches(0.6), CARD)
        tf = row.text_frame
        tf.margin_left = Inches(0.15)
        tf.margin_right = Inches(0.15)
        p = tf.paragraphs[0]
        p.text = f'{icon}  {c.get("text", "")}' 
        p.font.size = Pt(16)
        p.font.color.rgb = TEXT
        p.font.name = FONT
        p.alignment = PP_ALIGN.LEFT
        badge_w = Inches(1.0)
        badge = rounded_rect(slide, Inches(12.6) - badge_w, y + Inches(0.1), badge_w, Inches(0.4), status_color)
        set_text(badge, status.upper(), font_size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        y += Inches(0.75)

    slide = prs.slides.add_slide(blank)
    set_slide_bg(slide, BG)
    textbox(slide, Inches(0.6), Inches(0.4), Inches(12), Inches(0.7), labels['next'], font_size=32, bold=True, color=TEXT)
    lead = rounded_rect(slide, Inches(0.6), Inches(1.2), Inches(12), Inches(0.9), CARD)
    set_text(lead, summary.get('next_takeaway', ''), font_size=18, color=TEXT)
    y = Inches(2.35)
    for s in summary.get('next_steps', []):
        row = rounded_rect(slide, Inches(0.6), y, Inches(12), Inches(0.55), CARD)
        set_text(row, f'→  {s}', font_size=16, color=TEXT)
        y += Inches(0.7)

    prs.save(str(pptx_path))

def screenshot_tabs(html_path: Path, slug: str, width: int = 900, height: int = 650):
    tab_ids = ['main', 'format', 'built', 'next']
    paths = []
    with sync_playwright() as p:
        browser = p.chromium.launch(headless=True)
        page = browser.new_page(viewport={'width': width, 'height': height})
        page.goto(f"file:///{html_path.as_posix()}")
        for idx, tab_id in enumerate(tab_ids, start=1):
            page.locator(f"[data-tab='{tab_id}']").click()
            page.wait_for_timeout(150)
            png_path = html_path.parent / f'{slug}_tab{idx}.png'
            page.screenshot(path=str(png_path), full_page=False)
            paths.append(str(png_path))
        browser.close()
    return paths

def screenshot_html(html_path: Path, png_path: Path, width: int = 900, height: int = 650):
    with sync_playwright() as p:
        browser = p.chromium.launch(headless=True)
        page = browser.new_page(viewport={'width': width, 'height': height})
        page.goto(f"file:///{html_path.as_posix()}")
        page.screenshot(path=str(png_path), full_page=False)
        browser.close()

def slugify(text: str) -> str:
    text = re.sub(r"[^\w\s-]", "", text).strip().lower()
    return re.sub(r"[-\s]+", "-", text) or "summary"

def main():
    parser = argparse.ArgumentParser(description='Render ClawVision summary visual.')
    parser.add_argument('--summary', '-s', help='Path to JSON summary or - for stdin')
    parser.add_argument('--output', '-o', required=True, help='Output directory')
    parser.add_argument('--slug', help='Output slug (default: derived from title)')
    parser.add_argument('--lang', default='zh', help='Language code for UI labels'),
    parser.add_argument('--preset', default='minimal', help='Aesthetic preset: minimal, editorial, retro, luxury, playful')
    parser.add_argument('--png', action='store_true', help='Render PNG screenshots')
    parser.add_argument('--md', action='store_true', help='Render Markdown summary')
    parser.add_argument('--pptx', action='store_true', help='Render PowerPoint deck')
    args = parser.parse_args()

    if args.summary == '-' or args.summary is None:
        raw = sys.stdin.read()
    else:
        raw = Path(args.summary).read_text(encoding='utf-8')
    summary = json.loads(raw)

    out_dir = Path(args.output)
    out_dir.mkdir(parents=True, exist_ok=True)
    slug = args.slug or slugify(summary.get('title', 'summary'))
    md_file = f'{slug}.md'
    pptx_file = f'{slug}.pptx'
    html_path = out_dir / f'{slug}.html'
    html_path.write_text(render_html(summary, slug, md_file, pptx_file, lang=args.lang, preset_name=args.preset), encoding='utf-8')

    result = {"html": str(html_path), "preset": args.preset}
    if args.md:
        md_path = out_dir / md_file
        md_path.write_text(render_md(summary, lang=args.lang, preset_name=args.preset), encoding='utf-8')
        result["md"] = str(md_path)
    if args.pptx:
        pptx_path = out_dir / pptx_file
        render_pptx(summary, pptx_path, lang=args.lang, preset_name=args.preset)
        result["pptx"] = str(pptx_path)
    if args.png:
        tab_paths = screenshot_tabs(html_path, slug)
        result["png_tabs"] = tab_paths
        png_path = out_dir / f'{slug}.png'
        screenshot_html(html_path, png_path)
        result["png"] = str(png_path)

    print(json.dumps(result, ensure_ascii=False, indent=2))

if __name__ == '__main__':
    main()
