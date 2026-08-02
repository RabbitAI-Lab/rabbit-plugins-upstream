"""
ahkb_extract_pptx.py — PPTX 解析器
提取：每页文字、嵌入图片、嵌入音视频、全页截图（纯 Python，无外部依赖）
输出：chunks 嵌套结构，资源归属明确
"""
import os, json, zipfile
from pathlib import Path
from xml.etree import ElementTree as ET

# ─── 媒体文件类型检测 ───

IMAGE_EXTS = {'.png', '.jpg', '.jpeg', '.gif', '.bmp', '.tiff', '.wmf', '.emf'}
VIDEO_EXTS = {'.mp4', '.avi', '.mov', '.wmv', '.m4v', '.mpg', '.mpeg', '.flv', '.webm'}
AUDIO_EXTS = {'.mp3', '.wav', '.wma', '.aac', '.ogg', '.flac', '.m4a', '.mid', '.midi'}
# 所有媒体扩展名
MEDIA_EXTS = IMAGE_EXTS | VIDEO_EXTS | AUDIO_EXTS


def _get_media_type(ext):
    """根据扩展名判断媒体类型。"""
    ext = ext.lower()
    if ext in IMAGE_EXTS:
        return "image"
    elif ext in VIDEO_EXTS:
        return "video"
    elif ext in AUDIO_EXTS:
        return "audio"
    return "other"


def extract_pptx(filepath, workspace):
    """Extract text, images, and media from PPTX. Returns structured dict with chunks."""
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    prs = Presentation(filepath)
    base = Path(filepath).stem
    safe_base = "".join(c if c.isalnum() or c in '-_ ' else '_' for c in base)

    # 目录结构
    img_dir = Path(workspace) / "图片及其他资源" / "images"
    video_dir = Path(workspace) / "图片及其他资源" / "videos"
    audio_dir = Path(workspace) / "图片及其他资源" / "audios"
    other_dir = Path(workspace) / "图片及其他资源" / "others"
    for d in [img_dir, video_dir, audio_dir, other_dir]:
        d.mkdir(parents=True, exist_ok=True)

    result = {
        "file": str(filepath),
        "type": "pptx",
        "metadata": {
            "slide_count": len(prs.slides),
            "slide_width": prs.slide_width,
            "slide_height": prs.slide_height,
        },
        "chunks": [],
        "full_text": "",
        # 扁平列表向后兼容 + 方便 .ctx 生成
        "resources_flat": [],
    }

    img_counter = 0

    # ── 预先获取 ZIP 中的媒体文件清单（视频/音频等）──
    # 从 ppt/media/ 中找出所有媒体文件，按文件名索引
    zip_media_map = {}  # filename (e.g. "media1.mp4") → {"path": "ppt/media/media1.mp4", "bytes": ...}
    try:
        with zipfile.ZipFile(filepath, 'r') as z:
            for name in z.namelist():
                if name.startswith("ppt/media/"):
                    fname = Path(name).name
                    if fname.lower().endswith(tuple(MEDIA_EXTS)):
                        zip_media_map[fname] = {
                            "zip_path": name,
                            "bytes": z.read(name),
                            "ext": Path(name).suffix[1:].lower(),
                        }
    except Exception:
        pass

    # ── 尝试从 XML 构建关系映射（rId → media 文件名）──
    # 用于将媒体匹配到幻灯片
    slide_rels_map = {}  # slide_num → {rId: media_filename}
    try:
        with zipfile.ZipFile(filepath, 'r') as z:
            # 读取演示文稿的 rels
            pres_rels_path = "ppt/_rels/presentation.xml.rels"
            if pres_rels_path in z.namelist():
                pres_rels = ET.parse(z.open(pres_rels_path)).getroot()
                ns = {'r': 'http://schemas.openxmlformats.org/package/2006/relationships'}
                # 为每张幻灯片建立 rels 映射
                for slide_num in range(1, len(prs.slides) + 1):
                    slide_rels_path = f"ppt/slides/_rels/slide{slide_num}.xml.rels"
                    if slide_rels_path in z.namelist():
                        slide_media = {}
                        rels_tree = ET.parse(z.open(slide_rels_path)).getroot()
                        for rel in rels_tree:
                            rid = rel.get('Id')
                            target = rel.get('Target', '')
                            # 媒体文件引用
                            if target.startswith('../media/'):
                                media_fname = Path(target).name
                                slide_media[rid] = media_fname
                        slide_rels_map[slide_num] = slide_media
    except Exception:
        pass

    # ── 逐幻灯片处理 ──
    for slide_num, slide in enumerate(prs.slides, 1):
        slide_texts = []
        slide_images = []
        slide_media_refs = []  # 非图片媒体（视频/音频）
        matched_media = set()  # 记录已匹配的 ZIP 媒体文件名

        for shape in slide.shapes:
            # 提取文字
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    t = para.text.strip()
                    if t:
                        slide_texts.append(t)

            # 提取嵌入图片
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                try:
                    image = shape.image
                    img_bytes = image.blob
                    content_type = image.content_type
                    ext_map = {
                        "image/png": "png", "image/jpeg": "jpg",
                        "image/gif": "gif", "image/bmp": "bmp",
                        "image/tiff": "tiff",
                    }
                    ext = ext_map.get(content_type, "png")
                    img_counter += 1
                    fname = f"{safe_base}-s{slide_num:03d}-img{img_counter:02d}.{ext}"
                    with open(img_dir / fname, "wb") as f:
                        f.write(img_bytes)
                    slide_images.append({
                        "type": "image",
                        "filename": fname,
                        "ext": ext,
                        "source_ref": f"slide {slide_num} - Picture shape",
                    })
                except Exception:
                    pass

            # 递归处理组合图形
            if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                _extract_group_images(shape, slide_num, img_dir, safe_base,
                                      img_counter, slide_images, matched_media)

        # ── 从 ZIP 中匹配该幻灯片的视频/音频 ──
        slide_rels = slide_rels_map.get(slide_num, {})
        for rid, media_fname in slide_rels.items():
            if media_fname in zip_media_map and media_fname not in matched_media:
                media_info = zip_media_map[media_fname]
                mtype = _get_media_type(Path(media_fname).suffix)
                if mtype in ("video", "audio", "other"):
                    # 保存到对应目录
                    target_dir = {"video": video_dir, "audio": audio_dir, "other": other_dir}[mtype]
                    fname = f"{safe_base}-s{slide_num:03d}-{media_fname}"
                    with open(target_dir / fname, "wb") as f:
                        f.write(media_info["bytes"])
                    slide_media_refs.append({
                        "type": mtype,
                        "filename": fname,
                        "ext": media_info["ext"],
                        "source_ref": f"slide {slide_num} - Media shape ({rid})",
                    })
                    matched_media.add(media_fname)

        # 构建该幻灯片的内容文本
        slide_content = "\n".join(slide_texts)

        # 合并所有资源
        chunk_resources = slide_images + slide_media_refs

        # 对每个资源补充 context_text（该幻灯片的全部文本）
        for r in chunk_resources:
            r["context_text"] = slide_content

        chunk = {
            "id": f"slide-{slide_num:03d}",
            "heading": slide_texts[0] if slide_texts else "",
            "source_position": f"slide {slide_num}",
            "type": "slide",
            "text": slide_content,
            "resources": chunk_resources,
        }
        result["chunks"].append(chunk)
        result["full_text"] += f"\n\n--- 幻灯片{slide_num} ---\n\n{slide_content}"

    # ── 处理未匹配到具体幻灯片的 ZIP 媒体文件 ──
    all_matched = set()
    for slide_rels in slide_rels_map.values():
        all_matched.update(slide_rels.values())
    unmatched_media = [f for f in zip_media_map if f not in all_matched]
    if unmatched_media:
        orphan_resources = []
        for media_fname in unmatched_media:
            media_info = zip_media_map[media_fname]
            mtype = _get_media_type(Path(media_fname).suffix)
            if mtype in ("video", "audio", "other"):
                target_dir = {"video": video_dir, "audio": audio_dir, "other": other_dir}[mtype]
                fname = f"{safe_base}-{media_fname}"
                with open(target_dir / fname, "wb") as f:
                    f.write(media_info["bytes"])
                orphan_resources.append({
                    "type": mtype,
                    "filename": fname,
                    "ext": media_info["ext"],
                    "source_ref": f"unmatched media in ppt/media/",
                    "context_text": result["full_text"][:1000],
                })
        if orphan_resources:
            result["chunks"].append({
                "id": "media-orphan",
                "heading": "(未定位媒体资源)",
                "source_position": "ppt/media/ (unmatched)",
                "type": "media_bundle",
                "text": "",
                "resources": orphan_resources,
            })

    # ── 全页截图（后续由 ahkb.py 调用 render_pptx_full_slides 生成）──
    result["metadata"]["full_slides_generated"] = 0

    # ── 构建扁平资源列表（便利 .ctx 生成）──
    flat = []
    for chunk in result["chunks"]:
        for r in chunk["resources"]:
            r_copy = dict(r)
            r_copy["belongs_to_chunk"] = chunk["id"]
            r_copy["chunk_heading"] = chunk["heading"]
            r_copy["chunk_text"] = chunk["text"]
            flat.append(r_copy)
    result["resources_flat"] = flat

    return result


def _extract_group_images(shape, slide_num, img_dir, safe_base,
                          img_counter, slide_images, matched_media):
    """递归提取组合图形中的图片"""
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    for child in shape.shapes:
        if child.shape_type == MSO_SHAPE_TYPE.PICTURE:
            try:
                image = child.image
                ext_map = {"image/png": "png", "image/jpeg": "jpg", "image/gif": "gif"}
                ext = ext_map.get(image.content_type, "png")
                img_counter += 1
                fname = f"{safe_base}-s{slide_num:03d}-img{img_counter:02d}.{ext}"
                with open(img_dir / fname, "wb") as f:
                    f.write(image.blob)
                slide_images.append({
                    "type": "image",
                    "filename": fname,
                    "ext": ext,
                    "source_ref": f"slide {slide_num} - Group Picture",
                })
            except Exception:
                pass
        elif child.shape_type == MSO_SHAPE_TYPE.GROUP:
            _extract_group_images(child, slide_num, img_dir, safe_base,
                                  img_counter, slide_images, matched_media)


def render_pptx_full_slides(pptx_path, workspace, safe_base):
    """生成 PPT 全页截图，优先 PowerPoint → WPS → LibreOffice"""
    img_dir = Path(workspace) / "图片及其他资源" / "images"
    img_dir.mkdir(parents=True, exist_ok=True)
    full_slides = []

    def export_via_com(com_name):
        try:
            import win32com.client, pythoncom, time
            pythoncom.CoInitialize()
            app = win32com.client.Dispatch(com_name)
            try: app.Visible = False
            except: app.Visible = True
            try: app.WindowState = 2
            except: pass
            pres = app.Presentations.Open(str(Path(pptx_path).resolve()), WithWindow=False)
            time.sleep(0.5)
            result = []
            for i in range(1, pres.Slides.Count + 1):
                fname = f"slide-{safe_base}-{i:03d}.png"
                pres.Slides(i).Export(str((img_dir / fname).resolve()), "PNG", 1920, 1080)
                result.append({"number": i, "filename": fname})
            pres.Close(); app.Quit(); pythoncom.CoUninitialize()
            return result
        except:
            try: app.Quit()
            except: pass
            try: pythoncom.CoUninitialize()
            except: pass
            return None

    s = export_via_com("PowerPoint.Application")
    if s: return s
    for cn in ["Kwpp.Application", "WPP.Application", "WPS.Application"]:
        s = export_via_com(cn)
        if s: return s

    import shutil
    if shutil.which("soffice"):
        try:
            import subprocess, tempfile, pymupdf
            with tempfile.TemporaryDirectory() as td:
                r = subprocess.run(["soffice", "--headless", "--convert-to", "pdf",
                    "--outdir", td, str(Path(pptx_path).resolve())],
                    capture_output=True, text=True, timeout=60)
                if r.returncode == 0:
                    for pf in Path(td).glob("*.pdf"):
                        doc = pymupdf.open(str(pf))
                        for i in range(len(doc)):
                            pix = doc[i].get_pixmap(dpi=150)
                            fn = f"slide-{safe_base}-{i+1:03d}.png"
                            pix.save(str(img_dir / fn))
                            full_slides.append({"number": i+1, "filename": fn})
                        doc.close()
        except Exception as e:
            import sys as _s
            _s.stderr.write("[ahkb] LibreOffice error: " + str(e) + chr(10))

    if not full_slides:
        import sys as _s
        _s.stderr.write("[ahkb] PPT全页截图不可用（需安装PowerPoint/WPS/LibreOffice）" + chr(10))
    return full_slides


    def export_via_com(com_name):
        try:
            import win32com.client, pythoncom, time
            pythoncom.CoInitialize()
            app = win32com.client.Dispatch(com_name)
            try: app.Visible = False
            except: app.Visible = True
            try: app.WindowState = 2
            except: pass
            pres = app.Presentations.Open(str(Path(pptx_path).resolve()), WithWindow=False)
            time.sleep(0.5)
            result = []
            for i in range(1, pres.Slides.Count + 1):
                fname = f"slide-{safe_base}-{i:03d}.png"
                pres.Slides(i).Export(str((img_dir / fname).resolve()), "PNG", 1920, 1080)
                result.append({"number": i, "filename": fname})
            pres.Close(); app.Quit(); pythoncom.CoUninitialize()
            return result
        except:
            try: app.Quit()
            except: pass
            try: pythoncom.CoUninitialize()
            except: pass
            return None

    s = export_via_com("PowerPoint.Application")
    if s: return s
    for cn in ["Kwpp.Application", "WPP.Application", "WPS.Application"]:
        s = export_via_com(cn)
        if s: return s

    import shutil
    if shutil.which("soffice"):
        try:
            import subprocess, tempfile, pymupdf
            with tempfile.TemporaryDirectory() as td:
                r = subprocess.run(["soffice", "--headless", "--convert-to", "pdf",
                    "--outdir", td, str(Path(pptx_path).resolve())],
                    capture_output=True, text=True, timeout=60)
                if r.returncode == 0:
                    for pf in Path(td).glob("*.pdf"):
                        doc = pymupdf.open(str(pf))
                        for i in range(len(doc)):
                            pix = doc[i].get_pixmap(dpi=150)
                            fn = f"slide-{safe_base}-{i+1:03d}.png"
                            pix.save(str(img_dir / fn))
                            full_slides.append({"number": i+1, "filename": fn})
                        doc.close()
        except Exception as e:
            import sys as _s
        _s.stderr.write("[ahkb] LibreOffice error: " + str(e) + chr(10))

    if not full_slides:
        import sys as _s
        _s.stderr.write("[ahkb] PPT full slide unavailable (need PowerPoint/WPS/LibreOffice)\n")

    return full_slides

