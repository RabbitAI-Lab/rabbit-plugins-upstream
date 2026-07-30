#!/usr/bin/env python3
"""
ahkb.py — AHKB Knowledge Base Constructor 主工具 (v0.1.0)

自动以当前工作目录为知识库 Vault 根目录（即 VS Code 打开的工作空间）。
也可通过 --workspace <path> 指定其他目录。

Usage:
    python ahkb.py <command> [args]             # 自动使用当前目录
    python ahkb.py <command> [args] --workspace <path>  # 指定目录

Commands:
    scan                    Scan for new unprocessed documents
    extract <file>          Extract text, images and media from a document
    mark-processed          Mark a file as processed
    find-root               Find root knowledge graph .md file
    init                    Initialize workspace folder structure
    build-graph             Rebuild root knowledge graph .md
    build-kg-html           Build D3.js interactive knowledge graph HTML
    stats                   Show knowledge base statistics
    check-env               Check Python environment and dependencies
    set-graph-colors        Auto-configure Obsidian graph.json color groups
    audit                   Quality audit for knowledge base
    cross-link              Auto link: 知识元↔知识元 ([[链接]]) + 知识元↔资源 (.ctx)
    regenerate-ctx          Regenerate .ctx metadata files
    setup-permissions       Set Claude Code auto-permissions
    auto-install            Auto install missing Python dependencies
    set-kb-name <name>      Set or update the knowledge base name in settings
    clear-units-for-file    Clear knowledge units for a file (keep chunks/resources)
    list-processed          List processed documents with unit/chunk counts
    check-maintain-status   Check crosslink/maintain process status (for LLM polling)
"""
import sys, os, json, hashlib, datetime, glob, io, re, platform, subprocess, webbrowser, shutil
from pathlib import Path
from PIL import Image
from urllib.parse import urlparse
from ahkb_chunks import save_chunks, purge_all_chunks, get_chunk_stats
from ahkb_trash import _trash_file, _trash_dir

VERSION = "0.1.0"

# ─── 弹窗模式标记 ───
# 当脚本以 --popup-child 运行时，_POPUP_ACTIVE = True
# 用于避免在弹窗子进程中再次弹出新窗口
_POPUP_ACTIVE = "--popup-child" in sys.argv
_IN_GUI = False  # GUI弹窗递归保护锁
_DETACHED = "--detached" in sys.argv  # 解绑模式：启动独立子进程后立即返回

# Windows 控制台 UTF-8 输出支持
if sys.stdout.encoding != 'utf-8':
    sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8', errors='replace')

# ─── 子命令注册 ───
COMMANDS = {}

def command(name):
    """Decorator to register sub-commands."""
    def wrapper(fn):
        COMMANDS[name] = fn
        return fn
    return wrapper


# ─── 工具函数 ───

def _is_pid_alive(pid):
    """跨平台检查 PID 是否存活。"""
    try:
        if platform.system() == "Windows":
            import ctypes
            kernel32 = ctypes.windll.kernel32
            handle = kernel32.OpenProcess(0x0400, False, pid)  # PROCESS_QUERY_INFORMATION
            if handle:
                kernel32.CloseHandle(handle)
                return True
        else:
            os.kill(pid, 0)
            return True
    except (OSError, Exception):
        pass
    return False

def get_workspace(args):
    """Get workspace path from args or auto-detect.

    ⚠️ 拒绝将 skill 目录本身作为工作空间。
    """
    ws = None
    for i, a in enumerate(args):
        if a == "--workspace" and i + 1 < len(args):
            ws = Path(args[i + 1])
            break

    if ws:
        # --workspace 显式指定
        if not ws.exists():
            print(f"\n🔴 错误：工作空间不存在: {ws}\n")
            sys.exit(1)
        ws = ws.resolve()
    else:
        # 自动检测：从 cwd 向上找含 知识元/ 的目录
        cwd = Path.cwd().resolve()
        for parent in [cwd] + list(cwd.parents):
            if (parent / "知识元").exists():
                ws = parent
                break
        if not ws:
            # 未找到知识元/ 时回退到当前目录（适用于全新的、未初始化的知识库）
            ws = cwd

    # 🔴 禁止将 skill 目录或其任何父/子目录作为工作空间
    skill_dir = Path(__file__).resolve().parent.parent
    ws_resolved = ws.resolve()
    skill_resolved = skill_dir.resolve()
    # 收集所有需要排除的目录：skill目录本身 + skills/集群 + .claude/根
    _banned = {skill_resolved}
    _p = skill_resolved.parent
    while _p.name in ("skills", ".claude") or _p.name.startswith("."):
        _banned.add(_p)
        _p = _p.parent
    for _bp in _banned:
        if ws_resolved == _bp or _bp in ws_resolved.parents:
            print(f"\n🔴 错误：不允许将 skill 目录作为工作空间！")
            print(f"  检测到工作空间路径: {ws}")
            print(f"  Skill 目录路径: {skill_dir}")
            print(f"")
            print(f"  请将工作空间设置为您的知识库 Vault 目录（如 D:\\My Documents\\AHKB-CPS），")
            print(f"  而不是 skill 安装目录的任何父/子目录。")
            print(f"  可使用 --workspace <vault-path> 指定正确的知识库根目录。\n")
            sys.exit(1)

    # ✅ 明确告知用户当前工作空间路径
    print(f"📂 当前知识库工作空间: {ws_resolved}")

    return ws


def load_json(path):
    """Load JSON file, return default if not found."""
    if path.exists():
        try:
            with open(path, "r", encoding="utf-8") as f:
                return json.load(f)
        except (json.JSONDecodeError, OSError):
            return {}
    return {}


def save_json(path, data):
    """Save data to JSON file."""
    path.parent.mkdir(parents=True, exist_ok=True)
    with open(path, "w", encoding="utf-8") as f:
        json.dump(data, f, ensure_ascii=False, indent=2)


def detect_workspace_structure(workspace):
    """Check if the workspace has AHKB folder structure."""
    folders = ["原始文件", "知识元", "图片及其他资源", "临时工作文件"]
    existing = {}
    for f in folders:
        p = workspace / f
        existing[f] = p.is_dir()
    return existing


SUPPORTED_EXTS = {'.pptx', '.ppt', '.docx', '.doc', '.xlsx', '.xls',
                  '.pdf', '.md', '.html', '.htm', '.txt', '.csv'}

IGNORE_PREFIXES = {'.', '_'}
IGNORE_FILES = {'工作日记.md', '工作日记.html'}


def _is_root_node(filepath):
    """快速检查一个 .md 文件是否为根节点。"""
    try:
        with open(filepath, 'r', encoding='utf-8', errors='ignore') as f:
            head = f.read(500)
        return 'root_node: true' in head
    except Exception:
        return False


def get_manifest_path(workspace):
    return workspace / "原始文件" / "_processed_docs.json"


def _quick_count_text(filepath):
    """快速估算文件的有效文字量（不计空格、不计 HTML 标签、不计格式标记）。

    返回 int：有效中英文字符数，估算失败返回 0。
    """
    import re as _re

    ext = Path(filepath).suffix.lower()
    try:
        if ext in ('.md', '.txt', '.csv'):
            raw = Path(filepath).read_text(encoding='utf-8', errors='ignore')
            # 去掉 frontmatter（---...---）
            raw = _re.sub(r'^---.*?---\s*', '', raw, count=1, flags=_re.DOTALL)
            # 去掉 markdown 链接/图片语法、代码块
            raw = _re.sub(r'!?\[([^\]]*)\]\([^)]*\)', r'\1', raw)
            raw = _re.sub(r'```.*?```', '', raw, flags=_re.DOTALL)
            return len(_re.sub(r'\s+', '', raw))

        if ext in ('.html', '.htm'):
            raw = Path(filepath).read_text(encoding='utf-8', errors='ignore')
            clean = _re.sub(r'<[^>]+>', '', raw)
            clean = _re.sub(r'&[a-zA-Z#]+;', '', clean)
            return len(_re.sub(r'\s+', '', clean))

        if ext in ('.docx', '.doc'):
            from docx import Document
            doc = Document(str(filepath))
            text = ' '.join(p.text for p in doc.paragraphs)
            return len(_re.sub(r'\s+', '', text))

        if ext in ('.pptx', '.ppt'):
            from pptx import Presentation
            prs = Presentation(str(filepath))
            total = 0
            for slide in prs.slides:
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        for para in shape.text_frame.paragraphs:
                            total += len(_re.sub(r'\s+', '', para.text))
            return total

        if ext == '.pdf':
            import pymupdf
            doc = pymupdf.open(str(filepath))
            total = 0
            for page in doc:
                text = page.get_text()
                total += len(_re.sub(r'\s+', '', text))
            doc.close()
            return total

        if ext in ('.xlsx', '.xls'):
            from openpyxl import load_workbook
            wb = load_workbook(str(filepath), read_only=True, data_only=True)
            total = 0
            for sheet in wb.worksheets:
                for row in sheet.iter_rows():
                    for cell in row:
                        if cell.value and isinstance(cell.value, str):
                            total += len(_re.sub(r'\s+', '', cell.value))
            wb.close()
            return total
    except Exception:
        pass
    return 0


# ─── 权重设置文件 ───

DEFAULT_WEIGHTS = {
    "cContext": 0.5, "cTags": 0.5, "cStars": 0.5,
    "cEdited": 0.5, "cGranularity": 0.5, "cTextAmount": 0.5,
    "cLinksNum": 0.5, "cLinksDensity": 0.5,
}


def ensure_settings_file(workspace):
    """检查权重设置文件是否存在，如不存在则用默认值自动创建。"""
    p = workspace / "系统设置" / "project_settings.json"
    if not p.exists():
        p.parent.mkdir(parents=True, exist_ok=True)
        save_json(p, {"weights": dict(DEFAULT_WEIGHTS)})
        return {"created": True, "path": str(p)}
    return {"created": False, "path": str(p)}


def get_kb_name(workspace):
    """获取知识库名称：优先 project_settings.json，其次 _processed_docs.json，然后从根节点文件名推断。
    如果都找不到，输出 JSON 提示 AI 询问用户。"""
    settings_path = workspace / "系统设置" / "project_settings.json"
    if settings_path.exists():
        try:
            data = load_json(settings_path)
            name = data.get("kb_name", "").strip()
            if name:
                return name
        except Exception:
            pass
    manifest = load_json(get_manifest_path(workspace))
    name = manifest.get("knowledge_base", "").strip()
    if name:
        return name
    # 从根节点文件名推断（确保 HTML 文件名与根节点一致）
    for f in workspace.iterdir():
        if f.suffix.lower() == ".md":
            try:
                if "root_node: true" in f.read_text(500, errors="ignore"):
                    stem = f.stem
                    # 去掉尾部可能存在的"(根)"标记，返回纯净的知识库名称
                    if stem.endswith("(根)"):
                        stem = stem[:-3]
                    return stem
            except Exception:
                pass
    # 找不到名称 → 输出提示，让 AI 询问用户
    print(json.dumps({"warning": "kb_name_missing",
                       "message": "知识库名称未设置。请用户输入知识库名称，然后执行 ahkb.py set-kb-name <名称>"}),
          file=sys.stderr)
    return "知识库"


def set_kb_name(workspace, name):
    """将知识库名称同时保存到 project_settings.json 和 _processed_docs.json。"""
    name = name.strip()
    if not name:
        return False
    # 保存到 settings
    settings_path = workspace / "系统设置" / "project_settings.json"
    settings = load_json(settings_path)
    settings["kb_name"] = name
    save_json(settings_path, settings)
    # 保存到 manifest
    manifest_path = get_manifest_path(workspace)
    manifest = load_json(manifest_path)
    manifest["knowledge_base"] = name
    save_json(manifest_path, manifest)
    return True


# ─── 依赖检查状态 ───

def get_deps_status_path(workspace):
    """获取依赖检查状态文件路径。"""
    return workspace / "系统设置" / ".deps_status.json"


def load_deps_status(workspace):
    """读取依赖检查状态，不存在则返回 None。"""
    p = get_deps_status_path(workspace)
    if p.exists():
        try:
            return json.loads(p.read_text(encoding="utf-8"))
        except (json.JSONDecodeError, OSError):
            return None
    return None


def save_deps_status(workspace, data):
    """保存依赖检查状态。"""
    p = get_deps_status_path(workspace)
    p.parent.mkdir(parents=True, exist_ok=True)
    data["checked_at"] = datetime.datetime.now().isoformat()
    save_json(p, data)


# ─── 弹窗辅助函数 ───

def _get_result_file_path(args, workspace):
    """从 args 中提取 --result-file 路径，没有则返回 None。"""
    for i, a in enumerate(args):
        if a == "--result-file" and i + 1 < len(args):
            return Path(args[i + 1])
    return None


def _should_show_popup():
    """判断是否应该自动打开独立终端窗口（弹窗模式）。

    同时满足以下条件时弹窗：
    1. 不在弹窗子进程中（--popup-child 未设置）
    2. 无 --no-popup 标志
    3. stdout 不是 TTY（即被 Claude Code Bash 工具调用，非直接在终端运行）
    """
    return (not _POPUP_ACTIVE and not _IN_GUI
            and "--no-popup" not in sys.argv
            and not sys.stdout.isatty())


def _run_gui_window(cmd_name, cmd_func_name, args):
    import tkinter as tk
    from tkinter import scrolledtext, messagebox
    import threading
    import io as _io
    import sys as _sys
    global _GUI_RESULT
    _GUI_RESULT = {"ok": True}

    root = tk.Tk()
    root.title(f"AHKB - {cmd_name}")
    root.geometry("860x640")
    root.configure(bg='#1e1e1e')
    root.lift()
    root.attributes('-topmost', True)
    root.after(100, lambda: root.attributes('-topmost', False))

    # 关闭确认状态
    _processing_done = [False]

    def _on_closing():
        """关闭窗口时的确认提示"""
        if not _processing_done[0]:
            ret = messagebox.askokcancel(
                "确认关闭",
                "程序正在执行关联，强制关闭将导致知识库不完整、数据不一致。\n\n确定要关闭吗？",
                default="cancel",
                icon="warning"
            )
            if not ret:
                return
            os._exit(1)
        else:
            root.destroy()
    root.protocol("WM_DELETE_WINDOW", _on_closing)

    text_area = scrolledtext.ScrolledText(
        root, bg='#1e1e1e', fg='#d4d4d4', insertbackground='white',
        font=('Consolas', 10), wrap=tk.WORD, relief=tk.FLAT, borderwidth=0
    )
    text_area.pack(fill=tk.BOTH, expand=True, padx=5, pady=5)

    text_area.tag_configure("green", foreground="#4ec94e")
    text_area.tag_configure("blue", foreground="#569cd6")
    text_area.tag_configure("yellow", foreground="#dcdcaa")
    text_area.tag_configure("cyan", foreground="#26c6da")
    text_area.tag_configure("red", foreground="#f44336")
    text_area.tag_configure("bold", font=('Consolas', 10, 'bold'))

    # 颜色标签在运行时由 ansi_to_tag 动态应用

    def write_msg(text, tag=None):
        text_area.insert(tk.END, text, tag)
        text_area.see(tk.END)

    def ansi_to_tag(text):
        """解析 ANSI 转义码，返回 (clean_text, tag) 分段列表。"""
        segs = []
        i = 0
        cur = None
        buf = []
        while i < len(text):
            if text[i] == chr(27):
                if buf:
                    segs.append(("".join(buf), cur))
                    buf = []
                i += 1
                if i < len(text) and text[i] == "[":
                    i += 1
                    code = ""
                    while i < len(text) and (text[i].isdigit() or text[i] == ";"):
                        code += text[i]
                        i += 1
                    if i < len(text) and text[i] == "m":
                        i += 1
                    if code == "32": cur = "green"
                    elif code == "94": cur = "blue"
                    elif code == "93": cur = "yellow"
                    elif code == "96": cur = "cyan"
                    elif code == "91": cur = "red"
                    elif code == "0" or code == "": cur = None
                    else: cur = None
                continue
            buf.append(text[i])
            i += 1
        if buf:
            segs.append(("".join(buf), cur))
        return segs
    class StderrRedirect:
        def __init__(self, widget):
            self.widget = widget
            self.orig = _sys.stderr

        def write(self, msg):
            if not msg:
                return
            segs = ansi_to_tag(msg)
            # 一次性将所有段写入，避免逐段排队
            def do_write():
                for clean, tag in segs:
                    if clean:
                        text_area.insert(tk.END, clean, tag)
                    elif not tag and segs:
                        pass  # 空段跳过
                text_area.see(tk.END)
            root.after(0, do_write)

        def flush(self):
            pass

    redirect = StderrRedirect(text_area)

    def run_cmd_thread():
        nonlocal root
        nonlocal _processing_done
        try:
            global _IN_GUI
            global _GUI_RESULT
            _IN_GUI = True
            # 授权 cross_link 引擎：仅 GUI 线程可以解锁
            import ahkb_crosslink as _cl
            _cl._GUI_AUTHORIZED = True
            _sys.stderr = redirect
            _result_cmd = COMMANDS[cmd_func_name](args)
            if _result_cmd is not None:
                _GUI_RESULT = _result_cmd
            _sys.stderr = redirect.orig
            from ahkb_crosslink import _PAUSE_FLAG
            _PAUSE_FLAG[0] = False
            root.after(0, lambda: btn_pause.config(state=tk.DISABLED))
            root.after(0, write_msg, "知识库重构完毕\n", "green")
            _processing_done[0] = True
            if "--auto-close" in args:
                root.after(0, write_msg, "3 秒后自动关闭窗口...\n", "green")
                root.after(3000, root.destroy)
            else:
                root.after(0, write_msg, "\n本窗口可安全关闭。请点击「关闭窗口」按钮退出。\n", "green")
        except Exception as e:
            _sys.stderr = redirect.orig
            from ahkb_crosslink import _PAUSE_FLAG
            _PAUSE_FLAG[0] = False
            root.after(0, lambda: btn_pause.config(state=tk.DISABLED))
            _processing_done[0] = True
            root.after(0, write_msg, f"错误: {e}\n", "red")

    btn_frame = tk.Frame(root, bg='#1e1e1e')
    btn_frame.pack(fill=tk.X, padx=5, pady=(0, 5))

    # 暂停控制
    _PAUSE_BTN_STATE = [False]
    def _toggle_pause():
        from ahkb_crosslink import _PAUSE_FLAG
        _PAUSE_BTN_STATE[0] = not _PAUSE_BTN_STATE[0]
        if _PAUSE_BTN_STATE[0]:
            _PAUSE_FLAG[0] = True
            btn_pause.config(text="▶ 继续", bg='#4ec9b0')
        else:
            _PAUSE_FLAG[0] = False
            btn_pause.config(text="⏸ 暂停", bg='#dcdcaa')

    btn_pause = tk.Button(btn_frame, text="⏸ 暂停",
                           font=("Microsoft YaHei", 9), bg='#dcdcaa', fg='#1e1e1e',
                           activebackground='#e5e07a', relief=tk.FLAT, padx=15,
                           command=_toggle_pause)
    btn_pause.pack(side=tk.RIGHT, padx=(0, 5))

    tk.Button(btn_frame, text="关闭窗口", command=_on_closing,
              bg='#3c3c3c', fg='white', relief=tk.FLAT, padx=20).pack(side=tk.RIGHT)

    t = threading.Thread(target=run_cmd_thread, daemon=True)
    t.start()
    root.mainloop()
    return _GUI_RESULT


def _run_in_terminal(cmd_name, args, workspace):
    print(f"\n\\U0001f527 正在打开独立窗口执行：{cmd_name}")
    print(f"   请在新窗口中查看处理过程...\n")

    if platform.system() == "Windows":
        import uuid, tempfile, time as _time, os as _os
        script_path = Path(__file__).resolve()
        # 唯一信号文件：子进程执行完后创建，主进程检测到即继续
        sig_file = _os.path.join(tempfile.gettempdir(),
                                 f"ahkb_{uuid.uuid4().hex[:8]}.sig")
        # 子进程：新控制台窗口，--no-popup 防止递归弹窗
        child_args = [str(sys.executable), str(script_path),
                       cmd_name, "--no-popup", "--sig-file", sig_file] + list(args)
        subprocess.Popen(child_args, cwd=str(workspace),
                         creationflags=subprocess.CREATE_NEW_CONSOLE)
        # ⏳ 等待子进程执行完毕（信号文件出现），不等待窗口关闭
        while not _os.path.exists(sig_file):
            _time.sleep(0.3)
        try:
            _os.remove(sig_file)
        except Exception:
            pass
        return {"ok": True, "popup": True}
    else:
        import shutil, uuid, tempfile, time as _time, os as _os
        script_path = Path(__file__).resolve()
        sig_file = _os.path.join(tempfile.gettempdir(),
                                 f"ahkb_{uuid.uuid4().hex[:8]}.sig")
        child_args = [str(sys.executable), str(script_path),
                       cmd_name, "--no-popup", "--sig-file", sig_file] + list(args)
        launched = False
        def _q(s):
            s = str(s)
            return "'" + s.replace("'", "'\\''") + "'" if " " in s else s
        cmd_str = " ".join(_q(a) for a in child_args)
        if shutil.which("xterm"):
            subprocess.Popen(["xterm", "-title", f"AHKB - {cmd_name}",
                              "-e", cmd_str], cwd=str(workspace))
            launched = True
        elif shutil.which("gnome-terminal"):
            subprocess.Popen(["gnome-terminal", "--", "bash", "-c", cmd_str],
                             cwd=str(workspace))
            launched = True
        elif sys.platform == "darwin":
            ascript = f'tell app "Terminal" to do script "{_q(cmd_str)}"'
            subprocess.Popen(["osascript", "-e", ascript])
            launched = True
        if not launched:
            subprocess.Popen(child_args, cwd=str(workspace))
        while not _os.path.exists(sig_file):
            _time.sleep(0.3)
        try:
            _os.remove(sig_file)
        except Exception:
            pass
        return {"ok": True, "popup": True}

# ─── 资源目录映射 ───

RESOURCE_DIRS = {
    "image": "images",
    "video": "videos",
    "audio": "audios",
    "full_slide_capture": "images",
    "other": "others",
}


# ─── .ctx 文件生成 ───


def get_resource_dir(workspace, resource_type):
    """获取资源类型对应的目录路径。"""
    subdir = RESOURCE_DIRS.get(resource_type, "others")
    return workspace / "图片及其他资源" / subdir


def generate_ctx_file(workspace, resource_info, source_path, force=False):
    """为单个资源生成 .ctx 元数据文件。

    resource_info: dict with type, filename (or url), context_text, source_ref,
                   belongs_to_chunk, chunk_heading, chunk_text
    source_path: str — 源文档的相对路径
    force: bool — 若 True，即使文件存在也重新生成（保留已有的 belongs_to）
    """
    rtype = resource_info.get("type", "other")
    resource_dir = get_resource_dir(workspace, rtype)
    resource_dir.mkdir(parents=True, exist_ok=True)

    # 确定是本地文件还是远程URL
    has_local_file = "filename" in resource_info and resource_info["filename"]
    is_remote = "url" in resource_info and resource_info["url"].startswith("http")

    # 远程资源没有本地文件 → 不生成 .ctx（无用）
    if is_remote and not has_local_file:
        return None

    if has_local_file:
        stem = Path(resource_info["filename"]).stem
    elif is_remote:
        stem = Path(urlparse(resource_info["url"]).path).stem or "remote_resource"
    else:
        stem = "unnamed_resource"

    ctx_name = f"{stem}.ctx"
    ctx_path = resource_dir / ctx_name

    # 如果已存在：保留 belongs_to，但重新生成正文（确保上下文最新）
    existing_belongs_to = []
    existing_body = ""
    if ctx_path.exists():
        with open(ctx_path, "r", encoding="utf-8") as f:
            existing_content = f.read()
        bt_match = re.search(r'belongs_to:\s*\n((?:\s+-.*\n?)*)', existing_content)
        if bt_match:
            existing_belongs_to = [line.strip() for line in bt_match.group(1).split('\n') if line.strip()]
        # 也提取已有的 body 上下文（保留手工添加的说明文字）
        parts = existing_content.split("---", 2)
        if len(parts) >= 3:
            existing_body = parts[2]
        if not force:
            return ctx_name

    # 上下文文本
    context_text = resource_info.get("context_text", "")
    chunk_text = resource_info.get("chunk_text", "")
    full_context = context_text or chunk_text

    # 从上下文提取关键词作为初始标签（精简）
    tags = _extract_key_tags(full_context, max_tags=5)
    tags.append(rtype)

    # 自动计算 importance（根据图片尺寸）
    imp = 3
    if has_local_file and rtype in ("image", "full_slide_capture"):
        try:
            with Image.open(resource_dir / resource_info["filename"]) as _im:
                _s = int((_im.width * _im.height) ** 0.5)
            if _s < 100: imp = 1
            elif _s < 400: imp = 2
            elif _s < 1000: imp = 3
            else: imp = 4
        except:
            pass
    # 构建 YAML frontmatter（★ 严格按此字段顺序）
    fm_lines = ["---"]
    fm_lines.append("type: resource")
    fm_lines.append(f"resource_type: {rtype}")
    fm_lines.append(f"importance: {imp}")
    fm_lines.append("user_edited: false")
    now_str = datetime.datetime.now().strftime("%Y-%m-%d %H:%M:%S")
    fm_lines.append(f"last_edited_time: {now_str}")
    fm_lines.append(f'source: "{source_path}"')
    fm_lines.append(f"belongs_to_chunk: \"{resource_info.get('belongs_to_chunk', '')}\"")
    if resource_info.get("chunk_heading"):
        fm_lines.append(f'chunk_heading: "{resource_info["chunk_heading"]}"')
    if has_local_file:
        fm_lines.append(f'resource_file: "{resource_info["filename"]}"')
    if is_remote:
        fm_lines.append(f'remote_url: "{resource_info["url"]}"')
    fm_lines.append(f"tags: [{', '.join(tags)}]")
    # belongs_to 初始为空，由后续填写（★ 必须是最后一个字段）
    if existing_belongs_to:
        fm_lines.append("belongs_to:")
        for bt in existing_belongs_to:
            fm_lines.append(f"  - {bt}")
    else:
        fm_lines.append("belongs_to:")
    fm_lines.append("---")
    fm_lines.append("")

    # 引用资源
    body_lines = []
    if has_local_file:
        body_lines.append(f"![[{resource_info['filename']}]]")
    elif is_remote:
        body_lines.append(f"远程资源：{resource_info['url']}")
    body_lines.append("")

    # 上下文正文（精简，最多 300 字）
    if full_context:
        context_body = full_context[:300]
        body_lines.append("> " + context_body.replace("\n", "\n> "))
        if len(full_context) > 300:
            body_lines.append(">")
            body_lines.append("> _(上下文已截断)_")
    body_lines.append("")

    # 保留已存在的 body 中的补充说明（非引用行、非空行）
    if existing_body:
        extra = [l for l in existing_body.split("\n") if l.strip() and not l.strip().startswith("![[") and not l.strip().startswith(">") and not l.strip().startswith("远程资源")]
        if extra:
            body_lines.append("")
            body_lines.extend(extra)
            body_lines.append("")

    # 写入 .ctx 文件
    content = "\n".join(fm_lines + body_lines)
    with open(ctx_path, "w", encoding="utf-8") as f:
        f.write(content)
    return ctx_name


def _extract_key_tags(text, max_tags=5):
    """从文本中提取关键词作为标签（2-5 字短词，避免长句）。"""
    if not text or not text.strip():
        return []
    # 按标点拆分后取 2-5 字中文词或 3-8 字符英文词
    segments = re.split(r'[，。、；：！？\s,.;:!?\n\r]+', text)
    seen = set()
    result = []
    for seg in segments:
        cn = re.findall(r'[一-鿿]{2,7}', seg)
        en = re.findall(r'[A-Za-z]{3,8}', seg)
        for w in cn + en:
            if w not in seen:
                seen.add(w)
                result.append(w)
                if len(result) >= max_tags:
                    return result
    return result


# ─── 子命令: check-env ───

@command("check-env")
def cmd_check_env(args):
    """Check Python environment and required libraries (cached after first check)."""
    workspace = get_workspace(args)

    # 如果已有检查记录且全部就绪，跳过完整检查
    cached = load_deps_status(workspace)
    if cached and cached.get("all_ok"):
        print(json.dumps({
            "deps_checked": True,
            "skipped": True,
            "python": cached.get("python", sys.version),
            "libraries": cached.get("libraries", {}),
            "obsidian": cached.get("obsidian", {}),
            "checked_at": cached.get("checked_at", ""),
        }, ensure_ascii=False))
        return

    result = {"python": sys.version, "libraries": {}}
    libs = [
        ("python-pptx", "pptx"),
        ("python-docx", "docx"),
        ("openpyxl", "openpyxl"),
        ("pymupdf", "pymupdf"),
        ("jieba", "jieba"),
    ]
    for name, mod in libs:
        try:
            __import__(mod.replace("-", "_"))
            result["libraries"][name] = "OK"
        except ImportError:
            result["libraries"][name] = "MISSING"

    # ── Obsidian 检测（含便携版支持） ──
    obsidian_found = False
    obsidian_paths = []
    obsidian_methods = []

    # ── 方法1：Windows 注册表 App Paths（仅限安装版） ──
    if platform.system() == "Windows":
        try:
            import winreg
            for hkey, subkey in [
                (winreg.HKEY_LOCAL_MACHINE,
                 r"SOFTWARE\Microsoft\Windows\CurrentVersion\App Paths\Obsidian.exe"),
                (winreg.HKEY_LOCAL_MACHINE,
                 r"SOFTWARE\WOW6432Node\Microsoft\Windows\CurrentVersion\Uninstall\Obsidian"),
                (winreg.HKEY_CURRENT_USER,
                 r"SOFTWARE\Microsoft\Windows\CurrentVersion\Uninstall\Obsidian"),
            ]:
                if obsidian_found:
                    break
                try:
                    with winreg.OpenKey(hkey, subkey) as key:
                        for field in ["", "InstallLocation", "DisplayIcon"]:
                            try:
                                val, _ = winreg.QueryValueEx(key, field)
                                if val:
                                    exe = val if val.endswith(".exe") else os.path.join(val, "Obsidian.exe")
                                    obsidian_paths.append(exe)
                                    if os.path.isfile(exe):
                                        obsidian_found = True
                                        obsidian_methods.append("注册表")
                                        break
                            except OSError:
                                continue
                except OSError:
                    continue
        except ImportError:
            pass

    # ── 方法2：obsidian:// URL 协议（便携版也会注册）──
    if not obsidian_found and platform.system() == "Windows":
        try:
            import winreg
            for hkey in [winreg.HKEY_CLASSES_ROOT,
                          winreg.HKEY_CURRENT_USER,
                          winreg.HKEY_LOCAL_MACHINE]:
                try:
                    with winreg.OpenKey(hkey, r"obsidian\shell\open\command") as key:
                        val, _ = winreg.QueryValueEx(key, "")
                        if val:
                            # 格式: "D:\path\to\Obsidian.exe" "%1"
                            m = re.search(r'^"([^"]+\.exe)"', val)
                            if m and os.path.isfile(m[1]):
                                obsidian_paths.append(m[1])
                                obsidian_found = True
                                obsidian_methods.append("URL协议")
                                break
                except OSError:
                    continue
        except ImportError:
            pass

    # ── 方法3：标准路径 + 便携版常见路径快速扫描 ──
    if not obsidian_found:
        search_targets = []
        # 标准安装路径
        for env_var in ["LOCALAPPDATA", "PROGRAMFILES", "PROGRAMFILES(X86)"]:
            base = os.environ.get(env_var, "")
            if base:
                search_targets.append(os.path.join(base, "Obsidian", "Obsidian.exe"))
        # 便携版：只扫描实际存在的盘符下的常用路径
        user = os.environ.get("USERNAME", "")
        for drive_letter in "CDEFGH":
            drive = drive_letter + ":\\"
            if not os.path.exists(drive):
                continue
            for base_dir in [
                drive + "Obsidian",
                drive + "Portable",
                drive + "Apps",
                drive + "Programs",
                drive + "Tools",
                os.path.join(drive, f"Users", user, "AppData", "Local"),
            ]:
                search_targets.append(os.path.join(base_dir, "Obsidian.exe"))
        # 逐个快速检测（os.path.isfile 对不存在的文件极快，纳秒级返回）
        for p in search_targets:
            if os.path.isfile(p):
                obsidian_paths.append(p)
                obsidian_found = True
                obsidian_methods.append("路径搜索")
                break

    # ── 方法4：正在运行的进程（便携版即使没注册也能检出）──
    if not obsidian_found and platform.system() == "Windows":
        try:
            out = subprocess.check_output(
                ["tasklist", "/fi", "imagename eq Obsidian.exe", "/nh"],
                timeout=5, stderr=subprocess.DEVNULL
            ).decode("utf-8", errors="ignore")
            if "Obsidian" in out:
                obsidian_found = True
                obsidian_methods.append("运行进程")
                # 通过 PowerShell 获取完整路径
                try:
                    ps_out = subprocess.check_output(
                        ["powershell", "-NoProfile", "-Command",
                         "Get-Process Obsidian | Select-Object -ExpandProperty Path"],
                        timeout=5, stderr=subprocess.DEVNULL
                    ).decode("utf-8", errors="ignore").strip()
                    if ps_out and os.path.isfile(ps_out):
                        obsidian_paths.insert(0, ps_out)
                except Exception:
                    pass
        except Exception:
            pass

    # ── 方法5：where 命令快速兜底 ──
    if not obsidian_found and platform.system() == "Windows":
        try:
            out = subprocess.check_output(
                "where obsidian 2>nul", shell=True, timeout=5
            ).decode("utf-8", errors="ignore").strip()
            if out:
                for line in out.split("\n"):
                    p = line.strip()
                    if p and os.path.isfile(p):
                        obsidian_paths.append(p)
                        obsidian_found = True
                        obsidian_methods.append("where命令")
                        break
        except Exception:
            pass

    elif platform.system() == "Darwin":
        p = "/Applications/Obsidian.app"
        obsidian_paths.append(p)
        if os.path.isdir(p):
            obsidian_found = True
            obsidian_methods.append("标准路径")
    else:  # Linux
        p = shutil.which("obsidian")
        if p:
            obsidian_paths.append(p)
            obsidian_found = True
            obsidian_methods.append("which")
        else:
            for try_path in [
                "/usr/bin/obsidian", "/usr/local/bin/obsidian",
                "/opt/Obsidian/obsidian", os.path.expanduser("~/.local/bin/obsidian"),
            ]:
                obsidian_paths.append(try_path)
                if os.path.isfile(try_path):
                    obsidian_found = True
                    obsidian_methods.append("路径扫描")
                    break

    result["obsidian"] = {
        "installed": obsidian_found,
        "detected_by": obsidian_methods,
        "paths_checked": obsidian_paths,
    }

    # 检查权重设置文件
    result["settings_file"] = ensure_settings_file(workspace)

    # 判断是否全部就绪，写入缓存
    all_libs_ok = all(v == "OK" for v in result["libraries"].values())
    all_ok = all_libs_ok  # Obsidian 非必需，不纳入 all_ok 判断
    result["all_ok"] = all_ok

    # ── 自动安装缺失的依赖（无需用户干预）──
    if not all_libs_ok:
        import subprocess as _sp
        missing = [n for n, v in result["libraries"].items() if v == "MISSING"]
        import sys as _sys
        _sys.stderr.write(f"  检测到缺失依赖: {missing}，自动安装中...\n")
        try:
            _sp.check_call(
                [_sys.executable, "-m", "pip", "install"] + missing,
                stdout=_sp.DEVNULL, stderr=_sp.DEVNULL, timeout=120
            )
            for n in missing:
                result["libraries"][n] = "OK"
            result["all_ok"] = True
            _sys.stderr.write(f"  ✅ 安装成功: {missing}\n")
        except Exception as _e:
            _sys.stderr.write(f"  ❌ 安装失败: {_e}\n")

    save_deps_status(workspace, result)
    print(json.dumps(result, ensure_ascii=False))


# ─── 子命令: setup-permissions ───

@command("setup-permissions")
def cmd_setup_permissions(args):
    """设置 Claude Code 完全自动权限。"""
    workspace = get_workspace(args)
    claude_dir = workspace / ".claude"
    claude_dir.mkdir(parents=True, exist_ok=True)

    settings_path = claude_dir / "settings.json"
    full_perms = {
        "permissions": {
            "allow": [
                "Bash(*)", "Read(*)", "Write(*)", "Edit(*)",
                "Glob(*)", "Grep(*)", "Agent(*)", "NotebookEdit(*)",
                "Skill(*)", "AskUserQuestion(*)", "WebFetch(*)", "WebSearch(*)",
                "TodoWrite(*)", "Workflow(*)", "CronCreate(*)", "CronDelete(*)",
                "CronList(*)", "EnterPlanMode(*)", "ExitPlanMode(*)",
                "EnterWorktree(*)", "ExitWorktree(*)", "ScheduleWakeup(*)",
                "TaskOutput(*)", "TaskStop(*)",
            ]
        }
    }
    save_json(settings_path, full_perms)
    print(json.dumps({"ok": True, "settings": str(settings_path)}, ensure_ascii=False))


# ─── 子命令: auto-install ───

@command("auto-install")
def cmd_auto_install(args):
    """自动安装缺失的 Python 依赖库。安装成功后更新依赖检查标记。"""
    workspace = get_workspace(args)
    libs = [
        ("python-pptx", "pptx"),
        ("python-docx", "docx"),
        ("openpyxl", "openpyxl"),
        ("pymupdf", "pymupdf"),
        ("jieba", "jieba"),
        ("Pillow", "PIL"),
    ]
    missing = []
    for name, mod in libs:
        try:
            __import__(mod.replace("-", "_"))
        except ImportError:
            missing.append(name)

    if not missing:
        # 全部已就绪 → 更新标记文件
        save_deps_status(workspace, {
            "python": sys.version,
            "libraries": {n: "OK" for n, _ in libs},
            "obsidian": {},
            "all_ok": True,
        })
        print(json.dumps({"ok": True, "installed": [], "note": "所有依赖已就绪"}))
        return

    import subprocess
    pip_list = [n for n in missing if n != "Pillow"] + (["Pillow"] if "Pillow" in missing else [])
    print(json.dumps({"ok": False, "installing": pip_list}), file=sys.stderr)
    try:
        subprocess.check_call(
            [sys.executable, "-m", "pip", "install"] + pip_list,
            stdout=subprocess.DEVNULL, stderr=subprocess.DEVNULL, timeout=120
        )
        result = {"ok": True, "installed": pip_list}
        # 安装成功后更新标记
        libs_status = {n: "OK" for n, _ in libs}
        save_deps_status(workspace, {
            "python": sys.version,
            "libraries": libs_status,
            "obsidian": {},
            "all_ok": True,
        })
    except Exception as e:
        result = {"ok": False, "error": str(e), "installed": []}

    print(json.dumps(result, ensure_ascii=False))


# ─── 子命令: init ───

@command("init")
def cmd_init(args):
    """Initialize workspace: check env, create folders and root node.
    支持 --name <知识库名称> 参数，将名称写入 project_settings.json。
    """
    workspace = get_workspace(args)

    # ── 首次初始化时检查依赖并缓存结果 ──
    _cached = load_deps_status(workspace)
    if not _cached or not _cached.get("all_ok"):
        _libs = [
            ("python-pptx", "pptx"), ("python-docx", "docx"),
            ("openpyxl", "openpyxl"), ("pymupdf", "pymupdf"),
            ("jieba", "jieba"), ("Pillow", "PIL"),
        ]
        _miss = []
        for _n, _m in _libs:
            try:
                __import__(_m.replace("-", "_"))
            except ImportError:
                _miss.append(_n)
        if _miss:
            import subprocess as _sp
            import sys as _sys2
            _sys2.stderr.write(f"  ⚠ 安装缺失依赖: {_miss}...\n")
            _sp.check_call([sys.executable, "-m", "pip", "install"] + _miss, timeout=120)
        save_deps_status(workspace, {"all_ok": True, "libraries": {n: "OK" for n, _ in _libs}})

    # 提取 --name 参数
    kb_name = ""
    for i, a in enumerate(args):
        if a == "--name" and i + 1 < len(args):
            kb_name = args[i + 1].strip()
            break

    folders = [
        "原始文件", "知识元",
        "图片及其他资源/images",
        "图片及其他资源/videos",
        "图片及其他资源/audios",
        "图片及其他资源/others",
        "chunks",
        "临时工作文件",
        "产品成果",
    ]
    created = []
    for folder in folders:
        p = workspace / folder
        if not p.exists():
            p.mkdir(parents=True, exist_ok=True)
            created.append(folder)

    manifest = get_manifest_path(workspace)
    if not manifest.exists():
        save_json(manifest, {"version": "1.1", "knowledge_base": "", "processed": []})

    # 自动创建权重设置文件（如不存在）
    settings_result = ensure_settings_file(workspace)

    # 如有 --name，保存到 settings 和 manifest
    if kb_name:
        set_kb_name(workspace, kb_name)

    # 初始化概念索引
    save_concept_index(workspace)

    print(json.dumps({"ok": True, "created_folders": created, "workspace": str(workspace),
                       "kb_name": kb_name or "知识库", "settings": settings_result},
                     ensure_ascii=False))


# ─── 子命令: scan ───

@command("scan")
def cmd_scan(args):
    """Scan workspace for new unprocessed documents (recursive, with category detection)."""
    workspace = get_workspace(args)
    manifest = load_json(get_manifest_path(workspace))
    processed_files = {p["source_path"] for p in manifest.get("processed", [])}
    processed_names = {p["file"] for p in manifest.get("processed", [])}

    # AHKB 系统目录（递归扫描时跳过）— 注意：原始文件/ 不在其中，
    # 用户可能将未处理的文档直接放入原始文件/，这些应被当作新文件扫描。
    # 是否已处理仅以 _processed_docs.json 为准。
    SYSTEM_DIRS = {"知识元", "图片及其他资源", "临时工作文件", "系统设置", "产品成果", "chunks", "回收站"}
    # 备份目录名称（任何层级出现即跳过）
    BACKUP_DIR_NAMES = {"备份", "backup", "Backup", "BACKUP", "_backup", "_备份", ".backup", "archive", "存档"}

    new_files = []
    processed_list = []

    # 递归扫描整个工作空间（跳过系统/隐藏/备份目录）
    for f in workspace.rglob("*"):
        if not f.is_file():
            continue
        if f.suffix.lower() not in SUPPORTED_EXTS:
            continue
        if f.stat().st_size == 0:
            continue
        if f.name in IGNORE_FILES:
            continue
        if any(f.name.startswith(p) for p in IGNORE_PREFIXES):
            continue

        rel = f.relative_to(workspace)
        parts = rel.parts

        # 跳过 AHKB 系统目录
        if parts[0] in SYSTEM_DIRS:
            continue

        # 跳过隐藏目录（任何层级以 . 开头的目录）
        if any(part.startswith(".") for part in parts[:-1]):
            continue

        # 跳过备份目录（任何层级的目录名匹配备份关键词）
        if any(part in BACKUP_DIR_NAMES for part in parts[:-1]):
            continue

        # 跳过根节点 .md/.html
        if f.suffix.lower() == '.md' and _is_root_node(f):
            continue
        if f.suffix.lower() == '.html':
            # 排除知识地图 HTML（<名称>-知识地图.html）
            if f.stem.endswith('-知识地图'):
                continue
            # 排除根节点对应的 HTML（<名称>.html，对应 <名称>.md 是根节点）
            md_version = f.with_suffix('.md')
            if md_version.exists() and _is_root_node(md_version):
                continue

        rel_path = str(rel).replace("\\", "/")

        # 确定分类：子目录路径即为用户的分类
        category = str(Path(*parts[:-1])).replace("\\", "/") if len(parts) > 1 else ""

        entry = {
            "name": f.name,
            "path": rel_path,
            "category": category,
            "size": f.stat().st_size,
            "ext": f.suffix.lower()[1:],
            "text_chars": _quick_count_text(f),
        }

        if rel_path in processed_files or f.name in processed_names:
            entry["processed"] = True
            processed_list.append(entry)
        else:
            entry["processed"] = False
            new_files.append(entry)

    # ── 检查原始文件/ 中遗漏的未处理文件（合并到 new_files）──
    pending_files = []
    source_dir = workspace / "原始文件"
    # 建立 new_files 路径索引用于去重
    new_paths = {e["path"] for e in new_files}
    BACKUP_DIR_NAMES = {"备份", "backup", "Backup", "BACKUP", "_backup", "_备份", ".backup", "archive", "存档"}
    if source_dir.exists():
        for f in source_dir.rglob("*"):
            if not f.is_file() or f.suffix.lower() not in SUPPORTED_EXTS:
                continue
            if f.name in IGNORE_FILES:
                continue
            # 跳过隐藏目录和备份目录中的文件
            rel_parts = f.relative_to(workspace).parts
            if any(p.startswith(".") for p in rel_parts[1:-1]):
                continue
            if any(p in BACKUP_DIR_NAMES for p in rel_parts[1:-1]):
                continue
            rel_path = str(f.relative_to(workspace)).replace("\\", "/")
            if rel_path in processed_files or f.name in processed_names:
                continue  # 已处理完成
            if rel_path in new_paths:
                continue  # 已通过主扫描纳入 new_files，避免重复
            # 找到了！在原始文件/中但不在已处理清单中且未被主扫描捕获
            parts = f.relative_to(workspace).parts
            category = str(Path(*parts[1:-1])).replace("\\", "/") if len(parts) > 2 else ""
            parts = f.relative_to(workspace).parts
            category = str(Path(*parts[1:-1])).replace("\\", "/") if len(parts) > 2 else ""
            new_files.append({
                "name": f.name,
                "path": rel_path,
                "category": category,
                "size": f.stat().st_size,
                "ext": f.suffix.lower()[1:],
                "status": "extracted_not_processed",
                "text_chars": _quick_count_text(f),
            })

    # ── 检查 chunks/ 中已提取但未标记为已处理的文件 ──
    chunks_dir = workspace / "chunks"
    chunk_index_path = chunks_dir / "index.json"
    if chunk_index_path.exists():
        try:
            chunk_index = load_json(chunk_index_path)
            chunked_files = chunk_index.get("files", [])
            for src_path in chunked_files:
                # 已记录为已处理 → 跳过
                if src_path in processed_files:
                    continue
                # 已在 new_files 中 → 跳过（主扫描已捕获）
                if src_path in new_paths:
                    continue
                # 已在 pending_files 中 → 跳过
                if any(e["path"] == src_path for e in pending_files):
                    continue
                # 检查原始文件是否还存在
                src_file = workspace / src_path
                if not src_file.exists():
                    continue
                if src_file.suffix.lower() not in SUPPORTED_EXTS:
                    continue
                if src_file.name in IGNORE_FILES:
                    continue
                # 确认为"已提取但未标记"的文件
                pending_files.append({
                    "name": src_file.name,
                    "path": src_path,
                    "category": str(Path(src_path).parent).replace("\\", "/") if str(Path(src_path).parent) != "." else "",
                    "size": src_file.stat().st_size,
                    "ext": src_file.suffix.lower()[1:],
                    "status": "chunked_not_finalized",
                    "text_chars": _quick_count_text(src_file),
                })
        except Exception:
            pass  # chunks 索引损坏时静默跳过

    pending_files.sort(key=lambda x: x["path"])

    # ── 估算处理时间 ──
    total_min = 0.0
    for f_entry in new_files:
        sz = f_entry.get("size", 0)
        if sz < 50 * 1024:          # <50KB → 1min
            total_min += 1
        elif sz < 500 * 1024:       # 50~500KB → 3min
            total_min += 3
        elif sz < 5 * 1024 * 1024:  # 500KB~5MB → 5min
            total_min += 5
        else:                        # >5MB → 10min
            total_min += 10
    # pending_files 已有 chunks，只需 mark-processed，每文件 +0.5min
    for _ in pending_files:
        total_min += 0.5
    # Phase 2 maintain 固定 +2min
    if new_files or pending_files:
        total_min += 2
    est_optimistic = max(1, round(total_min * 0.8))
    est_pessimistic = max(2, round(total_min * 1.5))

    result = {
        "workspace": str(workspace),
        "new_files": new_files,
        "pending_files": pending_files,
        "processed_files": processed_list,
        "new_count": len(new_files),
        "pending_count": len(pending_files),
        "processed_count": len(processed_list),
        "estimated_minutes": {
            "optimistic": est_optimistic,
            "pessimistic": est_pessimistic,
        },
    }
    print(json.dumps(result, ensure_ascii=False))


# ─── 子命令: extract ───

@command("extract")
def cmd_extract(args):
    """Extract content from a document."""
    if len(args) < 1:
        print(json.dumps({"error": "Usage: ahkb.py extract <filepath> [--category <分类>] [--workspace <path>]"}))
        sys.exit(1)

    filepath = Path(args[0])
    if not filepath.is_absolute():
        filepath = get_workspace(args) / filepath
    if not filepath.exists():
        print(json.dumps({"error": f"File not found: {filepath}"}))
        sys.exit(1)

    workspace = get_workspace(args)

    # 解析 --category 参数（用户指定的分类，用于根目录文件）
    category = None
    for i, a in enumerate(args):
        if a == "--category" and i + 1 < len(args):
            category = args[i + 1]
            break

    ext = filepath.suffix.lower()

    # 确保资源目录存在
    for d in ["images", "videos", "audios", "others"]:
        (workspace / "图片及其他资源" / d).mkdir(parents=True, exist_ok=True)

    try:
        if ext == ".pdf":
            from ahkb_extract_pdf import extract_pdf
            result = extract_pdf(str(filepath), workspace)
        elif ext in (".pptx", ".ppt"):
            from ahkb_extract_pptx import extract_pptx
            result = extract_pptx(str(filepath), workspace)
            # 生成全页截图
            try:
                safe_base = "".join(c if c.isalnum() or c in '-_ ' else '_' for c in filepath.stem)
                from ahkb_extract_pptx import render_pptx_full_slides
                full_slides = render_pptx_full_slides(str(filepath), workspace, safe_base)
                if full_slides:
                    # 全页截图作为 full_slide_capture 资源注入到对应 chunk
                    for fs in full_slides:
                        sn = fs["number"]
                        fname = fs["filename"]
                        for chunk in result.get("chunks", []):
                            if chunk["id"] == f"slide-{sn:03d}":
                                # 检查是否已存在该资源
                                existing = [r for r in chunk["resources"]
                                            if r.get("type") == "full_slide_capture"
                                            and r.get("filename") == fname]
                                if not existing:
                                    chunk["resources"].append({
                                        "type": "full_slide_capture",
                                        "filename": fname,
                                        "ext": "png",
                                        "source_ref": f"slide {sn} - full slide render",
                                        "context_text": chunk["text"],
                                    })
                                break
                    result["metadata"]["full_slides_generated"] = len(full_slides)
                # 全页截图也加入 resources_flat（确保 .ctx 生成覆盖）
                for fs in full_slides:
                    fsn = fs["number"]
                    ffname = fs["filename"]
                    for chunk in result.get("chunks", []):
                        if chunk["id"] == f"slide-{fsn:03d}":
                            for r in chunk["resources"]:
                                if r.get("type") == "full_slide_capture" and r.get("filename") == ffname:
                                    r_copy = dict(r)
                                    r_copy["belongs_to_chunk"] = chunk["id"]
                                    r_copy["chunk_heading"] = chunk["heading"]
                                    r_copy["chunk_text"] = chunk["text"]
                                    result.setdefault("resources_flat", []).append(r_copy)
                            break
            except Exception as e:
                import sys as _sys; _sys.stderr.write(f"[ahkb] Full slide render error: {e}\n")
        elif ext in (".docx", ".doc"):
            from ahkb_extract_docx import extract_docx
            result = extract_docx(str(filepath), workspace)
        elif ext in (".xlsx", ".xls"):
            from ahkb_extract_xlsx import extract_xlsx
            result = extract_xlsx(str(filepath), workspace)
        elif ext in (".md", ".html", ".htm", ".txt", ".csv"):
            from ahkb_extract_md import extract_md
            result = extract_md(str(filepath), workspace)
        else:
            print(json.dumps({"error": f"Unsupported file type: {ext}"}))
            sys.exit(1)

        # 统一使用相对路径
        try:
            current_rel = filepath.relative_to(workspace).as_posix()
        except ValueError:
            current_rel = filepath.as_posix()

        # 添加文件哈希（使用原始路径）
        result["sha256"] = file_hash(str(filepath))

        # ── 注入当前知识库参数，确保 AI 使用的参数始终最新 ──
        settings_p = workspace / "系统设置" / "project_settings.json"
        if settings_p.exists():
            try:
                sd = json.loads(settings_p.read_text(encoding="utf-8"))
                result["weights"] = sd.get("weights", {})
            except Exception:
                result["weights"] = {}
        else:
            result["weights"] = {}

        # ── 注入完整的知识元生成指导（已处理，AI 可直接使用）──
        result["guidance"] = _compute_guidance(result["weights"])

        # ── 注入已有知识元概念索引（供AI去重前置检测）──
        result["concept_index"] = load_concept_index(workspace)

        # ═══════════════════════════════════════════════
        # 为所有资源生成 .ctx 元数据文件
        # ═══════════════════════════════════════════════
        ctx_count = 0

        for r in result.get("resources_flat", []):
            generate_ctx_file(workspace, r, current_rel)
            ctx_count += 1

        result["ctx_count"] = ctx_count

        # 提取成功后，自动将文件移入原始文件目录（按分类子目录保存）
        source_dir = workspace / "原始文件"
        source_dir.mkdir(parents=True, exist_ok=True)
        try:
            file_rel = filepath.relative_to(workspace)
            is_in_source = str(file_rel).startswith("原始文件")
        except ValueError:
            is_in_source = False
            file_rel = Path(filepath.name)

        # ═══════════════════════════════════════════════
        # 保存 chunk 到 chunks/ 目录（在文件移动前，使用原始路径）
        # ═══════════════════════════════════════════════
        try:
            chunk_result = save_chunks(result, current_rel, workspace)
            result["chunks_saved"] = chunk_result["chunk_count"]
        except Exception as e:
            result["chunks_saved"] = 0
            result["chunks_error"] = str(e)

        if not is_in_source:
            # 确定分类目录：子目录路径=用户分类，根目录文件需--category或归入"未分类"
            if str(file_rel.parent) != ".":
                # 文件已在子目录中 → 继承该子目录路径作为分类
                category_dir = str(file_rel.parent).replace("\\", "/")
            elif category:
                # 用户通过 --category 指定分类
                category_dir = category
            else:
                # 根目录文件且无分类指定 → 归入"未分类"
                category_dir = "未分类"

            dest_dir = source_dir / category_dir
            dest_dir.mkdir(parents=True, exist_ok=True)
            dest = dest_dir / filepath.name
            if dest.exists():
                stem = dest.stem
                suffix = dest.suffix
                dest = dest_dir / f"{stem}_{datetime.datetime.now().strftime('%Y%m%d_%H%M%S')}{suffix}"
            filepath.rename(dest)
            # 移动成功后，若原文件夹为空则删除，避免遗留空目录让用户误会
            try:
                orig_parent = filepath.parent
                if orig_parent.resolve() != workspace.resolve():
                    # 忽略系统残留文件（Windows 缩略图/桌面设置等）
                    remaining = [f for f in orig_parent.iterdir()
                                 if f.name not in ('Thumbs.db', 'desktop.ini', 'Desktop.ini')
                                 and not f.name.startswith('.')]
                    if not remaining:
                        orig_parent.rmdir()
            except Exception:
                pass
            result["moved_to"] = str(dest.relative_to(workspace))
            result["file"] = str(dest.relative_to(workspace))
        else:
            result["moved_to"] = None
            # 统一 result["file"] 为相对路径
            result["file"] = current_rel

        print(json.dumps(result, ensure_ascii=False, default=str))

    except Exception as e:
        print(json.dumps({"error": str(e), "file": str(filepath)}, ensure_ascii=False))
        sys.exit(1)


# ─── 子命令: mark-processed ───

@command("mark-processed")
def cmd_mark_processed(args):
    """Mark a file as processed in the manifest.

    🔴 知识元完整性检查（自动）：
    - 扫描 知识元/*.md，检查 frontmatter 中 source 是否匹配该文件
    - 零匹配且无 --force → ❌ 拒绝标记，输出 warning
    - 有匹配 → ✅ 正常标记
    """
    import sys as _sys

    # 从 args 中提取文件路径（跳过 -- 开头的 flag）
    filepath_arg = None
    for a in args:
        if not a.startswith("--"):
            filepath_arg = a
            break
    if not filepath_arg:
        print(json.dumps({"error": "Usage: ahkb.py mark-processed <filepath> [--force] [--workspace <path>]"}))
        sys.exit(1)

    workspace = get_workspace(args)
    force_mode = "--force" in args

    filepath = Path(filepath_arg)
    if not filepath.is_absolute():
        filepath = workspace / filepath
    try:
        rel_path = str(filepath.relative_to(workspace)).replace("\\", "/")
    except ValueError:
        rel_path = str(filepath).replace("\\", "/")

    # ═══════════════════════════════════════════════════════════
    # 🔴 知识元完整性检查（--force 可跳过）
    # ═══════════════════════════════════════════════════════════
    if not force_mode:
        knowledge_dir = workspace / "知识元"
        matched_units = []
        if knowledge_dir.exists():
            for f in sorted(knowledge_dir.glob("*.md")):
                try:
                    content = f.read_text(encoding="utf-8")
                    fm = parse_frontmatter(content)
                    if fm:
                        source = (fm.get("source", "") or "").replace("\\", "/").strip('"')
                        # Windows 上大小写不敏感比较
                        if source and source.lower() == rel_path.lower():
                            matched_units.append(f.stem)
                except Exception:
                    continue

        if not matched_units:
            print(json.dumps({
                "warning": True,
                "type": "no_knowledge_units",
                "source_path": rel_path,
                "message": (
                    f"⚠️ 知识元完整性检查未通过：在 知识元/ 中未找到任何 "
                    f"source 为「{rel_path}」的知识元。\n"
                    f"   请先执行 [C] 分析和 [D] 写入步骤，然后重新运行 mark-processed。\n"
                    f"   如需强制标记，请使用 --force 参数。"
                ),
            }, ensure_ascii=False))
            return

        # 有匹配 → 报告
        _sys.stderr.write(f"\n  ✅ 知识元完整性检查通过：{len(matched_units)} 个知识元匹配\n")

    # ★ 仅标记已处理（资源关联统一由 maintain 完成）
    manifest = load_json(get_manifest_path(workspace))
    if "processed" not in manifest:
        manifest["processed"] = []

    for p in manifest["processed"]:
        if p.get("source_path") and p["source_path"].replace("\\", "/").lower() == rel_path.lower():
            print(json.dumps({"ok": True, "already_exists": True}))
            return

    manifest["processed"].append({
        "file": filepath.name,
        "source_path": rel_path,
        "processed_date": datetime.datetime.now().isoformat(),
        "sha256": file_hash(str(filepath)),
    })
    save_json(get_manifest_path(workspace), manifest)
    print(json.dumps({"ok": True, "source_path": rel_path}))
    if not force_mode:
        _sys.stderr.write(f"  ✅ 已标记为已处理\n")


# ─── 子命令: clear-units-for-file ───

@command("clear-units-for-file")
def cmd_clear_units_for_file(args):
    """清除指定文档的知识元并从 manifest 移除标记（保留 chunks 和资源文件）。

    用法：ahkb.py clear-units-for-file <source_path> [--workspace <path>]
    """
    if len(args) < 1 or args[0].startswith("--"):
        print(json.dumps({"error": "Usage: ahkb.py clear-units-for-file <source_path> [--workspace <path>]"}))
        sys.exit(1)

    workspace = get_workspace(args)
    src_path = args[0].replace("\\", "/")
    knowledge_dir = workspace / "知识元"

    # 1. 删除匹配的知识元 .md 文件
    deleted_units = []
    deleted_names = []
    if knowledge_dir.exists():
        for f in list(knowledge_dir.glob("*.md")):
            try:
                content = f.read_text(encoding="utf-8", errors="ignore")
                fm = parse_frontmatter(content)
                if fm:
                    source = (fm.get("source", "") or "").replace("\\", "/").strip('"')
                    if source and source.lower() == src_path.lower():
                        _trash_file(f, workspace)
                        deleted_units.append(str(f.relative_to(workspace)))
                        deleted_names.append(f.stem)
            except Exception:
                continue

    # 2. 从已处理清单中移除
    manifest_path = get_manifest_path(workspace)
    manifest = load_json(manifest_path)
    processed = manifest.get("processed", [])
    new_processed = [
        p for p in processed
        if p.get("source_path", "").replace("\\", "/").lower() != src_path.lower()
    ]
    removed_count = len(processed) - len(new_processed)
    manifest["processed"] = new_processed
    save_json(manifest_path, manifest)

    # 3. 不动 chunks/ 和 图片及其他资源/ 下的任何文件

    print(json.dumps({
        "ok": True,
        "source_path": src_path,
        "deleted_units": len(deleted_units),
        "unit_names": deleted_names,
        "removed_from_manifest": removed_count > 0,
    }, ensure_ascii=False))
    import sys as _sys
    _sys.stderr.write(f"\n  ✅ 已清除 {len(deleted_units)} 个知识元，保留 chunks 和资源文件\n")
    if deleted_units:
        _sys.stderr.write(f"  📄 已删除: {', '.join(deleted_names)}\n")
    _sys.stderr.write(f"  ℹ️  请执行 [C] 分析步骤重新生成知识元\n")


@command("clear-manifest")
def cmd_clear_manifest(args):
    """Clear the processed files manifest (for rebuild)."""
    workspace = get_workspace(args)
    mpath = get_manifest_path(workspace)
    _trash_file(mpath, workspace)
    print(json.dumps({"ok": True, "cleared": True}))


# ─── 子命令: purge-kb ───

@command("purge-kb")
def cmd_purge_kb(args):
    """清除全部知识元和资源文件（绕过权限拦截，一步完成）。"""
    import subprocess as _sp, sys as _sys
    script = Path(__file__).parent / "ahkb_purge.py"
    ws = get_workspace(args)
    if ws:
        cmd = [_sys.executable, str(script), "--workspace", str(ws)]
    else:
        cmd = [_sys.executable, str(script)]
    r = _sp.run(cmd, cwd=Path.cwd())
    _sys.exit(r.returncode)


# ─── 子命令: find-root ───

@command("find-root")
def cmd_find_root(args):
    """Find the root knowledge graph .md file."""
    workspace = get_workspace(args)
    for f in workspace.iterdir():
        if f.suffix.lower() == ".md" and f.is_file():
            try:
                with open(f, "r", encoding="utf-8") as fh:
                    content = fh.read(500)
                if "root_node: true" in content:
                    print(json.dumps({"found": True, "path": str(f.relative_to(workspace)),
                                      "name": f.stem}))
                    return
            except Exception:
                pass
    print(json.dumps({"found": False}))


# ─── 根节点写入工具 ───

def _write_root_node(workspace, root_file, kb_name, units, tag_groups):
    """将知识库数据写入根节点 .md 文件。
    build-graph 调用此函数，无需 AI 手工介入。
    """
    today = datetime.datetime.now().strftime("%Y-%m-%d")

    # 如果是更新已有根节点，保留原始创建日期
    created = today
    if root_file.exists():
        try:
            old_fm = parse_frontmatter(root_file.read_text(encoding="utf-8"))
            if old_fm and old_fm.get("created"):
                created = old_fm["created"]
        except Exception:
            pass

    # 统计标签
    all_tags = []
    all_tags_set = set()
    for u in units:
        for t in u.get("tags", []):
            if t and t not in all_tags_set:
                all_tags_set.add(t)
                all_tags.append(t)

    # 统计资源数
    resource_count = 0
    resource_base = workspace / "图片及其他资源"
    if resource_base.exists():
        for sd in ["images", "videos", "audios", "others"]:
            sd_path = resource_base / sd
            if sd_path.exists():
                for f in sd_path.iterdir():
                    if not f.name.endswith(".ctx"):
                        resource_count += 1

    # 统计原始文件数
    manifest = load_json(get_manifest_path(workspace))
    source_count = len(manifest.get("processed", [])) if isinstance(manifest, dict) else 0

    # 概述
    overview = f"本知识库包含 {len(units)} 个知识元"
    if all_tags:
        top_tags = all_tags[:5]
        overview += f"，涵盖 {'、'.join(top_tags)}"
        if len(all_tags) > 5:
            overview += f"等 {len(all_tags)} 个领域"
    overview += "。"

    # 构建内容
    content = "---\n"
    content += f"tags: [{', '.join(all_tags) if all_tags else '知识库'}, 知识库, 根节点]\n"
    content += "root_node: true\n"
    content += f"created: {created}\n"
    content += f"updated: {today}\n"
    content += "---\n\n"
    content += f"# {kb_name} 知识库\n\n"
    content += "> 基于 AHKB-CPS v0.1.0 构建的全息知识库\n\n"
    content += "## 📖 知识库概述\n\n"
    content += f"{overview}\n\n"
    content += "## 📊 统计信息\n\n"
    content += "| 项目 | 数量 |\n"
    content += "|------|:----:|\n"
    content += f"| 原始文件 | {source_count} 个 |\n"
    content += f"| 知识元 | {len(units)} 个 |\n"
    content += f"| 资源文件 | {resource_count} 个 |\n"
    content += f"| 最后更新 | {today} |\n\n"
    content += "## 🗺️ 知识地图\n\n"
    content += f"> [打开交互式知识地图]({kb_name}-知识地图.html)\n\n"
    content += f"## 📋 知识元列表（{len(units)} 个）\n\n"
    content += "| 知识元 | 标签 | 来源 |\n"
    content += "|--------|------|------|\n"
    for u in units:
        name = u.get("title", u["file"].replace(".md", ""))
        tags = ", ".join(u.get("tags", [])) if u.get("tags") else "—"
        source = u.get("source", "—")
        content += f"| [[{name}]] | {tags} | {source} |\n"

    # 处理清单
    content += "\n## 📋 处理清单\n\n"
    content += "| 文件 | 状态 | 知识元数 |\n"
    content += "|------|:----:|:--------:|\n"
    if isinstance(manifest, dict):
        from collections import Counter
        source_counts = Counter()
        for u in units:
            s = u.get("source", "")
            if s:
                source_counts[s] = source_counts.get(s, 0) + 1
        for entry in manifest.get("processed", []):
            src_file = entry.get("source_path") or entry.get("file", "—")
            src_count = source_counts.get(src_file, 0)
            content += f"| {src_file} | ✅ 已处理 | {src_count} |\n"
    else:
        content += "| （无处理记录） | — | — |\n"

    content += "\n## 🔗 关联模块\n\n"
    content += "- 📊 [[幻灯片生成]] — 基于知识库生成演示文稿\n"
    content += "- ✍️ [[文章生成]] — 生成结构化文章\n"
    content += "- 🧠 [[全息脑图生成]] — 生成阿色全息脑图\n"

    try:
        root_file.write_text(content, encoding="utf-8")
    except Exception:
        pass


# ─── 子命令: build-graph ───

@command("build-graph")
def cmd_build_graph(args):
    """Rebuild the root knowledge graph .md file."""
    workspace = get_workspace(args)
    knowledge_dir = workspace / "知识元"

    if not knowledge_dir.exists():
        print(json.dumps({"error": "知识元/ directory not found"}))
        sys.exit(1)

    # 读取所有知识单元的 frontmatter
    units = []
    for f in sorted(knowledge_dir.glob("*.md")):
        try:
            with open(f, "r", encoding="utf-8") as fh:
                content = fh.read()
            fm = parse_frontmatter(content)
            if fm:
                resources = fm.get("resources", [])
                # resources 可能是字符串列表或 YAML 列表
                if isinstance(resources, str):
                    resources = [resources]
                units.append({
                    "file": f.name,
                    "title": fm.get("title", f.stem),
                    "tags": fm.get("tags", []),
                    "summary": fm.get("summary", ""),
                    "source": fm.get("source", ""),
                    "resources": resources,
                })
        except Exception:
            pass

    # 按标签归类
    tag_groups = {}
    for u in units:
        for tag in u.get("tags", []):
            tag_groups.setdefault(tag, []).append(u["file"])

    # 查找或创建根节点文件
    kb_name = get_kb_name(workspace)
    root_file = None
    for f in workspace.iterdir():
        if f.suffix.lower() == ".md" and f.is_file():
            try:
                with open(f, "r", encoding="utf-8") as fh:
                    if "root_node: true" in fh.read(500):
                        root_file = f
                        break
            except Exception:
                pass

    if root_file is None:
        root_file = workspace / f"{kb_name}(根).md"

    # ── 自动写入根节点 .md ──
    _write_root_node(workspace, root_file, kb_name, units, tag_groups)

    print(json.dumps({
        "ok": True,
        "root_file": str(root_file.relative_to(workspace)) if root_file else "",
        "units_count": len(units),
        "units": units,
        "tag_groups": {k: len(v) for k, v in tag_groups.items()},
    }, ensure_ascii=False))

    # ── 同时更新概念索引 ──
    save_concept_index(workspace)


# ─── 子命令: build-kg-html ───

@command("build-kg-html")
def cmd_build_kg_html(args):
    """Build D3.js interactive knowledge graph HTML from 知识元/."""
    # 委托给独立脚本
    script_dir = Path(__file__).parent
    script_path = script_dir / "ahkb_build_kg_html.py"
    if not script_path.exists():
        print(json.dumps({"error": f"脚本不存在: {script_path}"}))
        sys.exit(1)

    workspace = get_workspace(args)
    kb_name = get_kb_name(workspace)
    result = subprocess.run(
        [sys.executable, str(script_path),
         "--workspace", str(workspace),
         "--kb-name", kb_name],
        capture_output=True
    )
    # 手动解码，兼容 Windows 编码问题
    stderr_text = result.stderr.decode("utf-8", errors="replace")
    stdout_text = result.stdout.decode("utf-8", errors="replace")
    # 打印 stderr（进度信息）
    if stderr_text.strip():
        print(stderr_text, file=sys.stderr, end="")
    # 打印 stdout（JSON 结果）
    if stdout_text.strip():
        try:
            data = json.loads(stdout_text.strip())
            # 自动在本地浏览器中打开生成的 HTML
            if data.get("ok") and data.get("html_path"):
                html_abs = (workspace / data["html_path"]).resolve()
                opened = False
                try:
                    webbrowser.open(str(html_abs))
                    opened = True
                except Exception:
                    pass
                if not opened and platform.system() == "Windows":
                    try:
                        os.startfile(str(html_abs))
                        opened = True
                    except Exception:
                        pass
                data["opened_in_browser"] = opened
                if opened:
                    print(f"\n🌐 已在浏览器中打开: {data['html_path']}", file=sys.stderr)
                else:
                    print(f"\n📂 HTML 已生成但无法自动打开，请手动打开: {html_abs}", file=sys.stderr)
            print(json.dumps(data, ensure_ascii=False))
        except json.JSONDecodeError:
            print(stdout_text)


# ─── 子命令: stats ───

@command("stats")
def cmd_stats(args):
    """Show knowledge base statistics."""
    workspace = get_workspace(args)
    knowledge_dir = workspace / "知识元"
    resource_base = workspace / "图片及其他资源"
    source_dir = workspace / "原始文件"
    manifest = load_json(get_manifest_path(workspace))

    units_count = len(list(knowledge_dir.glob("*.md"))) if knowledge_dir.exists() else 0

    # 统计各类资源
    resource_counts = {}
    ctx_counts = {}
    for subdir in ["images", "videos", "audios", "others"]:
        d = resource_base / subdir
        if d.exists():
            # 统计媒体文件（非 .ctx）
            media_files = [f for f in d.iterdir() if f.is_file() and f.suffix.lower() != '.ctx']
            ctx_files = [f for f in d.iterdir() if f.is_file() and f.suffix.lower() == '.ctx']
            resource_counts[subdir] = len(media_files)
            ctx_counts[subdir] = len(ctx_files)

    processed_count = len(manifest.get("processed", []))
    # 统计原始文件（含子目录分类）
    source_count = 0
    if source_dir.exists():
        for f in source_dir.rglob("*"):
            if f.is_file() and f.suffix.lower() in SUPPORTED_EXTS:
                source_count += 1

    # 收集所有标签
    all_tags = {}
    if knowledge_dir.exists():
        for f in knowledge_dir.glob("*.md"):
            try:
                with open(f, "r", encoding="utf-8") as fh:
                    content = fh.read()
                fm = parse_frontmatter(content)
                if fm and "tags" in fm:
                    for tag in fm["tags"]:
                        all_tags[tag] = all_tags.get(tag, 0) + 1
            except Exception:
                pass

    top_tags = sorted(all_tags.items(), key=lambda x: -x[1])[:30]

    result = {
        "knowledge_units": units_count,
        "resources": resource_counts,
        "resource_ctx_files": ctx_counts,
        "processed_documents": processed_count,
        "source_documents": source_count,
        "top_tags": [{"tag": t, "count": c} for t, c in top_tags],
    }
    print(json.dumps(result, ensure_ascii=False))


# ─── 子命令: list-processed ───

@command("list-processed")
def cmd_list_processed(args):
    """列出已处理文档及其知识元数、chunk 数。"""
    workspace = get_workspace(args)
    manifest = load_json(get_manifest_path(workspace))
    processed = manifest.get("processed", [])
    knowledge_dir = workspace / "知识元"

    result = []
    for p in processed:
        src = p.get("source_path", "").replace("\\", "/")
        # 统计该文件对应的知识元数
        unit_count = 0
        if knowledge_dir.exists():
            for f in knowledge_dir.glob("*.md"):
                try:
                    content = f.read_text(encoding="utf-8", errors="ignore")
                    fm = parse_frontmatter(content)
                    if fm:
                        source = (fm.get("source", "") or "").replace("\\", "/").strip('"')
                        if source and source.lower() == src.lower():
                            unit_count += 1
                except Exception:
                    continue
        # 统计 chunk 数
        from ahkb_chunks import load_chunks_for_file
        chunk_data = load_chunks_for_file(src, workspace)
        chunks = chunk_data.get("chunks", [])
        chunk_count = len(chunks)

        result.append({
            "file": p.get("file", ""),
            "source_path": src,
            "units_count": unit_count,
            "chunk_count": chunk_count,
            "processed_date": p.get("processed_date", ""),
        })

    print(json.dumps({
        "ok": True,
        "total": len(result),
        "documents": result,
    }, ensure_ascii=False))


# ─── 共享：计算知识元生成指导 ───

def _compute_guidance(weights):
    """根据 weights 计算知识元生成指导（三档离散模式），返回 dict。

    废除旧公式 8*(1-cG)^2+0.5 和 300*cTA^2+150*cTA+50，
    改为三档离散行为映射。
    cGranularity: 0→精细 / 0.5→均衡 / 1→粗粒
    cTextAmount:  0→摘要 / 0.5→适中 / 1→详细
    """
    # ── cGranularity 三档映射（允许 0.25/0.75 作为切换点）──
    raw_g = float(weights.get("cGranularity", 0.5))
    if raw_g <= 0.25:
        g_level = 0  # 精细
    elif raw_g >= 0.75:
        g_level = 2  # 粗粒
    else:
        g_level = 1  # 均衡（默认）

    g_names = {0: "精细模式", 1: "均衡模式", 2: "粗粒模式"}
    g_name = g_names[g_level]

    # ── cTextAmount 三档映射 ──
    raw_t = float(weights.get("cTextAmount", 0.5))
    if raw_t <= 0.25:
        t_level = 0  # 摘要
    elif raw_t >= 0.75:
        t_level = 2  # 详细
    else:
        t_level = 1  # 适中（默认）

    t_names = {0: "摘要模式", 1: "适中模式", 2: "详细模式"}
    t_name = t_names[t_level]

    # ── 各档位参数 ──
    # 字数底线
    min_chars = {0: 50, 1: 60, 2: 120}[g_level]
    # 连续短节合并阈值（每节字数低于此则合并）
    merge_threshold = {0: 30, 1: 30, 2: 100}[g_level]

    # ── 展开风格描述 ──
    text_styles = {
        0: "【摘要】只写定义加 2~3 个核心要点，不展开论证。类似摘要卡片。",
        1: "【适中】写一个完整段落，含定义、关键说明、要点。可适当润色语序。",
        2: "【详细】充分展开，含定义、说明、论证依据、举例、意义或应用。可润色扩写但不得编造。",
    }
    text_style = text_styles[t_level]

    # ── 拆分策略描述 ──
    split_rules = {
        0: (
            "⚡ 精细拆分：概念有独立性即拆分，宁可细不要粗。"
            "每个独立概念/数据点/事实都必须独立成元。"
        ),
        1: "均衡拆分：有明确概念边界且可脱离父级独立存在的才独立。属于整体组成部分的合入父级。",
        2: (
            "⚡ 粗粒合并：只按大主题拆分，子概念合入父级。"
            "同一主题下只有人物和定义必须独立，其他类型尽量合并。"
        ),
    }
    split_rule = split_rules[g_level]

    # ── 强制独立规则（按档位）──
    force_rules = {
        0: "人物/定义强制独立。其他概念只要有独立倾向就拆，不要犹豫。",
        1: "人物/定义强制独立。其他概念按正常边界判断。",
        2: "人物/定义强制独立。其他概念只有完全独立才拆，边界模糊一律合并。",
    }

    # ═══ 护城河原则（任何时候不违反）═══
    principles = {
        "no_force_split_or_merge": force_rules[g_level],
        "no_fabrication": (
            "所有内容必须有源文档原文依据，绝对禁止编造、推测、混入大模型自身知识。"
            "cTextAmount较低时：只写原文核心内容，严格控制字数。"
            "cTextAmount较高时：允许通过润色、调整语序、增加连接词来扩写正文，"
            "但不得增加原文没有的事实信息、数据、人名、定义。"
            "所有关键信息点必须在原文中有对应依据。准确性高于一切。"
        ),
        "no_empty_unit": (
            "每个知识元的正文内容不得为空，必须有实质文字。"
            "标题必须有内容支撑，如果一个概念只有标题没有可写的内容，"
            "则不单独成元，合入上级知识元。"
        ),
        "low_temperature": (
            "🔴 低温模式：以事实准确性为最高标准。"
            "每个知识元写入前，对照源文档原文 chunk 逐句回溯检查关键事实点在原文中是否有对应。"
        ),
    }

    # ── 组合 instruction（AI 直接看到此字段）──
    combined = (
        f"🔴 人物和定义在任何档位都必须独立成元，不得合并！"
        f" 🔴 与已有知识元名称高度相似（同一概念的不同表述、全称与缩写、多余修饰词的差异），不新建，追加关联到已有知识元的 related_files！"
        f' 🔴 概念依赖性：如果一个概念是另一个概念的组成部分（如「访问层」之于「AI数智油气田参考架构」），脱离父级含义不完整，则不独立成元，合入父级知识元展开。'
        f" 拆分策略→【{g_name}】{split_rule} "
        f"详细程度→{text_style} "
        f"字数底线：每知识元正文≥{min_chars}字，"
        f"连续多节每节<{merge_threshold}字的合并为一节。"
    )

    return {
        "cGranularity": raw_g,
        "cGranularity_level": g_level,
        "cGranularity_name": g_name,
        "cTextAmount": raw_t,
        "cTextAmount_level": t_level,
        "cTextAmount_name": t_name,
        "min_chars_per_unit": min_chars,
        "split_level": g_name,
        "instruction": combined,
        "principles": principles,
        "concept_dedup": True,  # 标记：需要前置去重检测
    }


# ─── 概念索引（concept_index）函数 ───

def build_concept_index(workspace):
    """从知识元目录构建概念索引，返回概念列表。

    从每个知识元 .md 文件的 frontmatter 提取 [name, summary, tags] 三元组。
    返回格式: {"concepts": [{"name": ..., "summary": ..., "tags": [...]}, ...]}
    """
    knowledge_dir = Path(workspace) / "知识元"
    concepts = []
    if knowledge_dir.exists():
        for f in sorted(knowledge_dir.glob("*.md")):
            try:
                content = f.read_text(encoding="utf-8", errors="ignore")
                fm = parse_frontmatter(content)
                if fm:
                    # 确保 tags 始终为列表（parse_frontmatter 可能返回字符串）
                    raw_tags = fm.get("tags", [])
                    if isinstance(raw_tags, str):
                        raw_tags = [t.strip() for t in raw_tags.split(",") if t.strip()]
                    elif not isinstance(raw_tags, list):
                        raw_tags = []
                    concepts.append({
                        "name": f.stem,
                        "summary": fm.get("summary", ""),
                        "tags": raw_tags,
                    })
            except Exception:
                continue
    return {"concepts": concepts}


def save_concept_index(workspace):
    """构建并保存概念索引到 临时工作文件/concept_index.json"""
    index = build_concept_index(workspace)
    index["updated"] = datetime.datetime.now().isoformat(timespec="seconds")
    p = Path(workspace) / "临时工作文件" / "concept_index.json"
    p.parent.mkdir(parents=True, exist_ok=True)
    p.write_text(json.dumps(index, ensure_ascii=False, indent=2), encoding="utf-8")
    return index


def load_concept_index(workspace):
    """加载概念索引，不存在则返回空索引。"""
    p = Path(workspace) / "临时工作文件" / "concept_index.json"
    if p.exists():
        try:
            return json.loads(p.read_text(encoding="utf-8"))
        except Exception:
            pass
    return {"concepts": []}


# ─── 子命令: update-concept-index ───

@command("update-concept-index")
def cmd_update_concept_index(args):
    """更新概念索引（从知识元目录重建）。"""
    workspace = get_workspace(args)
    result = save_concept_index(workspace)
    count = len(result.get("concepts", []))
    print(json.dumps({
        "ok": True,
        "concepts_count": count,
        "updated": result.get("updated", ""),
    }, ensure_ascii=False))
    import sys as _sys
    _sys.stderr.write(f"\n📇 概念索引已更新：{count} 个知识元\n")


# ─── 子命令: knowledge-gen-guidance ───

@command("knowledge-gen-guidance")
def cmd_knowledge_gen_guidance(args):
    """根据 project_settings.json 中的参数，输出知识元生成指导给 AI 使用。"""
    workspace = get_workspace(args)
    p = workspace / "系统设置" / "project_settings.json"
    weights = {}
    if p.exists():
        try:
            data = json.loads(p.read_text(encoding="utf-8"))
            weights = data.get("weights", {})
        except Exception:
            pass

    result = _compute_guidance(weights)
    result["accuracy_requirement"] = "🔴🔴🔴 最高标准：准确性！每个知识元必须与源文档原文严格一致，禁止编造、禁止推测、禁止混入大模型自身知识。无原文依据的内容不得写入知识元正文。写入前必须逐条自检准确性清单。"
    result["proper_noun_rule"] = "🔴 专有名词（人名/地名/术语名）必须与源文档逐字一致，严禁润色或修正。原文写什么就写什么，即使看起来像错别字也不得擅自更改。同一个名称在全库中必须保持写法完全统一。🚫 严禁使用汉语拼音处理中文人名，必须按字形（汉字）逐字匹配，禁止通过拼音转换，拼音会导致同音字混淆错误。"
    print(json.dumps(result, ensure_ascii=False))


# ─── 子命令: regenerate-ctx ───

@command("regenerate-ctx")
def cmd_regenerate_ctx(args):
    """Regenerate all .ctx files (preserves belongs_to links, rebuilds metadata).

    Without --force: only fixes broken/unset metadata (missing source, blank chunk_heading, etc).
    With --force: fully rebuilds all .ctx files from their embedded content.
    """
    workspace = get_workspace(args)
    resource_base = workspace / "图片及其他资源"

    if not resource_base.exists():
        print(json.dumps({"error": "Resource directory not found"}))
        return

    force_mode = args_has(args, "--force")
    total = 0
    updated = 0
    skipped = 0
    errors = []

    for subdir in ["images", "videos", "audios", "others"]:
        d = resource_base / subdir
        if not d.exists():
            continue
        for f in d.iterdir():
            if not f.is_file() or f.suffix.lower() != '.ctx':
                continue
            total += 1
            try:
                with open(f, "r", encoding="utf-8") as fh:
                    content = fh.read()

                parts = content.split("---", 2)
                if len(parts) < 3:
                    errors.append(f"{f.name}: invalid format")
                    continue

                fm_text = parts[1]
                body_text = parts[2] if len(parts) > 2 else ""

                def get_fm(key, default=""):
                    m = re.search(r'^' + key + r':\s*"(.+)"', fm_text, re.M)
                    return m.group(1) if m else default

                def get_fm_list(key):
                    m = re.search(r'^' + key + r':\s*\[(.+)\]', fm_text, re.M)
                    if m:
                        return [x.strip().strip('"').strip("'") for x in m.group(1).split(",")]
                    return []

                rtype = get_fm("resource_type", "other")
                source = get_fm("source", "")
                chunk_heading = get_fm("chunk_heading", "")
                belongs_to_chunk = get_fm("belongs_to_chunk", "")
                resource_file = get_fm("resource_file", "")
                remote_url = get_fm("remote_url", "")

                # 保留已有的 belongs_to
                belongs_to_list = []
                bt_match = re.search(r'belongs_to:\s*\n((?:\s+-.*\n?)*)', content)
                if bt_match:
                    belongs_to_list = [line.strip() for line in bt_match.group(1).split('\n') if line.strip()]
                tags = get_fm_list("tags")

                # 判断是否需要更新
                needs_update = force_mode
                if not needs_update:
                    # 没有强制模式时，只修复明确有问题的字段
                    if not source:  # source 为空
                        needs_update = True
                    elif not chunk_heading and not belongs_to_chunk:  # 无标题
                        needs_update = True
                    elif not tags:  # 标签为空
                        needs_update = True

                if not needs_update:
                    skipped += 1
                    continue

                # 重建 frontmatter（★ 严格按此字段顺序）
                new_fm_lines = ["---"]
                new_fm_lines.append("type: resource")
                new_fm_lines.append(f"resource_type: {rtype}")
                if "importance:" not in fm_text:
                    new_fm_lines.append("importance: 3")
                else:
                    imp_m = re.search(r'^importance:\s*(\d+)', fm_text, re.MULTILINE)
                    new_fm_lines.append(f"importance: {imp_m.group(1) if imp_m else 3}")
                if "user_edited:" not in fm_text:
                    new_fm_lines.append("user_edited: false")
                else:
                    ue_m = re.search(r'^user_edited:\s*(true|false)', fm_text, re.MULTILINE)
                    new_fm_lines.append(f"user_edited: {ue_m.group(1) if ue_m else 'false'}")
                now_str = datetime.datetime.now().strftime("%Y-%m-%d %H:%M:%S")
                le_m = re.search(r'^last_edited_time:\s*(.*)', fm_text, re.MULTILINE)
                if le_m and le_m.group(1).strip():
                    new_fm_lines.append(f"last_edited_time: {le_m.group(1).strip()}")
                else:
                    new_fm_lines.append(f"last_edited_time: {now_str}")
                new_fm_lines.append(f'source: "{source}"')
                if belongs_to_chunk:
                    new_fm_lines.append(f'belongs_to_chunk: "{belongs_to_chunk}"')
                if chunk_heading:
                    new_fm_lines.append(f'chunk_heading: "{chunk_heading}"')
                if resource_file:
                    new_fm_lines.append(f'resource_file: "{resource_file}"')
                if remote_url:
                    new_fm_lines.append(f'remote_url: "{remote_url}"')
                new_fm_lines.append(f"tags: [{', '.join(tags)}]")
                # belongs_to 必须是最后一个字段
                if belongs_to_list:
                    new_fm_lines.append("belongs_to:")
                    for bt in belongs_to_list:
                        new_fm_lines.append(f"  - {bt}")
                else:
                    new_fm_lines.append("belongs_to:")
                new_fm_lines.append("---")
                new_fm_lines.append("")

                # 重建 body：保留已有正文，不丢失上下文
                new_body_lines = []
                if resource_file:
                    new_body_lines.append(f"![[{resource_file}]]")
                elif remote_url:
                    new_body_lines.append(f"远程资源：{remote_url}")
                new_body_lines.append("")

                # body 中已有关键上下文（> 引用的文字）
                body_lines = body_text.strip().split("\n")
                ctx_lines = [l for l in body_lines if l.strip().startswith(">")]
                if ctx_lines:
                    new_body_lines.extend(ctx_lines)
                else:
                    new_body_lines.append(f"> （资源属于「{chunk_heading or belongs_to_chunk or '未分类'}」）")
                new_body_lines.append("")

                new_content = "\n".join(new_fm_lines + new_body_lines)
                with open(f, "w", encoding="utf-8") as fh:
                    fh.write(new_content)
                updated += 1

            except Exception as e:
                errors.append(f"{f.name}: {str(e)}")

    result = {
        "ok": True,
        "total_ctx": total,
        "updated": updated,
        "skipped": skipped,
        "errors": errors[:10],
    }
    print(json.dumps(result, ensure_ascii=False))


def args_has(args, flag):
    """Check if a flag is present in args list."""
    return flag in args


# ─── Phase 1 完成度检查（供 maintain / cross-link 使用）───

def _check_phase1_complete(workspace):
    """检查 Phase 1（A→F 文档处理流水线）是否全部完成。

    返回 dict:
        {"complete": bool, "new_count": int, "pending_count": int, "total_processed": int,
         "unprocessed_files": [...], "error": str|None}

    仅当 new_count == 0 且 pending_count == 0 时，complete 才为 True。
    此函数是 maintain 的硬关卡——知识元未全部抽取完毕时禁止启动关联。
    """
    # 读取已处理清单
    manifest = load_json(get_manifest_path(workspace))
    processed_files = {p["source_path"] for p in manifest.get("processed", [])}
    processed_names = {p["file"] for p in manifest.get("processed", [])}

    # 系统目录和忽略文件（与 scan 保持一致）
    SYSTEM_DIRS = {"知识元", "图片及其他资源", "临时工作文件", "系统设置", "产品成果", "chunks", "回收站"}
    BACKUP_DIR_NAMES = {"备份", "backup", "Backup", "BACKUP", "_backup", "_备份", ".backup", "archive", "存档"}
    IGNORE_FILES = set()
    IGNORE_PREFIXES = ("~", ".~")

    SUPPORTED_EXTS = {".pptx", ".docx", ".xlsx", ".pdf", ".md", ".html", ".htm", ".txt", ".csv", ".tsv"}

    new_files = []
    pending_files = []
    new_paths = set()

    # ── 主扫描：遍历工作空间中所有支持的文件 ──
    for f in workspace.rglob("*"):
        if not f.is_file():
            continue
        if f.suffix.lower() not in SUPPORTED_EXTS:
            continue
        if f.stat().st_size == 0:
            continue
        if f.name in IGNORE_FILES:
            continue
        if any(f.name.startswith(p) for p in IGNORE_PREFIXES):
            continue

        rel = f.relative_to(workspace)
        parts = rel.parts
        if parts[0] in SYSTEM_DIRS:
            continue
        if any(part.startswith(".") for part in parts[:-1]):
            continue
        if any(part in BACKUP_DIR_NAMES for part in parts[:-1]):
            continue
        # 跳过根节点和知识地图
        if f.suffix.lower() == '.md' and _is_root_node(f):
            continue
        if f.suffix.lower() == '.html':
            if f.stem.endswith('-知识地图'):
                continue
            md_version = f.with_suffix('.md')
            if md_version.exists() and _is_root_node(md_version):
                continue

        rel_path = str(rel).replace("\\", "/")
        if rel_path in processed_files or f.name in processed_names:
            continue  # 已处理
        new_paths.add(rel_path)
        new_files.append(f.name)

    # ── 原始文件/ 中遗漏的未处理文件 ──
    source_dir = workspace / "原始文件"
    if source_dir.exists():
        for f in source_dir.rglob("*"):
            if not f.is_file() or f.suffix.lower() not in SUPPORTED_EXTS:
                continue
            if f.name in IGNORE_FILES:
                continue
            rel_parts = f.relative_to(workspace).parts
            if any(p.startswith(".") for p in rel_parts[1:-1]):
                continue
            if any(p in BACKUP_DIR_NAMES for p in rel_parts[1:-1]):
                continue
            rel_path = str(f.relative_to(workspace)).replace("\\", "/")
            if rel_path in processed_files or f.name in processed_names:
                continue
            if rel_path in new_paths:
                continue
            new_paths.add(rel_path)
            new_files.append(f.name)

    # ── chunks/ 中已提取但未标记为已处理的文件 ──
    chunks_dir = workspace / "chunks"
    chunk_index_path = chunks_dir / "index.json"
    if chunk_index_path.exists():
        try:
            chunk_index = load_json(chunk_index_path)
            chunked_files = chunk_index.get("files", [])
            for src_path in chunked_files:
                if src_path in processed_files:
                    continue
                if src_path in new_paths:
                    continue
                if any(src_path == e.get("path", "") for e in pending_files):
                    continue
                src_file = workspace / src_path
                if not src_file.exists():
                    continue
                if src_file.suffix.lower() not in SUPPORTED_EXTS:
                    continue
                if src_file.name in IGNORE_FILES:
                    continue
                new_paths.add(src_path)
                pending_files.append({
                    "name": src_file.name,
                    "path": src_path,
                })
        except Exception:
            pass

    total_processed = len(processed_files)
    new_count = len(new_files)
    pending_count = len(pending_files)
    complete = (new_count == 0 and pending_count == 0)

    result = {
        "complete": complete,
        "new_count": new_count,
        "pending_count": pending_count,
        "total_processed": total_processed,
        "unprocessed_samples": new_files[:10] + [p["name"] for p in pending_files[:10]],
    }

    if not complete:
        result["error"] = (
            f"Phase 1 尚未完成：还有 {new_count} 个文件未处理、"
            f"{pending_count} 个文件已提取但未完成关联。"
            f"请先完成所有文档的 A→F 流水线后再执行 maintain。"
        )

    return result


# ─── 子命令: cross-link ───

@command("cross-link")
def cmd_cross_link(args):
    """Cross-link: 知识元↔知识元 ([[链接]]) + 知识元↔资源 (.ctx) across all documents."""
    # 🔴 禁止使用--headless无头模式，必须弹出 tkinter GUI 窗口让用户可见
    # 如果已经在 GUI 窗口内（被 run_cmd_thread 调用），直接运行不递归弹窗
    if _IN_GUI:
        workspace = get_workspace(args)
        from ahkb_crosslink import cross_link
        result = cross_link(workspace)
        print(json.dumps(result, ensure_ascii=False))
        return result

    # ─── detached 模式：启动独立子进程，立即返回 ───
    if _DETACHED:
        workspace = get_workspace(args)
        ws_path = Path(workspace).resolve()
        tmp_dir = ws_path / "临时工作文件"
        tmp_dir.mkdir(parents=True, exist_ok=True)
        pid_file = tmp_dir / "_crosslink.pid"
        result_file = tmp_dir / "_crosslink_result.json"
        progress_file = tmp_dir / "_crosslink_progress.json"

        # PID 冲突检查
        if pid_file.exists():
            try:
                raw = pid_file.read_text(encoding="utf-8").strip()
                if not raw:
                    # 空文件 = 上轮残留，移入回收站
                    _trash_file(pid_file, ws_path)
                else:
                    old_data = json.loads(raw)
                    old_pid = old_data.get("pid", 0)
                    if _is_pid_alive(old_pid):
                        print(json.dumps({"status": "already_running", "pid": old_pid,
                            "message": "知识链构建任务已在运行中（PID: %d），请等待完成或手动终止。" % old_pid},
                            ensure_ascii=False))
                        return {"status": "already_running", "pid": old_pid}
                    else:
                        _trash_file(pid_file, ws_path)
            except (json.JSONDecodeError, ValueError):
                # JSON 解析失败（损坏），移入回收站
                _trash_file(pid_file, ws_path)

        # 🔴🔴 Phase 1 完成度硬检查（阻止知识元尚未全部抽取即启动关联）
        phase1_status = _check_phase1_complete(ws_path)
        if not phase1_status["complete"]:
            msg = {
                "status": "phase1_incomplete",
                "error": phase1_status["error"],
                "new_count": phase1_status["new_count"],
                "pending_count": phase1_status["pending_count"],
                "total_processed": phase1_status["total_processed"],
                "hint": "请先完成所有文档的 A→F 流水线（extract → 分析 → 写入 → verify-units → mark-processed），全部完成后才能执行 maintain。"
            }
            print(json.dumps(msg, ensure_ascii=False))
            return msg

        # 清理旧进度/结果文件（移入回收站）
        _trash_file(progress_file, ws_path)
        _trash_file(result_file, ws_path)

        # 🔴 自动更新概念索引（确保 crosslink 使用最新的知识元清单）
        _pre_index = save_concept_index(ws_path)
        sys.stderr.write(
            f"📇 概念索引已自动更新：{_pre_index.get('concepts_count', len(_pre_index.get('concepts', [])))} 个知识元\n"
        )

        # 构建子进程命令：直接调用 ahkb_crosslink.py
        script = Path(__file__).resolve().parent / "ahkb_crosslink.py"
        # 使用 pythonw.exe 避免控制台窗口（Windows）
        if platform.system() == "Windows":
            _pyw = Path(sys.executable).with_name("pythonw.exe")
            _py_exe = str(_pyw) if _pyw.exists() else sys.executable
        else:
            _py_exe = sys.executable
        cmd = [_py_exe, str(script), "--detached-child", "--workspace", str(ws_path)]
        if args_has(args, "--dry-run"):
            cmd.append("--dry-run")

        # 跨平台解绑启动
        if platform.system() == "Windows":
            CREATE_NEW_PROCESS_GROUP = 0x00000200
            proc = subprocess.Popen(
                cmd,
                creationflags=CREATE_NEW_PROCESS_GROUP,
                close_fds=True,
                stdin=subprocess.DEVNULL,
                stdout=subprocess.DEVNULL,
                stderr=subprocess.DEVNULL
            )
        else:
            proc = subprocess.Popen(
                cmd,
                start_new_session=True,
                close_fds=True,
                stdin=subprocess.DEVNULL,
                stdout=subprocess.DEVNULL,
                stderr=subprocess.DEVNULL
            )

        result = {
            "status": "started",
            "pid": proc.pid,
            "progress_file": str(progress_file),
            "result_file": str(result_file)
        }
        print(json.dumps(result, ensure_ascii=False))
        return result

    # 弹出 tkinter GUI 窗口（显示运行进度）
    result = _run_gui_window("知识链构建引擎", "cross-link", args)
    # GUI 关闭后把结果打印到 stdout，AI 才能检测到完成状态
    print(json.dumps(result, ensure_ascii=False))
    return result


# ─── 辅助函数 ───

def file_hash(filepath):
    """Compute SHA256 hash of a file."""
    h = hashlib.sha256()
    try:
        with open(filepath, "rb") as f:
            for chunk in iter(lambda: f.read(65536), b""):
                h.update(chunk)
        return h.hexdigest()[:16]
    except Exception:
        return ""


def parse_frontmatter(content):
    """Parse YAML frontmatter from markdown content (simple version)."""
    if not content.startswith("---"):
        return None
    parts = content.split("---", 2)
    if len(parts) < 3:
        return None
    fm = {}
    for line in parts[1].strip().split("\n"):
        line = line.strip()
        if not line:
            continue
        if ":" in line:
            k, v = line.split(":", 1)
            k = k.strip()
            v = v.strip()
            if v.startswith("[") and v.endswith("]"):
                try:
                    v = json.loads(v.replace("'", '"'))
                except json.JSONDecodeError:
                    v = [x.strip().strip('"').strip("'") for x in v[1:-1].split(",")]
            elif v.startswith('"') and v.endswith('"'):
                v = v[1:-1]
            elif v.startswith("'") and v.endswith("'"):
                v = v[1:-1]
            fm[k] = v
    return fm


# ─── 子命令: set-graph-colors ───

# 预定义颜色方案（按标签名特征自动匹配色系）
_COLOR_PALETTE = [
    # (color_name, rgb_int, keywords_in_tag)
    ("红色", 12595883, ["系统", "理论", "核心", "基础"]),
    ("深紫", 9324717,  ["大系统", "BSV"]),
    ("蓝色", 2715833,  ["AI", "智能", "人工", "机器", "大模型"]),
    ("紫色", 10185142, ["哲学", "思维", "思想"]),
    ("橙色", 15100450, ["方法", "实践", "工具", "规划"]),
    ("青色", 1754268,  ["全息", "HOST", "有机"]),
    ("金色", 13934503, ["道德经", "经典", "传统"]),
    ("绿色", 2601568,  ["自组织", "组织", "耗散", "生态"]),
    ("粉色", 15278115, ["美学", "美", "艺术", "审美"]),
    ("灰蓝", 6320523,  ["技术", "代码", "软件", "Node"]),
    ("深蓝", 2899536,  ["政治", "马克思", "习近平", "党"]),
    ("深红", 9576225,  ["马克思主义", "毛泽东"]),
    ("亮蓝", 48227,    ["智能", "意识", "涌现"]),
    ("翠绿", 3068113,  ["健康", "数据", "指标"]),
    ("淡紫", 11496133, ["脑图", "AHMM", "思维工具"]),
    ("天蓝", 40447,    ["英文", "English"]),
    ("灰色", 11184810, ["其他"]),
]

@command("set-graph-colors")
def cmd_set_graph_colors(args):
    """Auto-configure Obsidian graph.json with tag-based color groups."""
    workspace = get_workspace(args)
    knowledge_dir = workspace / "知识元"
    graph_path = workspace / ".obsidian" / "graph.json"

    if not knowledge_dir.exists():
        print(json.dumps({"error": "知识元目录不存在"}, ensure_ascii=False))
        sys.exit(1)

    # 统计所有知识元的标签频率
    from collections import Counter
    tag_counter = Counter()

    for md_file in knowledge_dir.glob("*.md"):
        content = md_file.read_text(encoding="utf-8", errors="ignore")
        # 从 frontmatter 提取 tags
        m = re.search(r'tags:\s*\[([^\]]+)\]', content)
        if m:
            tag_str = m.group(1)
            tags = [t.strip().strip('"').strip("'") for t in tag_str.split(",")]
            for t in tags:
                t = t.strip()
                if t and not t.startswith("来") and not t.startswith("#"):
                    tag_counter[t] += 1

    if not tag_counter:
        print(json.dumps({"error": "未找到任何标签"}, ensure_ascii=False))
        sys.exit(1)

    # 取前 20 个高频标签
    top_tags = [t for t, _ in tag_counter.most_common(20)]

    # 为每个标签分配颜色
    color_groups = []
    used_colors = set()

    for tag in top_tags:
        best_color = None
        best_score = 0
        # 按关键字匹配色系
        for cname, rgb, keywords in _COLOR_PALETTE:
            score = 0
            for kw in keywords:
                if kw in tag:
                    score += 1
            if score > best_score:
                best_score = score
                best_color = rgb

        if best_color is None:
            # 未匹配，使用灰色
            best_color = 11184810

        # 防止同一颜色重复太多
        color_key = best_color
        if color_key in used_colors:
            continue  # 跳过完全重复的颜色

        used_colors.add(color_key)
        color_groups.append({
            "query": f"tag:#{tag}",
            "color": {"a": 1, "rgb": best_color}
        })

    # 构建 graph.json
    graph_config = {
        "collapse-filter": False,
        "search": "",
        "showTags": True,
        "showAttachments": False,
        "hideUnresolved": False,
        "showOrphans": False,
        "collapse-color-groups": False,
        "colorGroups": color_groups,
        "collapse-display": False,
        "showArrow": False,
        "textFadeMultiplier": -2.5,
        "nodeSizeMultiplier": 1.68,
        "lineSizeMultiplier": 0.97,
        "collapse-forces": True,
        "centerStrength": 0.52,
        "repelStrength": 10,
        "linkStrength": 1,
        "linkDistance": 250,
        "scale": 0.26,
        "close": False
    }

    # 写入
    graph_path.parent.mkdir(parents=True, exist_ok=True)
    graph_path.write_text(json.dumps(graph_config, ensure_ascii=False, indent=2), encoding="utf-8")

    result = {
        "ok": True,
        "color_groups": len(color_groups),
        "tags_scanned": len(tag_counter),
        "file": str(graph_path.relative_to(workspace)),
        "message": "配置完成！请重启 Obsidian 或刷新关系图谱面板以使颜色生效。"
    }

    print(json.dumps(result, ensure_ascii=False))

    # 额外在 stderr 输出人类可读的提示
    print("\n🎨 关系图谱颜色已配置！请重启 Obsidian 或刷新关系图谱面板以使颜色生效。", file=sys.stderr)


# ─── 子命令: audit ───

@command("audit")
def cmd_audit(args):
    """Audit knowledge base quality: unit count, link density, resource attachment rate."""
    workspace = get_workspace(args)
    knowledge_dir = workspace / "知识元"
    resource_base = workspace / "图片及其他资源"
    manifest = load_json(get_manifest_path(workspace))
    source_dir = workspace / "原始文件"

    result = {
        "knowledge_units": {"total": 0, "new_in_session": 0},
        "link_density": {"with_outgoing_links": 0, "total": 0, "avg_links": 0.0, "zero_link_units": []},
        "resource_stats": {"total_ctx": 0, "with_belongs_to": 0},
        "unit_chunk_ratio": 0.0,
        "warnings": [],
    }

    # ── 知识元统计 ──
    if knowledge_dir.exists():
        unit_files = sorted(knowledge_dir.glob("*.md"))
        result["knowledge_units"]["total"] = len(unit_files)

        link_total = 0
        for uf in unit_files:
            content = uf.read_text(encoding="utf-8", errors="ignore")
            # 统计 [[出链]]
            links = re.findall(r'\[\[([^\]]+)\]\]', content)
            # 排除 resources 中的 .ctx 引用
            info_links = [l for l in links if not l.endswith('.ctx') and not l.startswith('KD-')]
            if info_links:
                result["link_density"]["with_outgoing_links"] += 1
                link_total += len(info_links)
            else:
                result["link_density"]["zero_link_units"].append(uf.stem)

        if len(unit_files) > 0:
            result["link_density"]["total"] = len(unit_files)
            result["link_density"]["avg_links"] = round(link_total / len(unit_files), 2)

    # 资源统计
    total_ctx = 0
    with_bt = 0
    if resource_base.exists():
        for subdir in ["images", "videos", "audios", "others"]:
            d = resource_base / subdir
            if not d.exists():
                continue
            for f in d.glob("*.ctx"):
                total_ctx += 1
                bt_pattern = re.compile(r"belongs_to:\s*\n\s+-")
                content = f.read_text(encoding="utf-8", errors="ignore")
                if bt_pattern.search(content):
                    with_bt += 1
    result["resource_stats"]["total_ctx"] = total_ctx
    result["resource_stats"]["with_belongs_to"] = with_bt

    # ── chunk 比例 ──
    processed_docs = len(manifest.get("processed", []))
    if processed_docs > 0:
        unit_chunk_approx = result["knowledge_units"]["total"] / processed_docs
        result["unit_chunk_ratio"] = round(unit_chunk_approx, 1)

    # ── 生成警告 ──
    if result["link_density"]["zero_link_units"]:
        result["warnings"].append(f"{len(result['link_density']['zero_link_units'])} 个知识元零关联")
    total_units = result["knowledge_units"]["total"]
    # 根据 cGranularity 设定比例阈值
    settings_path = workspace / "系统设置" / "project_settings.json"
    _raw_g = 0.5
    if settings_path.exists():
        try:
            _data = json.loads(settings_path.read_text(encoding="utf-8"))
            _raw_g = float(_data.get("weights", {}).get("cGranularity", 0.5))
        except Exception:
            pass
    _ratio_threshold = 5 if _raw_g <= 0.25 else (1.5 if _raw_g >= 0.75 else 3)
    if total_units > 0 and processed_docs > 0 and total_units / processed_docs < _ratio_threshold:
        result["warnings"].append(f"知识元总数({total_units})相对文档数({processed_docs})偏低(比例={total_units/processed_docs:.1f})，低于档位阈值({_ratio_threshold})")

    print(json.dumps(result, ensure_ascii=False))


# ─── 子命令: verify-units ───

@command("verify-units")
def cmd_verify_units(args):
    """验证文档的知识元产出是否符合参数设定（参数感知检查）。

    用法：
      ahkb.py verify-units <文件路径>                 # 检查单个文件
      ahkb.py verify-units --all                      # 检查全部已处理文件
      ahkb.py verify-units --pending                  # 检查待处理文件（chunked 但未标记）

    自动读取 project_settings.json 中的 cGranularity 设定阈值：
      - cGranularity ≤ 0.25 → 精细模式（每文档 ≥3 单元，chunk:单元 ≤ 3:1）
      - cGranularity 0.26~0.74 → 均衡模式（每文档 ≥2 单元，chunk:单元 ≤ 5:1）
      - cGranularity ≥ 0.75 → 粗粒模式（每文档 ≥1 单元，chunk:单元 ≤ 8:1）
    """
    workspace = get_workspace(args)

    # ── 读取参数设置 ──
    settings_path = workspace / "系统设置" / "project_settings.json"
    weights = {"cGranularity": 0.5, "cTextAmount": 0.5}
    if settings_path.exists():
        try:
            data = json.loads(settings_path.read_text(encoding="utf-8"))
            weights = data.get("weights", weights)
        except Exception:
            pass

    # ── 确定 cGranularity 档位 ──
    raw_g = float(weights.get("cGranularity", 0.5))
    if raw_g <= 0.25:
        g_level = 0
        g_name = "精细模式"
    elif raw_g >= 0.75:
        g_level = 2
        g_name = "粗粒模式"
    else:
        g_level = 1
        g_name = "均衡模式"

    # ── 三档阈值表 ──
    THRESHOLDS = {
        0: {"min_units_per_doc": 3, "max_chunk_per_unit": 3, "min_chars": 50,  "min_kb_ratio": 5},
        1: {"min_units_per_doc": 2, "max_chunk_per_unit": 5, "min_chars": 60,  "min_kb_ratio": 3},
        2: {"min_units_per_doc": 1, "max_chunk_per_unit": 8, "min_chars": 120, "min_kb_ratio": 1.5},
    }
    thr = THRESHOLDS[g_level]

    # ── 确定检查目标 ──
    targets = []
    if "--all" in args:
        manifest = load_json(get_manifest_path(workspace))
        targets = [p["source_path"] for p in manifest.get("processed", [])]
    elif "--pending" in args:
        from ahkb_chunks import load_chunk_index
        chunk_idx = load_chunk_index(workspace)
        manifest = load_json(get_manifest_path(workspace))
        processed_paths = {p["source_path"] for p in manifest.get("processed", [])}
        targets = [f for f in chunk_idx.get("files", []) if f not in processed_paths]
    elif len(args) >= 1 and not args[0].startswith("--"):
        filepath = Path(args[0])
        if not filepath.is_absolute():
            filepath = workspace / filepath
        try:
            targets.append(str(filepath.relative_to(workspace)).replace("\\", "/"))
        except ValueError:
            targets.append(str(filepath).replace("\\", "/"))
    else:
        manifest = load_json(get_manifest_path(workspace))
        targets = [p["source_path"] for p in manifest.get("processed", [])]

    if not targets:
        print(json.dumps({"ok": True, "status": "pass",
                          "message": "没有需要检查的文件", "results": []},
                         ensure_ascii=False))
        return

    # ── 逐个检查 ──
    results = []
    overall_status = "pass"
    knowledge_dir = workspace / "知识元"

    for src_path in targets:
        # 统计匹配的知识元
        matched_units = []
        if knowledge_dir.exists():
            for f in sorted(knowledge_dir.glob("*.md")):
                try:
                    fm = parse_frontmatter(f.read_text(encoding="utf-8"))
                    source = (fm.get("source", "") or "").replace("\\", "/").strip('"') if fm else ""
                    # Windows 上大小写不敏感比较
                    if source and source.lower() == src_path.lower():
                        matched_units.append(f.stem)
                except Exception:
                    continue

        unit_count = len(matched_units)

        # 统计该文件的 chunk 数（支持新旧路径两种方式）
        chunk_count = 0
        from ahkb_chunks import load_chunks_for_file, load_chunk_index
        chunk_data = load_chunks_for_file(src_path, workspace)
        chunk_count = chunk_data.get("chunk_count", 0)
        if chunk_count == 0:
            # 尝试按文件名匹配（文件可能已被移入 原始文件/）
            chunk_idx = load_chunk_index(workspace)
            for cf in chunk_idx.get("files", []):
                if Path(cf).name == Path(src_path).name:
                    chunk_data = load_chunks_for_file(cf, workspace)
                    chunk_count = chunk_data.get("chunk_count", 0)
                    break

        # 逐项判定
        file_warnings = []
        file_status = "pass"

        if unit_count == 0:
            file_warnings.append("知识元数为 0，未完成 [C]/[D] 步骤")
            file_status = "fail"
        elif unit_count < thr["min_units_per_doc"]:
            file_warnings.append(
                f"知识元数({unit_count}) < 档位最低要求({thr['min_units_per_doc']})"
            )
            file_status = "warning"

        if chunk_count > 0:
            ratio = chunk_count / max(unit_count, 1)
            if ratio > thr["max_chunk_per_unit"]:
                file_warnings.append(
                    f"chunk/知识元比({ratio:.1f}:1) > 档位上限({thr['max_chunk_per_unit']}:1)，"
                    f"可能遗漏了概念"
                )
                if file_status == "pass":
                    file_status = "warning"

        if file_status == "fail":
            overall_status = "fail"
        elif file_status == "warning" and overall_status != "fail":
            overall_status = "warning"

        results.append({
            "file": src_path,
            "unit_count": unit_count,
            "chunk_count": chunk_count,
            "thresholds_applied": {
                "granularity": g_name,
                "min_units_per_doc": thr["min_units_per_doc"],
                "max_chunk_per_unit": thr["max_chunk_per_unit"],
            },
            "status": file_status,
            "warnings": file_warnings,
        })

    # ── 全库整体统计 ──
    total_units = 0
    if knowledge_dir.exists():
        total_units = len(list(knowledge_dir.glob("*.md")))
    total_docs = len(targets)
    kb_ratio = total_units / max(total_docs, 1)
    if kb_ratio < thr["min_kb_ratio"]:
        overall_status = "warning" if overall_status != "fail" else "fail"

    print(json.dumps({
        "ok": True,
        "status": overall_status,
        "granularity": g_name,
        "thresholds": {
            "min_units_per_doc": thr["min_units_per_doc"],
            "max_chunk_per_unit": thr["max_chunk_per_unit"],
            "min_chars": thr["min_chars"],
            "min_kb_ratio": thr["min_kb_ratio"],
        },
        "kb_unit_count": total_units,
        "kb_doc_count": total_docs,
        "kb_unit_doc_ratio": round(kb_ratio, 2),
        "total_checked": len(results),
        "results": results,
    }, ensure_ascii=False))


# ─── 子命令: cleanup-stray-units ───

@command("cleanup-stray-units")
def cmd_cleanup_stray_units(args):
    """扫描工作空间根目录，将误放到根目录的 .md 移入知识元/。"""
    workspace = get_workspace(args)
    knowledge_dir = workspace / "知识元"
    import sys as _sys
    moved = []
    skipped = []
    if not knowledge_dir.exists():
        knowledge_dir.mkdir(parents=True, exist_ok=True)
    for f in sorted(workspace.glob("*.md")):
        if f.stem in ("AHKB知识库", "README"):
            skipped.append(str(f.relative_to(workspace)))
            continue
        target = knowledge_dir / f.name
        if target.exists():
            skipped.append(str(f.relative_to(workspace)))
            continue
        f.rename(target)
        # 移动后清理空文件夹
        try:
            orig_parent = f.parent
            if orig_parent.resolve() != workspace.resolve():
                remaining = [x for x in orig_parent.iterdir()
                             if x.name not in ('Thumbs.db', 'desktop.ini', 'Desktop.ini')
                             and not x.name.startswith('.')]
                if not remaining:
                    orig_parent.rmdir()
        except Exception:
            pass
        moved.append(str(f.relative_to(workspace)))
    if moved:
        print(f"📦 移入知识元/: {len(moved)} 个文件", file=_sys.stderr)
        for m in moved:
            print(f"  - {m}", file=_sys.stderr)
    else:
        _sys.stderr.write(chr(10) + "  " + chr(27) + "[92m" + "✅ 无游离知识元需处理" + chr(27) + "[0m" + chr(10))

# ─── 子命令: attach-resources ───

@command("attach-resources")
def cmd_attach_resources(args):
    """调用 ahkb_crosslink 进行双向资源挂接。"""
    workspace = get_workspace(args)

    # 弹窗模式：打开 Python 图形窗口
    if _should_show_popup():
        return _run_gui_window("知识链构建引擎", "attach-resources", args)

    dry_run = args_has(args, "--dry-run")
    from ahkb_crosslink import cross_link
    result = cross_link(workspace, verbose=True, dry_run=dry_run)

    # 子进程模式：写入结果文件
    result_file = _get_result_file_path(args, workspace)
    if result_file and _POPUP_ACTIVE:
        result_file.write_text(json.dumps(result, ensure_ascii=False), encoding="utf-8")
    elif "--sig-file" in sys.argv:
        pass  # 弹窗子进程，JSON 已在窗口中展示，不重复打印
    else:
        print(json.dumps(result, ensure_ascii=False))


@command("maintain")
def cmd_maintain(args):
    """自动维护知识库：知识元↔知识元 [[链接]] + 知识元↔资源双向关联。"""
    # 直接调用 cross-link 的 GUI 弹窗（与 menu 4 统一）
    # 蜂鸣提示已移入 ahkb_crosslink.py run_task()，在 cross_link 实际完成后播放
    result = cmd_cross_link(args)
    return result


# ─── 子命令: check-maintain-status ───

@command("check-maintain-status")
def cmd_check_maintain_status(args):
    """查询 maintain / crosslink 进程的当前状态（用于 LLM 轮询）。

    这是 LLM 查询 crosslink 状态的**唯一权威来源**。
    返回统一的状态 JSON，所有菜单的轮询代码必须使用此命令。

    状态码：
        "not_started" — maintain 从未启动
        "starting"     — 子进程已启动，但尚未写入进度文件
        "running"      — 正在运行，返回进度详情
        "complete"     — 已完成，写入 _maintain_completed.marker 标记
        "cancelled"    — 用户手动关闭了 GUI 窗口
        "error"        — 子进程异常退出（PID 文件存在但进程已死）
        "stale"        — 上一个任务的残留状态文件（PID 和进程都不在），已自动清理

    当状态为 "complete" 时，会写入 _maintain_completed.marker 标记文件。
    """
    workspace = get_workspace(args)
    ws_path = Path(workspace).resolve()
    tm_dir = ws_path / "临时工作文件"
    result_file = tm_dir / "_crosslink_result.json"
    progress_file = tm_dir / "_crosslink_progress.json"
    pid_file = tm_dir / "_crosslink.pid"
    marker_file = tm_dir / "_maintain_completed.marker"

    # ── 1. 检查结果文件（最高优先级） ──
    if result_file.exists():
        try:
            raw = result_file.read_text(encoding="utf-8").strip()
            if raw:
                r = json.loads(raw)
                status = r.get("status", "")
                if status == "complete":
                    # 写入完成标记（供 build-graph 交叉检查）
                    marker_data = {
                        "completed_at": datetime.datetime.now().isoformat(timespec="seconds"),
                        "units_updated": r.get("units_updated", 0),
                        "ctx_updated": r.get("ctx_updated", 0),
                        "matches_found": r.get("matches_found", 0),
                        "elapsed_seconds": r.get("elapsed_seconds", 0),
                    }
                    marker_file.write_text(json.dumps(marker_data, ensure_ascii=False), encoding="utf-8")
                    # 结果文件不再需要，清理
                    _trash_file(result_file, ws_path)
                    if progress_file.exists():
                        _trash_file(progress_file, ws_path)
                    if pid_file.exists():
                        _trash_file(pid_file, ws_path)
                    print(json.dumps({
                        "status": "complete",
                        "result": r,
                        "_display": (
                            f"[完成] 更新 {r.get('units_updated', 0)} 知识元, "
                            f"{r.get('ctx_updated', 0)} 资源, "
                            f"{r.get('matches_found', 0)} 匹配, "
                            f"耗时 {r.get('elapsed_seconds', 0) // 60} 分钟"
                        ),
                    }, ensure_ascii=False))
                    return
                elif status == "cancelled":
                    # 用户手动关闭，清理状态文件
                    _trash_file(result_file, ws_path)
                    _trash_file(progress_file, ws_path) if progress_file.exists() else None
                    _trash_file(pid_file, ws_path) if pid_file.exists() else None
                    print(json.dumps({
                        "status": "cancelled",
                        "_display": "[终止] 用户手动关闭了窗口",
                    }, ensure_ascii=False))
                    return
                else:
                    # 未知状态，当作异常处理，清理后返回
                    _trash_file(result_file, ws_path)
                    print(json.dumps({
                        "status": "error",
                        "detail": f"未知的结果状态: {status}",
                        "raw_status": status,
                        "_display": f"[异常] 未知状态: {status}",
                    }, ensure_ascii=False))
                    return
        except (json.JSONDecodeError, ValueError):
            # JSON 解析失败 → 文件损坏，清理
            _trash_file(result_file, ws_path)

    # ── 2. 检查 PID 是否存活 ──
    pid_alive = False
    pid_value = 0
    if pid_file.exists():
        try:
            raw = pid_file.read_text(encoding="utf-8").strip()
            if raw:
                pid_data = json.loads(raw)
                pid_value = pid_data.get("pid", 0)
                pid_alive = _is_pid_alive(pid_value)
        except (json.JSONDecodeError, ValueError):
            pass

    # ── 3. 检查进度文件 ──
    if progress_file.exists():
        try:
            raw = progress_file.read_text(encoding="utf-8").strip()
            if raw:
                p = json.loads(raw)
                if pid_alive:
                    # 进程正常 → 返回进度
                    print(json.dumps({
                        "status": "running",
                        "progress": p,
                        "done": p.get("done", 0),
                        "total": p.get("total", 1),
                        "phase": p.get("phase", ""),
                        "phase_name": p.get("phase_name", ""),
                        "elapsed_seconds": p.get("elapsed_seconds", 0),
                        "_display": (
                            f"[{p.get('done', 0)}/{p.get('total', 1)}] "
                            f"{p.get('phase_name', '')}，"
                            f"已用时 {p.get('elapsed_seconds', 0) // 60} 分钟"
                        ),
                    }, ensure_ascii=False))
                    return
                else:
                    # 进度文件存在但进程已死 → 异常退出
                    _trash_file(progress_file, ws_path)
                    _trash_file(pid_file, ws_path) if pid_file.exists() else None
                    print(json.dumps({
                        "status": "error",
                        "detail": "crosslink 进程异常退出（进度文件存在但 PID 已死）",
                        "_display": "[异常] crosslink 进程异常退出",
                    }, ensure_ascii=False))
                    return
        except (json.JSONDecodeError, ValueError):
            _trash_file(progress_file, ws_path)

    # ── 4. PID 存活但无进度文件 → 正在初始化 ──
    if pid_alive:
        print(json.dumps({
            "status": "starting",
            "pid": pid_value,
            "_display": "[等待] 进程正在初始化...",
        }, ensure_ascii=False))
        return

    # ── 5. PID 文件存在（已死）但无进度/结果 → 僵尸残留，清理 ──
    if pid_file.exists():
        _trash_file(pid_file, ws_path)
        print(json.dumps({
            "status": "stale",
            "detail": "检测到上一个任务的残留状态文件，已自动清理",
            "_display": "[清理] 上一个任务的残留状态已清除",
        }, ensure_ascii=False))
        return

    # ── 6. 没有任何线索 → 从未启动 ──
    print(json.dumps({
        "status": "not_started",
        "_display": "[未启动] maintain 从未运行",
    }, ensure_ascii=False))


# ─── 子命令: manage-resources ───

@command("manage-resources")
def cmd_manage_resources(args):
    """启动资源管理器 GUI（detached 子进程，与 maintain 方式一致）。"""
    workspace = get_workspace(args)
    ws_path = Path(workspace).resolve()
    tm_dir = ws_path / "临时工作文件"
    tm_dir.mkdir(parents=True, exist_ok=True)
    result_file = tm_dir / "_resource_manager_result.json"

    # 清理旧状态文件
    _trash_file(result_file, ws_path)

    script = Path(__file__).resolve().parent / "ahkb_manage_resources.py"
    # 直接使用 sys.executable（python.exe），不换 pythonw.exe
    cmd = [sys.executable, str(script), "--workspace", str(ws_path)]

    # 跨平台解绑启动（与 maintain 完全相同的方式）
    if platform.system() == "Windows":
        CREATE_NEW_PROCESS_GROUP = 0x00000200
        proc = subprocess.Popen(
            cmd,
            creationflags=CREATE_NEW_PROCESS_GROUP,
            close_fds=True,
            stdin=subprocess.DEVNULL,
            stdout=subprocess.DEVNULL,
            stderr=subprocess.DEVNULL
        )
    else:
        proc = subprocess.Popen(
            cmd,
            start_new_session=True,
            close_fds=True,
            stdin=subprocess.DEVNULL,
            stdout=subprocess.DEVNULL,
            stderr=subprocess.DEVNULL
        )
    print(json.dumps({
        "ok": True,
        "status": "detached",
        "result_file": str(result_file),
        "message": "资源管理器窗口已启动，操作完成后关闭窗口即可"
    }, ensure_ascii=False))


# ─── 入口 ───

@command("set-kb-name")
def cmd_set_kb_name(args):
    """设置知识库名称。用法：ahkb.py set-kb-name <名称> [--workspace <path>]"""
    workspace = get_workspace(args)
    if len(args) < 1 or args[0].startswith("--"):
        print(json.dumps({"error": "Usage: ahkb.py set-kb-name <知识库名称> [--workspace <path>]"}))
        sys.exit(1)
    name = args[0]
    set_kb_name(workspace, name)
    print(json.dumps({"ok": True, "kb_name": name, "saved_to": str(workspace / "系统设置" / "project_settings.json")}, ensure_ascii=False))


# ─── 子命令: create-root-node ───

@command("create-root-node")
def cmd_create_root_node(args):
    """创建根节点 .md 文件（仅含名称框架，无知识元数据）。用法：ahkb.py create-root-node [--workspace <path>]

    在重建知识库时，设置完名称后立即调用此命令。
    后续 maintain 完成后，由 build-graph + AI 手写入完整数据。"""
    workspace = get_workspace(args)
    settings_path = workspace / "系统设置" / "project_settings.json"
    kb_name = "知识库"
    try:
        settings = json.loads(settings_path.read_text(encoding="utf-8"))
        kb_name = settings.get("kb_name", "知识库")
    except Exception:
        pass
    root_name = f"{kb_name}(根)"
    root_path = workspace / f"{root_name}.md"
    today = datetime.datetime.now().strftime("%Y-%m-%d")
    content = "---\n"
    content += "tags: [知识库, 根节点]\n"
    content += "root_node: true\n"
    content += f"created: {today}\n"
    content += f"updated: {today}\n"
    content += "---\n\n"
    content += f"# {kb_name} 知识库\n\n"
    content += "> 基于 AHKB-CPS v0.1.0 构建的全息知识库\n\n"
    content += "## 📖 知识库概述\n\n"
    content += "（待补充）\n\n"
    content += "## 📊 统计信息\n\n"
    content += "| 项目 | 数量 |\n"
    content += "|------|:----:|\n"
    content += "| 知识元 | 0 个 |\n"
    content += f"| 最后更新 | {today} |\n\n"
    content += "## 🗺️ 知识地图\n\n"
    content += f"> [打开交互式知识地图]({kb_name}-知识地图.html)\n"
    root_path.write_text(content, encoding="utf-8")
    print(json.dumps({"ok": True, "root_node": root_name, "path": str(root_path)}, ensure_ascii=False))

def _get_sig_file_from_args():
    """从 sys.argv 提取 --sig-file 的值（弹窗子进程用）"""
    for i, a in enumerate(sys.argv):
        if a == "--sig-file" and i + 1 < len(sys.argv):
            return sys.argv[i + 1]
    return None

def main():
    # 修复 Windows 终端 GBK 编码导致中文乱码的问题
    import platform as _platform
    if _platform.system() == "Windows":
        try:
            import os as _os
            _os.system("chcp 65001 >nul")
        except Exception:
            pass
    for _s in (sys.stdout, sys.stderr):
        try:
            _s.reconfigure(encoding='utf-8')
        except Exception:
            pass

    if len(sys.argv) < 2:
        print(__doc__)
        sys.exit(1)

    cmd = sys.argv[1]
    if cmd in ("-h", "--help"):
        print(__doc__)
        return

    script_dir = Path(__file__).parent
    if str(script_dir) not in sys.path:
        sys.path.insert(0, str(script_dir))

    # ── 按需环境检查：仅在 init / check-env 时检查 ──
    # 日常命令不检查，如果缺库在执行时自然会报 ImportError
    # 用户可随时运行 check-env 查看/修复环境状态

    if cmd in COMMANDS:
        try:
            COMMANDS[cmd](sys.argv[2:])
        except Exception as e:
            print(f"Error: {e}")
            import traceback
            traceback.print_exc()
    else:
        print(f"Unknown command: {cmd}")
        print(f"Available: {', '.join(sorted(COMMANDS.keys()))}")

    # ─── 弹窗子进程收尾 ───
    # 如果指定了 --sig-file，说明是弹窗子进程
    # 创建信号文件通知父进程
    sig_file = _get_sig_file_from_args()
    if sig_file:
        try:
            open(sig_file, 'w').close()
        except Exception:
            pass
        import time as _time
        C = chr(27) + "[96m"  # cyan
        R = chr(27) + "[0m"
        sys.stderr.write(C + "\n ✔ 知识库已更新完毕。\n" + R)
        if "--auto-close" in sys.argv:
            sys.stderr.write(C + "\n ✖ 窗口将于 10 秒后自动关闭，您也可以手动关闭。\n" + R)
            _time.sleep(10)
        else:
            sys.stderr.write(C + "\n 请检查输出后按 Enter 键关闭窗口...\n" + R)
            try:
                input()
            except EOFError:
                _time.sleep(3)


if __name__ == "__main__":
    if "--version" in sys.argv or "-v" in sys.argv:
        print(json.dumps({"skill": "ahkb-cps", "version": VERSION}))
        sys.exit(0)
    main()

