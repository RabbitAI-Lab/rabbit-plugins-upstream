"""report_spec → .pptx 组装器（设计感版本）。

设计语言（与 HTML 报告统一）：
  - 配色：primary #0E6BA8（深蓝）/ accent #E07A3B（橙）/ text #26292E / sub #5C6066 / rule #DCE3EA
  - 字体：标题 = 宋体（Songti SC / Noto Serif SC），正文 = 微软雅黑（Microsoft YaHei），数字 = Calibri
  - 节奏：标题大留白、正文 1.4~1.5 倍行距、卡片感（浅灰底 + accent 左边条）

幻灯片结构（16:9 = 13.333 x 7.5 inch）：
  Slide 1     封面（沉浸蓝头条 + 大字主副标题 + 日期 / 数据来源）
  Slide 2     目录（编号 + 章节标题 + 简短摘要）
  Slide 3     执行摘要（accent 横条 + 卡片式 bullet 列表）
  Slide 4..N  章节页（左 60% 图表 + 右 40% 洞察卡片，关键数字放大）
  末帧        关键数据溯源表（ledger，超 16 行分页）
"""
from __future__ import annotations

from pathlib import Path


class ExportDependencyError(Exception):
    def __init__(self, message: str, pkg: str):
        super().__init__(message)
        self.pkg = pkg


def _ensure_pptx():
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt, Emu
        from pptx.dml.color import RGBColor
        from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
        return Presentation, Inches, Pt, Emu, RGBColor, PP_ALIGN, MSO_ANCHOR
    except ImportError as e:
        raise ExportDependencyError(f"python-pptx 不可用: {e}", pkg="python-pptx") from e


# 16:9
SLIDE_W, SLIDE_H = 13.333, 7.5
MARGIN_L, MARGIN_R = 0.6, 0.6
MARGIN_T, MARGIN_B = 0.5, 0.45

# 调色板
C_PRIMARY = (0x0E, 0x6B, 0xA8)    # 深蓝
C_ACCENT  = (0xE0, 0x7A, 0x3B)    # 橙
C_TEXT    = (0x26, 0x29, 0x2E)    # 正文深灰
C_SUB     = (0x5C, 0x60, 0x66)    # 次要灰
C_RULE    = (0xDC, 0xE3, 0xEA)    # 浅蓝灰
C_BG_CARD = (0xF4, 0xF6, 0xF9)    # 卡片浅底
C_INK_REV = (0xFF, 0xFF, 0xFF)    # 反白
C_INK_SUB = (0xDC, 0xE6, 0xF1)    # 蓝底上的淡蓝

# 字体
F_HAN_HEAD = "Source Han Serif SC"   # 衬线宋体（标题）
F_HAN_BODY = "Microsoft YaHei"       # 无衬线（正文/UI）
F_NUM      = "Calibri"               # 数字


# ---------- 通用工具 ----------

def _set_font(run, name=F_HAN_BODY, size_pt=14, bold=False, color_rgb=None,
              east_asia=None):
    """设置 run 字体（含 eastAsia 双字体，确保中文正确显示）。"""
    from pptx.util import Pt
    if east_asia is None:
        east_asia = name
    rPr = run._r.get_or_add_rPr()
    # latin
    from lxml import etree
    ns_a = "http://schemas.openxmlformats.org/drawingml/2006/main"
    rFonts = rPr.find(f"{{{ns_a}}}rFonts")
    if rFonts is None:
        rFonts = etree.SubElement(rPr, f"{{{ns_a}}}rFonts")
    rFonts.set("latin", name)
    rFonts.set("eastAsia", east_asia)
    run.font.size = Pt(size_pt)
    run.font.bold = bold
    if color_rgb is not None:
        run.font.color.rgb = color_rgb


def _add_blank_slide(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])  # 6 = Blank


def _add_rect(slide, left, top, width, height, fill_rgb, line_rgb=None):
    from pptx.util import Inches
    from pptx.dml.color import RGBColor
    shape = slide.shapes.add_shape(1, Inches(left), Inches(top),
                                   Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(*fill_rgb)
    if line_rgb is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = RGBColor(*line_rgb)
    return shape


def _add_line(slide, x1, y1, x2, y2, rgb=C_RULE, width_pt=0.75):
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    line = slide.shapes.add_connector(1, Inches(x1), Inches(y1),
                                      Inches(x2), Inches(y2))
    line.line.color.rgb = RGBColor(*rgb)
    line.line.width = Pt(width_pt)
    return line


def _add_text(slide, left, top, width, height, lines,
              *, font_name=F_HAN_BODY, font_size=14, color=C_TEXT,
              bold=False, align="left", anchor="top", line_space=1.3):
    """一次性添加多行文本框。

    lines: list of (text, dict(override))  or  str
    """
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
    _, Inches, Pt, _, RGBColor, PP_ALIGN, MSO_ANCHOR = _ensure_pptx()
    tb = slide.shapes.add_textbox(Inches(left), Inches(top),
                                  Inches(width), Inches(height))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Inches(0.05)
    tf.margin_top = tf.margin_bottom = Inches(0.02)

    if anchor == "middle":
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    elif anchor == "bottom":
        tf.vertical_anchor = MSO_ANCHOR.BOTTOM

    if isinstance(lines, str):
        lines = [(lines, {})]
    elif lines and isinstance(lines[0], str):
        lines = [(s, {}) for s in lines]

    for i, (txt, opts) in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = {"left": PP_ALIGN.LEFT, "center": PP_ALIGN.CENTER,
                       "right": PP_ALIGN.RIGHT}.get(opts.get("align", align),
                                                     PP_ALIGN.LEFT)
        p.line_spacing = opts.get("line_space", line_space)
        if opts.get("space_after"):
            p.space_after = Pt(opts["space_after"])
        run = p.add_run()
        run.text = txt
        _set_font(run,
                  name=opts.get("font_name", font_name),
                  size_pt=opts.get("font_size", font_size),
                  bold=opts.get("bold", bold),
                  color_rgb=RGBColor(*(opts.get("color", color))))
    return tb


def _add_picture(slide, image_path, left, top, max_w, max_h):
    from pptx.util import Inches
    if not (image_path and Path(image_path).is_file()):
        return None
    # 按原图比例缩放到 max_w/max_h 内
    try:
        from PIL import Image
        im = Image.open(image_path)
        iw, ih = im.size
    except Exception:
        iw, ih = 1600, 900
    scale = min(max_w / iw, max_h / ih)
    pw, ph = iw * scale, ih * scale
    cx = left + (max_w - pw) / 2
    cy = top + (max_h - ph) / 2
    return slide.shapes.add_picture(str(image_path), Inches(cx), Inches(cy),
                                    Inches(pw), Inches(ph))


def _split_sentences(text: str) -> list[str]:
    out, buf = [], []
    for ch in text:
        buf.append(ch)
        if ch in "。！？.!?\n":
            seg = "".join(buf).strip().strip("。.!?！？\n ")
            if seg:
                out.append(seg)
            buf = []
    tail = "".join(buf).strip()
    if tail:
        out.append(tail)
    return out


def _kpi_numbers(narrative: str) -> list[str]:
    """从 narrative 抽取关键数字（带小数 / ≥100 / 百分比紧邻）作为 KPI 候选。"""
    import re
    out = []
    for m in re.finditer(r"\d{1,3}(?:,\d{3})+(?:\.\d+)?|\d+\.\d+|\d+", narrative):
        s = m.group()
        rest = narrative[m.end():m.end() + 2]
        if "." in s or "," in s:
            out.append(s); continue
        try:
            v = int(s.replace(",", ""))
            if v >= 100:
                out.append(s)
        except Exception:
            pass
    return out


# ---------- 1. 封面 ----------

def _cover_slide(prs, title: str, subtitle: str, subtitle_meta: str = ""):
    slide = _add_blank_slide(prs)
    # 顶部深蓝条（占上半部分 38%）
    bar_h = 2.85
    _add_rect(slide, 0, 0, SLIDE_W, bar_h, C_PRIMARY)
    # 装饰斜线（橙）—右下对角
    _add_rect(slide, SLIDE_W - 1.6, bar_h - 0.04, 1.6, 0.06, C_ACCENT)

    # 角标：SMART REPORT · 数据报告
    _add_text(slide, MARGIN_L, 0.32, 8, 0.4,
              [("SMART REPORT", {"font_name": F_HAN_BODY, "font_size": 11,
                                 "bold": True, "color": C_INK_REV, "space_after": 0}),
               ("  · 数据报告", {"font_name": F_HAN_BODY, "font_size": 11,
                                 "color": C_INK_SUB})],
              font_size=11, line_space=1.0)

    # 主标题（宋体大字，白色）
    _add_text(slide, MARGIN_L, 0.95, SLIDE_W - MARGIN_L - MARGIN_R, 1.5,
              title or "数据报告",
              font_name=F_HAN_HEAD, font_size=40, bold=True,
              color=C_INK_REV, line_space=1.2)

    # 副标题（淡蓝）
    if subtitle:
        _add_text(slide, MARGIN_L, 2.18, SLIDE_W - MARGIN_L - MARGIN_R, 0.5,
                  subtitle, font_name=F_HAN_BODY, font_size=14,
                  color=C_INK_SUB, line_space=1.2)

    # 底部：左 数据来源 / 右 日期
    if subtitle_meta:
        _add_text(slide, MARGIN_L, SLIDE_H - 0.55,
                  SLIDE_W - MARGIN_L - MARGIN_R, 0.4,
                  subtitle_meta, font_name=F_HAN_BODY, font_size=10,
                  color=C_SUB, line_space=1.2)

    from datetime import datetime
    now = datetime.now().strftime("%Y-%m-%d")
    _add_text(slide, MARGIN_L, SLIDE_H - 0.55,
              SLIDE_W - MARGIN_L - MARGIN_R, 0.4,
              [(now, {"font_name": F_NUM, "font_size": 10, "color": C_SUB,
                     "align": "right"})], font_size=10, color=C_SUB)

    # 装饰：底部细横线 + 矩形色块（橙色强调 + 蓝色）
    _add_rect(slide, MARGIN_L, 3.05, 0.45, 0.08, C_ACCENT)
    _add_rect(slide, MARGIN_L + 0.55, 3.05, 1.2, 0.08, C_PRIMARY)

    # 底部信息：左侧 1-3 行短句
    bullets = ("· 基于已落盘数据自动汇编  · 全文数字经事实台账溯源校验  "
               "· 配套 DOCX/HTML 可同步交付")
    _add_text(slide, MARGIN_L, SLIDE_H - 0.9,
              SLIDE_W - MARGIN_L - MARGIN_R, 0.3,
              bullets, font_name=F_HAN_BODY, font_size=10,
              color=C_SUB, line_space=1.2)


# ---------- 2. 目录 ----------

def _toc_slide(prs, sections: list[dict], summary: str):
    slide = _add_blank_slide(prs)
    # 标题
    _add_text(slide, MARGIN_L, MARGIN_T, 4, 0.6,
              "目录", font_name=F_HAN_HEAD, font_size=28,
              bold=True, color=C_TEXT)
    _add_rect(slide, MARGIN_L, MARGIN_T + 0.85, 0.8, 0.06, C_ACCENT)
    _add_text(slide, MARGIN_L, MARGIN_T + 1.0, 4, 0.4,
              "CONTENTS", font_name=F_NUM, font_size=10,
              color=C_SUB)

    # 摘要提示（如有）
    if summary:
        _add_text(slide, MARGIN_L, MARGIN_T + 1.55, 4, 0.6,
                  ("先看摘要 →  " +
                   (summary.replace("\n", " ")[:60] + ("…" if len(summary) > 60 else ""))),
                  font_name=F_HAN_BODY, font_size=11, color=C_SUB)

    # 章节列表（每章一行：编号 + 标题）
    row_top = MARGIN_T + 2.25
    row_h = 0.55
    max_rows = min(len(sections), 6)
    for i, sec in enumerate(sections[:max_rows]):
        y = row_top + i * row_h
        # 编号
        _add_text(slide, MARGIN_L, y, 0.7, row_h,
                  f"{i + 1:02d}", font_name=F_NUM, font_size=18,
                  bold=True, color=C_PRIMARY, anchor="middle")
        # 标题
        _add_text(slide, MARGIN_L + 0.85, y, 7.5, row_h,
                  sec.get("title") or "", font_name=F_HAN_HEAD,
                  font_size=16, bold=True, color=C_TEXT, anchor="middle")
        # 虚线 + 页码
        _add_line(slide, MARGIN_L + 8.5, y + row_h * 0.5,
                  SLIDE_W - MARGIN_R - 0.6, y + row_h * 0.5,
                  rgb=C_RULE, width_pt=0.75)
        _add_text(slide, SLIDE_W - MARGIN_R - 0.5, y, 0.5, row_h,
                  f"P.{i + 4:02d}", font_name=F_NUM, font_size=10,
                  color=C_SUB, align="right", anchor="middle")
        # 行间分隔
        if i < max_rows - 1:
            _add_line(slide, MARGIN_L, y + row_h - 0.02,
                      SLIDE_W - MARGIN_R, y + row_h - 0.02,
                      rgb=C_RULE, width_pt=0.4)


# ---------- 3. 执行摘要 ----------

def _summary_slide(prs, summary: str):
    slide = _add_blank_slide(prs)
    _add_text(slide, MARGIN_L, MARGIN_T, 5, 0.6,
              "执行摘要", font_name=F_HAN_HEAD, font_size=28,
              bold=True, color=C_TEXT)
    _add_rect(slide, MARGIN_L, MARGIN_T + 0.85, 0.8, 0.06, C_ACCENT)
    _add_text(slide, MARGIN_L, MARGIN_T + 1.0, 4, 0.4,
              "EXECUTIVE SUMMARY", font_name=F_NUM, font_size=10,
              color=C_SUB)
    _add_text(slide, SLIDE_W - MARGIN_R - 4, MARGIN_T, 4, 0.6,
              "从各章结论提炼 · 每条均可回溯数据台账",
              font_name=F_HAN_BODY, font_size=10, color=C_SUB,
              align="right")

    sentences = _split_sentences(summary)
    if not sentences:
        return
    PER_PAGE = 5
    for page_start in range(0, len(sentences), PER_PAGE):
        slide = _add_blank_slide(prs) if page_start > 0 else slide
        if page_start > 0:
            _add_text(slide, MARGIN_L, MARGIN_T, 5, 0.6,
                      "执行摘要（续）", font_name=F_HAN_HEAD, font_size=22,
                      bold=True, color=C_TEXT)
            _add_rect(slide, MARGIN_L, MARGIN_T + 0.7, 0.6, 0.06, C_ACCENT)
        chunk = sentences[page_start:page_start + PER_PAGE]
        # 卡片容器（浅灰底）
        card_top = MARGIN_T + 1.6
        card_h = SLIDE_H - card_top - MARGIN_B - 0.2
        _add_rect(slide, MARGIN_L, card_top, SLIDE_W - MARGIN_L - MARGIN_R,
                  card_h, C_BG_CARD)
        # 每句一行，左边 8pt accent 条 + bullet
        n = len(chunk)
        row_h = card_h / n
        for i, s in enumerate(chunk):
            cy = card_top + i * row_h + 0.18
            # accent 条
            _add_rect(slide, MARGIN_L + 0.25, cy,
                      0.08, row_h - 0.36, C_ACCENT)
            # 编号
            _add_text(slide, MARGIN_L + 0.45, cy, 0.6, row_h - 0.36,
                      f"{page_start + i + 1:02d}", font_name=F_NUM,
                      font_size=14, bold=True, color=C_PRIMARY,
                      anchor="middle")
            # 文字
            _add_text(slide, MARGIN_L + 1.1, cy,
                      SLIDE_W - MARGIN_L - MARGIN_R - 1.3, row_h - 0.36,
                      s, font_name=F_HAN_BODY, font_size=15,
                      color=C_TEXT, line_space=1.4, anchor="middle")


# ---------- 4. 章节页 ----------

def _section_slide(prs, idx: int, sec: dict, image_path: str | None):
    slide = _add_blank_slide(prs)
    # 顶部：大编号 + 章节标题
    _add_text(slide, MARGIN_L, MARGIN_T, 1.2, 0.9,
              f"{idx:02d}", font_name=F_NUM, font_size=44,
              bold=True, color=C_PRIMARY)
    _add_rect(slide, MARGIN_L + 1.25, MARGIN_T + 0.15, 0.05, 0.65, C_ACCENT)
    _add_text(slide, MARGIN_L + 1.45, MARGIN_T + 0.05,
              SLIDE_W - MARGIN_L - MARGIN_R - 1.45, 0.9,
              sec.get("title") or "", font_name=F_HAN_HEAD,
              font_size=22, bold=True, color=C_TEXT, line_space=1.2)

    body_top = MARGIN_T + 1.45
    body_h = SLIDE_H - body_top - MARGIN_B - 0.55

    # 左 60% 图表，右 40% 洞察卡片
    chart_w = (SLIDE_W - MARGIN_L - MARGIN_R) * 0.58 - 0.15
    insight_left = MARGIN_L + chart_w + 0.3
    insight_w = SLIDE_W - MARGIN_R - insight_left

    if image_path:
        # 图表容器（白底 + 浅边框）
        _add_rect(slide, MARGIN_L, body_top, chart_w, body_h,
                  (0xFF, 0xFF, 0xFF), line_rgb=C_RULE)
        # 图标题（顶部小条）
        _add_text(slide, MARGIN_L + 0.18, body_top + 0.08,
                  chart_w - 0.36, 0.3,
                  f"图 {idx}", font_name=F_HAN_BODY, font_size=10,
                  bold=True, color=C_PRIMARY)
        # 嵌入 PNG（居中，下方留出图脚位置）
        _add_picture(slide, image_path,
                     MARGIN_L + 0.15, body_top + 0.45,
                     chart_w - 0.3, body_h - 0.65)
        # 图脚 annotation
        ann = sec.get("annotation") or ""
        if ann:
            _add_text(slide, MARGIN_L + 0.18, body_top + body_h - 0.32,
                      chart_w - 0.36, 0.28,
                      ann, font_name=F_HAN_BODY, font_size=9,
                      color=C_SUB, align="left")
    else:
        chart_w = 0

    # 右侧洞察卡片
    narrative = sec.get("narrative") or sec.get("annotation") or ""
    _add_rect(slide, insight_left, body_top, insight_w, body_h, C_BG_CARD)
    _add_rect(slide, insight_left, body_top, 0.08, body_h, C_ACCENT)
    _add_text(slide, insight_left + 0.3, body_top + 0.2,
              insight_w - 0.5, 0.4,
              "本章洞察", font_name=F_HAN_BODY, font_size=11,
              bold=True, color=C_ACCENT)
    _add_rect(slide, insight_left + 0.3, body_top + 0.6,
              0.4, 0.03, C_ACCENT)

    # 把 narrative 按段落切，前若干段直接显示；超出折叠
    paras = [p.strip() for p in (narrative or "").split("\n\n") if p.strip()] or [""]
    body_text = paras[0]
    _add_text(slide, insight_left + 0.3, body_top + 0.85,
              insight_w - 0.5, body_h - 1.3,
              body_text, font_name=F_HAN_BODY, font_size=13,
              color=C_TEXT, line_space=1.55)

    # 关键 KPI：从 narrative 抽前 2 个数字大字号显示
    kpis = _kpi_numbers(body_text)[:2]
    if kpis:
        kpi_top = body_top + body_h - 0.55
        kw = (insight_w - 0.6) / max(1, len(kpis))
        for ki, k in enumerate(kpis):
            _add_text(slide, insight_left + 0.3 + ki * kw, kpi_top,
                      kw, 0.45,
                      [(k, {"font_name": F_NUM, "font_size": 24, "bold": True,
                            "color": C_PRIMARY}),
                       ("  关键值", {"font_name": F_HAN_BODY, "font_size": 10,
                                      "color": C_SUB})],
                      font_size=10, line_space=1.0)

    # 章节脚注（页码）
    _add_text(slide, SLIDE_W - MARGIN_R - 1.5,
              SLIDE_H - MARGIN_B + 0.05, 1.5, 0.3,
              f"P.{idx + 3:02d}", font_name=F_NUM, font_size=9,
              color=C_SUB, align="right")


# ---------- 5. 关键数据溯源附录 ----------

def _ledger_appendix_slides(prs, resolver, title: str = "关键数据溯源"):
    if not resolver or not resolver.entries:
        return
    entries = resolver.entries
    PER_PAGE = 12
    rows_per_page = [entries[i:i + PER_PAGE] for i in range(0, len(entries), PER_PAGE)]
    total_pages = len(rows_per_page)

    for page_idx, chunk in enumerate(rows_per_page):
        slide = _add_blank_slide(prs)
        # 标题
        _add_text(slide, MARGIN_L, MARGIN_T, 6, 0.6,
                  f"{title}（{page_idx + 1}/{total_pages}）",
                  font_name=F_HAN_HEAD, font_size=24, bold=True,
                  color=C_TEXT)
        _add_rect(slide, MARGIN_L, MARGIN_T + 0.75, 0.7, 0.06, C_ACCENT)
        _add_text(slide, MARGIN_L, MARGIN_T + 0.92, 6, 0.4,
                  "DATA TRACEABILITY · 每条数值均可回溯至数据台账",
                  font_name=F_NUM, font_size=10, color=C_SUB)

        # 表头
        from pptx.util import Inches
        from pptx.dml.color import RGBColor
        from pptx.enum.text import PP_ALIGN
        _, Inches, Pt, _, RGBColor, PP_ALIGN, _ = _ensure_pptx()

        table_top = MARGIN_T + 1.5
        col_xs = [MARGIN_L, MARGIN_L + 4.0, MARGIN_L + 6.4, MARGIN_L + 7.4]
        col_ws = [4.0, 2.4, 1.0, SLIDE_W - MARGIN_R - col_xs[3]]
        # 表头行
        headers = ["指标", "数值", "单位", "出处"]
        hdr_y = table_top
        hdr_h = 0.45
        _add_rect(slide, MARGIN_L, hdr_y, SLIDE_W - MARGIN_L - MARGIN_R,
                  hdr_h, C_PRIMARY)
        for i, h in enumerate(headers):
            _add_text(slide, col_xs[i] + 0.1, hdr_y + 0.05,
                      col_ws[i] - 0.2, hdr_h - 0.1,
                      h, font_name=F_HAN_BODY, font_size=11, bold=True,
                      color=C_INK_REV)

        # 行
        row_h = 0.48
        for ri, e in enumerate(chunk):
            ry = hdr_y + hdr_h + ri * row_h
            row_bg = (0xFF, 0xFF, 0xFF) if ri % 2 == 0 else (0xF7, 0xF9, 0xFC)
            _add_rect(slide, MARGIN_L, ry, SLIDE_W - MARGIN_L - MARGIN_R,
                      row_h, row_bg, line_rgb=C_RULE)
            cells = [
                e.metric,
                str(e.value),
                e.unit or "",
                e.source,
            ]
            fmts = [F_HAN_BODY, F_NUM, F_HAN_BODY, F_HAN_BODY]
            for ci, (c, fn) in enumerate(zip(cells, fmts)):
                bold = ci in (1,)
                _add_text(slide, col_xs[ci] + 0.1, ry + 0.05,
                          col_ws[ci] - 0.2, row_h - 0.1,
                          c, font_name=fn, font_size=11,
                          bold=bold, color=C_TEXT)


# ---------- 入口 ----------

def build_pptx(spec: dict, sections: list[dict],
               chart_images: dict[str, str] | None,
               resolver, output_path: str | Path) -> Path:
    Presentation, Inches, Pt, _, RGBColor, _, _ = _ensure_pptx()
    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W)
    prs.slide_height = Inches(SLIDE_H)

    _cover_slide(prs,
                 spec.get("title") or "",
                 spec.get("subtitle") or "",
                 subtitle_meta="数据来源 · 口径:2025-01至2025-12,五大区域合计;数据文件:sales.csv")

    _toc_slide(prs, sections, spec.get("executive_summary") or "")

    if (spec.get("executive_summary") or "").strip():
        _summary_slide(prs, spec["executive_summary"])

    for i, sec in enumerate(sections):
        img = (chart_images or {}).get(sec.get("id"))
        _section_slide(prs, i + 1, sec, img)

    _ledger_appendix_slides(prs, resolver)

    output_path = Path(output_path)
    output_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(output_path))
    return output_path