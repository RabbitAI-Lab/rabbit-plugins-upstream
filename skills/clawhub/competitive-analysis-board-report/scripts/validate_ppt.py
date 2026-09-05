# -*- coding: utf-8 -*-
# PPT 校验：①旧战略表述残留 ②决策页旧口径 ③新表述覆盖 ④零越界 ⑤总页数
# 用法：python validate_ppt.py  （与本文件同目录下的 xxx.pptx 自动定位）
import os
from pptx import Presentation
from pptx.util import Emu, Inches

HERE = os.path.dirname(os.path.abspath(__file__))
PATTERNS = ("雨轩食品竞品分析与战略决策(董事会汇报版).pptx",
            "竞品分析与战略决策(董事会汇报版).pptx")
PATH = None
for p in PATTERNS:
    cand = os.path.join(HERE, p)
    if os.path.exists(cand):
        PATH = cand
        break
if PATH is None:
    PATH = os.path.join(HERE, PATTERNS[0])

prs = Presentation(PATH)
SW, SH = prs.slide_width, prs.slide_height  # EMU
EMU_IN = 914400


def collect_text(shapes, acc):
    for sh in shapes:
        if sh.has_text_frame:
            for p in sh.text_frame.paragraphs:
                t = "".join(r.text for r in p.runs)
                if t.strip():
                    acc.append(t)
        if sh.shape_type == 6:  # group
            collect_text(sh.shapes, acc)


slides_txt = []
for i, sl in enumerate(prs.slides, 1):
    acc = []
    collect_text(sl.shapes, acc)
    slides_txt.append((i, "\n".join(acc)))

# 1) leftover old terms check
old_terms = ["抓出口", "做财税", "财税筹划", "出口战略", "第一增长极（", "第一增长极"]
leftover_hits = []
for i, txt in slides_txt:
    for t in old_terms:
        if t in txt:
            leftover_hits.append((i, t))

# 2) new terms must be present somewhere
new_terms = ["全链路数字化", "制度系统化", "强数字化", "建制度", "SOP", "风控/合规三道防线", "AI 赋能"]
missing_new = [t for t in new_terms if not any(t in txt for _, txt in slides_txt)]

# 3) overflow check
overflow = []
for i, sl in enumerate(prs.slides, 1):
    for sh in sl.shapes:
        try:
            l, t, w, h = sh.left, sh.top, sh.width, sh.height
        except Exception:
            continue
        if l is None or t is None or w is None or h is None:
            continue
        r = l + w
        b = t + h
        tol = 0.02 * EMU_IN
        if l < -tol or t < -tol or r > SW + tol or b > SH + tol:
            overflow.append((i, sh.shape_type, round(l / EMU_IN, 2), round(t / EMU_IN, 2),
                             round(r / EMU_IN, 2), round(b / EMU_IN, 2)))

print("=== 1) 旧战略表述残留检查 ===")
print("残留命中(应为空):", leftover_hits if leftover_hits else "无 ✅")
# also specifically scan decision slides 2,8,9,10,11,13 for the four forbidden decision phrases
print("\n=== 2) 决策页(2/8/9/10/11/13)旧口径扫描 ===")
forbidden = ["抓出口", "做财税", "财税筹划", "出口战略"]
for i in (2, 8, 9, 10, 11, 13):
    txt = slides_txt[i - 1][1]
    bad = [f for f in forbidden if f in txt]
    print(f"  P{i}: {'含旧口径 ' + str(bad) if bad else '干净 ✅'}")

print("\n=== 3) 新战略表述覆盖检查 ===")
print("缺失新表述(应为空):", missing_new if missing_new else "无 ✅")

print("\n=== 4) 零越界检查 (slide 13.333x7.5in) ===")
print("越界形状(应为空):", overflow if overflow else "无 ✅")

print("\n=== 5) 总页数 ===", len(prs.slides._sldIdLst))
