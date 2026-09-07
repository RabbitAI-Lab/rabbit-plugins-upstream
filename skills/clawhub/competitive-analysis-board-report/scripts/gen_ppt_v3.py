# -*- coding: utf-8 -*-
"""董事会汇报版 PPT v3：决策脊柱型结构（定位→对手地图→九宫格对标→3大战略动作→节奏→监测）"""
# 本脚本为"雨轩食品"实例。适配新公司：改封面 KPI、对手地图/九宫格/品牌榜复函、P2–P13 文字、决策表，
# 复用 Page.sec/bul/tbl/box/view/note/finish 与组件 kpi_row/scale_bars/month_chart/two_col/matrix_9box/rival_map。
# 禁止手动画死坐标，所有内容走流式 y 游标（见 references/engine_notes.md）。out 已相对化，开箱可跑。
import os
import math
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

HERE = os.path.dirname(os.path.abspath(__file__))
C_DARK = RGBColor(0x1F, 0x3B, 0x63)
C_BLUE = RGBColor(0x18, 0x5F, 0xA5)
C_LIGHT = RGBColor(0xE6, 0xF1, 0xFB)
C_RED = RGBColor(0xA3, 0x2D, 0x2D)
C_AMBER = RGBColor(0xBA, 0x75, 0x17)
C_GREEN = RGBColor(0x0F, 0x6E, 0x56)
C_GRAY = RGBColor(0x5F, 0x5E, 0x5A)
C_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
C_TXT = RGBColor(0x2C, 0x2C, 0x2A)
FONT = "微软雅黑"

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]
PAGE = {"n": 0}
LOG = []


def vlen(text):
    w = 0.0
    for ch in text:
        w += 1.0 if ord(ch) > 0x2E80 else 0.55
    return w


def text_h(text, width, size, min_h=0.0):
    cw = size / 72.0
    per_line = max(1.0, (width - 0.12) / cw)
    lines = max(1, math.ceil(vlen(text) / per_line))
    return max(min_h, lines * (size * 1.42 / 72.0) + 0.045)


class Page:
    def __init__(self, title_text, size=30):
        PAGE["n"] += 1
        self.s = prs.slides.add_slide(BLANK)
        self.no = PAGE["n"]
        tb = self.s.shapes.add_textbox(Inches(0.42), Inches(0.18), Inches(12.5), Inches(0.62))
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title_text
        p.font.size = Pt(size)
        p.font.bold = True
        p.font.name = FONT
        p.font.color.rgb = C_DARK
        ln = self.s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.42), Inches(0.84), Inches(12.5), Pt(2.5))
        ln.fill.solid()
        ln.fill.fore_color.rgb = C_BLUE
        ln.line.fill.background()
        ln.shadow.inherit = False
        self.y = 0.99

    def sec(self, text, size=18, color=C_BLUE, gap=0.14, width=12.4):
        self.y += gap
        h = text_h(text, width, size)
        tb = self.s.shapes.add_textbox(Inches(0.45), Inches(self.y), Inches(width), Inches(h))
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(size)
        p.font.bold = True
        p.font.name = FONT
        p.font.color.rgb = color
        self.y += h
        return self

    def bul(self, items, size=16, width=12.3, gap=0.09, item_gap=0.055, color=C_TXT, left=0.5):
        self.y += gap
        for it in items:
            t = "• " + it
            h = text_h(t, width - 0.2, size)
            tb = self.s.shapes.add_textbox(Inches(left), Inches(self.y), Inches(width), Inches(h))
            tf = tb.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = t
            p.font.size = Pt(size)
            p.font.name = FONT
            p.font.color.rgb = color
            self.y += h + item_gap
        self.y -= item_gap
        return self

    def tbl(self, headers, rows, size=13.5, head_size=None, row_pad=0.075,
            col_w=None, left=0.45, width=12.4, gap=0.09):
        self.y += gap
        hs = head_size or size
        hs_max = text_h(headers[0], (col_w[0] if col_w else width / len(headers)) - 0.14, hs) + row_pad
        for i, htxt in enumerate(headers):
            cwid = (col_w[i] if col_w else width / len(headers)) - 0.14
            hs_max = max(hs_max, text_h(htxt, cwid, hs) + row_pad)
        row_hs = []
        for row in rows:
            rh = 0.0
            for i, cell in enumerate(row):
                cwid = (col_w[i] if col_w else width / len(rows[0])) - 0.14
                cell_h = 0.0
                for ln in str(cell).split("\n"):
                    cell_h += text_h(ln, cwid, size)
                rh = max(rh, cell_h + row_pad * 2)
            row_hs.append(rh)
        total = hs_max + sum(row_hs)
        tbl_shape = self.s.shapes.add_table(len(rows) + 1, len(headers), Inches(left),
                                            Inches(self.y), Inches(width), Inches(total))
        tbl = tbl_shape.table
        if col_w:
            for i, w in enumerate(col_w):
                tbl.columns[i].width = Inches(w)
        tbl.rows[0].height = Inches(hs_max)
        for i, htxt in enumerate(headers):
            c = tbl.cell(0, i)
            c.text = htxt
            c.fill.solid()
            c.fill.fore_color.rgb = C_DARK
            c.vertical_anchor = MSO_ANCHOR.MIDDLE
            c.margin_left = Inches(0.07)
            c.margin_top = Inches(0.015)
            c.margin_bottom = Inches(0.015)
            p = c.text_frame.paragraphs[0]
            p.font.size = Pt(hs)
            p.font.bold = True
            p.font.name = FONT
            p.font.color.rgb = C_WHITE
        for r_i, row in enumerate(rows, start=1):
            tbl.rows[r_i].height = Inches(row_hs[r_i - 1])
            for c_i, v in enumerate(row):
                c = tbl.cell(r_i, c_i)
                c.text = str(v)
                c.fill.solid()
                c.fill.fore_color.rgb = C_WHITE if r_i % 2 else C_LIGHT
                c.vertical_anchor = MSO_ANCHOR.MIDDLE
                c.margin_left = Inches(0.07)
                c.margin_top = Inches(0.015)
                c.margin_bottom = Inches(0.015)
                lines = str(v).split("\n")
                for li, ln in enumerate(lines):
                    p = c.text_frame.paragraphs[0] if li == 0 else c.text_frame.add_paragraph()
                    p.text = ln
                    p.font.size = Pt(size)
                    p.font.name = FONT
                    p.font.color.rgb = C_TXT
        self.y += total
        return self

    def box(self, height, draw_fn, gap=0.12, left=0.45, width=12.4):
        self.y += gap
        draw_fn(self.s, self.y, left, width)
        self.y += height
        return self

    def view(self, text, size=15, color=C_RED, fill=RGBColor(0xFC, 0xF0, 0xF0), align_bottom=True):
        h = text_h("◆ " + text, 12.1, size) + 0.12
        top = min(self.y + 0.18, 6.98 - h) if align_bottom else self.y + 0.18
        bx = self.s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.45), Inches(top),
                                     Inches(12.4), Inches(h))
        bx.fill.solid()
        bx.fill.fore_color.rgb = fill
        bx.line.color.rgb = color
        bx.line.width = Pt(1)
        bx.shadow.inherit = False
        tf = bx.text_frame
        tf.word_wrap = True
        tf.margin_left = Inches(0.12)
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.text = "◆ " + text
        p.font.size = Pt(size)
        p.font.bold = True
        p.font.name = FONT
        p.font.color.rgb = color
        self.y = top + h
        return self

    def note(self, text, size=13.5, gap=0.1, color=C_TXT):
        self.y += gap
        h = text_h(text, 12.3, size)
        tb = self.s.shapes.add_textbox(Inches(0.45), Inches(self.y), Inches(12.4), Inches(h))
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(size)
        p.font.name = FONT
        p.font.color.rgb = color
        self.y += h
        return self

    def finish(self, src="数据来源：国家统计局 / 海关总署 / 农业农村部 / 上市公司年报 / 政府通稿 / 竞品情报第1期"):
        tb = self.s.shapes.add_textbox(Inches(0.45), Inches(7.08), Inches(9.6), Inches(0.28))
        tf = tb.text_frame
        p = tf.paragraphs[0]
        p.text = src
        p.font.size = Pt(9)
        p.font.name = FONT
        p.font.color.rgb = C_GRAY
        tb2 = self.s.shapes.add_textbox(Inches(12.15), Inches(7.08), Inches(0.75), Inches(0.28))
        p2 = tb2.text_frame.paragraphs[0]
        p2.text = str(self.no)
        p2.font.size = Pt(10)
        p2.font.bold = True
        p2.font.name = FONT
        p2.font.color.rgb = C_DARK
        p2.alignment = PP_ALIGN.RIGHT
        LOG.append((self.no, round(self.y, 2)))
        return self.s


# ============ 复用组件 ============
def kpi_row(sl, items, top, left=0.8, width=11.8, height=1.45, size=13):
    n = len(items)
    gap = 0.16
    w = (width - gap * (n - 1)) / n
    for i, (val, label, color) in enumerate(items):
        x = left + i * (w + gap)
        box = sl.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(top), Inches(w), Inches(height))
        box.fill.solid()
        box.fill.fore_color.rgb = C_WHITE
        box.line.color.rgb = color
        box.line.width = Pt(1.5)
        box.shadow.inherit = False
        tf = box.text_frame
        tf.word_wrap = True
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.text = val
        p.font.size = Pt(25)
        p.font.bold = True
        p.font.name = FONT
        p.font.color.rgb = color
        p.alignment = PP_ALIGN.CENTER
        p2 = tf.add_paragraph()
        p2.text = label
        p2.font.size = Pt(size)
        p2.font.name = FONT
        p2.font.color.rgb = C_GRAY
        p2.alignment = PP_ALIGN.CENTER


def scale_bars(sl, top, left, width):
    data = [("皓月（牛，东北）", 144.28, C_GRAY),
            ("大庄园（羊，重叠最高）", 50.33, C_RED),
            ("恒都（牛，河南有基地）", 32.24, C_AMBER),
            ("雨轩（羊，本品）", 24, C_DARK),
            ("科尔沁（牛）", 12.5, C_BLUE),
            ("蒙都（羊）", 4.69, C_BLUE),
            ("额尔敦（羊）", 4.33, C_BLUE)]
    tv = 145
    row_h = 0.46
    label_w = 2.7
    bar_max = width - label_w - 1.5
    for i, (lab, val, color) in enumerate(data):
        y = top + i * row_h
        tb = sl.shapes.add_textbox(Inches(left), Inches(y), Inches(label_w), Inches(row_h - 0.04))
        tf = tb.text_frame
        tf.word_wrap = True
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.text = lab
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.name = FONT
        p.font.color.rgb = C_DARK
        bw = max(0.12, bar_max * val / tv)
        bar = sl.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(left + label_w), Inches(y + 0.06), Inches(bw), Inches(row_h - 0.16))
        bar.fill.solid()
        bar.fill.fore_color.rgb = color
        bar.line.fill.background()
        bar.shadow.inherit = False
        vt = sl.shapes.add_textbox(Inches(left + label_w + bw + 0.08), Inches(y), Inches(1.4), Inches(row_h - 0.04))
        vtf = vt.text_frame
        vtf.vertical_anchor = MSO_ANCHOR.MIDDLE
        vp = vtf.paragraphs[0]
        vp.text = f"{val} 亿"
        vp.font.size = Pt(13)
        vp.font.bold = True
        vp.font.name = FONT
        vp.font.color.rgb = color


def month_chart(sl, top, left, width):
    months = [("1月", 14, C_RED), ("2月", 9, C_AMBER), ("3月", 5, C_BLUE), ("4月", 4.5, C_BLUE),
              ("5月", 5, C_BLUE), ("6月", 6.5, C_BLUE), ("7月", 8, C_BLUE), ("8月", 7.5, C_BLUE),
              ("9月", 8.5, C_BLUE), ("10月", 9.5, C_AMBER), ("11月", 10.5, C_AMBER), ("12月", 12, C_RED)]
    bar_h = 1.4
    bw = 0.78
    gap = 0.24
    for i, (m, v, color) in enumerate(months):
        x = left + i * (bw + gap)
        h = bar_h * v / 14
        bar = sl.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(top + 0.28 + bar_h - h), Inches(bw), Inches(h))
        bar.fill.solid()
        bar.fill.fore_color.rgb = color
        bar.line.fill.background()
        bar.shadow.inherit = False
        pv = sl.shapes.add_textbox(Inches(x - 0.05), Inches(top + bar_h - h - 0.02), Inches(bw + 0.1), Inches(0.24))
        pvf = pv.text_frame
        pp = pvf.paragraphs[0]
        pp.text = f"{v}%"
        pp.font.size = Pt(12)
        pp.font.bold = True
        pp.font.name = FONT
        pp.font.color.rgb = color
        pp.alignment = PP_ALIGN.CENTER
        vt = sl.shapes.add_textbox(Inches(x - 0.05), Inches(top + 0.28 + bar_h + 0.04), Inches(bw + 0.1), Inches(0.24))
        vtf = vt.text_frame
        p2 = vtf.paragraphs[0]
        p2.text = m
        p2.font.size = Pt(12)
        p2.font.bold = True
        p2.font.name = FONT
        p2.font.color.rgb = C_DARK
        p2.alignment = PP_ALIGN.CENTER


def two_col(sl, top, left, width):
    cols = [
        ("✅ 可以做的（按确定性排序）", C_GREEN, RGBColor(0xE1, 0xF5, 0xEE), [
            "供应链：11 月前锁定春节档羊源（活羊 +6.8%、75% 外采，最大成本风险）；近地化 25% → 40%",
            "数字化：Q3 启动全链路数字化试点（智能排产/损耗预警），Q4 验证 ROI；GM 直管专班",
            "营销：8 月定春节档盘子，9 月锁中秋企业团购，主打真羊肉 / 0 掺假 / 出口级标准",
            "竞品：绑死巴奴 / 锅圈年度框架协议；冲击胖东来；正面研究茂源",
            "制度化：核心流程 SOP 化，建风控/合规三道防线，降低经营波动、承接渠道扩张",
        ]),
        ("❌ 不可以做的（五条红线）", C_RED, RGBColor(0xFC, 0xEB, 0xEB), [
            "不盲目扩建重资产屠宰产能 —— 设计产能 150 万只/年，伊赛 27 亿 → 破产的教训",
            "不参与线上价格战 —— 成本结构撑不住，打信任牌不打价格牌",
            "不为进胖东来牺牲毛利底线 —— 背书重要，但赔本买卖不做",
            "不重仓牛肉辅线 —— 配额耗尽后进口涨价，且雨轩核心在羊",
            "出口不单押中东 —— 美以伊冲突、霍尔木兹封锁，必须国别多元化",
            "不把数字化做成『买一套系统』的一次性工程 —— 重数据贯通与持续迭代",
        ]),
    ]
    w = 6.1
    h = 4.35
    for i, (name, color, fill, items) in enumerate(cols):
        x = left + i * (w + 0.15)
        box = sl.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(top), Inches(w), Inches(h))
        box.fill.solid()
        box.fill.fore_color.rgb = fill
        box.line.color.rgb = color
        box.line.width = Pt(2)
        box.shadow.inherit = False
        tf = box.text_frame
        tf.word_wrap = True
        tf.margin_left = Inches(0.14)
        tf.margin_top = Inches(0.08)
        p = tf.paragraphs[0]
        p.text = name
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.name = FONT
        p.font.color.rgb = color
        for it in items:
            pp = tf.add_paragraph()
            pp.text = "• " + it
            pp.font.size = Pt(13.5)
            pp.font.name = FONT
            pp.font.color.rgb = C_TXT
            pp.space_after = Pt(7)


# ============ 新增：九宫格对标矩阵 ============
def matrix_9box(sl, top, left, width):
    dims = ["规模", "品牌", "线上", "线下渠道", "价格", "品控", "出口", "利润韧性"]
    cols = ["雨轩", "全国巨头", "区域龙头", "垂直新锐"]
    # 强=2 中=1 弱=0 ；颜色：绿/黄/红
    M = {
        "规模":       [1, 2, 1, 0],
        "品牌":       [1, 2, 1, 1],
        "线上":       [0, 1, 0, 2],
        "线下渠道":   [2, 2, 1, 0],
        "价格":       [1, 1, 1, 0],
        "品控":       [2, 2, 1, 1],
        "出口":       [2, 1, 0, 0],
        "利润韧性":   [1, 2, 0, 0],
    }
    label_w = 1.85
    col_w = (width - label_w) / 4
    row_h = 0.34
    head_h = 0.42
    # 表头
    hx = left + label_w
    for j, c in enumerate(cols):
        x = hx + j * col_w
        box = sl.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(top), Inches(col_w - 0.06), Inches(head_h))
        box.fill.solid()
        box.fill.fore_color.rgb = C_DARK
        box.line.fill.background()
        box.shadow.inherit = False
        tf = box.text_frame
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.text = c
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.name = FONT
        p.font.color.rgb = C_WHITE
        p.alignment = PP_ALIGN.CENTER
    # 行
    for i, d in enumerate(dims):
        y = top + head_h + i * row_h
        # 维度标签
        lb = sl.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(left), Inches(y), Inches(label_w - 0.06), Inches(row_h - 0.04))
        lb.fill.solid()
        lb.fill.fore_color.rgb = RGBColor(0xD3, 0xD1, 0xC7)
        lb.line.fill.background()
        lb.shadow.inherit = False
        tf = lb.text_frame
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.text = d
        p.font.size = Pt(13)
        p.font.bold = True
        p.font.name = FONT
        p.font.color.rgb = C_TXT
        p.alignment = PP_ALIGN.CENTER
        for j, v in enumerate(M[d]):
            x = hx + j * col_w
            col = [RGBColor(0xFC, 0xEB, 0xEB), RGBColor(0xFA, 0xEE, 0xDA), RGBColor(0xE1, 0xF5, 0xEE)][v]
            edge = [C_RED, C_AMBER, C_GREEN][v]
            txt = ["弱", "中", "强"][v]
            cell = sl.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(col_w - 0.06), Inches(row_h - 0.04))
            cell.fill.solid()
            cell.fill.fore_color.rgb = col
            cell.line.color.rgb = edge
            cell.line.width = Pt(0.75)
            cell.shadow.inherit = False
            tf = cell.text_frame
            tf.vertical_anchor = MSO_ANCHOR.MIDDLE
            p = tf.paragraphs[0]
            p.text = txt
            p.font.size = Pt(12.5)
            p.font.bold = True
            p.font.name = FONT
            p.font.color.rgb = edge
            p.alignment = PP_ALIGN.CENTER


# ============ 新增：对手地图分层卡片 ============
def rival_map(sl, top, left, width):
    tiers = [
        ("① 全国巨头", "皓月 144 亿 / 大庄园 50 亿 / 恒都 32 亿",
         "规模与品牌碾压；大庄园与雨轩业务线全面重叠，最该深研", C_GRAY),
        ("② 区域龙头", "恒都（泌阳基地 + 北上）、科尔沁系",
         "区域渠道贴身竞争；恒都新设呼和浩特子公司=扩张信号", C_AMBER),
        ("③ 垂直新锐", "宁羴源（2024 线上 1.8 亿·抖音羊肉榜 TOP1 三年）、大希地（GMV 30 亿）、小牛凯西",
         "线上打法颠覆传统；抢家庭餐桌预算，须重点对标", C_BLUE),
        ("④ 同省同赛道", "茂源肉业（鹿邑，锅圈系清真深加工）",
         "家门口直接竞争；同在锅圈体系，须正面对标防替代", C_RED),
    ]
    w = (width - 0.45) / 4
    gap = 0.15
    for i, (t, sub, desc, color) in enumerate(tiers):
        x = left + i * (w + gap)
        box = sl.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(top), Inches(w), Inches(2.25))
        box.fill.solid()
        box.fill.fore_color.rgb = C_WHITE
        box.line.color.rgb = color
        box.line.width = Pt(2)
        box.shadow.inherit = False
        tf = box.text_frame
        tf.word_wrap = True
        tf.margin_left = Inches(0.12)
        tf.margin_top = Inches(0.08)
        p = tf.paragraphs[0]
        p.text = t
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.name = FONT
        p.font.color.rgb = color
        p2 = tf.add_paragraph()
        p2.text = sub
        p2.font.size = Pt(12)
        p2.font.bold = True
        p2.font.name = FONT
        p2.font.color.rgb = C_DARK
        p2.space_after = Pt(4)
        p3 = tf.add_paragraph()
        p3.text = desc
        p3.font.size = Pt(12)
        p3.font.name = FONT
        p3.font.color.rgb = C_GRAY


# ============ P1 封面 ============
PAGE["n"] += 1
s = prs.slides.add_slide(BLANK)
bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(2.3))
bg.fill.solid()
bg.fill.fore_color.rgb = C_DARK
bg.line.fill.background()
bg.shadow.inherit = False
tb = s.shapes.add_textbox(Inches(0.8), Inches(0.6), Inches(11.8), Inches(1.15))
tf = tb.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "雨轩食品 · 竞品分析与战略决策"
p.font.size = Pt(40)
p.font.bold = True
p.font.name = FONT
p.font.color.rgb = C_WHITE
p2 = tf.add_paragraph()
p2.text = "决策脊柱版 · 定位 → 对手地图 → 九宫格对标 → 三大战略动作"
p2.font.size = Pt(20)
p2.font.name = FONT
p2.font.color.rgb = RGBColor(0xBD, 0xD7, 0xEE)
tb = s.shapes.add_textbox(Inches(0.8), Inches(2.68), Inches(11.8), Inches(0.45))
tf = tb.text_frame
p = tf.paragraphs[0]
p.text = "新乡市雨轩清真食品股份有限公司 · 董事会汇报"
p.font.size = Pt(22)
p.font.bold = True
p.font.name = FONT
p.font.color.rgb = C_DARK
tb = s.shapes.add_textbox(Inches(0.8), Inches(3.28), Inches(11.8), Inches(0.32))
tf = tb.text_frame
p = tf.paragraphs[0]
p.text = "2026 年 8 月 30 日   总经理办公室"
p.font.size = Pt(17)
p.font.name = FONT
p.font.color.rgb = C_GRAY

kpi_row(s, [("15.7 亿", "2024 营收（2022 峰值 26 亿）\n三年连续下滑", C_DARK),
            ("3–4 倍", "领先直接竞品\n额尔敦 4.33 亿", C_BLUE),
            ("60%+", "黄河滩羊占河南\n羊肉制品市场", C_GREEN),
            ("第 7", "全国羊肉品牌榜\n线上仅个位数", C_RED)], top=3.85, height=1.55)
bx = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(5.72), Inches(11.8), Inches(1.05))
bx.fill.solid()
bx.fill.fore_color.rgb = RGBColor(0xFC, 0xF0, 0xF0)
bx.line.color.rgb = C_RED
bx.shadow.inherit = False
tf = bx.text_frame
tf.word_wrap = True
tf.vertical_anchor = MSO_ANCHOR.MIDDLE
p = tf.paragraphs[0]
p.text = "◆ 总纲：拐点已至但由供给驱动、需求仍弱 —— 锁成本、强数字化、抢旺季、建制度；以销定产、现金为王。不赌拐点，但绝不错过窗口。"
p.font.size = Pt(17)
p.font.bold = True
p.font.name = FONT
p.font.color.rgb = C_RED
p.alignment = PP_ALIGN.CENTER

# ============ P2 一页纸结论 ============
pg = Page("一页纸结论：董事会要拍板的 3 件事")
pg.sec("结论先行 —— 行业拐点是真，但由供给驱动、需求仍弱", 18, C_BLUE, gap=0.06)
pg.bul([
    "周期：2026 羊肉连年减产 + 存栏 6 年最低 + 进口三年连降 + 价格全面转正 = 拐点确认；但官方预测消费仍 -1.7%，属供给驱动型",
    "定位：雨轩 2024 营收 15.7 亿（2022 峰值 26 亿、连续下滑），约为直接竞品额尔敦（4.33 亿）的 3–4 倍；更大对手是皓月/大庄园/恒都，大庄园重叠度最高",
    "护城河：国产冷鲜 + 可溯源 + 清真 + AEO 出口 + 河南规上屠宰 85% + 已绑定巴奴/锅圈；资质含农业产业化国家重点龙头、国家级肉羊产业集群链主、中国肉类协会会长单位；短板在线上（约 3%）与全国品牌（第 7），且 2025E 净利率仅约 1.25%、盈利极薄",
], size=16, item_gap=0.08)
pg.sec("请董事会决策的三项事项", 18, C_DARK, gap=0.22)
pg.tbl(["#", "决策事项", "时限", "一句话理由"],
       [["①", "2026–2027 春节档总盘子与预算", "9.30 前", "错过 8 月定盘、10 月出礼盒两节点则春节只能吃残羹"],
        ["②", "全链路数字化与 AI 赋能专项", "Q3 启动 / Q4 验证", "从羊源到餐桌全链数据贯通，用 AI 降损耗、提排产与定价效率，是存量竞争下的核心能力"],
        ["③", "制度系统化与风控合规体系建设", "Q4 前框架落地", "把依赖个人的经验转化为 SOP/风控/合规体系，降低波动、承接专项整治与渠道扩张"]],
       size=14, col_w=[0.6, 4.6, 1.6, 5.6], gap=0.06, row_pad=0.05)
pg.view("一句话：用确定的动作抓确定的窗口 —— 锁成本、强数字化、抢旺季、建制度；三件确定性最高的事现在就能做", 15)
pg.finish()

# ============ P3 雨轩定位 ============
pg = Page("① 雨轩竞争定位：羊肉赛道头部，全国全口径有更大对手")
pg.sec("规模与基本面（公开可查口径）", 18, C_DARK, gap=0.06)
pg.tbl(["指标", "数值", "信号"],
       [["营收（合并口径）", "2022:26 / 2023:21 / 2024:15.7 / 2025E:16 亿（三年连降）", "羊肉主业第一梯队"],
        ["屠宰量", "设计产能 150 万只/年；累计（2019 投产）超 300 万只；2024 屠宰 69 万只", "河南规上屠宰 85%+"],
        ["线上（预制菜）", "2023 线上破 8,000 万；爆款『香辣羊蝎子』20 分钟 3 万单 / 400 万", "占整体约 3%+，最大短板"],
        ["出口", "2024.9 首发中东，2025 翻 3 倍、目标 2,000 万美元；AEO 认证", "已布局双通道"],
        ["品牌", "全国榜第 7（chinabgao 2025）；『黄河滩羊』占河南羊肉制品 60%+", "全国弱、区域绝对主导"]],
       size=14, col_w=[2.6, 6.4, 3.4], gap=0.08, row_pad=0.05)
pg.sec("护城河 vs 软肋", 18, C_GREEN, gap=0.2)
pg.bul([
    "护城河：全产业链闭环（养-屠-深加-冷链-外贸-保税仓）+ 国产冷鲜可溯源 + 清真/AEO 出口资质 + 已绑定巴奴/锅圈；资质含农业产业化国家重点龙头、国家级肉羊产业集群链主、中国肉类协会会长单位",
    "软肋：营收三年连降（26→21→15.7 亿）、2025E 净利率仅约 1.25%（净利 0.2 亿）盈利极薄、线上占比低、全国品牌第 7、羊源约 75% 外采",
], size=15, item_gap=0.08)
pg.view("对外表述建议：雨轩是中原及全国羊肉深加工领域头部企业，全国牛羊肉综合实力位列前五（公开可查口径）", 15)
pg.finish()

# ============ P4 对手地图（分层）============
pg = Page("② 对手地图：四类对手，威胁性质完全不同")
pg.sec("竞品三层界定（按威胁来源分层）", 18, C_BLUE, gap=0.06)
pg.box(2.25, rival_map, gap=0.08, width=12.4)
pg.sec("直接 / 间接竞品清单（含羊肉品牌榜名次）", 18, C_DARK, gap=0.16)
pg.tbl(["层级", "代表企业（品牌榜名次）", "对雨轩的意义"],
       [["直接（羊肉为主）", "大庄园(#1)、额尔敦(#3)、蒙都、草原宏宝、小巴依、盐池滩羊系、茂源", "同赛道贴身竞争，体量小但打法新"],
        ["间接 / 潜在", "大希地、小牛凯西（抢家庭预算）、恒都泌阳、中原皓月（开封）", "分流需求与渠道"],
        ["备注", "雨轩斋居品牌榜第 7（chinabgao 2025·共 10 强），河南唯一上榜", "品牌弱于北方产区老牌，高于区域小厂"]],
       size=13, col_w=[2.2, 6.4, 3.8], gap=0.06, row_pad=0.05)
pg.view("结论：雨轩体量约为直接竞品额尔敦（4.33 亿）的 3–4 倍，同赛道无同量级对手；真威胁来自品类外（进口/低价替代）与地域（大庄园、恒都泌阳、茂源）", 15)
pg.finish()

# ============ P4b 品牌力全景：羊肉品牌榜 TOP10 ============
pg = Page("⑤ 品牌力全景：羊肉品牌榜 TOP10（雨轩斋第 7）")
pg.sec("中国报告大厅 2025 羊肉十大品牌（公开榜单，可核验）", 18, C_BLUE, gap=0.06)
pg.tbl(["排名", "品牌", "企业", "省份"],
       [["1", "大庄园", "大庄园肉业集团", "黑龙江"],
        ["2", "涝河桥", "宁夏涝河桥肉食品", "宁夏"],
        ["3", "额尔敦羊业", "内蒙古额尔敦羊业", "内蒙古"],
        ["4", "吉羊羊", "锡林郭勒盟羊羊牧业", "内蒙古"],
        ["5", "青青草原", "内蒙古青青草原牧业", "内蒙古"],
        ["6", "宁鑫", "盐池县鑫海食品", "宁夏"],
        ["7", "雨轩斋 ★", "新乡市雨轩清真食品", "河南"],
        ["8", "盐池滩羊", "盐池滩羊产业发展集团", "宁夏"],
        ["9", "天顺源", "黑龙江天顺源清真食品", "黑龙江"],
        ["10", "草之味", "苏尼特左旗满都拉图肉食品", "内蒙古"]],
       size=12.5, col_w=[0.9, 2.0, 6.3, 1.3], gap=0.06, row_pad=0.028)
pg.note("绿框=雨轩；榜单前 3 均为北方产区老牌（大庄园东北牛+羊、额尔敦内蒙羊），是品牌与规模的双标杆", size=13, gap=0.08, color=C_GRAY)
pg.view("读图：雨轩是河南唯一上榜企业（第 7），品牌认知弱于产区老牌，但高于多数区域小厂——『品牌』是较『规模/渠道』更优先的补强项", 15)
pg.finish()

# ============ P5 九宫格对标矩阵（核心）============
pg = Page("③ 竞争维度对标矩阵（九宫格 · 报告核心）")
pg.sec("8 维度 × 4 梯队，雨轩与对手逐一比对", 18, C_BLUE, gap=0.06)
pg.box(3.0, matrix_9box, gap=0.08, width=12.4)
pg.note("绿=强 黄=中 红=弱 · 九宫格已用公开核验信号校准（宁羴源 2024 线上 1.8 亿、电商占比 80%、抖音羊肉榜 TOP1 三年；雨轩线上数千万级占整体约 3%）", size=12.5, gap=0.08, color=C_GRAY)
pg.view("读图：雨轩在『线下渠道 / 品控 / 出口』三格占优；『线上』一格最弱且与垂直新锐（宁羴源）差距最大——这是战略资源应优先填补的缺口", 15)
pg.finish()

# ============ P6 差距与护城河 ============
pg = Page("④ 差距 → 护城河：把短板变成进攻支点")
pg.sec("关键差距与对应护城河打法", 18, C_DARK, gap=0.06)
pg.tbl(["差距（九宫格红/黄格）", "风险", "护城河打法"],
       [["线上弱（vs 宁羴源/大希地）", "家庭端增量被抢、品牌年轻化停滞", "复制自播矩阵 + 爆品扩 SKU，线上占比目标 15%"],
        ["全国品牌第 7 / 区域 60%+", "全国溢价弱、礼盒定价受抑；河南市场已占 60%+", "借专项整治打『真羊肉/0 掺假/全程溯源』信任牌，全国化靠千店计划"],
        ["羊源 75% 外采", "活羊 +6.8% 推高成本、毛利承压", "11 月前锁价 + 近地化 25%→40%"],
        ["全国巨头规模碾压", "大庄园/恒都正面重叠", "差异化：清真 + 出口 + 中原主场，不拼规模"],
        ["深加工率仅约 5%", "行业天花板低、增长靠 α", "预制菜/礼盒/出口高附加值四条线"]],
       size=13.5, col_w=[3.0, 4.0, 5.4], gap=0.08, row_pad=0.05)
pg.view("原则：不在巨头强项（规模/品牌）上硬碰，而在其弱项（信任/数字化/近地化/深加工）上建壁垒", 15)
pg.finish()

# ============ P7 战略动作一：锁羊源·修基本盘 ============
pg = Page("战略动作一：锁羊源 · 修基本盘（与周期同频）")
pg.sec("为什么是现在：蛛网周期节拍", 18, C_GREEN, gap=0.06)
pg.bul([
    "机理：母畜淘汰→存栏降(6–12 月)→出栏降(6–12 月)→价格升(3–6 月)→补栏(12–24 月)；羊繁殖一轮 12–18 月",
    "2024 谷底 → 2025 磨底（额尔敦营收 +10.7% 仍亏 755 万=铁证）→ 2026 拐点确认 → 2027 复苏上行",
    "官方交叉验证：2024 全国羊存栏 3 亿只（-6.8%，近 6 年最低）、羊肉产量 517.75 万吨（-2.54%）、进口 36.65 万吨（-15.49%）—— 与 2026 H1 信号一致，供给收缩确立",
    "窗口：2026–2027 供给紧张持续，但生猪 -23.1% 压制需求 —— 结构性机会 + 防御性操作",
], size=15, item_gap=0.08)
pg.sec("动作清单", 18, C_DARK, gap=0.2)
pg.bul([
    "11 月前锁定春节档羊源（活羊 +6.8%、75% 外采，最大成本风险）；推进近地化 25%→40%",
    "以销定产、不囤货；设计产能 150 万只/年，不盲目扩建重资产屠宰产能（伊赛教训）",
    "修复核心品类、重构基本盘、堵漏节流（呼应 2026 六大任务）",
], size=15, item_gap=0.08)
pg.view("判错止损线：若 Q4 价格未站稳 72–75、出栏降幅 <6%、进口 >35 万吨，则修正为更保守基调", 14)
pg.finish()

# ============ P8 战略动作二：全链路数字化·AI 赋能 ============
def kpi4(sl, top, left, width):
    kpi_row(sl, [("ERP/MES", "全链数据底座\n羊源→餐桌贯通", C_DARK),
                 ("AI 排产", "智能排产+需求预测\n降闲置、提周转", C_BLUE),
                 ("损耗预警", "冷链 IoT 实时监控\n直接压降成本", C_GREEN),
                 ("动态定价", "数据驱动毛利\n守住而非走量", C_AMBER)],
            top=top, left=left, width=width, height=1.35, size=13)

pg = Page("战略动作二：全链路数字化 · AI 赋能降本增效")
pg.box(1.5, kpi4, gap=0.06)
pg.sec("为什么数字化是当前最优抓手", 18, C_GREEN, gap=0.18)
pg.bul([
    "底座具备：ERP/MES 与冷链 IoT 可贯通，羊源—屠宰—深加工—渠道全链路数据打通条件已成熟",
    "降本最直接：AI 需求预测 + 智能排产降低产能闲置与库存损耗；损耗预警直接压成本（活羊 +6.8% 推高成本背景下尤其关键）",
    "定价更优：动态定价/毛利监控，在结构弱需求下守住毛利而非走量",
    "内部已认：公司《三年规划》已将『物联网溯源 + 数字化产销协同平台』列为科技支柱，方向一致，缺落地专班与 ROI 验证",
    "组织保障：GM 直管数字化专班，Q3 启动试点、Q4 验证 ROI，不做无业务价值的 IT 项目",
], size=16, item_gap=0.08)
pg.view("建议：Q3 启动全链路数字化试点（智能排产/损耗预警），Q4 验证 ROI；目标从羊源到餐桌数据贯通，把数字化做成持续能力而非一次性工程", 15)
pg.finish()

# ============ P9 战略动作三：制度化系统化·旺季与渠道 ============
pg = Page("战略动作三：制度化系统化 · 旺季与渠道（风控合规）")
pg.sec("渠道卡位：已占双雄，最大空白在胖东来", 18, C_DARK, gap=0.06)
pg.tbl(["渠道 / 对手", "关键数据", "雨轩卡位", "动作"],
       [["巴奴毛肚火锅", "150 店、年营收 25 亿+", "★ 已是客户", "升级年度框架协议，锁量锁价"],
        ["锅圈食汇", "万店规模", "★ 已是客户", "深化定制 SKU，防茂源替代"],
        ["胖东来", "2025 销售 235 亿（+38.7%）", "✗ 缺席（最大遗憾）", "全力冲击，但不牺牲毛利"],
        ["千店计划（自营）", "2026 试点50→2027 400→2028 1000+ 加盟店", "◐ 规划中", "冷鲜现切+调理+熟食，数字化/SOP 托底"],
        ["永辉 / 盒马 / 大张", "永辉调改 100→300 家", "区域铺货中", "借店中店/专柜扩全国网络"]],
       size=13.5, col_w=[2.6, 3.3, 2.7, 3.8], gap=0.06, row_pad=0.05)
pg.sec("制度系统化与风控合规", 18, C_AMBER, gap=0.2)
pg.bul([
    "制度系统化：把依赖个人的经验转化为 SOP 化核心流程，降低经营波动、承接渠道扩张与专项整治",
    "风控/合规三道防线：业务自查 → 职能风控 → 独立合规审计；建出口合规（RASFF/FSIS）与食品安全月度监控",
    "信任牌：借 2026 假牛羊肉专项整治（重庆 732 起）打『真羊肉/0 掺假/出口级』标准，转化品牌窗口",
], size=15, item_gap=0.08)
pg.view("旺季节奏：春节(12–2月)>中秋国庆(9–10月)>双十一(10–11月)>夏季烧烤(6–8月)；8 月定盘、10 月出礼盒、礼盒落位 150–400 元；制度框架 Q4 前落地，与旺季作战并行", 14)
pg.finish()

# ============ P10 DO / DON'T ============
pg = Page("2026 下半年：可以做的与不可以做的")
pg.box(4.5, two_col, gap=0.06)
pg.view("核心原则：用确定的动作抓确定的窗口 —— 锁成本、强数字化、抢旺季、建制度，三件事确定性最高；产能与线上重仓等 Q4 数据验证后再加码", 15)
pg.finish()

# ============ P11 出手节奏 ============
pg = Page("出手节奏：周期与旺季共同决定『何时动手』")
pg.sec("月度收入占比模型（行业估算 C 级，待内部数据校准）", 17, C_DARK, gap=0.06)
pg.box(1.9, month_chart, gap=0.06)
pg.sec("时间轴：与董事会决策节点对齐", 17, C_BLUE, gap=0.16)
pg.tbl(["时点", "动作", "对应决策"],
       [["8 月", "定春节档盘子、启动全链路数字化试点立项", "决策①②"],
        ["9 月", "锁中秋企业团购、数字化试点上线、制度框架启动", "决策②③"],
        ["10 月", "出礼盒、冲击胖东来、制度系统化框架落地", "决策③"],
        ["11 月", "锁定春节档羊源、备货、数字化试点 ROI 复盘", "决策①②"],
        ["Q4", "二次验证拐点，验证数字化 ROI，再决定加码", "复盘"]],
       size=14, col_w=[1.4, 7.2, 3.8], gap=0.06, row_pad=0.05)
pg.view("节奏逻辑：旺季前 2–3 个月必须完成备货与锁价；数字化试点（Q3 立项、Q4 验证）与制度框架（Q4 前）与旺季作战并行推进", 14)
pg.finish()

# ============ P12 持续监测 ============
pg = Page("④ 持续监测：竞品情报闭环（已出第 1 期）")
pg.sec("情报速递机制（L1 监控已落地）", 18, C_BLUE, gap=0.06)
pg.bul([
    "模板固化：每月 1 日市监总局 + 5 省通报；每周一价格周报；每月 15 日 RASFF/FSIS；每季度店铺 SKU/黑猫/招聘；双月海底捞 SRM 巡检",
    "第 1 期已覆盖：抽检黑名单（太原/桃源/天牧）、出口合规（RASFF 氨基脲）、海底捞年报（鲜切溯源门槛）、宁羴源全量对标、大希地投诉",
], size=15, item_gap=0.08)
pg.sec("第 1 期关键信号", 18, C_DARK, gap=0.2)
pg.tbl(["信号类型", "内容", "对雨轩"],
       [["监管", "2026 假牛羊肉专项整治年（重庆 732 起）", "信任牌窗口放大"],
        ["出口", "RASFF 通报中国咸羊肠氨基脲，欧盟线需提前布控", "建月度合规监控"],
        ["竞品", "宁羴源 2024 线上 1.8 亿、电商占比 80%、抖音羊肉榜 TOP1 三年、日峰值 226 万", "线上打法对标"],
        ["舆情", "大希地 2024 GMV 约 30 亿、月发 300 万单，但客单价 100 元级、投诉承压", "流量打法可借鉴、毛利结构需警惕"]],
       size=13.5, col_w=[1.8, 6.4, 4.2], gap=0.06, row_pad=0.05)
pg.view("闭环：情报速递 → 月度滚动更新 → 反哺九宫格与战略动作；建议设为每周一自动拉取价格 + 竞品动态", 14)
pg.finish()

# ============ P13 决策请求 ============
pg = Page("请董事会决策的三项事项")
items = [
    ("①", "2026–2027 春节档总盘子与预算", "9 月 30 日前批准",
     "含 SKU 四档矩阵（引流 50–100 / 利润 150–300 / 形象 300+ / 年夜饭组合）、渠道资源分配、礼盒备货量；错过 8 月定盘子、10 月出礼盒两个节点，春节档只能吃残羹", C_RED),
    ("②", "全链路数字化与 AI 赋能专项", "Q3 启动 / Q4 验证",
     "从羊源到餐桌全链数据贯通（ERP/MES/冷链 IoT 底座），用 AI 需求预测、智能排产、损耗预警、动态定价降本增效；GM 直管专班，Q3 试点、Q4 验证 ROI", C_GREEN),
    ("③", "制度系统化与风控合规体系建设", "Q4 前框架落地",
     "把依赖个人的经验转化为 SOP/风控/合规体系（业务自查→职能风控→独立合规三道防线）；承接专项整治与渠道扩张，建出口合规(RASFF/FSIS)与食品安全月度监控", C_BLUE),
]
for i, (num, t1, deadline, desc, color) in enumerate(items):
    h1 = text_h(f"{num}  {t1}     【{deadline}】", 12.1, 19)
    h2 = text_h(desc, 12.1, 14.5)
    bh = h1 + h2 + 0.30
    pg.y += 0.20 if i else 0.08
    box = pg.s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.45), Inches(pg.y), Inches(12.4), Inches(bh))
    box.fill.solid()
    box.fill.fore_color.rgb = C_WHITE
    box.line.color.rgb = color
    box.line.width = Pt(2)
    box.shadow.inherit = False
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.16)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.text = f"{num}  {t1}     【{deadline}】"
    p.font.size = Pt(19)
    p.font.bold = True
    p.font.name = FONT
    p.font.color.rgb = color
    p2 = tf.add_paragraph()
    p2.text = desc
    p2.font.size = Pt(14)
    p2.font.name = FONT
    p2.font.color.rgb = C_TXT
    pg.y += bh
pg.view("结语：行业拐点由供给驱动、需求仍弱 —— 锁成本、强数字化、抢旺季、建制度；以销定产、现金为王。不赌拐点，但绝不错过窗口。", 16)
pg.finish()

out = os.path.join(HERE, "雨轩食品竞品分析与战略决策(董事会汇报版).pptx")
prs.save(out)
print("saved | slides:", len(prs.slides._sldIdLst))
print("每页内容底部 y：", LOG)
