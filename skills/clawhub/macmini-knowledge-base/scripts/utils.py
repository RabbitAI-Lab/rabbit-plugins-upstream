#!/usr/bin/env python3
"""
知识库共享工具函数 v2.1
2026-08-12 修复：
- 修复 OCR 路径不工作（--skip-text → --force-ocr）
- 新增 CMap 残缺度自检（is_cmap_broken, 阈值 0.03）
- 新增 PDFExtractError 异常类（含文件路径）
- OCR 超时 60s → 600s（支持大文档）
"""
import os
import json
import shutil
import subprocess
import re as re_module
import uuid
import tempfile
import pymupdf
import docx
import openpyxl
import pptx

try:
    from kreuzberg import extract_file_sync as kreuzberg_extract
    KREUZBERG_AVAILABLE = True
except Exception:
    KREUZBERG_AVAILABLE = False

_DEFAULT_KNOWLEDGE_DIR = os.path.expanduser("~/.openclaw/workspace/knowledge")
_STATE_FILE    = os.path.join(_DEFAULT_KNOWLEDGE_DIR, ".analysis/analysis_state.json")
_CACHE_FILE    = os.path.join(_DEFAULT_KNOWLEDGE_DIR, ".analysis/.catalog_cache.json")
_PROGRESS_FILE = os.path.join(_DEFAULT_KNOWLEDGE_DIR, ".analysis/.catalog_progress.json")

# ==== 异常类 ====
class PDFExtractError(Exception):
    """PDF 提取失败异常（含文件路径）"""
    def __init__(self, filepath, reason):
        self.filepath = filepath
        self.reason = reason
        super().__init__(f"PDF 提取失败: {filepath} - {reason}")


# v2.1: 通用文件提取失败异常
class ExtractError(Exception):
    """通用文件提取失败异常（含文件路径）"""
    def __init__(self, filepath, reason):
        self.filepath = filepath
        self.reason = reason
        super().__init__(f"文件提取失败: {filepath} - {reason}")

# ==== CMap 残缺度检测（C 方案核心）====
CMAP_BAD_RATIO_THRESHOLD = 0.03   # 异常字符比例阈值（用户指定）
CMAP_CID_COUNT_THRESHOLD = 10      # (cid:xxxx) 出现次数阈值
CMAP_SAMPLE_CHARS_MIN = 50         # 文本长度低于此值不做自检（可能是图像PDF）

# v2.2: 文本提取上限（用户要求"完整阅读"，原 [:8000] 太小）
# 50 万字 ≈ 1MB，覆盖 99.9% 正常文档；超大 Excel/Word 也不会无限制膨胀
MAX_EXTRACT_LEN = 500_000

def is_cmap_broken(text, threshold=CMAP_BAD_RATIO_THRESHOLD):
    """检测文本是否含异常字符（CMap残缺/PUA污染/未映射CID字面值）

    Returns:
        bool: True 表示需要走 OCR
    """
    if not text or len(text.strip()) < CMAP_SAMPLE_CHARS_MIN:
        return False

    total = len(text)
    # PUA 私用区 U+E000-F8FF
    pua_count = sum(1 for c in text if 0xE000 <= ord(c) <= 0xF8FF)
    # CJK 扩展 B/C/D
    cjk_ext = sum(1 for c in text if 0x20000 <= ord(c) <= 0x2EBEF)
    # CJK 兼容区 U+F900-FAFF
    cjk_compat = sum(1 for c in text if 0xF900 <= ord(c) <= 0xFAFF)
    # (cid:xxxx) 字面值（pdfplumber 提取失败标志）
    cid_count = text.count('(cid:')

    bad_count = pua_count + cjk_ext + cjk_compat + cid_count
    bad_ratio = bad_count / total

    return bad_ratio > threshold or cid_count > CMAP_CID_COUNT_THRESHOLD


def load_json(path=None):
    if path is None:
        path = _STATE_FILE
    if os.path.exists(path):
        try:
            with open(path, 'r', encoding='utf-8') as f:
                return json.load(f)
        except (json.JSONDecodeError, IOError):
            pass
    return None

def save_json(path, data):
    if path is None:
        path = _STATE_FILE
    os.makedirs(os.path.dirname(path), exist_ok=True)
    with open(path, 'w', encoding='utf-8') as f:
        json.dump(data, f, ensure_ascii=False, indent=2)

def load_state():
    data = load_json(_STATE_FILE)
    return data if data else {}

def save_state(state):
    save_json(_STATE_FILE, state)

def load_cache():
    data = load_json(_CACHE_FILE)
    return data if data else {}

def save_cache(cache):
    save_json(_CACHE_FILE, cache)

def load_progress():
    data = load_json(_PROGRESS_FILE)
    return data if data else {}

def save_progress(progress):
    save_json(_PROGRESS_FILE, progress)

def is_gibberish(text, strict=False):
    if not text or len(text.strip()) < 20:
        return True
    alpha_count   = sum(1 for c in text if c.isalpha())
    upper_count   = sum(1 for c in text if c.isupper())
    space_count   = sum(1 for c in text if c.isspace())
    chinese_count = sum(1 for c in text if '\u4e00' <= c <= '\u9fff')
    chinese_ratio = chinese_count / len(text) if len(text) > 0 else 0
    if alpha_count > 10:
        upper_ratio = upper_count / alpha_count
        space_ratio = space_count / len(text)
        if upper_ratio > 0.45 and space_ratio < 0.05:
            return True
    if chinese_ratio > 0.3:
        return False
    cleaned = re_module.sub(r'[\s\n]', '', text)
    if len(cleaned) == 0:
        return True
    normal_count = len(re_module.findall(r'[A-Za-z0-9 .,;:\'"!?()-]', cleaned))
    ratio = normal_count / len(cleaned) if len(cleaned) > 0 else 0
    return ratio < (0.5 if strict else 0.3)

def extract_via_kreuzberg(filepath):
    if not KREUZBERG_AVAILABLE:
        return None
    try:
        result = kreuzberg_extract(filepath)
        if result and result.content and result.content.strip():
            return result.content
    except Exception:
        pass
    return None


def _quick_page_check(filepath):
    """快速预检前3页是否需要走OCR（文字量极少=编码PDF）"""
    try:
        doc = pymupdf.open(filepath)
        total_chars = 0
        for i in range(min(3, len(doc))):
            t = doc[i].get_text()
            total_chars += len(t.strip())
        doc.close()
        if total_chars < 200:
            return True
    except:
        pass
    return False

def extract_doc_via_antiword(filepath):
    try:
        result = subprocess.run(['antiword', filepath], capture_output=True, timeout=10)
        if result.returncode == 0:
            text = result.stdout.decode('utf-8', errors='replace')
            if len(text) > 100:
                return text
    except Exception:
        pass
    return None

def _kill_proc_tree(pid):
    try:
        # v1.4.3: 替换 os.system 为 subprocess.run（避免 shell 注入 + 抑制输出）
        subprocess.run(["pkill", "-P", str(pid)],
                       capture_output=True, timeout=5)
    except (subprocess.TimeoutExpired, Exception):
        pass
    try:
        os.kill(pid, 9)
    except:
        pass

def convert_old_office(filepath, ext):
    # v1.4.3 fix2: 拆分为内部函数（finally 清理 tmp_dir 前先把结果 cp 到稳定位置）
    tmp_dir = tempfile.mkdtemp(prefix="office_convert_")
    try:
        convert_map = {".doc": "docx", ".xls": "xlsx", ".ppt": "pptx"}
        new_ext = convert_map.get(ext)
        if not new_ext:
            return (False, "")
        expected_name = os.path.basename(filepath).replace(ext, '.' + new_ext)
        tmp_out = os.path.join(tmp_dir, expected_name)
        if os.path.exists(tmp_out):
            try:
                os.remove(tmp_out)
            except:
                pass
        proc = subprocess.Popen(
            ['soffice', '--headless', '--convert-to', new_ext, '--outdir', tmp_dir, filepath],
            stdout=subprocess.DEVNULL, stderr=subprocess.DEVNULL
        )
        try:
            outs, errs = proc.communicate(timeout=60)
            if proc.returncode == 0 and os.path.exists(tmp_out):
                # v1.4.3 fix2: 把结果 cp 到 tmp_dir 外的稳定路径（否则 finally 会删掉）
                # 用 tempfile.mkstemp（更安全，无 race condition）
                stable_fd, stable_out = tempfile.mkstemp(suffix='.' + new_ext)
                os.close(stable_fd)
                shutil.copy2(tmp_out, stable_out)
                return (True, stable_out)
        except subprocess.TimeoutExpired:
            try:
                proc.kill()
            except:
                pass
            _kill_proc_tree(proc.pid)
        return (False, "")
    finally:
        # v1.4.3: finally 清理私有目录
        shutil.rmtree(tmp_dir, ignore_errors=True)

def extract_doc_text(filepath):
    """提取 .doc 文本（v2.1: 加 OCR fallback）

    流程：
    1. antiword 提取（直接读 .doc 流）
    2. CMap 自检（is_cmap_broken）
    3. 自检失败 → soffice 转 PDF + OCRmyPDF 强制 OCR
    4. 失败 → 抛 ExtractError（含文件路径）
    """
    text = extract_doc_via_antiword(filepath)
    if text and not is_cmap_broken(text):
        return text

    # 兜底：antiword 提取为空或乱码 → 走 OCR
    success, converted = convert_old_office(filepath, '.doc')
    if success:
        try:
            doc = docx.Document(converted)
            text = "\n".join(p.text for p in doc.paragraphs if p.text.strip())
            if text.strip() and not is_cmap_broken(text):
                return text[:MAX_EXTRACT_LEN]
        finally:
            try:
                os.remove(converted)
            except:
                pass

    # 最终兜底：OCR
    ocr_text = ocr_office_via_ocr(filepath, '.doc')
    if ocr_text:
        return ocr_text
    raise ExtractError(filepath, ".doc 提取失败：antiword/soffice/OCR 都返回空或乱码")

def extract_pdf_text(filepath):
    """提取 PDF 文本（CMap 残缺度自检 v2.1）

    流程：
    1. _quick_page_check：前 3 页文字量少 → 直接 OCR
    2. kreuzberg 默认提取 + CMap 自检
    3. pymupdf 默认提取 + CMap 自检
    4. 自检失败或提取为空 → OCR（--force-ocr 已修复）
    5. OCR 失败 → 抛 PDFExtractError（含文件路径）
    """
    # 前 3 页文字量少 = 可能是扫描件/图像 PDF
    if _quick_page_check(filepath):
        result = extract_pdf_via_ocr(filepath)
        if result and not result.startswith("【"):
            return result
        raise PDFExtractError(filepath, "OCR 后仍无文本（可能是扫描件或损坏文件）")

    # kreuzberg 提取 + CMap 自检
    text = extract_via_kreuzberg(filepath)
    if text and not is_cmap_broken(text):
        return text

    # pymupdf 提取 + CMap 自检
    try:
        doc = pymupdf.open(filepath)
        text = ""
        for page_num in range(len(doc)):
            t = doc[page_num].get_text()
            if t.strip():
                text += t + "\n"
        doc.close()
        if text.strip() and not is_cmap_broken(text):
            return text
    except Exception:
        pass

    # 走到这里：默认提取失败或检测到乱码 → 走 OCR
    result = extract_pdf_via_ocr(filepath)
    if result and not result.startswith("【"):
        return result
    raise PDFExtractError(filepath, "OCR 失败：所有提取路径都返回空或乱码")

def extract_docx_text(filepath):
    text = extract_via_kreuzberg(filepath)
    if text:
        return text
    try:
        doc = docx.Document(filepath)
        text = "\n".join(p.text for p in doc.paragraphs if p.text.strip())
        return text[:MAX_EXTRACT_LEN] if text.strip() else "【Word文档无文字内容】"
    except Exception as e:
        return "【DOCX提取失败】" + str(e)

def extract_xlsx_text(filepath):
    text = extract_via_kreuzberg(filepath)
    if text:
        return text
    try:
        wb = openpyxl.load_workbook(filepath, data_only=True)
        sheets_text = []
        for sheet in wb.worksheets:
            for row in sheet.iter_rows(values_only=True):
                row_text = " ".join(str(c) for c in row if c is not None)
                if row_text.strip():
                    sheets_text.append(row_text)
        return "\n".join(sheets_text)[:MAX_EXTRACT_LEN] if sheets_text else "【Excel无内容】"
    except Exception as e:
        return "【XLSX提取失败】" + str(e)

def extract_pptx_text(filepath):
    text = extract_via_kreuzberg(filepath)
    if text:
        return text
    try:
        prs = pptx.Presentation(filepath)
        text_parts = []
        for i, slide in enumerate(prs.slides[:20]):
            slide_text = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_text.append(shape.text)
            if slide_text:
                text_parts.append(f"[Page {i+1}]\n" + "\n".join(slide_text))
        return "\n".join(text_parts)[:MAX_EXTRACT_LEN] if text_parts else "【PPT无文字内容】"
    except Exception as e:
        return "【PPTX提取失败】" + str(e)

def extract_pdf_via_ocr(filepath):
    """用 OCRmyPDF 强制 OCR 提取 PDF 文本（v2.1 已修复 --force-ocr）

    Returns:
        str: 提取的文本。如果失败，返回 "【XXX失败】..." 标记
    """
    # v1.4.3: 用 tempfile.mkdtemp() 创建 per-run 私有目录
    tmp_dir = tempfile.mkdtemp(prefix="ocrmypdf_")
    tmp_pdf = os.path.join(tmp_dir, uuid.uuid4().hex + ".pdf")
    tmp_out = os.path.join(tmp_dir, uuid.uuid4().hex + "_ocr.pdf")
    try:
        shutil.copy2(filepath, tmp_pdf)
        # v2.1 修复：用 --force-ocr 替代 --skip-text，确保 OCRmyPDF 强制 OCR 所有页面
        result = subprocess.run(
            ["ocrmypdf", "-l", "chi_sim+eng", "--force-ocr",
             "--pages", "1-999", tmp_pdf, tmp_out],
            capture_output=True, timeout=600   # v2.1: 超时从 60s 提升到 600s
        )
        if result.returncode == 0 and os.path.exists(tmp_out):
            doc = pymupdf.open(tmp_out)
        else:
            doc = pymupdf.open(filepath)
        text = ""
        for page_num in range(len(doc)):
            t = doc[page_num].get_text()
            if t.strip():
                text += t + "\n"
        doc.close()
        if text.strip():
            # 过滤封面页乱码片段
            lines = text.split('\n')
            cleaned_lines = []
            for line in lines:
                if is_gibberish(line) and len(line.strip()) < 100:
                    continue
                cleaned_lines.append(line)
            text = '\n'.join(cleaned_lines).strip()
            if text:
                return text
    except Exception as e:
        pass
    finally:
        # v1.4.3: finally 清理整个私有目录
        shutil.rmtree(tmp_dir, ignore_errors=True)
    return "【PDF文字提取失败】"

def ocr_office_via_ocr(filepath, source_ext):
    """通用 Office 文件 OCR 路径（v2.1 新增）

    流程：
    1. soffice headless 转 PDF
    2. OCRmyPDF --force-ocr
    3. pymupdf 读文本

    Returns:
        str: 提取的文本。失败返回 None。
    """
    # v1.4.3: 用 tempfile.mkdtemp() 创建 per-run 私有目录
    tmp_dir = tempfile.mkdtemp(prefix="office_ocr_")
    try:
        # 1. soffice 转 PDF
        proc = subprocess.Popen(
            ['soffice', '--headless', '--convert-to', 'pdf', '--outdir', tmp_dir, filepath],
            stdout=subprocess.DEVNULL, stderr=subprocess.DEVNULL
        )
        try:
            proc.communicate(timeout=120)
        except subprocess.TimeoutExpired:
            try:
                proc.kill()
            except:
                pass
            _kill_proc_tree(proc.pid)
        
        # 找生成的 PDF
        base = os.path.splitext(os.path.basename(filepath))[0]
        converted_pdf = os.path.join(tmp_dir, base + ".pdf")
        if not os.path.exists(converted_pdf):
            return None
        
        # 2. OCRmyPDF --force-ocr
        ocr_pdf = os.path.join(tmp_dir, uuid.uuid4().hex + "_ocr.pdf")
        result = subprocess.run(
            ["ocrmypdf", "-l", "chi_sim+eng", "--force-ocr",
             "--pages", "1-999", converted_pdf, ocr_pdf],
            capture_output=True, timeout=600
        )
        # 3. 读文本
        if result.returncode == 0 and os.path.exists(ocr_pdf):
            doc = pymupdf.open(ocr_pdf)
        else:
            doc = pymupdf.open(converted_pdf)
        text = ""
        for page_num in range(len(doc)):
            t = doc[page_num].get_text()
            if t.strip():
                text += t + "\n"
        doc.close()
        return text.strip() if text.strip() else None
    except Exception:
        return None
    finally:
        # v1.4.3: finally 清理整个私有目录
        shutil.rmtree(tmp_dir, ignore_errors=True)


def extract_xls_text(filepath):
    """Extract text from legacy .xls Excel files using xlrd"""
    try:
        import xlrd
        wb = xlrd.open_workbook(filepath)
        text_parts = []
        for sheet in wb.sheets():
            text_parts.append(f"[{sheet.name}]")
            for row in sheet.get_rows():
                row_data = []
                for cell in row:
                    if cell.ctype == xlrd.XL_CELL_TEXT:
                        row_data.append(cell.value)
                    elif cell.ctype == xlrd.XL_CELL_NUMBER:
                        row_data.append(str(cell.value))
                    elif cell.ctype == xlrd.XL_CELL_DATE:
                        try:
                            dt = xlrd.xldate_as_datetime(cell.value, wb.datemode)
                            row_data.append(dt.strftime('%Y-%m-%d'))
                        except:
                            row_data.append(str(cell.value))
                    elif cell.ctype == xlrd.XL_CELL_EMPTY:
                        pass
                if any(row_data):
                    text_parts.append(' | '.join(row_data))
        return '\n'.join(text_parts)
    except Exception as e:
        return None
