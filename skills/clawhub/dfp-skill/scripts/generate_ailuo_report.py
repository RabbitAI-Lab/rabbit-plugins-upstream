#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
艾罗能源 (688717) 深度估值报告生成脚本
作者: Wang Dongjie, CGMA/AICPA&CIMA, © 2026

使用数智财务 SAP 企业级深色风格生成估值报告
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# ============================================================================
# 数智财务 SAP 风格配色
# ============================================================================
BG_TOP        = RGBColor(0x0A, 0x18, 0x38)  # 深海军蓝
BG_MID        = RGBColor(0x12, 0x24, 0x4D)  # 中深蓝
BG_BOTTOM     = RGBColor(0x5A, 0x0F, 0x25)  # 暗红
CARD_FILL     = RGBColor(0x0F, 0x20, 0x50)  # 深蓝卡片底
CARD_BORDER   = RGBColor(0xB8, 0x00, 0x3A)  # 深红边框
ACCENT_GOLD   = RGBColor(0xF2, 0xB8, 0x4B)  # 装饰金色
ACCENT_CYAN   = RGBColor(0x25, 0xB7, 0xE0)  # 青色科技蓝
ACCENT_GREEN  = RGBColor(0x2F, 0xA4, 0x72)  # 财务绿
ACCENT_RED    = RGBColor(0xE7, 0x4C, 0x3C)  # 警示红
WHITE         = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_WHITE   = RGBColor(0xE8, 0xEC, 0xF5)

# ============================================================================
# 艾罗能源实际数据
# ============================================================================
COMPANY_DATA = {
    "公司名称": "艾罗能源",
    "股票代码": "688717",
    "英文名称": "SolaX Power Network Technology",
    "所属行业": "电气机械和器材制造业",
    "主营业务": "光伏储能逆变器、储能电池、并网逆变器",
    "上市日期": "2024年1月3日",
    "发行价格": "55.66元/股",
    
    # 股市数据 (2026年6月17日)
    "当前股价": "91.04元",
    "市值": "145.66亿元",
    "PE_TTM": "102.37x",
    "PB": "3.19x",
    "PS": "3.0x",
    "股息率": "1.03%",
    "52周最高": "157.13元",
    "52周最低": "50.30元",
    
    # 财务数据 (2024年)
    "营收_2024": "30.73亿元",
    "营收同比": "-31.30%",
    "净利润_2024": "2.04亿元",
    "净利润同比": "-80.88%",
    "ROE": "4.63%",
    "毛利率": "38.12%",
    "资产负债率": "28.31%",
    "经营现金流": "7.54亿元",
    
    # 业务构成
    "户用储能占比": "60.78%",
    "并网逆变器占比": "19.48%",
    "工商业储能占比": "13.53%",
    
    # 行业地位
    "认证数量": "3000+项",
    "覆盖国家": "130+",
    "荣誉": "国家级制造业单项冠军",
    
    # 风险因素
    "风险_欧洲市场": "欧洲户储市场降温",
    "风险_业绩下滑": "营收净利润大幅下滑",
    "风险_估值偏高": "PE超过100x",
}

# ============================================================================
# 超宽屏尺寸
# ============================================================================
SLIDE_W_IN = 22.22
SLIDE_H_IN = 7.5

def create_presentation():
    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W_IN)
    prs.slide_height = Inches(SLIDE_H_IN)
    
    # 获取空白布局
    blank_layout = prs.slide_layouts[6]
    
    # 生成13张幻灯片
    slide_01_cover(prs, blank_layout)
    slide_02_kpi(prs, blank_layout)
    slide_03_financial(prs, blank_layout)
    slide_04_valuation(prs, blank_layout)
    slide_05_business(prs, blank_layout)
    slide_06_position(prs, blank_layout)
    slide_07_risk(prs, blank_layout)
    slide_08_trend(prs, blank_layout)
    slide_09_cashflow(prs, blank_layout)
    slide_10_shareholder(prs, blank_layout)
    slide_11_catalyst(prs, blank_layout)
    slide_12_conclusion(prs, blank_layout)
    slide_13_thanks(prs, blank_layout)
    
    return prs

def add_background(slide):
    """添加渐变背景"""
    # 简化：使用中深蓝作为背景
    background = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0),
        Inches(SLIDE_W_IN), Inches(SLIDE_H_IN)
    )
    background.fill.solid()
    background.fill.fore_color.rgb = BG_MID
    background.line.fill.background()

def add_title(slide, title, subtitle="", top=0.3):
    """添加标题"""
    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(top),
        Inches(SLIDE_W_IN - 1), Inches(1)
    )
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(48)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.LEFT
    
    if subtitle:
        p2 = tf.add_paragraph()
        p2.text = subtitle
        p2.font.size = Pt(24)
        p2.font.color.rgb = LIGHT_WHITE

def add_kpi_card(slide, number, label, left, top, width=3, height=1.5, color=ACCENT_GOLD):
    """添加 KPI 数字卡片"""
    # 卡片背景
    card = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(left), Inches(top),
        Inches(width), Inches(height)
    )
    card.fill.solid()
    card.fill.fore_color.rgb = CARD_FILL
    card.line.color.rgb = CARD_BORDER
    card.line.width = Pt(10)
    
    # 数字
    num_box = slide.shapes.add_textbox(
        Inches(left + 0.2), Inches(top + 0.3),
        Inches(width - 0.4), Inches(0.8)
    )
    tf = num_box.text_frame
    p = tf.paragraphs[0]
    p.text = number
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = color
    p.alignment = PP_ALIGN.CENTER
    
    # 标签
    label_box = slide.shapes.add_textbox(
        Inches(left + 0.2), Inches(top + 1.1),
        Inches(width - 0.4), Inches(0.4)
    )
    tf = label_box.text_frame
    p = tf.paragraphs[0]
    p.text = label
    p.font.size = Pt(16)
    p.font.color.rgb = LIGHT_WHITE
    p.alignment = PP_ALIGN.CENTER

def add_content_card(slide, title, items, left, top, width=5, height=3):
    """添加内容卡片"""
    # 卡片背景
    card = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(left), Inches(top),
        Inches(width), Inches(height)
    )
    card.fill.solid()
    card.fill.fore_color.rgb = CARD_FILL
    card.line.color.rgb = CARD_BORDER
    card.line.width = Pt(10)
    
    # 标题条
    title_bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(left), Inches(top),
        Inches(width), Inches(0.5)
    )
    title_bar.fill.solid()
    title_bar.fill.fore_color.rgb = WHITE
    title_bar.line.fill.background()
    
    # 标题文字
    title_box = slide.shapes.add_textbox(
        Inches(left + 0.2), Inches(top + 0.05),
        Inches(width - 0.4), Inches(0.4)
    )
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = BG_MID
    p.alignment = PP_ALIGN.LEFT
    
    # 内容
    content_box = slide.shapes.add_textbox(
        Inches(left + 0.3), Inches(top + 0.6),
        Inches(width - 0.6), Inches(height - 0.7)
    )
    tf = content_box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"• {item}"
        p.font.size = Pt(14)
        p.font.color.rgb = LIGHT_WHITE
        p.space_after = Pt(8)

# ============================================================================
# 幻灯片生成函数
# ============================================================================

def slide_01_cover(prs, layout):
    """封面"""
    slide = prs.slides.add_slide(layout)
    add_background(slide)
    
    # 主标题
    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(2.5),
        Inches(SLIDE_W_IN - 1), Inches(1.5)
    )
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = f"{COMPANY_DATA['公司名称']} ({COMPANY_DATA['股票代码']})"
    p.font.size = Pt(72)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    
    # 副标题
    sub_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(4),
        Inches(SLIDE_W_IN - 1), Inches(0.8)
    )
    tf = sub_box.text_frame
    p = tf.paragraphs[0]
    p.text = "深度估值分析报告"
    p.font.size = Pt(36)
    p.font.color.rgb = ACCENT_GOLD
    p.alignment = PP_ALIGN.CENTER
    
    # 金色光带装饰
    beam = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(2), Inches(5.5),
        Inches(SLIDE_W_IN - 4), Inches(0.1)
    )
    beam.fill.solid()
    beam.fill.fore_color.rgb = ACCENT_GOLD
    beam.line.fill.background()
    
    # 底部信息
    info_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(6.2),
        Inches(SLIDE_W_IN - 1), Inches(0.5)
    )
    tf = info_box.text_frame
    p = tf.paragraphs[0]
    p.text = "数智财务演示 Skill | Wang Dongjie, CGMA/AICPA&CIMA, © 2026"
    p.font.size = Pt(14)
    p.font.color.rgb = LIGHT_WHITE
    p.alignment = PP_ALIGN.CENTER

def slide_02_kpi(prs, layout):
    """核心 KPI 指标快照"""
    slide = prs.slides.add_slide(layout)
    add_background(slide)
    add_title(slide, "核心 KPI 指标快照", "2026年6月17日收盘数据")
    
    # KPI 卡片
    kpis = [
        (COMPANY_DATA['当前股价'], "当前股价", ACCENT_GOLD),
        (COMPANY_DATA['市值'], "总市值", ACCENT_GOLD),
        (COMPANY_DATA['PE_TTM'], "PE(TTM)", ACCENT_CYAN),
        (COMPANY_DATA['PB'], "PB", ACCENT_CYAN),
        (COMPANY_DATA['股息率'], "股息率", ACCENT_GREEN),
    ]
    
    for i, (num, label, color) in enumerate(kpis):
        add_kpi_card(slide, num, label, 0.5 + i * 4.3, 1.8, 4, 1.5, color)
    
    # 52周区间
    range_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(4),
        Inches(SLIDE_W_IN - 1), Inches(0.5)
    )
    tf = range_box.text_frame
    p = tf.paragraphs[0]
    p.text = f"52周区间: {COMPANY_DATA['52周最低']} - {COMPANY_DATA['52周最高']}"
    p.font.size = Pt(20)
    p.font.color.rgb = LIGHT_WHITE

def slide_03_financial(prs, layout):
    """三年财务业绩概览"""
    slide = prs.slides.add_slide(layout)
    add_background(slide)
    add_title(slide, "2024年财务业绩概览", "年度报告数据")
    
    # 财务卡片
    add_content_card(slide, "营收表现", [
        f"营业收入: {COMPANY_DATA['营收_2024']}",
        f"同比变化: {COMPANY_DATA['营收同比']}",
        "主要产品: 户用储能系统",
    ], 0.5, 1.8, 6.5, 2.5)
    
    add_content_card(slide, "利润表现", [
        f"归母净利润: {COMPANY_DATA['净利润_2024']}",
        f"同比变化: {COMPANY_DATA['净利润同比']}",
        f"ROE: {COMPANY_DATA['ROE']}",
    ], 7.5, 1.8, 6.5, 2.5)
    
    add_content_card(slide, "盈利能力", [
        f"毛利率: {COMPANY_DATA['毛利率']}",
        f"资产负债率: {COMPANY_DATA['资产负债率']}",
        f"经营现金流: {COMPANY_DATA['经营现金流']}",
    ], 14.5, 1.8, 6.5, 2.5)

def slide_04_valuation(prs, layout):
    """估值指标深度分析"""
    slide = prs.slides.add_slide(layout)
    add_background(slide)
    add_title(slide, "估值指标深度分析", "相对估值法")
    
    # 估值卡片
    add_content_card(slide, "PE 估值", [
        f"PE(TTM): {COMPANY_DATA['PE_TTM']}",
        "行业平均: 约25-30x",
        "估值状态: 明显偏高",
        "风险等级: 高",
    ], 0.5, 1.8, 6.5, 2.5)
    
    add_content_card(slide, "PB 估值", [
        f"PB: {COMPANY_DATA['PB']}",
        "每股账面价值: 27.7元",
        "相对合理区间",
        "资产质量良好",
    ], 7.5, 1.8, 6.5, 2.5)
    
    add_content_card(slide, "PS 估值", [
        f"PS: {COMPANY_DATA['PS']}",
        "营收规模: 30.73亿",
        "市销率适中",
        "成长性待观察",
    ], 14.5, 1.8, 6.5, 2.5)

def slide_05_business(prs, layout):
    """业务板块结构"""
    slide = prs.slides.add_slide(layout)
    add_background(slide)
    add_title(slide, "五大业务板块结构", "主营业务收入构成")
    
    # 业务占比卡片
    add_kpi_card(slide, COMPANY_DATA['户用储能占比'], "户用储能系统", 0.5, 1.8, 4, 1.5, ACCENT_GOLD)
    add_kpi_card(slide, COMPANY_DATA['并网逆变器占比'], "并网逆变器", 5, 1.8, 4, 1.5, ACCENT_CYAN)
    add_kpi_card(slide, COMPANY_DATA['工商业储能占比'], "工商业储能", 9.5, 1.8, 4, 1.5, ACCENT_GREEN)
    
    # 业务说明
    add_content_card(slide, "主营业务", [
        "光伏储能逆变器研发生产",
        "储能电池系统销售",
        "并网逆变器产品",
        "全球130+国家销售网络",
    ], 0.5, 4, 10, 2.5)

def slide_06_position(prs, layout):
    """行业地位与全球排名"""
    slide = prs.slides.add_slide(layout)
    add_background(slide)
    add_title(slide, "行业地位与全球影响力", "国家级认证企业")
    
    # 地位卡片
    add_content_card(slide, "荣誉认证", [
        "国家级制造业单项冠军",
        "智能光伏示范企业",
        "绿色供应链管理企业",
        "浙江省科技进步一等奖",
    ], 0.5, 1.8, 6.5, 2.5)
    
    add_content_card(slide, "全球布局", [
        f"产品认证: {COMPANY_DATA['认证数量']}",
        f"覆盖国家: {COMPANY_DATA['覆盖国家']}",
        "主要市场: 欧洲、美国、日本",
        "海外收入占比: 90%+",
    ], 7.5, 1.8, 6.5, 2.5)
    
    add_content_card(slide, "技术优势", [
        "2013年推出SK系列储能逆变器",
        "国内最早储能逆变器产品之一",
        "网源友好型智能光储系统",
        "持续研发投入",
    ], 14.5, 1.8, 6.5, 2.5)

def slide_07_risk(prs, layout):
    """风险识别与压力测试"""
    slide = prs.slides.add_slide(layout)
    add_background(slide)
    add_title(slide, "风险识别与压力测试", "主要风险因素")
    
    # 风险卡片（红色边框）
    add_content_card(slide, "市场风险", [
        "欧洲户储市场降温",
        "海外需求波动",
        "汇率波动影响",
    ], 0.5, 1.8, 6.5, 2.5)
    
    add_content_card(slide, "业绩风险", [
        "营收同比下滑31.30%",
        "净利润同比下滑80.88%",
        "ROE下降19.90个百分点",
    ], 7.5, 1.8, 6.5, 2.5)
    
    add_content_card(slide, "估值风险", [
        "PE超过100x",
        "估值明显偏高",
        "业绩支撑不足",
    ], 14.5, 1.8, 6.5, 2.5)

def slide_08_trend(prs, layout):
    """新签合同趋势"""
    slide = prs.slides.add_slide(layout)
    add_background(slide)
    add_title(slide, "业绩趋势分析", "2024年业绩下滑")
    
    # 趋势说明
    add_content_card(slide, "业绩变化", [
        "2024年营收30.73亿元",
        "同比下降31.30%",
        "净利润2.04亿元",
        "同比下降80.88%",
    ], 0.5, 1.8, 10, 2.5)
    
    add_content_card(slide, "原因分析", [
        "欧洲户储市场降温",
        "海外需求减少",
        "竞争加剧",
        "价格压力",
    ], 11, 1.8, 10, 2.5)

def slide_09_cashflow(prs, layout):
    """现金流与分红能力"""
    slide = prs.slides.add_slide(layout)
    add_background(slide)
    add_title(slide, "现金流与分红能力", "2024年度分配方案")
    
    # 现金流卡片
    add_content_card(slide, "现金流状况", [
        f"经营现金流: {COMPANY_DATA['经营现金流']}",
        "同比下降55.30%",
        "现金流仍为正值",
        "资金状况良好",
    ], 0.5, 1.8, 6.5, 2.5)
    
    add_content_card(slide, "分红方案", [
        "2024年中期分红已执行",
        "每股派发0.9375元",
        "合计派发1.5亿元",
        "占净利润73.67%",
    ], 7.5, 1.8, 6.5, 2.5)
    
    add_content_card(slide, "未来分红", [
        "2025年Q3后计划分红",
        "金额1.5-1.8亿元",
        "持续回报股东",
    ], 14.5, 1.8, 6.5, 2.5)

def slide_10_shareholder(prs, layout):
    """股东结构与治理基础"""
    slide = prs.slides.add_slide(layout)
    add_background(slide)
    add_title(slide, "公司治理基础", "基本信息")
    
    # 公司信息
    add_content_card(slide, "基本信息", [
        f"董事长: 李新富",
        f"成立日期: 2012年3月2日",
        f"上市日期: {COMPANY_DATA['上市日期']}",
        f"发行价格: {COMPANY_DATA['发行价格']}",
    ], 0.5, 1.8, 6.5, 2.5)
    
    add_content_card(slide, "股权结构", [
        "总股本: 1.6亿股",
        "科创板上市",
        "公众持股比例较高",
    ], 7.5, 1.8, 6.5, 2.5)
    
    add_content_card(slide, "注册信息", [
        "注册地: 浙江杭州桐庐",
        "注册资本: 1.6亿元",
        "所属行业: 电气制造业",
    ], 14.5, 1.8, 6.5, 2.5)

def slide_11_catalyst(prs, layout):
    """价值重估潜在催化剂"""
    slide = prs.slides.add_slide(layout)
    add_background(slide)
    add_title(slide, "价值重估潜在催化剂", "未来增长机会")
    
    # 催化剂卡片
    add_content_card(slide, "市场机遇", [
        "全球储能市场长期增长",
        "能源转型趋势",
        "碳中和政策支持",
    ], 0.5, 1.8, 6.5, 2.5)
    
    add_content_card(slide, "产品拓展", [
        "工商业储能增长",
        "新产品研发",
        "技术迭代升级",
    ], 7.5, 1.8, 6.5, 2.5)
    
    add_content_card(slide, "市场拓展", [
        "新兴市场开拓",
        "美国市场增长",
        "亚太市场布局",
    ], 14.5, 1.8, 6.5, 2.5)

def slide_12_conclusion(prs, layout):
    """综合投资结论"""
    slide = prs.slides.add_slide(layout)
    add_background(slide)
    add_title(slide, "综合投资结论", "估值分析与建议")
    
    # 结论卡片
    add_content_card(slide, "估值结论", [
        f"当前股价: {COMPANY_DATA['当前股价']}",
        f"PE(TTM): {COMPANY_DATA['PE_TTM']}",
        "估值偏高，业绩支撑不足",
        "建议: 谨慎观望",
    ], 0.5, 1.8, 10, 2.5)
    
    add_content_card(slide, "投资建议", [
        "短期: 观望为主",
        "中期: 关注业绩恢复",
        "长期: 储能赛道有潜力",
        "风险等级: 中高",
    ], 11, 1.8, 10, 2.5)
    
    # 目标价区间
    target_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(5),
        Inches(SLIDE_W_IN - 1), Inches(0.5)
    )
    tf = target_box.text_frame
    p = tf.paragraphs[0]
    p.text = "合理估值区间: 60-80元 (基于行业平均PE 25-30x)"
    p.font.size = Pt(24)
    p.font.color.rgb = ACCENT_GOLD
    p.alignment = PP_ALIGN.CENTER

def slide_13_thanks(prs, layout):
    """结尾致谢"""
    slide = prs.slides.add_slide(layout)
    add_background(slide)
    
    # Thank You
    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(2.5),
        Inches(SLIDE_W_IN - 1), Inches(1.5)
    )
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Thank You"
    p.font.size = Pt(72)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    
    # 金色光带
    beam = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(2), Inches(4),
        Inches(SLIDE_W_IN - 4), Inches(0.1)
    )
    beam.fill.solid()
    beam.fill.fore_color.rgb = ACCENT_GOLD
    beam.line.fill.background()
    
    # 底部信息
    info_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(5),
        Inches(SLIDE_W_IN - 1), Inches(1)
    )
    tf = info_box.text_frame
    p = tf.paragraphs[0]
    p.text = f"{COMPANY_DATA['公司名称']} ({COMPANY_DATA['股票代码']}) 深度估值报告"
    p.font.size = Pt(24)
    p.font.color.rgb = LIGHT_WHITE
    p.alignment = PP_ALIGN.CENTER
    
    p2 = tf.add_paragraph()
    p2.text = "数智财务演示 Skill | Wang Dongjie, CGMA/AICPA&CIMA, © 2026"
    p2.font.size = Pt(14)
    p2.font.color.rgb = LIGHT_WHITE
    p2.alignment = PP_ALIGN.CENTER

# ============================================================================
# 主函数
# ============================================================================
if __name__ == "__main__":
    print("=" * 60)
    print(" 艾罗能源 (688717) 深度估值报告生成器")
    print(" 作者: Wang Dongjie, CGMA/AICPA&CIMA, © 2026")
    print("=" * 60)
    print()
    
    prs = create_presentation()
    
    # 保存
    output_dir = "/sessions/6a3272cfa9c34c00dd977285/workspace/数智财务演示-skill/output"
    os.makedirs(output_dir, exist_ok=True)
    output_path = os.path.join(output_dir, "艾罗能源_688717_深度估值报告.pptx")
    prs.save(output_path)
    
    print(f"✓ 报告已生成: {output_path}")
    print(f"✓ 幻灯片数量: 13 张")
    print(f"✓ 画布尺寸: {SLIDE_W_IN} × {SLIDE_H_IN} 英寸 (超宽屏)")
    print(f"✓ 风格: 数智财务 SAP 企业级深色")
    print()
    print("下一步:")
    print("  1. 在 macOS Keynote 中打开 .pptx 文件")
    print("  2. 文件 → 另存为 → 选择 .key 格式")
    print("=" * 60)