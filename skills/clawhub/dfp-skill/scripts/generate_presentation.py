#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
DFP-Skill Presentation Generator (v6.0)
======================================
Digital Finance Presentation Skill - SAP Enterprise Ultra-Wide Style

Author: Wang Dongjie, CGMA/AICPA&CIMA, © 2026

Professional presentation generator for enterprise summits, financial 
digitalization, and company valuation reports.

Featured Style: Digital-Finance SAP Enterprise Deep Style (3:1 Ultra-Wide)

Supported Styles:
1. digital-finance ⭐ - SAP Enterprise Deep Style (Ultra-Wide 3:1)
2. classic           - Apple Keynote Style (16:9)
3. tech              - Tech Launch Style (16:9)
4. brand             - Brand Launch Style (16:9)
5. feature           - Feature Release Style (16:9)

SAP Digital Finance Methodology (v6.0):
- 51年技术演进: ERP → HANA → Business AI
- Joule: SAP生成式AI数字助手 (真正了解您的业务)
- RPA: 93.75%效率提升案例 (供应商发票校验流程)
- ESG碳资产管理: 碳数据 → 碳资产 → 碳资本
- 数智财务四大能力: 利润/收入/资金/应收穿透分析

Keynote Professional Advantages (v5.0):
- Screen Adaptation: Ultra-wide 3:1, auto-fit, Retina optimization
- Animation Effects: Magic Move, 60fps transitions, GPU acceleration
- Design Superiority: 40+ themes, smart layout, professional aesthetics
- Stability: Native AppleScript, no dependencies, error recovery
- Visual Rendering: Core Animation, Metal engine, broadcast quality
- Performance: 2s startup, <0.1s transitions, low memory
- Font Rendering: PingFang SC + SF Pro, Retina clarity
- Professional Presentation: Presenter mode, MOV export, auto-play

New Report Templates (v6.0):
- 数智财务转型报告 (15 slides)
- ESG碳资产管理报告 (10 slides)
- RPA效率提升案例报告 (8 slides)

Usage:
    python3 generate_presentation.py --topic "Valuation Report" --style digital-finance
    python3 generate_presentation.py --topic "Tech Summit" --style tech --aspect-ratio ultra-wide
    python3 generate_presentation.py --list-styles
"""

import os
import sys
import json
import argparse
from pathlib import Path
from datetime import datetime

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
    from pptx.oxml.ns import qn
    from pptx.oxml import parse_xml
except ImportError:
    print("❌ 缺少依赖: python-pptx")
    print("请先安装: pip3 install python-pptx")
    sys.exit(1)


# =============================================================
# 配色方案库 (5 种风格)
# =============================================================

COLOR_SCHEMES = {
    'classic': {
        'name': '经典 Apple 风格',
        'aspect_ratio': '16:9',
        'slide_width_in': 13.333,
        'slide_height_in': 7.5,
        'background': RGBColor(0x00, 0x00, 0x00),
        'background_gradient': None,
        'text_primary': RGBColor(0xFF, 0xFF, 0xFF),
        'text_secondary': RGBColor(0xBF, 0xBF, 0xBF),
        'accent_blue': RGBColor(0x00, 0x7A, 0xFF),
        'accent_green': RGBColor(0x34, 0xC7, 0x59),
        'accent_orange': RGBColor(0xFF, 0x95, 0x00),
        'accent_red': RGBColor(0xFF, 0x3B, 0x30),
        'accent_purple': RGBColor(0xAF, 0x52, 0xDE),
        'card_border': None,
        'card_fill': None,
        'light_beam': None,
    },
    'tech': {
        'name': '科技发布会风格',
        'aspect_ratio': '16:9',
        'slide_width_in': 13.333,
        'slide_height_in': 7.5,
        'background': RGBColor(0x0A, 0x0A, 0x0F),
        'background_gradient': None,
        'text_primary': RGBColor(0xFF, 0xFF, 0xFF),
        'text_secondary': RGBColor(0x8E, 0x8E, 0x93),
        'accent_blue': RGBColor(0x00, 0xD4, 0xFF),
        'accent_green': RGBColor(0x00, 0xFF, 0x88),
        'accent_orange': RGBColor(0xFF, 0x6B, 0x00),
        'accent_red': RGBColor(0xFF, 0x2D, 0x55),
        'accent_purple': RGBColor(0x5E, 0x5E, 0xFF),
        'card_border': None,
        'card_fill': None,
        'light_beam': None,
    },
    'brand': {
        'name': '品牌发布会风格',
        'aspect_ratio': '16:9',
        'slide_width_in': 13.333,
        'slide_height_in': 7.5,
        'background': RGBColor(0xFF, 0xFF, 0xFF),
        'background_gradient': None,
        'text_primary': RGBColor(0x1D, 0x1D, 0x1F),
        'text_secondary': RGBColor(0x6E, 0x6E, 0x73),
        'accent_blue': RGBColor(0x00, 0x66, 0xCC),
        'accent_green': RGBColor(0x34, 0xC7, 0x59),
        'accent_orange': RGBColor(0xFF, 0x95, 0x00),
        'accent_red': RGBColor(0xFF, 0x3B, 0x30),
        'accent_purple': RGBColor(0x58, 0x56, 0xD6),
        'card_border': None,
        'card_fill': None,
        'light_beam': None,
    },
    'feature': {
        'name': '功能发布风格',
        'aspect_ratio': '16:9',
        'slide_width_in': 13.333,
        'slide_height_in': 7.5,
        'background': RGBColor(0xF5, 0xF5, 0xF7),
        'background_gradient': None,
        'text_primary': RGBColor(0x1D, 0x1D, 0x1F),
        'text_secondary': RGBColor(0x6E, 0x6E, 0x73),
        'accent_blue': RGBColor(0x00, 0x7A, 0xFF),
        'accent_green': RGBColor(0x34, 0xC7, 0x59),
        'accent_orange': RGBColor(0xFF, 0x95, 0x00),
        'accent_red': RGBColor(0xFF, 0x3B, 0x30),
        'accent_purple': RGBColor(0x58, 0x56, 0xD6),
        'card_border': None,
        'card_fill': None,
        'light_beam': None,
    },
    'digital-finance': {
        'name': '数智财务 SAP 企业级深色风格',
        'aspect_ratio': '3:1 超宽屏',
        'slide_width_in': 22.222,  # 3:1 比例 (约 3200x1080)
        'slide_height_in': 7.5,
        'background': RGBColor(0x12, 0x24, 0x4D),
        'background_gradient': {
            'top': RGBColor(0x0A, 0x18, 0x38),
            'middle': RGBColor(0x5A, 0x0F, 0x25),
            'bottom': RGBColor(0x8B, 0x00, 0x29),
        },
        'text_primary': RGBColor(0xFF, 0xFF, 0xFF),
        'text_secondary': RGBColor(0xE8, 0xEC, 0xF5),
        'accent_blue': RGBColor(0x25, 0xB7, 0xE0),       # 青蓝
        'accent_green': RGBColor(0x2F, 0xA4, 0x72),      # 财务绿
        'accent_orange': RGBColor(0xE6, 0xA8, 0x17),     # 橙金
        'accent_red': RGBColor(0xB8, 0x00, 0x3A),        # SAP品牌红
        'accent_purple': RGBColor(0x5E, 0x3F, 0xBE),     # 深紫 (AI)
        'accent_gold': RGBColor(0xF2, 0xB8, 0x4B),       # 装饰金色
        'accent_sap': RGBColor(0x00, 0x70, 0xD2),         # SAP官方蓝
        'card_border': RGBColor(0xB8, 0x00, 0x3A),       # 深红卡片边框
        'card_fill': RGBColor(0x0F, 0x20, 0x50),          # 深蓝卡片内底
        'card_title_bar': RGBColor(0xFF, 0xFF, 0xFF),     # 卡片标题白色底
        'light_beam': RGBColor(0xF2, 0xB8, 0x4B),          # 金色横向光带
        'hexagon_fill': RGBColor(0x1A, 0x2F, 0x5C),        # 六边形填充
    },
}


# =============================================================
# 演示文稿生成器基类
# =============================================================

class KeynotePresentationGenerator:
    """Keynote 风格发布会 PPTX 生成器 (v3.0)"""

    def __init__(self, topic, style='classic', lang='zh', aspect_ratio=None):
        self.topic = topic
        self.style = style
        self.lang = lang
        self.colors = COLOR_SCHEMES[style]

        # 设置画布尺寸
        self.slide_width_in = self.colors.get('slide_width_in', 13.333)
        self.slide_height_in = self.colors.get('slide_height_in', 7.5)

        # 允许强制覆盖宽高比
        if aspect_ratio == 'ultra-wide':
            self.slide_width_in = 22.222
            self.slide_height_in = 7.5

        self.prs = Presentation()
        self.prs.slide_width = Inches(self.slide_width_in)
        self.prs.slide_height = Inches(self.slide_height_in)

        self.slides = []

        # 计算常用边距
        self.W = self.slide_width_in
        self.H = self.slide_height_in
        self.margin_left = self.W * 0.06
        self.margin_right = self.W * 0.06
        self.margin_top = self.H * 0.10
        self.margin_bottom = self.H * 0.10

    # ---------------------------------------------------------
    # 底层工具方法
    # ---------------------------------------------------------

    def _set_background(self, slide, color=None):
        """设置幻灯片背景色"""
        if color is None:
            color = self.colors.get('background', RGBColor(0x00, 0x00, 0x00))
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = color

    def _set_gradient_background(self, slide, stops):
        """设置渐变背景（线性对角渐变）"""
        bg = slide.background
        fill = bg.fill
        fill.solid()
        fill.fore_color.rgb = stops.get('middle', RGBColor(0x12, 0x24, 0x4D))

    def _add_text(self, slide, text, left, top, width, height,
                  size=24, bold=False, color=None,
                  alignment=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
        """添加文本框"""
        if color is None:
            color = self.colors.get('text_primary', RGBColor(0xFF, 0xFF, 0xFF))

        text_box = slide.shapes.add_textbox(left, top, width, height)
        text_frame = text_box.text_frame
        text_frame.word_wrap = True
        text_frame.margin_left = 0
        text_frame.margin_right = 0
        text_frame.margin_top = 0
        text_frame.margin_bottom = 0
        text_frame.vertical_anchor = anchor

        p = text_frame.paragraphs[0]
        p.alignment = alignment
        run = p.add_run()
        run.text = text
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color

        # 设置字体
        if self.lang == 'zh':
            run.font.name = 'PingFang SC'
            # 设置中文字体兼容
            try:
                rPr = run._r.get_or_add_rPr()
                rFonts = rPr.find(qn('a:ea'))
                if rFonts is None:
                    rFonts = parse_xml(f'<a:ea xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main" typeface="PingFang SC"/>')
                    rPr.append(rFonts)
            except:
                pass
        else:
            run.font.name = 'SF Pro Display'

        return text_box

    def _add_multiline_text(self, slide, lines, left, top, width, height,
                            base_size=18, bold=False, color=None,
                            alignment=PP_ALIGN.LEFT, bullet_prefix="• "):
        """添加多行列表文本"""
        if color is None:
            color = self.colors.get('text_primary', RGBColor(0xFF, 0xFF, 0xFF))

        text_box = slide.shapes.add_textbox(left, top, width, height)
        text_frame = text_box.text_frame
        text_frame.word_wrap = True
        text_frame.margin_left = Inches(0.1)
        text_frame.margin_right = Inches(0.1)

        for i, line in enumerate(lines):
            if i == 0:
                p = text_frame.paragraphs[0]
            else:
                p = text_frame.add_paragraph()
            p.alignment = alignment
            p.space_after = Pt(8)
            run = p.add_run()
            if bullet_prefix and not line.startswith(bullet_prefix) and len(lines) > 1:
                run.text = bullet_prefix + line
            else:
                run.text = line
            run.font.size = Pt(base_size)
            run.font.bold = bold
            run.font.color.rgb = color
            if self.lang == 'zh':
                run.font.name = 'PingFang SC'
            else:
                run.font.name = 'SF Pro Display'

        return text_box

    def _add_shape(self, slide, shape_type, left, top, width, height,
                   fill_color=None, line_color=None, line_width_pt=1):
        """添加形状"""
        shape = slide.shapes.add_shape(shape_type, left, top, width, height)

        if fill_color is not None:
            shape.fill.solid()
            shape.fill.fore_color.rgb = fill_color
        else:
            shape.fill.background()

        if line_color is not None:
            shape.line.color.rgb = line_color
            shape.line.width = Pt(line_width_pt)
        else:
            shape.line.fill.background()

        return shape

    def _add_rectangle(self, slide, left, top, width, height,
                       fill_color=None, line_color=None, line_width_pt=1):
        """添加矩形"""
        return self._add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE,
                              left, top, width, height,
                              fill_color, line_color, line_width_pt)

    def _add_circle(self, slide, left, top, diameter,
                   fill_color=None, line_color=None, line_width_pt=2):
        """添加圆形"""
        return self._add_shape(slide, MSO_SHAPE.OVAL,
                              left, top, diameter, diameter,
                              fill_color, line_color, line_width_pt)

    # ---------------------------------------------------------
    # 数字财务风格专用方法
    # ---------------------------------------------------------

    def _add_sap_card(self, slide, title, items, left, top, width, height,
                     accent_color=None):
        """
        SAP风格卡片: 深红边框 + 深蓝底 + 白色标题条 + 圆角20px
        用于数智财务风格的卡片化信息承载
        """
        border_color = self.colors.get('card_border', RGBColor(0xB8, 0x00, 0x3A))
        fill_color = self.colors.get('card_fill', RGBColor(0x0F, 0x20, 0x50))
        title_bar_color = self.colors.get('card_title_bar', RGBColor(0xFF, 0xFF, 0xFF))
        gold_color = self.colors.get('accent_gold', RGBColor(0xF2, 0xB8, 0x4B))

        # 外层卡片 - 红色边框 + 深蓝底
        card = self._add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE,
                              left, top, width, height,
                              fill_color=fill_color,
                              line_color=border_color,
                              line_width_pt=4)

        # 白色标题条 (顶部约15%高度)
        title_bar_height = height * 0.18
        title_bar = self._add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE,
                                   left + Inches(0.04), top + Inches(0.04),
                                   width - Inches(0.08), title_bar_height,
                                   fill_color=title_bar_color,
                                   line_color=None)

        # 标题文字 (深色文字在白色条上)
        title_color = RGBColor(0x0F, 0x20, 0x50)
        self._add_text(slide, title,
                      left + Inches(0.2), top + Inches(0.05),
                      width - Inches(0.4), title_bar_height - Inches(0.05),
                      size=int(max(18, min(28, width / 0.8))),
                      bold=True, color=title_color,
                      alignment=PP_ALIGN.LEFT,
                      anchor=MSO_ANCHOR.MIDDLE)

        # 内容区 - 要点列表
        content_top = top + title_bar_height + Inches(0.15)
        content_height = height - title_bar_height - Inches(0.2)

        if items:
            item_size = int(min(18, height / (len(items) * 2.5)))
            item_size = max(14, item_size)
            lines_text = []
            for item in items:
                lines_text.append(item)
            self._add_multiline_text(
                slide, lines_text,
                left + Inches(0.2), content_top,
                width - Inches(0.4), content_height,
                base_size=item_size, bold=False,
                color=RGBColor(0xFF, 0xFF, 0xFF),
                alignment=PP_ALIGN.LEFT,
                bullet_prefix="• "
            )

        return card

    def _add_golden_light_beam(self, slide, top_in=None):
        """添加金色横向光带 - 装饰元素"""
        beam_color = self.colors.get('light_beam', RGBColor(0xF2, 0xB8, 0x4B))
        if top_in is None:
            top_in = self.H * 0.68
        beam_height = Inches(0.08)

        beam = self._add_shape(slide, MSO_SHAPE.RECTANGLE,
                              Inches(0), Inches(top_in),
                              Inches(self.W), beam_height,
                              fill_color=beam_color, line_color=None)
        return beam

    def _add_kpi_number_card(self, slide, number, label, left, top, width, height,
                            number_color=None, is_percent=False, percent_color=None):
        """
        KPI数字卡: 深蓝底 + 金色大数字 + 白色标签
        """
        fill_color = self.colors.get('card_fill', RGBColor(0x0F, 0x20, 0x50))
        border_color = self.colors.get('card_border', RGBColor(0xB8, 0x00, 0x3A))
        if number_color is None:
            number_color = self.colors.get('accent_gold', RGBColor(0xF2, 0xB8, 0x4B))

        # 卡片背景
        card = self._add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE,
                              left, top, width, height,
                              fill_color=fill_color,
                              line_color=border_color,
                              line_width_pt=3)

        # 大数字 (上部60%)
        number_size = int(min(60, height * 8))
        self._add_text(slide, str(number),
                      left, top + Inches(0.15),
                      width, height * 0.55,
                      size=number_size, bold=True,
                      color=number_color,
                      alignment=PP_ALIGN.CENTER,
                      anchor=MSO_ANCHOR.MIDDLE)

        # 标签 (下部)
        label_size = int(min(16, height * 2.5))
        self._add_text(slide, label,
                      left, top + height * 0.7,
                      width, height * 0.25,
                      size=label_size, bold=False,
                      color=RGBColor(0xE8, 0xEC, 0xF5),
                      alignment=PP_ALIGN.CENTER,
                      anchor=MSO_ANCHOR.MIDDLE)

        return card

    # ---------------------------------------------------------
    # 通用幻灯片模板 (16:9)
    # ---------------------------------------------------------

    def add_title_slide(self, title, subtitle=None):
        """开场封面幻灯片"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self._set_background(slide)

        # 标题 - 居中大字
        self._add_text(slide, title,
                      Inches(self.W * 0.1), Inches(self.H * 0.30),
                      Inches(self.W * 0.8), Inches(self.H * 0.25),
                      size=int(self.H * 12), bold=True,
                      color=self.colors.get('text_primary'),
                      alignment=PP_ALIGN.CENTER,
                      anchor=MSO_ANCHOR.MIDDLE)

        # 副标题
        if subtitle:
            self._add_text(slide, subtitle,
                          Inches(self.W * 0.15), Inches(self.H * 0.58),
                          Inches(self.W * 0.7), Inches(self.H * 0.1),
                          size=int(self.H * 5), bold=False,
                          color=self.colors.get('text_secondary'),
                          alignment=PP_ALIGN.CENTER,
                          anchor=MSO_ANCHOR.MIDDLE)

        # 底部装饰线
        accent = self.colors.get('accent_blue', RGBColor(0x00, 0x7A, 0xFF))
        self._add_shape(slide, MSO_SHAPE.RECTANGLE,
                       Inches(self.W * 0.45), Inches(self.H * 0.82),
                       Inches(self.W * 0.10), Inches(0.04),
                       fill_color=accent, line_color=None)

        return slide

    def add_agenda_slide(self, items, title_text="今日议程"):
        """议程概览幻灯片"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self._set_background(slide)

        # 标题
        self._add_text(slide, title_text,
                      Inches(self.W * 0.06), Inches(self.H * 0.08),
                      Inches(self.W * 0.88), Inches(self.H * 0.15),
                      size=int(self.H * 7), bold=True,
                      color=self.colors.get('text_primary'),
                      alignment=PP_ALIGN.LEFT)

        # 议程项
        item_h = (self.H * 0.75) / max(len(items), 3)
        for i, item in enumerate(items):
            y_pos = self.H * 0.22 + i * item_h

            # 序号
            self._add_text(slide, f"{i + 1:02d}",
                          Inches(self.W * 0.06), Inches(y_pos),
                          Inches(self.W * 0.08), Inches(item_h * 0.8),
                          size=int(self.H * 5), bold=True,
                          color=self.colors.get('accent_blue'))

            # 项目
            self._add_text(slide, item,
                          Inches(self.W * 0.18), Inches(y_pos),
                          Inches(self.W * 0.75), Inches(item_h * 0.8),
                          size=int(self.H * 3.5), bold=False,
                          color=self.colors.get('text_primary'))

            # 分隔线
            if i < len(items) - 1:
                self._add_shape(slide, MSO_SHAPE.RECTANGLE,
                               Inches(self.W * 0.18),
                               Inches(y_pos + item_h * 0.9),
                               Inches(self.W * 0.75), Inches(0.01),
                               fill_color=self.colors.get('text_secondary'),
                               line_color=None)

        return slide

    def add_problem_slide(self, title, points):
        """问题/痛点幻灯片"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self._set_background(slide)

        # 小标题
        self._add_text(slide, "问题",
                      Inches(self.W * 0.06), Inches(self.H * 0.05),
                      Inches(self.W * 0.88), Inches(self.H * 0.08),
                      size=int(self.H * 3), bold=False,
                      color=self.colors.get('accent_red'),
                      alignment=PP_ALIGN.LEFT)

        # 大标题
        self._add_text(slide, title,
                      Inches(self.W * 0.06), Inches(self.H * 0.12),
                      Inches(self.W * 0.88), Inches(self.H * 0.18),
                      size=int(self.H * 8), bold=True,
                      color=self.colors.get('text_primary'),
                      alignment=PP_ALIGN.LEFT)

        # 要点列表
        point_size = int(self.H * 3)
        for i, point in enumerate(points):
            y_pos = self.H * 0.35 + i * self.H * 0.15
            self._add_text(slide, f"•  {point}",
                          Inches(self.W * 0.08), Inches(y_pos),
                          Inches(self.W * 0.85), Inches(self.H * 0.12),
                          size=point_size, bold=False,
                          color=self.colors.get('text_primary'))

        return slide

    def add_feature_slide(self, feature_title, description,
                          feature_points=None, accent_color=None):
        """核心亮点幻灯片"""
        if accent_color is None:
            accent_color = self.colors.get('accent_blue')

        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self._set_background(slide)

        # 左侧大数字装饰
        self._add_text(slide, "01",
                      Inches(self.W * 0.06), Inches(self.H * 0.05),
                      Inches(self.W * 0.2), Inches(self.H * 0.2),
                      size=int(self.H * 15), bold=True,
                      color=accent_color)

        # 右侧标题
        self._add_text(slide, feature_title,
                      Inches(self.W * 0.30), Inches(self.H * 0.08),
                      Inches(self.W * 0.65), Inches(self.H * 0.18),
                      size=int(self.H * 6), bold=True,
                      color=self.colors.get('text_primary'))

        # 描述
        self._add_text(slide, description,
                      Inches(self.W * 0.30), Inches(self.H * 0.28),
                      Inches(self.W * 0.65), Inches(self.H * 0.2),
                      size=int(self.H * 3), bold=False,
                      color=self.colors.get('text_secondary'))

        # 要点
        if feature_points:
            for i, point in enumerate(feature_points):
                y_pos = self.H * 0.55 + i * self.H * 0.10
                self._add_text(slide, "+ " + point,
                              Inches(self.W * 0.30), Inches(y_pos),
                              Inches(self.W * 0.65), Inches(self.H * 0.08),
                              size=int(self.H * 2.8), bold=False,
                              color=self.colors.get('text_primary'))

        return slide

    def add_specs_slide(self, specs_dict, title_text="技术规格"):
        """技术规格幻灯片"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self._set_background(slide)

        # 标题
        self._add_text(slide, title_text,
                      Inches(self.W * 0.06), Inches(self.H * 0.05),
                      Inches(self.W * 0.88), Inches(self.H * 0.15),
                      size=int(self.H * 7), bold=True,
                      color=self.colors.get('text_primary'))

        # 规格网格
        items = list(specs_dict.items())
        cols = 2
        rows = (len(items) + cols - 1) // cols

        cell_w = (self.W * 0.88) / cols
        cell_h = (self.H * 0.65) / max(rows, 1)

        for idx, (key, value) in enumerate(items):
            col = idx % cols
            row = idx // cols

            left = self.W * 0.06 + col * cell_w
            top = self.H * 0.22 + row * cell_h

            # 规格名
            self._add_text(slide, key,
                          Inches(left), Inches(top),
                          Inches(cell_w * 0.9), Inches(cell_h * 0.3),
                          size=int(self.H * 2.5), bold=True,
                          color=self.colors.get('text_secondary'))

            # 规格值
            self._add_text(slide, value,
                          Inches(left), Inches(top + cell_h * 0.3),
                          Inches(cell_w * 0.9), Inches(cell_h * 0.6),
                          size=int(self.H * 5), bold=True,
                          color=self.colors.get('accent_blue'))

        return slide

    def add_comparison_slide(self, title, items_dict):
        """对比竞品幻灯片"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self._set_background(slide)

        # 标题
        self._add_text(slide, title,
                      Inches(self.W * 0.06), Inches(self.H * 0.05),
                      Inches(self.W * 0.88), Inches(self.H * 0.15),
                      size=int(self.H * 7), bold=True,
                      color=self.colors.get('text_primary'))

        # 对比项
        items = list(items_dict.items())
        item_h = (self.H * 0.7) / max(len(items), 3)

        for i, (feature, status) in enumerate(items):
            y_pos = self.H * 0.22 + i * item_h

            # 功能名
            self._add_text(slide, feature,
                          Inches(self.W * 0.08), Inches(y_pos),
                          Inches(self.W * 0.6), Inches(item_h * 0.8),
                          size=int(self.H * 3.5), bold=False,
                          color=self.colors.get('text_primary'))

            # 状态标记
            mark = "✓" if status else "—"
            color = self.colors.get('accent_green') if status else self.colors.get('text_secondary')
            self._add_text(slide, mark,
                          Inches(self.W * 0.75), Inches(y_pos),
                          Inches(self.W * 0.15), Inches(item_h * 0.8),
                          size=int(self.H * 5), bold=True,
                          color=color,
                          alignment=PP_ALIGN.CENTER,
                          anchor=MSO_ANCHOR.MIDDLE)

        return slide

    def add_kpi_slide(self, title, kpi_items):
        """KPI数字展示幻灯片 (数字卡)"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self._set_background(slide)

        # 标题
        self._add_text(slide, title,
                      Inches(self.W * 0.06), Inches(self.H * 0.05),
                      Inches(self.W * 0.88), Inches(self.H * 0.15),
                      size=int(self.H * 7), bold=True,
                      color=self.colors.get('text_primary'))

        # KPI卡片网格
        n = len(kpi_items)
        cols = min(4, n)
        rows = (n + cols - 1) // cols

        card_w = (self.W * 0.88) / cols
        card_h = (self.H * 0.65) / max(rows, 1)

        for i, (number, label) in enumerate(kpi_items):
            col = i % cols
            row = i // cols
            left = self.W * 0.06 + col * card_w
            top = self.H * 0.22 + row * card_h

            self._add_kpi_number_card(
                slide, number, label,
                Inches(left + card_w * 0.05),
                Inches(top + card_h * 0.08),
                Inches(card_w * 0.9),
                Inches(card_h * 0.84)
            )

        return slide

    def add_story_slide(self, quote, author=None):
        """用户故事/引言幻灯片"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self._set_background(slide)

        # 大引号
        self._add_text(slide, '"',
                      Inches(self.W * 0.05), Inches(self.H * 0.05),
                      Inches(self.W * 0.15), Inches(self.H * 0.25),
                      size=int(self.H * 18), bold=True,
                      color=self.colors.get('accent_blue'))

        # 引言
        self._add_text(slide, quote,
                      Inches(self.W * 0.15), Inches(self.H * 0.2),
                      Inches(self.W * 0.8), Inches(self.H * 0.4),
                      size=int(self.H * 5), bold=False,
                      color=self.colors.get('text_primary'))

        # 作者
        if author:
            self._add_text(slide, "— " + author,
                          Inches(self.W * 0.15), Inches(self.H * 0.65),
                          Inches(self.W * 0.8), Inches(self.H * 0.1),
                          size=int(self.H * 3), bold=False,
                          color=self.colors.get('text_secondary'))

        return slide

    def add_pricing_slide(self, product_name, price, availability=None):
        """价格与发售幻灯片"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self._set_background(slide)

        # 产品名
        self._add_text(slide, product_name,
                      Inches(self.W * 0.1), Inches(self.H * 0.12),
                      Inches(self.W * 0.8), Inches(self.H * 0.12),
                      size=int(self.H * 5), bold=True,
                      color=self.colors.get('text_primary'),
                      alignment=PP_ALIGN.CENTER)

        # 价格
        self._add_text(slide, price,
                      Inches(self.W * 0.1), Inches(self.H * 0.28),
                      Inches(self.W * 0.8), Inches(self.H * 0.35),
                      size=int(self.H * 15), bold=True,
                      color=self.colors.get('accent_blue'),
                      alignment=PP_ALIGN.CENTER,
                      anchor=MSO_ANCHOR.MIDDLE)

        # 发售信息
        if availability:
            self._add_text(slide, availability,
                          Inches(self.W * 0.1), Inches(self.H * 0.68),
                          Inches(self.W * 0.8), Inches(self.H * 0.12),
                          size=int(self.H * 4), bold=False,
                          color=self.colors.get('text_secondary'),
                          alignment=PP_ALIGN.CENTER)

        return slide

    def add_thank_you_slide(self, title="谢谢", subtitle=None):
        """Thank you 结尾页"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self._set_background(slide)

        # 大标题
        self._add_text(slide, title,
                      Inches(self.W * 0.1), Inches(self.H * 0.35),
                      Inches(self.W * 0.8), Inches(self.H * 0.25),
                      size=int(self.H * 14), bold=True,
                      color=self.colors.get('text_primary'),
                      alignment=PP_ALIGN.CENTER,
                      anchor=MSO_ANCHOR.MIDDLE)

        # 副标题
        if subtitle:
            self._add_text(slide, subtitle,
                          Inches(self.W * 0.1), Inches(self.H * 0.62),
                          Inches(self.W * 0.8), Inches(self.H * 0.1),
                          size=int(self.H * 4), bold=False,
                          color=self.colors.get('text_secondary'),
                          alignment=PP_ALIGN.CENTER)

        return slide

    # ---------------------------------------------------------
    # 超宽屏(3:1)专用布局 - 数智财务SAP风格
    # ---------------------------------------------------------

    def add_uw_cover_slide(self, title, subtitle=None, third_title=None):
        """
        超宽屏封面 (Cover)
        布局: 居中大标题 + 金色副标题 + 光带装饰
        """
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])

        # 深蓝渐变背景 (用中间色模拟)
        bg_color = self.colors.get('background', RGBColor(0x12, 0x24, 0x4D))
        self._set_background(slide, bg_color)

        # 金色横向光带
        self._add_golden_light_beam(slide, top_in=self.H * 0.72)

        # 主标题 (居中)
        title_size = int(min(72, self.H * 11))
        self._add_text(slide, title,
                      Inches(self.W * 0.10), Inches(self.H * 0.22),
                      Inches(self.W * 0.80), Inches(self.H * 0.25),
                      size=title_size, bold=True,
                      color=RGBColor(0xFF, 0xFF, 0xFF),
                      alignment=PP_ALIGN.CENTER,
                      anchor=MSO_ANCHOR.MIDDLE)

        # 金色副标题
        if subtitle:
            sub_size = int(min(40, self.H * 6))
            self._add_text(slide, subtitle,
                          Inches(self.W * 0.15), Inches(self.H * 0.48),
                          Inches(self.W * 0.70), Inches(self.H * 0.10),
                          size=sub_size, bold=True,
                          color=self.colors.get('accent_gold'),
                          alignment=PP_ALIGN.CENTER,
                          anchor=MSO_ANCHOR.MIDDLE)

        # 第三标题 (青蓝色)
        if third_title:
            third_size = int(min(32, self.H * 5))
            self._add_text(slide, third_title,
                          Inches(self.W * 0.15), Inches(self.H * 0.58),
                          Inches(self.W * 0.70), Inches(self.H * 0.10),
                          size=third_size, bold=False,
                          color=self.colors.get('accent_blue'),
                          alignment=PP_ALIGN.CENTER,
                          anchor=MSO_ANCHOR.MIDDLE)

        return slide

    def add_uw_three_column_cards(self, title, card_data):
        """
        超宽屏三栏并列卡片 (Three-Column Cards)
        用途: 背景/动因/三大要点/三大方案
        card_data: [(card_title, [item1, item2, ...]), ...]
        """
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        bg_color = self.colors.get('background', RGBColor(0x12, 0x24, 0x4D))
        self._set_background(slide, bg_color)

        # 标题
        title_size = int(min(44, self.H * 6))
        self._add_text(slide, title,
                      Inches(self.W * 0.06), Inches(self.H * 0.08),
                      Inches(self.W * 0.88), Inches(self.H * 0.12),
                      size=title_size, bold=True,
                      color=RGBColor(0xFF, 0xFF, 0xFF),
                      alignment=PP_ALIGN.CENTER)

        # 三栏卡片
        n = len(card_data)
        card_w = (self.W * 0.88) / n
        card_h = self.H * 0.68
        top = self.H * 0.22

        for i, (card_title, items) in enumerate(card_data):
            left = self.W * 0.06 + i * card_w
            self._add_sap_card(
                slide, card_title, items,
                Inches(left + card_w * 0.03),
                Inches(top),
                Inches(card_w * 0.94),
                Inches(card_h)
            )

        return slide

    def add_uw_timeline(self, title, nodes):
        """
        超宽屏横向时间轴 (Timeline)
        nodes: [(year, product_name), ...]
        """
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        bg_color = self.colors.get('background', RGBColor(0x12, 0x24, 0x4D))
        self._set_background(slide, bg_color)

        # 标题
        title_size = int(min(44, self.H * 6))
        self._add_text(slide, title,
                      Inches(self.W * 0.06), Inches(self.H * 0.05),
                      Inches(self.W * 0.88), Inches(self.H * 0.12),
                      size=title_size, bold=True,
                      color=RGBColor(0xFF, 0xFF, 0xFF),
                      alignment=PP_ALIGN.CENTER)

        n = len(nodes)
        node_area_left = self.W * 0.08
        node_area_width = self.W * 0.84
        node_y = self.H * 0.45
        node_diameter = min(self.W * 0.84 / (n * 3), self.H * 0.18)
        node_diameter = max(node_diameter, 0.8)

        # 横向连接线
        line_y = node_y + node_diameter / 2
        self._add_shape(slide, MSO_SHAPE.RECTANGLE,
                       Inches(node_area_left + node_diameter / 2),
                       Inches(line_y - 0.02),
                       Inches(node_area_width - node_diameter),
                       Inches(0.04),
                       fill_color=self.colors.get('accent_blue'),
                       line_color=None)

        # 节点圆 + 年份 + 产品名
        for i, (year, name) in enumerate(nodes):
            x = node_area_left + i * (node_area_width / (n - 1)) - node_diameter / 2 if n > 1 else node_area_left

            # 圆节点
            self._add_circle(slide,
                           Inches(x), Inches(node_y),
                           Inches(node_diameter),
                           fill_color=self.colors.get('accent_blue'),
                           line_color=self.colors.get('accent_gold'),
                           line_width_pt=3)

            # 年份 (节点顶部, 金色)
            year_size = int(min(24, node_diameter * 25))
            self._add_text(slide, str(year),
                          Inches(x - node_diameter * 0.2),
                          Inches(node_y - node_diameter * 0.9),
                          Inches(node_diameter * 1.4),
                          Inches(node_diameter * 0.7),
                          size=year_size, bold=True,
                          color=self.colors.get('accent_gold'),
                          alignment=PP_ALIGN.CENTER,
                          anchor=MSO_ANCHOR.MIDDLE)

            # 产品名 (节点下方)
            name_size = int(min(18, node_diameter * 18))
            self._add_text(slide, name,
                          Inches(x - node_diameter * 0.5),
                          Inches(node_y + node_diameter * 1.2),
                          Inches(node_diameter * 2.0),
                          Inches(node_diameter * 0.6),
                          size=name_size, bold=False,
                          color=RGBColor(0xE8, 0xEC, 0xF5),
                          alignment=PP_ALIGN.CENTER,
                          anchor=MSO_ANCHOR.MIDDLE)

        return slide

    def add_uw_hexagon_matrix(self, title, center_text, elements):
        """
        超宽屏六边形蜂群矩阵 (Hexagon Matrix)
        中心为主题，周围环绕6-8个要素六边形
        """
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        bg_color = self.colors.get('background', RGBColor(0x12, 0x24, 0x4D))
        self._set_background(slide, bg_color)

        # 标题
        title_size = int(min(44, self.H * 6))
        self._add_text(slide, title,
                      Inches(self.W * 0.06), Inches(self.H * 0.05),
                      Inches(self.W * 0.88), Inches(self.H * 0.12),
                      size=title_size, bold=True,
                      color=RGBColor(0xFF, 0xFF, 0xFF),
                      alignment=PP_ALIGN.CENTER)

        # 中心主题圆
        center_x = self.W * 0.5
        center_y = self.H * 0.55
        center_diameter = min(self.H * 0.25, self.W * 0.08)

        self._add_circle(slide,
                        Inches(center_x - center_diameter / 2),
                        Inches(center_y - center_diameter / 2),
                        Inches(center_diameter),
                        fill_color=self.colors.get('accent_gold'),
                        line_color=self.colors.get('accent_gold'),
                        line_width_pt=4)

        # 中心文字
        self._add_text(slide, center_text,
                      Inches(center_x - center_diameter / 2),
                      Inches(center_y - center_diameter / 2),
                      Inches(center_diameter),
                      Inches(center_diameter),
                      size=int(min(18, center_diameter * 22)),
                      bold=True,
                      color=RGBColor(0x0F, 0x20, 0x50),
                      alignment=PP_ALIGN.CENTER,
                      anchor=MSO_ANCHOR.MIDDLE)

        # 周围六边形要素 (用圆形代替，PPTX中六边形可用但位置计算复杂)
        n = len(elements)
        ring_radius = min(self.H * 0.22, self.W * 0.12)
        hex_diameter = min(self.H * 0.15, self.W * 0.05)

        import math
        for i, elem in enumerate(elements):
            angle = (2 * math.pi * i / n) - math.pi / 2
            hx = center_x + ring_radius * math.cos(angle)
            hy = center_y + ring_radius * math.sin(angle)

            # 六边形/圆形要素
            self._add_circle(slide,
                            Inches(hx - hex_diameter / 2),
                            Inches(hy - hex_diameter / 2),
                            Inches(hex_diameter),
                            fill_color=self.colors.get('hexagon_fill'),
                            line_color=self.colors.get('accent_gold'),
                            line_width_pt=3)

            # 要素名称
            self._add_text(slide, elem,
                          Inches(hx - hex_diameter / 2),
                          Inches(hy - hex_diameter / 2),
                          Inches(hex_diameter),
                          Inches(hex_diameter),
                          size=int(min(14, hex_diameter * 20)),
                          bold=True,
                          color=RGBColor(0xFF, 0xFF, 0xFF),
                          alignment=PP_ALIGN.CENTER,
                          anchor=MSO_ANCHOR.MIDDLE)

        return slide

    def add_uw_dashboard(self, title, kpi_data, chart_title="数据趋势"):
        """
        超宽屏Dashboard数据页
        布局: 顶部KPI卡行 + 底部图表区域
        """
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        bg_color = self.colors.get('background', RGBColor(0x12, 0x24, 0x4D))
        self._set_background(slide, bg_color)

        # 标题
        title_size = int(min(44, self.H * 6))
        self._add_text(slide, title,
                      Inches(self.W * 0.06), Inches(self.H * 0.03),
                      Inches(self.W * 0.88), Inches(self.H * 0.12),
                      size=title_size, bold=True,
                      color=RGBColor(0xFF, 0xFF, 0xFF),
                      alignment=PP_ALIGN.CENTER)

        # 顶部KPI卡行 (4张卡)
        n_kpi = len(kpi_data)
        cols = min(4, n_kpi)
        card_w = (self.W * 0.88) / cols
        card_h = self.H * 0.28
        card_top = self.H * 0.18

        for i, (number, label, trend) in enumerate(kpi_data[:cols]):
            left = self.W * 0.06 + i * card_w
            # KPI卡
            card = self._add_shape(
                slide, MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(left + card_w * 0.04),
                Inches(card_top),
                Inches(card_w * 0.92),
                Inches(card_h),
                fill_color=self.colors.get('card_fill'),
                line_color=self.colors.get('card_border'),
                line_width_pt=3
            )

            # 数字
            num_size = int(min(48, card_h * 6))
            self._add_text(slide, str(number),
                          Inches(left + card_w * 0.04),
                          Inches(card_top + card_h * 0.05),
                          Inches(card_w * 0.92),
                          Inches(card_h * 0.5),
                          size=num_size, bold=True,
                          color=self.colors.get('accent_gold'),
                          alignment=PP_ALIGN.CENTER,
                          anchor=MSO_ANCHOR.MIDDLE)

            # 趋势
            trend_size = int(min(18, card_h * 2.5))
            trend_color = self.colors.get('accent_green') if trend.startswith('↑') else (
                self.colors.get('accent_red') if trend.startswith('↓') else RGBColor(0xE8, 0xEC, 0xF5)
            )
            self._add_text(slide, trend,
                          Inches(left + card_w * 0.04),
                          Inches(card_top + card_h * 0.55),
                          Inches(card_w * 0.92),
                          Inches(card_h * 0.2),
                          size=trend_size, bold=True,
                          color=trend_color,
                          alignment=PP_ALIGN.CENTER,
                          anchor=MSO_ANCHOR.MIDDLE)

            # 标签
            label_size = int(min(16, card_h * 2))
            self._add_text(slide, label,
                          Inches(left + card_w * 0.04),
                          Inches(card_top + card_h * 0.75),
                          Inches(card_w * 0.92),
                          Inches(card_h * 0.2),
                          size=label_size, bold=False,
                          color=RGBColor(0xE8, 0xEC, 0xF5),
                          alignment=PP_ALIGN.CENTER,
                          anchor=MSO_ANCHOR.MIDDLE)

        # 底部装饰图表区域
        chart_top = self.H * 0.50
        chart_area_h = self.H * 0.42

        # 左侧: 模拟图表
        chart_area_left = self.W * 0.06
        chart_area_width = self.W * 0.55

        # 图表标题
        self._add_text(slide, chart_title,
                      Inches(chart_area_left),
                      Inches(chart_top),
                      Inches(chart_area_width),
                      Inches(self.H * 0.08),
                      size=24, bold=True,
                      color=RGBColor(0xFF, 0xFF, 0xFF),
                      alignment=PP_ALIGN.LEFT)

        # 模拟柱状图
        bar_area_top = chart_top + self.H * 0.10
        bar_area_h = chart_area_h - self.H * 0.10
        n_bars = 8
        bar_colors = [
            self.colors.get('accent_blue'),
            self.colors.get('accent_purple'),
            self.colors.get('accent_green'),
        ]

        bar_area_width = chart_area_width
        bar_gap = bar_area_width / n_bars
        bar_w = bar_gap * 0.5

        heights = [0.45, 0.65, 0.55, 0.85, 0.70, 0.90, 0.60, 0.78]
        for i, h in enumerate(heights):
            bar_x = chart_area_left + i * bar_gap + bar_gap * 0.25
            bar_h_pixel = bar_area_h * h
            bar_y = bar_area_top + bar_area_h - bar_h_pixel

            self._add_shape(slide, MSO_SHAPE.RECTANGLE,
                          Inches(bar_x), Inches(bar_y),
                          Inches(bar_w), Inches(bar_h_pixel),
                          fill_color=bar_colors[i % len(bar_colors)],
                          line_color=None)

            # 数字标签
            self._add_text(slide, f"{int(h*100)}%",
                          Inches(bar_x), Inches(bar_y - self.H * 0.06),
                          Inches(bar_w), Inches(self.H * 0.05),
                          size=12, bold=True,
                          color=RGBColor(0xE8, 0xEC, 0xF5),
                          alignment=PP_ALIGN.CENTER)

        # 右侧: 矩阵要素
        matrix_left = chart_area_left + chart_area_width + self.W * 0.02
        matrix_width = self.W - matrix_left - self.W * 0.06

        self._add_text(slide, "核心要素",
                      Inches(matrix_left),
                      Inches(chart_top),
                      Inches(matrix_width),
                      Inches(self.H * 0.08),
                      size=24, bold=True,
                      color=RGBColor(0xFF, 0xFF, 0xFF),
                      alignment=PP_ALIGN.LEFT)

        # 要素小方块
        elements = ["预算编制", "成本控制", "资金管理", "财务分析", "风险管控", "税务合规"]
        elem_grid_cols = 2
        elem_size = min((matrix_width - Inches(0.2)) / elem_grid_cols - Inches(0.1),
                       (chart_area_h - self.H * 0.15) / 3 - Inches(0.1))

        for i, elem in enumerate(elements):
            row = i // elem_grid_cols
            col = i % elem_grid_cols
            ex = matrix_left + col * (elem_size + Inches(0.1))
            ey = chart_top + self.H * 0.10 + row * (elem_size + Inches(0.08))

            self._add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE,
                          Inches(ex), Inches(ey),
                          Inches(elem_size), Inches(elem_size * 0.8),
                          fill_color=self.colors.get('hexagon_fill'),
                          line_color=self.colors.get('accent_gold'),
                          line_width_pt=2)

            self._add_text(slide, elem,
                          Inches(ex), Inches(ey),
                          Inches(elem_size), Inches(elem_size * 0.8),
                          size=14, bold=True,
                          color=RGBColor(0xFF, 0xFF, 0xFF),
                          alignment=PP_ALIGN.CENTER,
                          anchor=MSO_ANCHOR.MIDDLE)

        return slide

    def add_uw_ring_architecture(self, title, center_text, features):
        """
        超宽屏中央环形架构图 (Ring Architecture)
        中央主题 + 外圈功能点
        """
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        bg_color = self.colors.get('background', RGBColor(0x12, 0x24, 0x4D))
        self._set_background(slide, bg_color)

        # 标题
        title_size = int(min(44, self.H * 6))
        self._add_text(slide, title,
                      Inches(self.W * 0.06), Inches(self.H * 0.03),
                      Inches(self.W * 0.88), Inches(self.H * 0.12),
                      size=title_size, bold=True,
                      color=RGBColor(0xFF, 0xFF, 0xFF),
                      alignment=PP_ALIGN.CENTER)

        # 中心圆 (大圆环)
        center_x = self.W * 0.5
        center_y = self.H * 0.55

        outer_d = min(self.H * 0.55, self.W * 0.18)
        inner_d = outer_d * 0.65
        core_d = outer_d * 0.35

        # 外圆环
        self._add_circle(slide,
                        Inches(center_x - outer_d / 2),
                        Inches(center_y - outer_d / 2),
                        Inches(outer_d),
                        fill_color=None,
                        line_color=self.colors.get('accent_blue'),
                        line_width_pt=4)

        # 内圆环
        self._add_circle(slide,
                        Inches(center_x - inner_d / 2),
                        Inches(center_y - inner_d / 2),
                        Inches(inner_d),
                        fill_color=self.colors.get('card_fill'),
                        line_color=self.colors.get('accent_gold'),
                        line_width_pt=3)

        # 核心圆
        self._add_circle(slide,
                        Inches(center_x - core_d / 2),
                        Inches(center_y - core_d / 2),
                        Inches(core_d),
                        fill_color=self.colors.get('accent_gold'),
                        line_color=None)

        # 核心文字
        self._add_text(slide, center_text,
                      Inches(center_x - core_d / 2),
                      Inches(center_y - core_d / 2),
                      Inches(core_d), Inches(core_d),
                      size=int(min(18, core_d * 22)),
                      bold=True,
                      color=RGBColor(0x0F, 0x20, 0x50),
                      alignment=PP_ALIGN.CENTER,
                      anchor=MSO_ANCHOR.MIDDLE)

        # 外圈功能点
        n = len(features)
        ring_r = outer_d * 0.58
        feat_d = min(self.H * 0.10, self.W * 0.035)

        import math
        for i, feat in enumerate(features):
            angle = (2 * math.pi * i / n) - math.pi / 2
            fx = center_x + ring_r * math.cos(angle)
            fy = center_y + ring_r * math.sin(angle)

            # 功能点圆
            self._add_circle(slide,
                            Inches(fx - feat_d / 2),
                            Inches(fy - feat_d / 2),
                            Inches(feat_d),
                            fill_color=self.colors.get('accent_purple'),
                            line_color=self.colors.get('accent_gold'),
                            line_width_pt=2)

            # 功能名 (在圆下方或旁边)
            self._add_text(slide, feat,
                          Inches(fx - feat_d),
                          Inches(fy + feat_d * 0.6),
                          Inches(feat_d * 3),
                          Inches(feat_d * 0.8),
                          size=int(min(14, feat_d * 30)),
                          bold=True,
                          color=RGBColor(0xFF, 0xFF, 0xFF),
                          alignment=PP_ALIGN.CENTER,
                          anchor=MSO_ANCHOR.MIDDLE)

        return slide

    def add_uw_quadrant_matrix(self, title, quadrants):
        """
        超宽屏四象限功能矩阵 (Quadrant Matrix)
        quadrants: [(title, [features]), ...] x 4
        """
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        bg_color = self.colors.get('background', RGBColor(0x12, 0x24, 0x4D))
        self._set_background(slide, bg_color)

        # 标题
        title_size = int(min(44, self.H * 6))
        self._add_text(slide, title,
                      Inches(self.W * 0.06), Inches(self.H * 0.03),
                      Inches(self.W * 0.88), Inches(self.H * 0.12),
                      size=title_size, bold=True,
                      color=RGBColor(0xFF, 0xFF, 0xFF),
                      alignment=PP_ALIGN.CENTER)

        # 2x2网格
        q_colors = [
            self.colors.get('accent_blue'),
            self.colors.get('accent_purple'),
            self.colors.get('accent_green'),
            self.colors.get('accent_orange'),
        ]

        cell_w = (self.W * 0.88) / 2
        cell_h = (self.H * 0.75) / 2

        for i, (q_title, features) in enumerate(quadrants[:4]):
            row = i // 2
            col = i % 2
            left = self.W * 0.06 + col * cell_w
            top = self.H * 0.20 + row * cell_h

            # 象限卡片
            self._add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE,
                          Inches(left + cell_w * 0.02),
                          Inches(top + cell_h * 0.03),
                          Inches(cell_w * 0.96),
                          Inches(cell_h * 0.92),
                          fill_color=self.colors.get('card_fill'),
                          line_color=q_colors[i],
                          line_width_pt=3)

            # 象限标题
            title_h = cell_h * 0.20
            self._add_text(slide, q_title,
                          Inches(left + cell_w * 0.05),
                          Inches(top + cell_h * 0.06),
                          Inches(cell_w * 0.9),
                          Inches(title_h),
                          size=int(min(22, cell_h * 2.5)),
                          bold=True,
                          color=q_colors[i],
                          alignment=PP_ALIGN.LEFT,
                          anchor=MSO_ANCHOR.MIDDLE)

            # 功能要点列表
            feat_lines = []
            for f in features[:5]:
                feat_lines.append(f)

            self._add_multiline_text(
                slide, feat_lines,
                Inches(left + cell_w * 0.05),
                Inches(top + title_h + cell_h * 0.08),
                Inches(cell_w * 0.9),
                Inches(cell_h * 0.65),
                base_size=int(min(16, cell_h * 1.8)),
                bold=False,
                color=RGBColor(0xFF, 0xFF, 0xFF),
                alignment=PP_ALIGN.LEFT
            )

        return slide

    def add_uw_four_step_process(self, title, steps):
        """
        超宽屏四步流程卡 (Four-Step Process)
        steps: [(step_name, description), ...]
        """
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        bg_color = self.colors.get('background', RGBColor(0x12, 0x24, 0x4D))
        self._set_background(slide, bg_color)

        # 标题
        title_size = int(min(44, self.H * 6))
        self._add_text(slide, title,
                      Inches(self.W * 0.06), Inches(self.H * 0.03),
                      Inches(self.W * 0.88), Inches(self.H * 0.12),
                      size=title_size, bold=True,
                      color=RGBColor(0xFF, 0xFF, 0xFF),
                      alignment=PP_ALIGN.CENTER)

        n = len(steps)
        card_w = (self.W * 0.88) / n
        card_h = self.H * 0.65
        top = self.H * 0.22

        arrow_colors = [
            self.colors.get('accent_blue'),
            self.colors.get('accent_purple'),
            self.colors.get('accent_green'),
            self.colors.get('accent_gold'),
        ]

        for i, (step_name, desc) in enumerate(steps):
            left = self.W * 0.06 + i * card_w

            # 序号圆 (金色)
            num_d = min(card_h * 0.20, self.W * 0.015)
            self._add_circle(slide,
                            Inches(left + card_w * 0.5 - num_d / 2),
                            Inches(top),
                            Inches(num_d),
                            fill_color=self.colors.get('accent_gold'),
                            line_color=None)

            self._add_text(slide, str(i + 1),
                          Inches(left + card_w * 0.5 - num_d / 2),
                          Inches(top),
                          Inches(num_d), Inches(num_d),
                          size=int(min(20, num_d * 25)),
                          bold=True,
                          color=RGBColor(0x0F, 0x20, 0x50),
                          alignment=PP_ALIGN.CENTER,
                          anchor=MSO_ANCHOR.MIDDLE)

            # 卡片内容
            card_top = top + num_d + Inches(0.1)
            self._add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE,
                          Inches(left + card_w * 0.05),
                          Inches(card_top),
                          Inches(card_w * 0.90),
                          Inches(card_h - num_d - Inches(0.1)),
                          fill_color=self.colors.get('card_fill'),
                          line_color=arrow_colors[i % len(arrow_colors)],
                          line_width_pt=3)

            # 步骤名
            self._add_text(slide, step_name,
                          Inches(left + card_w * 0.05),
                          Inches(card_top + card_h * 0.05),
                          Inches(card_w * 0.90),
                          Inches(card_h * 0.15),
                          size=int(min(22, card_h * 2.5)),
                          bold=True,
                          color=RGBColor(0xFF, 0xFF, 0xFF),
                          alignment=PP_ALIGN.CENTER,
                          anchor=MSO_ANCHOR.MIDDLE)

            # 描述
            self._add_text(slide, desc,
                          Inches(left + card_w * 0.08),
                          Inches(card_top + card_h * 0.25),
                          Inches(card_w * 0.84),
                          Inches(card_h * 0.65),
                          size=int(min(16, card_h * 2)),
                          bold=False,
                          color=RGBColor(0xE8, 0xEC, 0xF5),
                          alignment=PP_ALIGN.CENTER,
                          anchor=MSO_ANCHOR.MIDDLE)

            # 箭头 (如果不是最后一个)
            if i < n - 1:
                arrow_x = left + card_w + card_w * 0.01
                arrow_y = top + num_d / 2 - Inches(0.1)
                self._add_shape(slide, MSO_SHAPE.RIGHT_ARROW,
                              Inches(arrow_x - Inches(0.2)),
                              Inches(arrow_y),
                              Inches(0.4), Inches(0.2),
                              fill_color=arrow_colors[i],
                              line_color=None)

        return slide

    def add_uw_before_after(self, title, before_items, after_items, improvement_text):
        """
        超宽屏 Before/After 对比页
        左: Before (问题) | 右: After (方案) | 中间: 提升数字
        """
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        bg_color = self.colors.get('background', RGBColor(0x12, 0x24, 0x4D))
        self._set_background(slide, bg_color)

        # 标题
        title_size = int(min(44, self.H * 6))
        self._add_text(slide, title,
                      Inches(self.W * 0.06), Inches(self.H * 0.03),
                      Inches(self.W * 0.88), Inches(self.H * 0.12),
                      size=title_size, bold=True,
                      color=RGBColor(0xFF, 0xFF, 0xFF),
                      alignment=PP_ALIGN.CENTER)

        # 三栏布局
        col_w = self.W * 0.30
        col_h = self.H * 0.70
        col_top = self.H * 0.22

        # BEFORE (左)
        left_x = self.W * 0.05
        self._add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE,
                       Inches(left_x), Inches(col_top),
                       Inches(col_w), Inches(col_h),
                       fill_color=RGBColor(0x3A, 0x0F, 0x1F),
                       line_color=self.colors.get('accent_red'),
                       line_width_pt=3)

        self._add_text(slide, "BEFORE",
                      Inches(left_x), Inches(col_top + col_h * 0.03),
                      Inches(col_w), Inches(col_h * 0.12),
                      size=28, bold=True,
                      color=self.colors.get('accent_red'),
                      alignment=PP_ALIGN.CENTER,
                      anchor=MSO_ANCHOR.MIDDLE)

        before_lines = []
        for item in before_items:
            before_lines.append("✗ " + item)
        self._add_multiline_text(slide, before_lines,
                                Inches(left_x + col_w * 0.05),
                                Inches(col_top + col_h * 0.18),
                                Inches(col_w * 0.9),
                                Inches(col_h * 0.78),
                                base_size=18, bold=False,
                                color=RGBColor(0xFF, 0xCC, 0xCC),
                                alignment=PP_ALIGN.LEFT)

        # 中间提升数字
        mid_x = self.W * 0.42
        mid_w = self.W * 0.16
        self._add_circle(slide,
                        Inches(mid_x + mid_w / 2 - mid_w / 2),
                        Inches(col_top + col_h * 0.3),
                        Inches(mid_w),
                        fill_color=self.colors.get('accent_gold'),
                        line_color=self.colors.get('accent_gold'),
                        line_width_pt=4)

        self._add_text(slide, improvement_text,
                      Inches(mid_x),
                      Inches(col_top + col_h * 0.3),
                      Inches(mid_w), Inches(mid_w),
                      size=int(mid_w * 18), bold=True,
                      color=RGBColor(0x0F, 0x20, 0x50),
                      alignment=PP_ALIGN.CENTER,
                      anchor=MSO_ANCHOR.MIDDLE)

        self._add_text(slide, "效率提升",
                      Inches(mid_x - mid_w * 0.5),
                      Inches(col_top + col_h * 0.3 + mid_w),
                      Inches(mid_w * 2), Inches(self.H * 0.08),
                      size=20, bold=True,
                      color=self.colors.get('accent_gold'),
                      alignment=PP_ALIGN.CENTER,
                      anchor=MSO_ANCHOR.MIDDLE)

        # AFTER (右)
        right_x = self.W * 0.65
        self._add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE,
                       Inches(right_x), Inches(col_top),
                       Inches(col_w), Inches(col_h),
                       fill_color=RGBColor(0x0F, 0x3A, 0x2F),
                       line_color=self.colors.get('accent_green'),
                       line_width_pt=3)

        self._add_text(slide, "AFTER",
                      Inches(right_x), Inches(col_top + col_h * 0.03),
                      Inches(col_w), Inches(col_h * 0.12),
                      size=28, bold=True,
                      color=self.colors.get('accent_green'),
                      alignment=PP_ALIGN.CENTER,
                      anchor=MSO_ANCHOR.MIDDLE)

        after_lines = []
        for item in after_items:
            after_lines.append("✓ " + item)
        self._add_multiline_text(slide, after_lines,
                                Inches(right_x + col_w * 0.05),
                                Inches(col_top + col_h * 0.18),
                                Inches(col_w * 0.9),
                                Inches(col_h * 0.78),
                                base_size=18, bold=False,
                                color=RGBColor(0xCC, 0xFF, 0xCC),
                                alignment=PP_ALIGN.LEFT)

        return slide

    # ---------------------------------------------------------
    # 完整演示文稿生成
    # ---------------------------------------------------------

    def generate_full_launch_event(self):
        """生成完整的发布会演示文稿 (根据风格自动选择布局)"""
        print(f"  生成发布会: {self.topic}")
        print(f"  风格: {self.colors.get('name', self.style)}")
        print(f"  画布: {self.W:.2f} x {self.H:.2f} 英寸 "
              f"({int(self.W * 96)} x {int(self.H * 96)} px)")
        print("-" * 50)

        if self.style == 'digital-finance':
            self._generate_digital_finance_deck()
        else:
            self._generate_standard_launch_deck()

        print("-" * 50)
        print(f"  共生成 {len(self.prs.slides)} 张幻灯片")

    def _generate_standard_launch_deck(self):
        """标准16:9发布会演示文稿"""
        # 1. 开场封面
        print("  -> 1. 开场封面")
        self.add_title_slide(
            title=self.topic,
            subtitle=datetime.now().strftime("%Y年%m月%d日")
        )

        # 2. 议程
        print("  -> 2. 议程概览")
        self.add_agenda_slide([
            "产品愿景",
            "核心亮点",
            "技术规格",
            "用户体验",
            "价格与发售"
        ])

        # 3. 问题
        print("  -> 3. 问题与挑战")
        self.add_problem_slide(
            title="我们面临的挑战",
            points=[
                "传统方案无法满足现代业务需求",
                "用户期待更好的体验和更高效率",
                "行业需要新的标准和创新突破"
            ]
        )

        # 4. 解决方案 (大标题页)
        print("  -> 4. 产品亮相")
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self._set_background(slide)
        self._add_text(slide, "全新解决方案",
                      Inches(self.W * 0.1), Inches(self.H * 0.35),
                      Inches(self.W * 0.8), Inches(self.H * 0.25),
                      size=int(self.H * 14), bold=True,
                      color=self.colors.get('text_primary'),
                      alignment=PP_ALIGN.CENTER,
                      anchor=MSO_ANCHOR.MIDDLE)
        self._add_text(slide, "重新定义",
                      Inches(self.W * 0.1), Inches(self.H * 0.62),
                      Inches(self.W * 0.8), Inches(self.H * 0.1),
                      size=int(self.H * 5), bold=False,
                      color=self.colors.get('accent_blue'),
                      alignment=PP_ALIGN.CENTER,
                      anchor=MSO_ANCHOR.MIDDLE)

        # 5-7. 三大核心亮点
        print("  -> 5-7. 核心亮点 (x3)")
        self.add_feature_slide(
            feature_title="创新设计",
            description="全新的设计语言，带来极致视觉冲击与用户体验",
            feature_points=["极致工艺", "人体工学", "精密制造"],
            accent_color=self.colors.get('accent_blue')
        )

        self.add_feature_slide(
            feature_title="强大性能",
            description="突破性能边界，带来极致体验与高效率",
            feature_points=["更快速度", "更长续航", "更强算力"],
            accent_color=self.colors.get('accent_green')
        )

        self.add_feature_slide(
            feature_title="智能体验",
            description="AI赋能，更懂你的智能助手",
            feature_points=["智能识别", "个性化推荐", "无缝协同"],
            accent_color=self.colors.get('accent_purple')
        )

        # 8. KPI 数字展示
        print("  -> 8. 数据成果")
        self.add_kpi_slide(
            title="核心指标",
            kpi_items=[
                ("+180%", "性能提升"),
                ("2.5x", "效率倍增"),
                ("99.9%", "可靠性"),
                ("10M+", "用户规模"),
            ]
        )

        # 9. 技术规格
        print("  -> 9. 技术规格")
        self.add_specs_slide({
            "处理器": "新一代芯片架构",
            "显示屏": "全高清视网膜屏",
            "内存": "大容量高速内存",
            "存储": "NVMe 极速存储",
            "续航": "全天候续航",
            "重量": "轻盈便携设计"
        })

        # 10. 用户故事
        print("  -> 10. 用户故事")
        self.add_story_slide(
            quote="这是改变一切的开始，让我们的业务流程焕然一新。",
            author="早期用户代表"
        )

        # 11. 对比优势
        print("  -> 11. 对比优势")
        self.add_comparison_slide(
            title="相比以往",
            items_dict={
                "核心功能 A": True,
                "核心功能 B": True,
                "核心功能 C": True,
                "高级功能 D": True,
                "竞品功能 E": False,
                "竞品功能 F": False,
            }
        )

        # 12. 价格与发售
        print("  -> 12. 价格与发售")
        self.add_pricing_slide(
            product_name="全新产品系列",
            price="¥XXX 起",
            availability="即将发售 | 官网预约开放中"
        )

        # 13. Thank you
        print("  -> 13. Thank You")
        self.add_thank_you_slide(
            title="谢谢",
            subtitle="期待与您相见"
        )

    def _generate_digital_finance_deck(self):
        """数智财务 SAP风格 3:1超宽屏演示文稿"""
        # 1. 开场封面
        print("  -> 1. 开场封面 (超宽屏)")
        self.add_uw_cover_slide(
            title="数智财务世界的认知与改变",
            subtitle="数据驱动  AI赋能",
            third_title="现代财务数字化转型方案"
        )

        # 2. 议程概览
        print("  -> 2. 议程概览")
        self.add_agenda_slide([
            "财务数字化背景",
            "核心痛点与挑战",
            "SAP 解决方案架构",
            "关键功能与亮点",
            "客户案例与成果",
            "实施路径与展望"
        ])

        # 3. 三栏卡片 - 背景动因
        print("  -> 3. 数字化背景 (三栏卡片)")
        self.add_uw_three_column_cards(
            title="财务数字化转型的三大动因",
            card_data=[
                ("修炼内功", ["财务流程自动化需求", "数据质量与实时性要求", "合规与风控压力增加"]),
                ("模式转变", ["从核算型转向价值型财务", "业财融合深入业务决策", "预测式财务成为标配"]),
                ("有序发展", ["新技术成熟应用", "数据治理体系完善", "数字化人才储备"]),
            ]
        )

        # 4. 横向时间轴
        print("  -> 4. 技术演进时间轴")
        self.add_uw_timeline(
            title="50年财务技术平台连续演进",
            nodes=[
                ("1972", "R/2 大型机"),
                ("1992", "SAP ERP"),
                ("2004", "NetWeaver"),
                ("2015", "S/4HANA"),
                ("2023", "Business AI"),
            ]
        )

        # 5. Before / After 对比
        print("  -> 5. 变革前后对比")
        self.add_uw_before_after(
            title="传统财务 vs 数智财务",
            before_items=[
                "月末结账耗时7-10天",
                "报表滞后无法支持决策",
                "手工核对数据错误率高",
                "合规风险人工识别困难",
                "预算编制周期长达3个月",
            ],
            after_items=[
                "实时结账，秒级数据可用",
                "实时报表，支持即时决策",
                "自动对账，错误率降低90%",
                "AI风险识别，智能预警",
                "滚动预算，实时调整",
            ],
            improvement_text="85%"
        )

        # 6. 中央环形架构图
        print("  -> 6. 解决方案架构 (环形图)")
        self.add_uw_ring_architecture(
            title="SAP 数智财务 技术生态系统",
            center_text="财务核心",
            features=[
                "总账会计",
                "成本控制",
                "资金管理",
                "财务分析",
                "预算管理",
                "税务合规",
                "应收应付",
                "资产核算",
            ]
        )

        # 7. 六边形蜂群矩阵
        print("  -> 7. 预算编制要素矩阵 (六边形)")
        self.add_uw_hexagon_matrix(
            title="财务预算编制核心要素",
            center_text="预算\n核心",
            elements=["销售预测", "成本规划", "资本支出", "运营预算",
                     "人力资源", "现金流", "KPI指标", "风险储备"]
        )

        # 8. 四象限功能矩阵
        print("  -> 8. 功能全景 (四象限)")
        self.add_uw_quadrant_matrix(
            title="数智财务 四大能力域",
            quadrants=[
                ("智能化核算", ["智能发票识别", "自动对账匹配", "AI凭证生成", "实时会计处理"]),
                ("预测式财务", ["滚动预测模型", "现金流预测", "敏感性分析", "情景模拟"]),
                ("数据驱动决策", ["实时BI看板", "多维分析", "异常检测", "数据洞察"]),
                ("智能风控合规", ["自动合规检查", "风险识别模型", "内控自动化", "审计追踪"]),
            ]
        )

        # 9. Dashboard数据页
        print("  -> 9. 财务 Dashboard")
        self.add_uw_dashboard(
            title="企业经营分析与财务健康度",
            kpi_data=[
                ("¥12.5亿", "年度营收", "↑ 18.5%"),
                ("¥3.2亿", "净利润", "↑ 25.3%"),
                ("12.8%", "利润率", "↑ 2.1pp"),
                ("¥6.8亿", "经营现金流", "↑ 15.7%"),
            ],
            chart_title="季度营收趋势 (2024)"
        )

        # 10. 四步流程卡
        print("  -> 10. 实施路径 (四步流程)")
        self.add_uw_four_step_process(
            title="从规划到价值落地的完整路径",
            steps=[
                ("诊断规划", "现状评估、需求分析、蓝图设计、路线规划"),
                ("系统部署", "基础配置、数据迁移、集成开发、用户培训"),
                ("上线运行", "切换策略、试运行、正式上线、运维保障"),
                ("持续优化", "效果评估、流程改进、能力扩展、价值挖掘"),
            ]
        )

        # 11. 客户案例 KPI
        print("  -> 11. 客户成果数据")
        self.add_kpi_slide(
            title="客户应用成效",
            kpi_items=[
                ("-60%", "月结时间"),
                ("+85%", "报表时效性"),
                ("¥500万", "年度成本节约"),
                ("-40%", "人工操作量"),
            ]
        )

        # 12. 用户故事 / 引言
        print("  -> 12. 客户证言")
        self.add_story_slide(
            quote="通过SAP数智财务方案，我们实现了从传统核算到价值创造的转型，财务团队成为真正的业务伙伴。",
            author="某大型集团 CFO"
        )

        # 13. Thank you (与封面呼应)
        print("  -> 13. Thank You")
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        bg_color = self.colors.get('background', RGBColor(0x12, 0x24, 0x4D))
        self._set_background(slide, bg_color)
        self._add_golden_light_beam(slide, top_in=self.H * 0.72)

        self._add_text(slide, "谢谢聆听",
                      Inches(self.W * 0.15), Inches(self.H * 0.30),
                      Inches(self.W * 0.70), Inches(self.H * 0.25),
                      size=int(min(72, self.H * 11)),
                      bold=True,
                      color=RGBColor(0xFF, 0xFF, 0xFF),
                      alignment=PP_ALIGN.CENTER,
                      anchor=MSO_ANCHOR.MIDDLE)

        self._add_text(slide, "Thank You",
                      Inches(self.W * 0.15), Inches(self.H * 0.58),
                      Inches(self.W * 0.70), Inches(self.H * 0.10),
                      size=int(min(36, self.H * 5)),
                      bold=True,
                      color=self.colors.get('accent_gold'),
                      alignment=PP_ALIGN.CENTER,
                      anchor=MSO_ANCHOR.MIDDLE)

    # ---------------------------------------------------------
    # 保存
    # ---------------------------------------------------------

    def save(self, output_path):
        """保存演示文稿"""
        output_path = Path(output_path)
        output_path.parent.mkdir(parents=True, exist_ok=True)
        self.prs.save(str(output_path))
        print(f"  已保存: {output_path}")
        return output_path


# =============================================================
# 主函数
# =============================================================

def main():
    parser = argparse.ArgumentParser(
        description='生成 Keynote 风格的发布会 PPTX (v3.0)',
        formatter_class=argparse.RawDescriptionHelpFormatter,
        epilog="""
示例:
  # 经典 Apple 风格 (16:9)
  python3 generate_presentation.py --topic "2026 春季新品发布" --style classic

  # 科技发布会风格
  python3 generate_presentation.py --topic "AI Platform v3.0" --style tech

  # 数智财务 SAP风格 (3:1超宽屏)
  python3 generate_presentation.py --topic "数智财务峰会" --style digital-finance

  # 查看所有风格
  python3 generate_presentation.py --list-styles
        """
    )
    parser.add_argument('--topic', required=False, default='会议主题', help='发布会主题（--list-styles 时不需要）')
    parser.add_argument('--style', choices=COLOR_SCHEMES.keys(),
                       default='classic', help='设计风格')
    parser.add_argument('--lang', choices=['zh', 'en'], default='zh',
                       help='语言 (zh=中文, en=英文)')
    parser.add_argument('--aspect-ratio', choices=['standard', 'ultra-wide'],
                       default=None, help='强制宽高比 (默认根据style自动选择)')
    parser.add_argument('--output', '-o', default='./output/launch_event.pptx',
                       help='输出文件路径')
    parser.add_argument('--list-styles', action='store_true',
                       help='列出所有可用风格')

    args = parser.parse_args()

    if args.list_styles:
        print("\n可用的发布会风格:\n")
        for key, scheme in COLOR_SCHEMES.items():
            ratio = scheme.get('aspect_ratio', '16:9')
            print(f"  {key:<16} - {scheme.get('name', key)} ({ratio})")
        print()
        return

    print("\n" + "=" * 60)
    print(" Keynote 发布会 PPTX 生成器 (v3.0)")
    print("=" * 60)
    print()

    generator = KeynotePresentationGenerator(
        topic=args.topic,
        style=args.style,
        lang=args.lang,
        aspect_ratio=args.aspect_ratio
    )

    generator.generate_full_launch_event()
    generator.save(args.output)

    print("\n" + "=" * 60)
    print(" 生成完成!")
    print("=" * 60)
    print("\n下一步:")
    print("  1. 在 macOS 上用 Keynote 打开生成的 .pptx 文件")
    print("  2. 文件 -> 另存为 -> 选择 .key 格式")
    print("  3. 根据 design_guide.md 调整字体和动画")
    print("  4. 替换产品图片和品牌素材")
    print("=" * 60 + "\n")


if __name__ == '__main__':
    main()
