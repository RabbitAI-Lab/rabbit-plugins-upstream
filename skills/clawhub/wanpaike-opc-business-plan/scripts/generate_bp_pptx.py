#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
商业计划书PPT生成脚本
基于python-pptx库，将商业计划书内容自动生成为专业PPT演示文稿

用法：
  python generate_bp_pptx.py --data <json_file> --output <output_path> [--style <style_name>]

参数：
  --data    : 商业计划书内容JSON文件路径（必需）
  --output  : 输出PPTX文件路径（默认：./商业计划书.pptx）
  --style   : PPT风格（tech_blue|elegant_gold|clean_green|classic_dark），默认tech_blue

数据JSON格式说明：
  {
    "project_name": "项目名称",
    "tagline": "一句话描述",
    "industry": "行业",
    "stage": "发展阶段",
    "financing_round": "融资轮次",
    "date": "日期",
    "executive_summary": {
      "pain_point": "痛点",
      "solution": "方案",
      "market_size": "市场规模",
      "team_highlight": "团队亮点",
      "financing_ask": "融资需求"
    },
    "pain_points": [
      {"title": "痛点标题", "description": "痛点描述"}
    ],
    "solution": {
      "core": "核心方案描述",
      "advantages": ["优势1", "优势2"]
    },
    "market": {
      "tam_sam_som": {"tam": "1000亿", "sam": "300亿", "som": "30亿"},
      "growth_rate": "25%/年",
      "competitors": [
        {"name": "竞品A", "strength": "优势", "weakness": "弱点", "pricing": "定价"}
      ]
    },
    "feasibility": {
      "tech": "高/中/低",
      "business": "高/中/低",
      "operations": "高/中/低",
      "finance": "高/中/低",
      "timing": "高/中/低"
    },
    "business_model": {
      "value_proposition": "价值主张",
      "customer_segments": "客户细分",
      "channels": "渠道通路",
      "customer_relationships": "客户关系",
      "revenue_streams": "收入来源",
      "key_resources": "核心资源",
      "key_activities": "关键业务",
      "key_partners": "重要伙伴",
      "cost_structure": "成本结构"
    },
    "financial": {
      "forecast": [
        {"year": "Y1", "revenue": "250万", "gross_profit": "175万", "operating_profit": "-50万", "margin": "-20%"}
      ],
      "key_metrics": {"CAC": "1万", "LTV": "17.5万", "LTV_CAC": "17.5", "Payback": "10月"},
      "unit_economics": {"CAC": "1万", "LTV": "17.5万", "LTV/CAC": "17.5"}
    },
    "financing": {
      "amount": "3000万",
      "valuation": "1.5亿",
      "use_of_funds": "研发40%/营销30%/团队20%/运营10%",
      "runway": "24个月"
    },
    "team": {
      "members": [
        {"name": "姓名", "role": "职位", "background": "背景描述"}
      ]
    },
    "milestones": [
      {"period": "0-6月", "title": "概念验证期", "goals": ["目标1", "目标2"]}
    ],
    "risks": [
      {"type": "市场风险", "description": "描述", "probability": "中", "impact": "高", "mitigation": "应对策略"}
    ],
    "valuation": {
      "tech_moat": "8/10",
      "market_momentum": "7/10",
      "team_factor": "9/10",
      "business_flywheel": "6/10",
      "exit_path": "7/10",
      "estimated_range": "1.2亿-2.5亿"
    },
    "contact": "联系方式"
  }
"""

import json
import argparse
import os
import sys

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
    from pptx.enum.shapes import MSO_SHAPE
except ImportError:
    print("正在安装python-pptx...")
    import subprocess
    subprocess.check_call([sys.executable, "-m", "pip", "install", "python-pptx", "-q"])
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
    from pptx.enum.shapes import MSO_SHAPE

# ============================================================
# 风格配置
# ============================================================
STYLES = {
    "tech_blue": {
        "name": "科技蓝",
        "primary": RGBColor(0x0B, 0x3D, 0x91),
        "secondary": RGBColor(0x00, 0x7B, 0xFF),
        "accent": RGBColor(0x00, 0xD4, 0xFF),
        "bg_dark": RGBColor(0x0A, 0x14, 0x28),
        "bg_light": RGBColor(0xF0, 0xF4, 0xF8),
        "text_dark": RGBColor(0x1A, 0x1A, 0x2E),
        "text_light": RGBColor(0xFF, 0xFF, 0xFF),
        "text_muted": RGBColor(0x6B, 0x7B, 0x8D),
        "success": RGBColor(0x00, 0xC4, 0x8C),
        "warning": RGBColor(0xFF, 0x8C, 0x00),
        "danger": RGBColor(0xFF, 0x3B, 0x3B),
        "chart_colors": [
            RGBColor(0x00, 0x7B, 0xFF), RGBColor(0x00, 0xD4, 0xFF),
            RGBColor(0x00, 0xC4, 0x8C), RGBColor(0xFF, 0x8C, 0x00),
            RGBColor(0x8B, 0x5C, 0xF6), RGBColor(0xFF, 0x3B, 0x3B),
        ],
    },
    "elegant_gold": {
        "name": "商务金",
        "primary": RGBColor(0x1A, 0x2A, 0x3A),
        "secondary": RGBColor(0xC5, 0xA5, 0x5A),
        "accent": RGBColor(0xE8, 0xD5, 0xA3),
        "bg_dark": RGBColor(0x0F, 0x1A, 0x2E),
        "bg_light": RGBColor(0xFA, 0xF7, 0xF0),
        "text_dark": RGBColor(0x1A, 0x1A, 0x1A),
        "text_light": RGBColor(0xFF, 0xFF, 0xFF),
        "text_muted": RGBColor(0x8B, 0x7E, 0x74),
        "success": RGBColor(0x2E, 0x8B, 0x57),
        "warning": RGBColor(0xD4, 0x8B, 0x2C),
        "danger": RGBColor(0xB2, 0x22, 0x22),
        "chart_colors": [
            RGBColor(0xC5, 0xA5, 0x5A), RGBColor(0x1A, 0x2A, 0x3A),
            RGBColor(0x2E, 0x8B, 0x57), RGBColor(0xD4, 0x8B, 0x2C),
            RGBColor(0x6A, 0x5A, 0xCD), RGBColor(0xB2, 0x22, 0x22),
        ],
    },
    "clean_green": {
        "name": "清新绿",
        "primary": RGBColor(0x0D, 0x6E, 0x4B),
        "secondary": RGBColor(0x10, 0xB9, 0x81),
        "accent": RGBColor(0x6E, 0xE7, 0xB7),
        "bg_dark": RGBColor(0x0A, 0x2E, 0x1F),
        "bg_light": RGBColor(0xF0, 0xFD, 0xF4),
        "text_dark": RGBColor(0x1A, 0x2E, 0x1A),
        "text_light": RGBColor(0xFF, 0xFF, 0xFF),
        "text_muted": RGBColor(0x6B, 0x8F, 0x71),
        "success": RGBColor(0x10, 0xB9, 0x81),
        "warning": RGBColor(0xF5, 0xA6, 0x23),
        "danger": RGBColor(0xEF, 0x44, 0x44),
        "chart_colors": [
            RGBColor(0x10, 0xB9, 0x81), RGBColor(0x0D, 0x6E, 0x4B),
            RGBColor(0xF5, 0xA6, 0x23), RGBColor(0x3B, 0x82, 0xF6),
            RGBColor(0x8B, 0x5C, 0xF6), RGBColor(0xEF, 0x44, 0x44),
        ],
    },
    "classic_dark": {
        "name": "经典暗色",
        "primary": RGBColor(0xE0, 0xE0, 0xE0),
        "secondary": RGBColor(0x00, 0xD4, 0xFF),
        "accent": RGBColor(0xFF, 0xD7, 0x00),
        "bg_dark": RGBColor(0x1A, 0x1A, 0x2E),
        "bg_light": RGBColor(0x2D, 0x2D, 0x3D),
        "text_dark": RGBColor(0x1A, 0x1A, 0x1A),
        "text_light": RGBColor(0xE0, 0xE0, 0xE0),
        "text_muted": RGBColor(0x80, 0x80, 0x90),
        "success": RGBColor(0x00, 0xFF, 0x88),
        "warning": RGBColor(0xFF, 0xD7, 0x00),
        "danger": RGBColor(0xFF, 0x44, 0x44),
        "chart_colors": [
            RGBColor(0x00, 0xD4, 0xFF), RGBColor(0xFF, 0xD7, 0x00),
            RGBColor(0x00, 0xFF, 0x88), RGBColor(0xFF, 0x44, 0x44),
            RGBColor(0x8B, 0x5C, 0xF6), RGBColor(0xFF, 0x8C, 0x00),
        ],
    },
}


# ============================================================
# 辅助函数
# ============================================================

def add_shape_with_text(slide, left, top, width, height, text, font_size=12,
                        font_color=None, bg_color=None, bold=False, alignment=PP_ALIGN.LEFT,
                        font_name="Microsoft YaHei", shape_type=MSO_SHAPE.ROUNDED_RECTANGLE):
    shape = slide.shapes.add_shape(shape_type, left, top, width, height)
    if bg_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = bg_color
    else:
        shape.fill.background()
    shape.line.fill.background()
    tf = shape.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = font_color or RGBColor(0x33, 0x33, 0x33)
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return shape


def add_textbox(slide, left, top, width, height, text, font_size=12,
                font_color=None, bold=False, alignment=PP_ALIGN.LEFT,
                font_name="Microsoft YaHei"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = font_color or RGBColor(0x33, 0x33, 0x33)
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def add_bullet_text(slide, left, top, width, height, items, font_size=14,
                    font_color=None, font_name="Microsoft YaHei", spacing=Pt(6)):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"• {item}"
        p.font.size = Pt(font_size)
        p.font.color.rgb = font_color or RGBColor(0x33, 0x33, 0x33)
        p.font.name = font_name
        p.space_after = spacing
    return txBox


def add_table(slide, left, top, width, height, headers, rows, style_cfg,
              header_color=None, font_size=10):
    n_rows = len(rows) + 1
    n_cols = len(headers)
    table_shape = slide.shapes.add_table(n_rows, n_cols, left, top, width, height)
    table = table_shape.table
    for j, h in enumerate(headers):
        cell = table.cell(0, j)
        cell.text = h
        for paragraph in cell.text_frame.paragraphs:
            paragraph.font.size = Pt(font_size)
            paragraph.font.bold = True
            paragraph.font.color.rgb = style_cfg["text_light"]
            paragraph.font.name = "Microsoft YaHei"
            paragraph.alignment = PP_ALIGN.CENTER
        cell.fill.solid()
        cell.fill.fore_color.rgb = header_color or style_cfg["primary"]
    for i, row in enumerate(rows):
        for j, val in enumerate(row):
            cell = table.cell(i + 1, j)
            cell.text = str(val)
            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.size = Pt(font_size)
                paragraph.font.color.rgb = style_cfg["text_dark"]
                paragraph.font.name = "Microsoft YaHei"
                paragraph.alignment = PP_ALIGN.CENTER
            if i % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = style_cfg["bg_light"]
    return table_shape


def add_section_header(slide, title, style):
    """统一的页面标题+分隔线"""
    add_textbox(slide, Inches(0.8), Inches(0.4), Inches(10), Inches(0.6),
                title, font_size=28, font_color=style["primary"], bold=True)
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.1), Inches(2), Inches(0.05))
    line.fill.solid()
    line.fill.fore_color.rgb = style["secondary"]
    line.line.fill.background()


# ============================================================
# 页面生成函数
# ============================================================

def create_cover(prs, data, style):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.background; fill = bg.fill; fill.solid(); fill.fore_color.rgb = style["bg_dark"]
    add_textbox(slide, Inches(1.2), Inches(1.5), Inches(9.6), Inches(1.2),
                data.get("project_name", "商业计划书"), font_size=44,
                font_color=style["text_light"], bold=True)
    add_textbox(slide, Inches(1.2), Inches(2.8), Inches(9.6), Inches(0.8),
                data.get("tagline", ""), font_size=20, font_color=style["accent"])
    info_parts = []
    if data.get("industry"): info_parts.append(f"行业：{data['industry']}")
    if data.get("stage"): info_parts.append(f"阶段：{data['stage']}")
    if data.get("financing_round"): info_parts.append(f"融资轮次：{data['financing_round']}")
    info_parts.append(f"日期：{data.get('date', '2025年')}")
    add_textbox(slide, Inches(1.2), Inches(4.0), Inches(9.6), Inches(0.5),
                "  |  ".join(info_parts), font_size=14, font_color=style["text_muted"])
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(6.8), Inches(13.3), Inches(0.2))
    bar.fill.solid(); bar.fill.fore_color.rgb = style["secondary"]; bar.line.fill.background()


def create_executive_summary(prs, data, style):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.background; fill = bg.fill; fill.solid(); fill.fore_color.rgb = style["bg_light"]
    add_section_header(slide, "执行摘要", style)
    summary = data.get("executive_summary", {})
    items = []
    if summary.get("pain_point"): items.append(f"痛点：{summary['pain_point']}")
    if summary.get("solution"): items.append(f"方案：{summary['solution']}")
    if summary.get("market_size"): items.append(f"市场：{summary['market_size']}")
    if summary.get("team_highlight"): items.append(f"团队：{summary['team_highlight']}")
    if summary.get("financing_ask"): items.append(f"融资：{summary['financing_ask']}")
    add_bullet_text(slide, Inches(0.8), Inches(1.4), Inches(11.5), Inches(4.5),
                    items, font_size=16, font_color=style["text_dark"], spacing=Pt(12))


def create_pain_point(prs, data, style):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.background; fill = bg.fill; fill.solid(); fill.fore_color.rgb = style["bg_light"]
    add_section_header(slide, "问题与痛点", style)
    pain_points = data.get("pain_points", [])
    for i, pp in enumerate(pain_points[:4]):
        row, col = i // 2, i % 2
        left, top = Inches(0.8 + col * 6.2), Inches(1.5 + row * 2.4)
        w, h = Inches(5.8), Inches(2.0)
        add_shape_with_text(slide, left, top, w, h, "", font_size=1, bg_color=RGBColor(0xFF, 0xFF, 0xFF))
        add_textbox(slide, left + Inches(0.3), top + Inches(0.2), w - Inches(0.6), Inches(0.5),
                    pp.get("title", f"痛点{i+1}"), font_size=16, font_color=style["danger"], bold=True)
        add_textbox(slide, left + Inches(0.3), top + Inches(0.8), w - Inches(0.6), Inches(1.0),
                    pp.get("description", ""), font_size=13, font_color=style["text_dark"])


def create_solution(prs, data, style):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.background; fill = bg.fill; fill.solid(); fill.fore_color.rgb = style["bg_light"]
    add_section_header(slide, "解决方案与差异化优势", style)
    solution = data.get("solution", {})
    add_textbox(slide, Inches(0.8), Inches(1.4), Inches(5.5), Inches(0.5),
                "核心方案", font_size=18, font_color=style["secondary"], bold=True)
    add_textbox(slide, Inches(0.8), Inches(2.0), Inches(5.5), Inches(2.0),
                solution.get("core", ""), font_size=14, font_color=style["text_dark"])
    add_textbox(slide, Inches(7.0), Inches(1.4), Inches(5.5), Inches(0.5),
                "差异化优势", font_size=18, font_color=style["secondary"], bold=True)
    add_bullet_text(slide, Inches(7.0), Inches(2.0), Inches(5.5), Inches(3.0),
                    solution.get("advantages", []), font_size=14, font_color=style["text_dark"])


def create_market_analysis(prs, data, style):
    # 页1：市场规模
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.background; fill = bg.fill; fill.solid(); fill.fore_color.rgb = style["bg_light"]
    add_section_header(slide, "市场规模（TAM/SAM/SOM）", style)
    market = data.get("market", {})
    tam_data = market.get("tam_sam_som", {})
    metrics = [
        ("TAM（潜在市场总量）", tam_data.get("tam", "待补充"), style["primary"]),
        ("SAM（可服务市场）", tam_data.get("sam", "待补充"), style["secondary"]),
        ("SOM（可获得市场）", tam_data.get("som", "待补充"), style["accent"]),
    ]
    for i, (label, value, color) in enumerate(metrics):
        left = Inches(0.8 + i * 4.1)
        add_shape_with_text(slide, left, Inches(1.5), Inches(3.8), Inches(2.5), "", font_size=1, bg_color=color)
        add_textbox(slide, left + Inches(0.3), Inches(1.8), Inches(3.2), Inches(0.5),
                    label, font_size=16, font_color=style["text_light"], bold=True)
        add_textbox(slide, left + Inches(0.3), Inches(2.5), Inches(3.2), Inches(1.2),
                    str(value), font_size=28, font_color=style["text_light"], bold=True, alignment=PP_ALIGN.CENTER)
    if market.get("growth_rate"):
        add_textbox(slide, Inches(0.8), Inches(4.5), Inches(11.5), Inches(0.5),
                    f"市场年均增长率：{market['growth_rate']}", font_size=16,
                    font_color=style["text_dark"], bold=True)

    # 页2：竞争格局
    slide2 = prs.slides.add_slide(prs.slide_layouts[6])
    bg2 = slide2.background; fill2 = bg2.fill; fill2.solid(); fill2.fore_color.rgb = style["bg_light"]
    add_section_header(slide2, "竞争格局分析", style)
    competitors = data.get("market", {}).get("competitors", [])
    if competitors:
        headers = ["竞争对手", "核心优势", "主要弱点", "定价策略"]
        rows = [[c.get("name", ""), c.get("strength", ""), c.get("weakness", ""), c.get("pricing", "")]
                for c in competitors[:5]]
        add_table(slide2, Inches(0.8), Inches(1.5), Inches(11.5), Inches(3.5), headers, rows, style, font_size=11)


def create_feasibility(prs, data, style):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.background; fill = bg.fill; fill.solid(); fill.fore_color.rgb = style["bg_light"]
    add_section_header(slide, "可行性分析 & OPC-BP可行性雷达", style)
    feasibility = data.get("feasibility", {})
    dims = [
        ("技术可行性", feasibility.get("tech", "中"), "技术路径与储备"),
        ("商业可行性", feasibility.get("business", "中"), "商业模式与需求"),
        ("运营可行性", feasibility.get("operations", "中"), "团队与资源匹配"),
        ("财务可行性", feasibility.get("finance", "中"), "盈利路径清晰度"),
        ("时机可行性", feasibility.get("timing", "中"), "市场窗口与节奏"),
    ]
    for i, (dim, level, desc) in enumerate(dims):
        left = Inches(0.8 + (i % 3) * 4.1)
        top = Inches(1.5 + (i // 3) * 2.5)
        w, h = Inches(3.8), Inches(2.0)
        color_map = {"高": style["success"], "中": style["warning"], "低": style["danger"]}
        dim_color = color_map.get(level, style["text_muted"])
        add_shape_with_text(slide, left, top, w, h, "", font_size=1, bg_color=RGBColor(0xFF, 0xFF, 0xFF))
        add_textbox(slide, left + Inches(0.2), top + Inches(0.15), w - Inches(0.4), Inches(0.4),
                    dim, font_size=14, font_color=style["text_dark"], bold=True)
        add_textbox(slide, left + Inches(0.2), top + Inches(0.6), w - Inches(0.4), Inches(0.5),
                    level, font_size=24, font_color=dim_color, bold=True, alignment=PP_ALIGN.CENTER)
        add_textbox(slide, left + Inches(0.2), top + Inches(1.3), w - Inches(0.4), Inches(0.5),
                    desc, font_size=11, font_color=style["text_muted"])


def create_business_model(prs, data, style):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.background; fill = bg.fill; fill.solid(); fill.fore_color.rgb = style["bg_light"]
    add_section_header(slide, "商业模式画布", style)
    bmc = data.get("business_model", {})
    bmc_items = [
        ("价值主张", bmc.get("value_proposition", "")),
        ("客户细分", bmc.get("customer_segments", "")),
        ("渠道通路", bmc.get("channels", "")),
        ("客户关系", bmc.get("customer_relationships", "")),
        ("收入来源", bmc.get("revenue_streams", "")),
        ("核心资源", bmc.get("key_resources", "")),
        ("关键业务", bmc.get("key_activities", "")),
        ("重要伙伴", bmc.get("key_partners", "")),
        ("成本结构", bmc.get("cost_structure", "")),
    ]
    for i, (label, value) in enumerate(bmc_items):
        row, col = i // 3, i % 3
        left = Inches(0.8 + col * 4.1)
        top = Inches(1.5 + row * 1.8)
        w, h = Inches(3.8), Inches(1.5)
        add_shape_with_text(slide, left, top, w, h, "", font_size=1, bg_color=RGBColor(0xFF, 0xFF, 0xFF))
        add_textbox(slide, left + Inches(0.2), top + Inches(0.1), w - Inches(0.4), Inches(0.35),
                    label, font_size=12, font_color=style["secondary"], bold=True)
        add_textbox(slide, left + Inches(0.2), top + Inches(0.5), w - Inches(0.4), Inches(0.9),
                    str(value)[:80], font_size=11, font_color=style["text_dark"])


def create_financial_forecast(prs, data, style):
    # 页1：收入与利润预测
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.background; fill = bg.fill; fill.solid(); fill.fore_color.rgb = style["bg_light"]
    add_section_header(slide, "财务预测（3-5年）", style)
    finance = data.get("financial", {})
    forecast = finance.get("forecast", [])
    if forecast:
        headers = ["年度", "收入", "毛利", "运营利润", "利润率"]
        rows = [[f.get("year", ""), f.get("revenue", ""), f.get("gross_profit", ""),
                 f.get("operating_profit", ""), f.get("margin", "")] for f in forecast]
        add_table(slide, Inches(0.8), Inches(1.5), Inches(11.5), Inches(3.5), headers, rows, style, font_size=11)
    metrics = finance.get("key_metrics", {})
    if metrics:
        add_bullet_text(slide, Inches(0.8), Inches(5.2), Inches(11.5), Inches(1.5),
                        [f"{k}：{v}" for k, v in metrics.items()], font_size=13, font_color=style["text_dark"])

    # 页2：单位经济 & 融资
    slide2 = prs.slides.add_slide(prs.slide_layouts[6])
    bg2 = slide2.background; fill2 = bg2.fill; fill2.solid(); fill2.fore_color.rgb = style["bg_light"]
    add_section_header(slide2, "单位经济模型 & 融资规划", style)
    unit_econ = finance.get("unit_economics", {})
    if unit_econ:
        add_bullet_text(slide2, Inches(0.8), Inches(1.4), Inches(5.5), Inches(3.5),
                        [f"{k}：{v}" for k, v in unit_econ.items()], font_size=14, font_color=style["text_dark"], spacing=Pt(10))
    financing = data.get("financing", {})
    if financing:
        add_textbox(slide2, Inches(7.0), Inches(1.4), Inches(5.5), Inches(0.5),
                    "融资规划", font_size=18, font_color=style["secondary"], bold=True)
        fin_items = []
        if financing.get("amount"): fin_items.append(f"融资金额：{financing['amount']}")
        if financing.get("valuation"): fin_items.append(f"估值：{financing['valuation']}")
        if financing.get("use_of_funds"): fin_items.append(f"资金用途：{financing['use_of_funds']}")
        if financing.get("runway"): fin_items.append(f"资金支撑：{financing['runway']}")
        add_bullet_text(slide2, Inches(7.0), Inches(2.0), Inches(5.5), Inches(3.0),
                        fin_items, font_size=14, font_color=style["text_dark"])


def create_team(prs, data, style):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.background; fill = bg.fill; fill.solid(); fill.fore_color.rgb = style["bg_light"]
    add_section_header(slide, "核心团队", style)
    team = data.get("team", {}).get("members", [])
    for i, member in enumerate(team[:5]):
        col, row = i % 3, i // 3
        left, top = Inches(0.8 + col * 4.1), Inches(1.5 + row * 2.5)
        w, h = Inches(3.8), Inches(2.0)
        add_shape_with_text(slide, left, top, w, h, "", font_size=1, bg_color=RGBColor(0xFF, 0xFF, 0xFF))
        add_textbox(slide, left + Inches(0.2), top + Inches(0.15), w - Inches(0.4), Inches(0.4),
                    member.get("name", ""), font_size=16, font_color=style["primary"], bold=True)
        add_textbox(slide, left + Inches(0.2), top + Inches(0.55), w - Inches(0.4), Inches(0.3),
                    member.get("role", ""), font_size=12, font_color=style["secondary"])
        add_textbox(slide, left + Inches(0.2), top + Inches(0.9), w - Inches(0.4), Inches(0.9),
                    member.get("background", ""), font_size=11, font_color=style["text_dark"])


def create_milestones(prs, data, style):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.background; fill = bg.fill; fill.solid(); fill.fore_color.rgb = style["bg_light"]
    add_section_header(slide, "发展里程碑", style)
    milestones = data.get("milestones", [])
    for i, ms in enumerate(milestones[:5]):
        left = Inches(0.8 + i * 2.4)
        dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, left + Inches(0.85), Inches(1.5), Inches(0.4), Inches(0.4))
        dot.fill.solid(); dot.fill.fore_color.rgb = style["secondary"]; dot.line.fill.background()
        if i < len(milestones) - 1:
            conn = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left + Inches(1.25), Inches(1.65), Inches(2.0), Inches(0.05))
            conn.fill.solid(); conn.fill.fore_color.rgb = style["secondary"]; conn.line.fill.background()
        add_shape_with_text(slide, left, Inches(2.0), Inches(2.2), Inches(3.5), "", font_size=1, bg_color=RGBColor(0xFF, 0xFF, 0xFF))
        add_textbox(slide, left + Inches(0.15), Inches(2.15), Inches(1.9), Inches(0.4),
                    ms.get("period", ""), font_size=14, font_color=style["secondary"], bold=True, alignment=PP_ALIGN.CENTER)
        add_textbox(slide, left + Inches(0.15), Inches(2.6), Inches(1.9), Inches(0.4),
                    ms.get("title", ""), font_size=13, font_color=style["primary"], bold=True)
        goals = ms.get("goals", [])
        if isinstance(goals, list):
            add_bullet_text(slide, left + Inches(0.15), Inches(3.1), Inches(1.9), Inches(2.0),
                            goals, font_size=10, font_color=style["text_dark"])


def create_risk_analysis(prs, data, style):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.background; fill = bg.fill; fill.solid(); fill.fore_color.rgb = style["bg_light"]
    add_section_header(slide, "风险分析与应对策略", style)
    risks = data.get("risks", [])
    if risks:
        headers = ["风险类型", "风险描述", "概率", "影响", "应对策略"]
        rows = [[r.get("type", ""), r.get("description", "")[:20], r.get("probability", ""),
                 r.get("impact", ""), r.get("mitigation", "")[:25]] for r in risks[:6]]
        add_table(slide, Inches(0.8), Inches(1.5), Inches(11.5), Inches(4.5), headers, rows, style, font_size=10)


def create_valuation_model(prs, data, style):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.background; fill = bg.fill; fill.solid(); fill.fore_color.rgb = style["bg_light"]
    add_section_header(slide, "OPC 科创企业估值模型", style)
    valuation = data.get("valuation", {})
    dims = [
        ("技术壁垒", valuation.get("tech_moat", "待评估")),
        ("市场势能", valuation.get("market_momentum", "待评估")),
        ("团队系数", valuation.get("team_factor", "待评估")),
        ("商业飞轮", valuation.get("business_flywheel", "待评估")),
        ("退出路径", valuation.get("exit_path", "待评估")),
    ]
    for i, (dim, val) in enumerate(dims):
        left = Inches(0.8 + i * 2.4)
        add_shape_with_text(slide, left, Inches(1.5), Inches(2.2), Inches(2.5),
                            "", font_size=1, bg_color=RGBColor(0xFF, 0xFF, 0xFF))
        add_textbox(slide, left + Inches(0.1), Inches(1.65), Inches(2.0), Inches(0.4),
                    dim, font_size=13, font_color=style["primary"], bold=True, alignment=PP_ALIGN.CENTER)
        add_textbox(slide, left + Inches(0.1), Inches(2.1), Inches(2.0), Inches(0.5),
                    str(val), font_size=20, font_color=style["secondary"], bold=True, alignment=PP_ALIGN.CENTER)
    if valuation.get("estimated_range"):
        add_shape_with_text(slide, Inches(0.8), Inches(4.3), Inches(11.5), Inches(1.5),
                            "", font_size=1, bg_color=style["primary"])
        add_textbox(slide, Inches(1.0), Inches(4.5), Inches(11.0), Inches(0.5),
                    "综合估值区间", font_size=16, font_color=style["text_light"], bold=True, alignment=PP_ALIGN.CENTER)
        add_textbox(slide, Inches(1.0), Inches(5.0), Inches(11.0), Inches(0.5),
                    str(valuation["estimated_range"]), font_size=28, font_color=style["accent"], bold=True, alignment=PP_ALIGN.CENTER)


def create_closing(prs, data, style):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.background; fill = bg.fill; fill.solid(); fill.fore_color.rgb = style["bg_dark"]
    add_textbox(slide, Inches(1.5), Inches(2.0), Inches(10.0), Inches(1.2),
                "Thank You", font_size=48, font_color=style["text_light"], bold=True, alignment=PP_ALIGN.CENTER)
    if data.get("contact"):
        add_textbox(slide, Inches(1.5), Inches(3.5), Inches(10.0), Inches(0.8),
                    data["contact"], font_size=18, font_color=style["text_muted"], alignment=PP_ALIGN.CENTER)
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(6.8), Inches(13.3), Inches(0.2))
    bar.fill.solid(); bar.fill.fore_color.rgb = style["secondary"]; bar.line.fill.background()


# ============================================================
# 主函数
# ============================================================

def generate_pptx(data, output_path, style_name="tech_blue"):
    style = STYLES.get(style_name, STYLES["tech_blue"])
    prs = Presentation()
    prs.slide_width = Inches(13.3)
    prs.slide_height = Inches(7.5)

    create_cover(prs, data, style)
    create_executive_summary(prs, data, style)
    create_pain_point(prs, data, style)
    create_solution(prs, data, style)
    create_market_analysis(prs, data, style)
    create_feasibility(prs, data, style)
    create_business_model(prs, data, style)
    create_financial_forecast(prs, data, style)
    create_team(prs, data, style)
    create_milestones(prs, data, style)
    create_risk_analysis(prs, data, style)
    create_valuation_model(prs, data, style)
    create_closing(prs, data, style)

    output_dir = os.path.dirname(output_path)
    if output_dir:
        os.makedirs(output_dir, exist_ok=True)
    prs.save(output_path)
    print(f"[成功] PPT已生成：{output_path}")
    print(f"[信息] 共 {len(prs.slides)} 页，风格：{style['name']}")
    return output_path


def validate_data(data):
    """验证JSON数据结构的完整性，确保关键字段存在"""
    required_fields = ["project_name", "tagline", "executive_summary"]
    missing = [f for f in required_fields if f not in data]
    if missing:
        warnings = f"[警告] 缺少必填字段：{', '.join(missing)}，部分页面可能显示不完整"
        print(warnings)
        return False
    return True


def main():
    parser = argparse.ArgumentParser(description="商业计划书PPT生成器")
    parser.add_argument("--data", required=True, help="商业计划书JSON数据文件路径")
    parser.add_argument("--output", default="./商业计划书.pptx", help="输出PPTX文件路径")
    parser.add_argument("--style", default="tech_blue", choices=list(STYLES.keys()), help="PPT风格")
    args = parser.parse_args()

    try:
        with open(args.data, "r", encoding="utf-8") as f:
            data = json.load(f)
    except FileNotFoundError:
        print(f"[错误] 数据文件不存在：{args.data}")
        print("[提示] 请检查 --data 参数指定的文件路径是否正确")
        sys.exit(1)
    except json.JSONDecodeError as e:
        print(f"[错误] JSON数据格式错误：{e}")
        print("[提示] 请确保数据文件是有效的JSON格式")
        sys.exit(1)

    print(f"[参数] 数据文件：{args.data}")
    print(f"[参数] 输出路径：{args.output}")
    print(f"[参数] 风格：{args.style}")

    validate_data(data)

    try:
        generate_pptx(data, args.output, args.style)
    except Exception as e:
        print(f"[错误] PPT生成失败：{e}")
        print("[提示] 请检查数据文件内容是否完整，或联系技术支持")
        sys.exit(1)


if __name__ == "__main__":
    main()
